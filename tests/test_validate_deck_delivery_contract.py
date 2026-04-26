import json
import sys
from pathlib import Path

from pptx import Presentation

from scripts import validate_deck_delivery_contract as contract_script


def _build_presentation(path: Path, slide_count: int) -> None:
    prs = Presentation()
    while len(prs.slides) < slide_count:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {len(prs.slides)}"
    prs.save(path)


def _director_sidecar() -> dict:
    return {
        "director": "Jesper Tyrer",
        "territory": "APAC",
        "open_land_deals": 8,
        "open_land_arr": 9_700_000,
        "q1_land_wins": 1,
        "q1_land_wins_arr": 1_900_000,
        "q1_land_lost": 14,
        "q1_land_lost_arr": 2_800_000,
        "q2_renewals": 1,
        "q2_renewals_acv": 455_000,
        "approved_2026": 2,
        "conditionally_approved": 1,
        "missing_stage3": 0,
    }


def test_main_validates_director_decks_and_exec_rollup(
    tmp_path: Path, monkeypatch
) -> None:
    workbooks_dir = tmp_path / "workbooks" / "2026-04-22"
    decks_dir = tmp_path / "decks" / "2026-04-22" / "land-only"
    workbooks_dir.mkdir(parents=True, exist_ok=True)
    decks_dir.mkdir(parents=True, exist_ok=True)
    (workbooks_dir / "jesper-tyrer.xlsx").write_text("", encoding="utf-8")

    _build_presentation(decks_dir / "jesper-tyrer-LAND.pptx", 6)
    (decks_dir / "jesper-tyrer-LAND.json").write_text(
        json.dumps(_director_sidecar()),
        encoding="utf-8",
    )
    _build_presentation(decks_dir / "Exec Rollup.pptx", 3)
    (decks_dir / "Exec Rollup.json").write_text(
        json.dumps({"deck": "Exec Rollup"}),
        encoding="utf-8",
    )

    monkeypatch.setattr(
        contract_script,
        "OUTPUT_ROOT",
        tmp_path / "output" / "deck_delivery_contract",
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "validate_deck_delivery_contract.py",
            "--date",
            "2026-04-22",
            "--workbooks-dir",
            str(workbooks_dir),
            "--decks-dir",
            str(decks_dir),
        ],
    )

    assert contract_script.main() == 0
    audit_path = (
        tmp_path
        / "output"
        / "deck_delivery_contract"
        / "2026-04-22"
        / "deck_delivery_contract_audit.json"
    )
    payload = json.loads(audit_path.read_text(encoding="utf-8"))
    assert payload["status"] == "ok"
    assert payload["expected_director_count"] == 1
    assert payload["validated_director_count"] == 1
    assert payload["directors"][0]["slug"] == "jesper-tyrer"
    assert payload["directors"][0]["slide_count"] == 6
    assert payload["exec_rollup"]["slide_count"] == 3


def test_main_fails_when_director_deck_is_missing(
    tmp_path: Path, monkeypatch
) -> None:
    workbooks_dir = tmp_path / "workbooks" / "2026-04-22"
    decks_dir = tmp_path / "decks" / "2026-04-22" / "land-only"
    workbooks_dir.mkdir(parents=True, exist_ok=True)
    decks_dir.mkdir(parents=True, exist_ok=True)
    (workbooks_dir / "jesper-tyrer.xlsx").write_text("", encoding="utf-8")

    _build_presentation(decks_dir / "Exec Rollup.pptx", 3)
    (decks_dir / "Exec Rollup.json").write_text(
        json.dumps({"deck": "Exec Rollup"}),
        encoding="utf-8",
    )

    monkeypatch.setattr(
        contract_script,
        "OUTPUT_ROOT",
        tmp_path / "output" / "deck_delivery_contract",
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "validate_deck_delivery_contract.py",
            "--date",
            "2026-04-22",
            "--workbooks-dir",
            str(workbooks_dir),
            "--decks-dir",
            str(decks_dir),
        ],
    )

    assert contract_script.main() == 1
    audit_path = (
        tmp_path
        / "output"
        / "deck_delivery_contract"
        / "2026-04-22"
        / "deck_delivery_contract_audit.json"
    )
    payload = json.loads(audit_path.read_text(encoding="utf-8"))
    assert payload["status"] == "failed"
    assert {
        "scope": "director",
        "slug": "jesper-tyrer",
        "issue": "missing_deck",
        "message": f"missing {decks_dir / 'jesper-tyrer-LAND.pptx'}",
    } in payload["failures"]
