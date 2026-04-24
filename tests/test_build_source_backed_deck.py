from __future__ import annotations

import json
from pathlib import Path
import sys

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.build_source_backed_deck import build_source_backed_deck  # noqa: E402


def _truth_packet(tmp_path: Path) -> Path:
    source_gate = tmp_path / "source_backed_publish_gate.json"
    source_gate.write_text(
        json.dumps(
            {
                "status": "ok",
                "source_run_dir": str(tmp_path / "sources" / "run-a"),
                "source_bundle_dir": str(tmp_path / "source_bundles" / "run-a"),
                "director_bundle_dir": str(tmp_path / "director_bundles" / "run-a"),
                "counts": {
                    "source_extract_count": 12,
                    "selected_source_count": 12,
                    "director_bundle_count": 2,
                },
            }
        )
        + "\n",
        encoding="utf-8",
    )
    source_bundle_dir = tmp_path / "source_bundles" / "run-a"
    source_bundle_dir.mkdir(parents=True)
    (source_bundle_dir / "source_bundle_manifest.json").write_text(
        json.dumps(
            {
                "source_run_id": "run-a",
                "summary": {
                    "forward_fallback_count": 1,
                    "territories": [
                        {
                            "director": "Jesper Tyrer",
                            "territory": "APAC",
                            "display_quarter_title": "Q2 2026",
                            "display_reason": "current_quarter",
                            "current_quarter_active_deals": 3,
                            "forward_quarter_active_deals": 1,
                        },
                        {
                            "director": "Megan Miceli",
                            "territory": "Canada",
                            "display_quarter_title": "Q3 2026",
                            "display_reason": "forward_quarter_fallback",
                            "current_quarter_active_deals": 0,
                            "forward_quarter_active_deals": 2,
                        },
                    ],
                },
            }
        )
        + "\n",
        encoding="utf-8",
    )
    truth_packet = tmp_path / "deck_truth_packet.json"
    truth_packet.write_text(
        json.dumps(
            {
                "snapshot_date": "2026-04-30",
                "status": "ok",
                "summary": {
                    "director_count": 2,
                    "metric_count": 12,
                    "claim_count": 12,
                    "high_blocker_count": 0,
                    "tieout_mismatch_count": 0,
                },
                "sources": {
                    "source_backed_publish_gate": str(source_gate),
                    "analyst_workbook_path": str(tmp_path / "run-a" / "source_backed_analyst_workbook.xlsx"),
                },
                "thinkcell": {
                    "recommended_element_names": [
                        "TruthStatus",
                        "RegionalRollupsTable",
                        "DirectorKpiTable",
                        "PublishBlockersTable",
                    ]
                },
                "regional_rollups": [
                    {
                        "region": "APAC",
                        "territories": ["APAC"],
                        "totals": {
                            "open_arr": 1_500_000,
                            "open_deals": 3,
                            "deal_risk_rows": 1,
                            "close_date_event_count": 0,
                        },
                    },
                    {
                        "region": "NAM",
                        "territories": ["Canada"],
                        "totals": {
                            "open_arr": 900_000,
                            "open_deals": 2,
                            "deal_risk_rows": 0,
                            "close_date_event_count": 0,
                        },
                    },
                ],
                "directors": [
                    {
                        "director": "Jesper Tyrer",
                        "territory": "APAC",
                        "open_arr": 1_500_000,
                        "open_deals": 3,
                        "deal_risk_rows": 1,
                        "tieout_mismatch_count": 0,
                        "bundle_issue_count": 0,
                    },
                    {
                        "director": "Megan Miceli",
                        "territory": "Canada",
                        "open_arr": 900_000,
                        "open_deals": 2,
                        "deal_risk_rows": 0,
                        "tieout_mismatch_count": 0,
                        "bundle_issue_count": 0,
                    },
                ],
            }
        )
        + "\n",
        encoding="utf-8",
    )
    return truth_packet


def test_build_source_backed_deck_creates_standard_package(tmp_path: Path) -> None:
    output_path = tmp_path / "deck.pptx"

    result = build_source_backed_deck(
        truth_packet_path=_truth_packet(tmp_path),
        output_path=output_path,
    )

    assert output_path.exists()
    assert result["status"] == "ok"
    assert result["source_run_id"] == "run-a"
    assert result["slide_count"] == 6

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 6
    assert (
        sum(
            1
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_table", False)
        )
        >= 3
    )
    assert (
        sum(
            1
            for slide in prs.slides
            for shape in slide.shapes
            if getattr(shape, "has_chart", False)
        )
        == 1
    )
    full_text = _deck_text(prs)
    assert "Truth Status: OK" in full_text
    assert "Quarter display logic" in full_text
    assert "Forward fallback" in full_text
    assert "think-cell" in full_text
    assert "€2.4M open ARR across 5 deals" in full_text
    assert "Top 2 books carry €2.4M (100%) of open ARR" in full_text
    assert "Use this packet as the standard deck production input" in full_text

    manifest_path = output_path.with_name("source_backed_deck_manifest.json")
    assert manifest_path.exists()
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert "Source-backed" in manifest["visual_contract"]["required_text"]
    assert "Truth Status: OK" in manifest["visual_contract"]["required_text"]


def _deck_text(prs: Presentation) -> str:
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.append(cell.text)
    return "\n".join(texts)
