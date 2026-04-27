"""Track F / F5 — render / overflow gates tests.

Positive control runs against the live anchor at
``~/Downloads/jesper-tyrer-LAND.pptx``; skipped if absent.

Negative controls synthesize a minimal .pptx that deliberately violates
each gate (off-slide table, title-region drift) and asserts the validator
catches it.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from pptx import Presentation
from pptx.util import Inches

from scripts.validate_deck_render import validate_render


REPO_ROOT = Path(__file__).resolve().parents[1]
APAC_DECK = Path("/Users/test/Downloads/jesper-tyrer-LAND.pptx")
TEMPLATE_PATH = REPO_ROOT / "assets" / "SimCorp_PPT_Template.pptx"


@pytest.mark.skipif(not APAC_DECK.exists(), reason="Live APAC pptx not present.")
def test_apac_anchor_render_passes():
    """The current production deck must clear all blocker gates.

    Footer-missing and legal-disclaimer-missing are warnings (forward-
    state debt) — acceptable. Blockers are: title placeholder missing,
    title drift outside region, and table off-slide.
    """
    report = validate_render(APAC_DECK)
    assert report["status"] == "pass", [
        f for f in report["findings"] if f["severity"] == "blocker"
    ]
    assert report["blocker_count"] == 0
    assert report["slide_count"] == 16


def _build_synth_deck(tmp_path: Path, *, off_slide_table: bool = False) -> Path:
    """Build a minimal one-slide deck. By default it passes; flags
    inject violations."""
    if not TEMPLATE_PATH.exists():
        pytest.skip("SimCorp template not present")
    prs = Presentation(str(TEMPLATE_PATH))
    # 'Title and Content' layout — has ph 144 + 145
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Find ph 144 and set its text
    for shape in slide.shapes:
        try:
            ph = shape.placeholder_format
            if ph and ph.idx == 144:
                shape.text_frame.text = "Synth Title"
        except (ValueError, AttributeError):
            continue

    if off_slide_table:
        # Place a 2x2 table that extends past the right edge.
        # Slide width 13.333" — start at 12" with width 5" -> ends at 17".
        slide.shapes.add_table(2, 2, Inches(12), Inches(2), Inches(5), Inches(1))

    out = tmp_path / "synth.pptx"
    prs.save(out)
    return out


@pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="SimCorp template not present.")
def test_off_slide_table_blocks(tmp_path):
    pptx = _build_synth_deck(tmp_path, off_slide_table=True)
    report = validate_render(pptx)
    assert report["status"] == "fail"
    assert any(f["code"] == "table_off_slide" for f in report["findings"])


@pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="SimCorp template not present.")
def test_title_drift_blocks_when_title_not_in_region(tmp_path):
    """Build a slide whose only text frames sit BELOW the title-region
    bottom; the validator should flag title_drift_outside_region."""
    prs = Presentation(str(TEMPLATE_PATH))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Find ph 144 and force its top way below the title region.
    for shape in slide.shapes:
        try:
            ph = shape.placeholder_format
            if ph and ph.idx == 144:
                shape.text_frame.text = "Synth Title"
                shape.top = Inches(5)  # well below the 1.5" region cap
        except (ValueError, AttributeError):
            continue
    out = tmp_path / "drift.pptx"
    prs.save(out)
    report = validate_render(out)
    # Note: this synth has only 1 slide, so it'll also fail the
    # "16-slide" expectation. We just check the title finding is there.
    assert any(f["code"] == "title_drift_outside_region" for f in report["findings"])


@pytest.mark.skipif(not APAC_DECK.exists(), reason="Live APAC pptx not present.")
def test_cover_slide_exempt_from_title_region():
    """Cover layout 'Title 1' has the title at ~4.2" by design.
    The validator must NOT flag the cover for title_drift_outside_region."""
    report = validate_render(APAC_DECK)
    cover_findings = [
        f
        for f in report["findings"]
        if f["path"].startswith("slides[1].")
        and f["code"] in ("title_drift_outside_region", "title_placeholder_missing")
    ]
    assert cover_findings == [], cover_findings
