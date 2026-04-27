"""Track F / F6 — golden visual regression tests.

Positive: rendering the live anchor against the committed baseline
produces no drift findings.

Negative: mutating the deck (adding a shape that overlaps the cover's
frozen region) triggers a frozen_region_drift blocker.

Both tests need libreoffice to run; skipped if the soffice binary
isn't available.
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest

from pptx import Presentation
from pptx.util import Inches

from scripts.monthly_platform import deck_contract
from scripts.validate_deck_visual_regression import (
    DEFAULT_BASELINE_PATH,
    _load_regions,
    DEFAULT_REGIONS_PATH,
    hash_deck,
    verify_against_baseline,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
APAC_DECK = Path("/Users/test/Downloads/jesper-tyrer-LAND.pptx")
SOFFICE_BIN = "/opt/homebrew/bin/soffice"


def _have_soffice() -> bool:
    return Path(SOFFICE_BIN).exists() or bool(shutil.which("soffice"))


@pytest.mark.skipif(not APAC_DECK.exists(), reason="Live APAC pptx not present.")
@pytest.mark.skipif(not _have_soffice(), reason="soffice (LibreOffice) not installed.")
@pytest.mark.skipif(
    not DEFAULT_BASELINE_PATH.exists(),
    reason="Baseline not captured; run validate_deck_visual_regression.py --mode capture",
)
def test_apac_anchor_visual_regression_passes():
    contract = deck_contract.load()
    regions = _load_regions(DEFAULT_REGIONS_PATH)
    snapshot = hash_deck(APAC_DECK, contract=contract, regions=regions, dpi=100)
    report = verify_against_baseline(snapshot, DEFAULT_BASELINE_PATH)
    assert report["status"] == "pass", report["findings"]
    assert report["blocker_count"] == 0
    assert report["slide_count"] == report["baseline_slide_count"]


@pytest.mark.skipif(not APAC_DECK.exists(), reason="Live APAC pptx not present.")
@pytest.mark.skipif(not _have_soffice(), reason="soffice (LibreOffice) not installed.")
@pytest.mark.skipif(
    not DEFAULT_BASELINE_PATH.exists(),
    reason="Baseline not captured.",
)
def test_drift_detected_when_cover_mutated(tmp_path):
    """Mutate the cover slide by adding a colored shape over the
    frozen-region area, then verify the regression validator surfaces
    a frozen_region_drift blocker."""
    mutated = tmp_path / "mutated.pptx"
    prs = Presentation(str(APAC_DECK))
    cover = prs.slides[0]
    # Drop a 4"x4" rectangle in the middle of the cover. This will
    # change the cover's full_slide hash regardless of where you put
    # it (cover declares full_slide as its frozen region).
    cover.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(4),
        Inches(2),
        Inches(4),
        Inches(4),
    )
    prs.save(str(mutated))

    contract = deck_contract.load()
    regions = _load_regions(DEFAULT_REGIONS_PATH)
    snapshot = hash_deck(mutated, contract=contract, regions=regions, dpi=100)
    report = verify_against_baseline(snapshot, DEFAULT_BASELINE_PATH)
    assert report["status"] == "fail"
    assert any(
        f["code"] == "frozen_region_drift" and "slides[1]" in f["path"]
        for f in report["findings"]
    )
