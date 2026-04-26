from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.build_sales_region_monthly_shell import build_shell_deck, region_component_books


def test_region_component_books_places_mea_under_emea() -> None:
    emea = region_component_books("EMEA")
    apac = region_component_books("APAC")
    assert any("Middle East & Africa - Mourad Essofi" == item for item in emea)
    assert all("Middle East & Africa" not in item for item in apac)


def test_build_shell_deck_creates_regional_template(tmp_path: Path) -> None:
    output_path = tmp_path / "emea-shell.pptx"
    result = build_shell_deck(
        region_name="EMEA",
        snapshot_date="2026-04-10",
        output_path=output_path,
    )
    assert output_path.exists()
    assert result["slide_count"] == 13
    prs = Presentation(str(output_path))
    assert len(prs.slides) == 13
    title_slide = prs.slides[0]
    texts = []
    for shape in title_slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text)
    assert "Sales Region Monthly" in texts
    assert "EMEA" in texts
