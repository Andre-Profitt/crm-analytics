from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))
TESTS_DIR = Path(__file__).resolve().parent
if str(TESTS_DIR) not in sys.path:
    sys.path.insert(0, str(TESTS_DIR))

from scripts.build_source_backed_deck import build_source_backed_deck  # noqa: E402
from scripts.polish_source_backed_deck_language import polish_deck_language  # noqa: E402
from scripts.validate_source_backed_deck_table_contract import (  # noqa: E402
    validate_table_contract,
)

from test_build_source_backed_deck import _truth_packet  # noqa: E402


def test_table_contract_passes_polished_source_backed_deck(tmp_path: Path) -> None:
    deck_path = _polished_deck(tmp_path)

    audit = validate_table_contract(
        deck_path=deck_path,
        snapshot_date="2026-04-30",
        source_run_id="run-a",
        output_path=tmp_path / "table-contract.json",
    )

    assert audit["status"] == "ok"
    assert audit["summary"]["finding_count"] == 0
    assert audit["checks"]["table_count"] == 5
    assert audit["checks"]["expected_table_count"] == 5
    assert [table["name"] for table in audit["checks"]["tables_checked"]] == [
        "publish_gate",
        "regional_rollup",
        "director_book",
        "quarter_policy",
        "production_handoff",
    ]
    assert Path(audit["output_path"]).exists()


def test_table_contract_blocks_header_and_style_drift(tmp_path: Path) -> None:
    deck_path = _polished_deck(tmp_path)
    presentation = Presentation(str(deck_path))
    table = next(shape.table for shape in presentation.slides[1].shapes if shape.has_table)
    table.cell(0, 0).text = "Gate drift"
    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = RGBColor(0, 0, 0)
    presentation.save(deck_path)

    audit = validate_table_contract(
        deck_path=deck_path,
        snapshot_date="2026-04-30",
        source_run_id="run-a",
        output_path=tmp_path / "table-contract.json",
    )

    assert audit["status"] == "blocked"
    issues = {finding["issue"] for finding in audit["findings"]}
    assert "table_headers_drifted" in issues
    assert "table_header_fill_drifted" in issues


def _polished_deck(tmp_path: Path) -> Path:
    deck_path = tmp_path / "deck.pptx"
    build_source_backed_deck(
        truth_packet_path=_truth_packet(tmp_path),
        output_path=deck_path,
    )
    polish_deck_language(
        deck_path=deck_path,
        manifest_path=deck_path.with_name("source_backed_deck_manifest.json"),
        snapshot_date="2026-04-30",
        source_run_id="run-a",
        output_path=tmp_path / "polish.json",
    )
    return deck_path
