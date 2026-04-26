from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.build_sales_region_gold_deck import build_gold_deck


def test_build_gold_deck_creates_populated_regional_example(tmp_path: Path) -> None:
    output_path = tmp_path / "emea-gold.pptx"
    result = build_gold_deck(
        region_snapshot_path=ROOT / "output" / "sales_region_snapshots" / "2026-04-10" / "emea.json",
        output_path=output_path,
    )
    assert output_path.exists()
    assert result["slide_count"] == 13

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 13
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    full_text = "\n".join(texts)
    assert "EMEA" in full_text
    assert "€31.3M" in full_text
    assert "€12.9M" in full_text
