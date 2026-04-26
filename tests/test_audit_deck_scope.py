import json
import sys
from pathlib import Path

from pptx import Presentation

from scripts import audit_deck_scope as scope_script


def test_main_writes_structured_audit_artifacts(
    tmp_path: Path, monkeypatch
) -> None:
    decks_dir = (
        tmp_path / "output" / "simcorp_director_decks" / "2026-04-22" / "land-only"
    )
    decks_dir.mkdir(parents=True, exist_ok=True)
    deck_path = decks_dir / "jesper-tyrer-LAND.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    if title is not None:
        title.text = "Executive Summary"
    prs.save(deck_path)
    (decks_dir / "jesper-tyrer-LAND.json").write_text(
        json.dumps(
            {
                "open_land_deals": 12,
                "open_land_arr": 1_200_000,
                "q1_land_lost": 2,
                "q1_land_lost_arr": 200_000,
                "q1_land_wins": 3,
                "q1_land_wins_arr": 300_000,
            }
        ),
        encoding="utf-8",
    )

    monkeypatch.setattr(scope_script, "ROOT", tmp_path)
    monkeypatch.setattr(scope_script, "VAULT", tmp_path / "obsidian")
    monkeypatch.setattr(scope_script, "OUTPUT_ROOT", tmp_path / "output" / "deck_scope_audit")
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "audit_deck_scope.py",
            "--date",
            "2026-04-22",
        ],
    )

    assert scope_script.main() == 0
    audit_path = (
        tmp_path
        / "output"
        / "deck_scope_audit"
        / "2026-04-22"
        / "deck_scope_audit.json"
    )
    payload = json.loads(audit_path.read_text(encoding="utf-8"))
    assert payload["status"] == "ok"
    assert payload["decks_audited"] == 1
    assert payload["flag_count"] == 0
    assert payload["results"][0]["slug"] == "jesper-tyrer"
    assert (
        tmp_path / "obsidian" / "Monthly" / "2026-04" / "scope-audit.md"
    ).exists()
