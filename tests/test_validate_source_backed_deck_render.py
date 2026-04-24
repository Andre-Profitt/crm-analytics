from __future__ import annotations

import subprocess
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation

from scripts import validate_source_backed_deck_render as render


def _write_deck(path: Path, *, slides: int = 2) -> None:
    prs = Presentation()
    for index in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(100_000, 100_000, 5_000_000, 600_000)
        textbox.text = f"Rendered slide {index + 1}"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def _write_slide(path: Path, *, label: str) -> None:
    image = Image.new("RGB", (960, 540), "white")
    image.putdata(
        [
            (
                (x * 17 + y * 3) % 255,
                (x * 11 + y * 7) % 255,
                (x * 5 + y * 13) % 255,
            )
            for y in range(540)
            for x in range(960)
        ]
    )
    draw = ImageDraw.Draw(image)
    draw.rectangle((80, 80, 880, 460), fill=(20, 85, 160))
    draw.text((120, 120), label, fill="white")
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def _fake_render_run(command: list[str], **kwargs: Any) -> subprocess.CompletedProcess[str]:
    command_text = " ".join(str(part) for part in command)
    if "--convert-to pdf" in command_text:
        outdir = Path(command[command.index("--outdir") + 1])
        deck_path = Path(command[-1])
        (outdir / f"{deck_path.stem}.pdf").write_bytes(b"x" * 25_000)
    elif "-png" in command:
        prefix = Path(command[-1])
        _write_slide(prefix.parent / "slide-1.png", label="Slide 1")
        _write_slide(prefix.parent / "slide-2.png", label="Slide 2")
    return subprocess.CompletedProcess(command, 0, stdout="ok", stderr="")


def test_validate_deck_render_passes_headless_pdf_png_outputs(
    tmp_path: Path,
    monkeypatch,
) -> None:
    deck_path = tmp_path / "deck.pptx"
    _write_deck(deck_path, slides=2)
    monkeypatch.setattr(render.shutil, "which", lambda name: f"/usr/bin/{name}")
    monkeypatch.setattr(render.subprocess, "run", _fake_render_run)

    audit = render.validate_deck_render(
        deck_path=deck_path,
        snapshot_date="2026-04-30",
        source_run_id="render-test",
        output_path=tmp_path / "render" / "audit.json",
    )

    assert audit["status"] == "ok"
    assert audit["checks"]["deck_slide_count"] == 2
    assert audit["checks"]["rendered_slide_count"] == 2
    assert audit["checks"]["rendered_png_checked"] is True
    assert Path(audit["pdf_path"]).exists()
    assert Path(audit["montage_path"]).exists()
    assert Path(audit["output_path"]).exists()


def test_validate_deck_render_blocks_slide_count_mismatch(
    tmp_path: Path,
    monkeypatch,
) -> None:
    deck_path = tmp_path / "deck.pptx"
    _write_deck(deck_path, slides=2)
    monkeypatch.setattr(render.shutil, "which", lambda name: f"/usr/bin/{name}")
    monkeypatch.setattr(render.subprocess, "run", _fake_render_run)

    audit = render.validate_deck_render(
        deck_path=deck_path,
        expected_slide_count=3,
        snapshot_date="2026-04-30",
        source_run_id="render-test",
        output_path=tmp_path / "render" / "audit.json",
    )

    assert audit["status"] == "blocked"
    assert any(
        finding["issue"] == "rendered_slide_count_mismatch"
        for finding in audit["findings"]
    )
