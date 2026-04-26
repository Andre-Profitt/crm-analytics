import json
import sys
from pathlib import Path

from pptx import Presentation

from scripts import audit_deck_fonts as audit_script


def _build_presentation(path: Path, slide_count: int) -> None:
    prs = Presentation()
    while len(prs.slides) < slide_count:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {len(prs.slides)}"
    prs.save(path)


def test_main_writes_warning_audit_when_font_issues_exist(
    tmp_path: Path, monkeypatch
) -> None:
    decks_dir = tmp_path / "decks" / "2026-04-22" / "land-only"
    decks_dir.mkdir(parents=True, exist_ok=True)
    _build_presentation(decks_dir / "jesper-tyrer-LAND.pptx", 3)
    _build_presentation(decks_dir / "Exec Rollup.pptx", 2)

    fake_detect = tmp_path / "fake_detect_font.py"
    fake_detect.write_text(
        """
import json, sys
name = sys.argv[-1]
if name.endswith('jesper-tyrer-LAND.pptx'):
    print(json.dumps({
        'font_missing_overall': ['calibri'],
        'font_missing_by_slide': {'2': ['calibri']},
        'font_substituted_overall': [],
        'font_substituted_by_slide': {}
    }))
else:
    print(json.dumps({
        'font_missing_overall': [],
        'font_missing_by_slide': {},
        'font_substituted_overall': [],
        'font_substituted_by_slide': {}
    }))
""".strip(),
        encoding="utf-8",
    )

    monkeypatch.setattr(audit_script, "OUTPUT_ROOT", tmp_path / "output" / "deck_font_audit")
    monkeypatch.setattr(audit_script, "DETECT_FONT_SCRIPT", fake_detect)
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "audit_deck_fonts.py",
            "--date",
            "2026-04-22",
            "--decks-dir",
            str(decks_dir),
        ],
    )

    assert audit_script.main() == 0
    audit_path = (
        tmp_path / "output" / "deck_font_audit" / "2026-04-22" / "deck_font_audit.json"
    )
    payload = json.loads(audit_path.read_text(encoding="utf-8"))
    assert payload["status"] == "warning"
    assert payload["deck_count"] == 2
    assert payload["decks_with_issues"] == 1
    assert payload["decks"][0]["deck"] == "Exec Rollup"
    assert payload["decks"][1]["deck"] == "jesper-tyrer-LAND"
    assert payload["decks"][1]["font_missing_overall"] == ["calibri"]

