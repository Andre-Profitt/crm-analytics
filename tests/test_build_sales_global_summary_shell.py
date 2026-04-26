from __future__ import annotations

import json
from pathlib import Path
import sys

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.build_sales_global_summary_shell import build_shell_deck


def test_build_shell_deck_creates_global_template(tmp_path: Path) -> None:
    output_path = tmp_path / "global-shell.pptx"
    result = build_shell_deck(
        snapshot_date="2026-04-10",
        output_path=output_path,
    )

    assert output_path.exists()
    assert result["slide_count"] == 8
    assert result["builder"] == "simcorp-native-shell"
    assert result["publish_safe"] is False

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 8
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text.append(shape.text)
    joined = "\n".join(all_text)
    assert "Sales Global Summary" in joined
    assert "Global Commercial Approval Overview" in joined
    assert "North America" in joined
    assert "Appendix and Guardrails" in joined


def test_build_shell_deck_renders_region_forward_quarter_footnote(tmp_path: Path) -> None:
    fill_payload_path = tmp_path / "fill-payload.json"
    fill_payload_path.write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "id": "apac-region-summary",
                        "slots": {
                            "region_name": "APAC",
                            "headline_pipeline_arr_q2": "€1.0M",
                            "q2_commit_arr": "€700K",
                            "q2_best_case_arr": "€300K",
                            "q2_omitted_arr": "€0",
                            "quarterly_pipeline_label": "Q3",
                            "quarterly_pipeline_title": "Q3 2026",
                            "quarterly_pipeline_display_reason": "forward_quarter_fallback",
                            "quarterly_pipeline_footnote": "No Q2 2026 in-scope pipeline; showing Q3 2026 forward-quarter outlook.",
                            "approval_rate_stage3_plus": "20.0%",
                            "renewal_open_acv": "€300K",
                            "top_risk": "APAC risk.",
                            "top_action": "APAC action.",
                        },
                    }
                ]
            }
        ),
        encoding="utf-8",
    )

    output_path = tmp_path / "global-shell-filled.pptx"
    build_shell_deck(
        snapshot_date="2026-04-30",
        output_path=output_path,
        fill_payload_path=fill_payload_path,
    )

    prs = Presentation(str(output_path))
    joined = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "Q3 active ARR" in joined
    assert (
        "No Q2 2026 in-scope pipeline; showing Q3 2026 forward-quarter outlook."
        in joined
    )
