from __future__ import annotations

import json
from pathlib import Path
import sys

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))
TESTS_DIR = Path(__file__).resolve().parent
if str(TESTS_DIR) not in sys.path:
    sys.path.insert(0, str(TESTS_DIR))

from scripts.build_source_backed_deck import build_source_backed_deck  # noqa: E402
from scripts.validate_source_backed_deck_visuals import validate_deck_visuals  # noqa: E402

from test_build_source_backed_deck import _truth_packet  # noqa: E402


def test_validate_source_backed_deck_visuals_passes_builder_output(tmp_path: Path) -> None:
    truth_packet_path = _truth_packet(tmp_path)
    deck_path = tmp_path / "deck.pptx"
    manifest = build_source_backed_deck(
        truth_packet_path=truth_packet_path,
        output_path=deck_path,
    )

    audit = validate_deck_visuals(
        deck_path=deck_path,
        truth_packet_path=truth_packet_path,
        manifest_path=Path(manifest["manifest_path"]),
        output_path=tmp_path / "audit.json",
    )

    assert audit["status"] == "ok"
    assert audit["summary"]["finding_count"] == 0
    assert audit["checks"]["slide_count"] == 6
    assert audit["checks"]["table_count"] >= 3
    assert audit["checks"]["chart_count"] == 1
    assert Path(audit["output_path"]).exists()


def test_validate_source_backed_deck_visuals_blocks_placeholders(tmp_path: Path) -> None:
    deck_path = tmp_path / "placeholder.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 6_000_000, 1_000_000)
    textbox.text = "TODO placeholder deck"
    presentation.save(deck_path)

    audit = validate_deck_visuals(
        deck_path=deck_path,
        min_slides=1,
        output_path=tmp_path / "placeholder-audit.json",
    )

    assert audit["status"] == "blocked"
    issues = {finding["issue"] for finding in audit["findings"]}
    assert "placeholder_text" in issues
    assert "table_count_below_minimum" in issues
    assert json.loads(Path(audit["output_path"]).read_text(encoding="utf-8"))["status"] == "blocked"
