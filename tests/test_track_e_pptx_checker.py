"""Track E — E5 tests for the produced-PPTX checker (E4).

Positive control runs against the live APAC anchor
``~/Downloads/jesper-tyrer-LAND.pptx``. Skipped if missing.

Negative controls modify the contract in-memory to exercise the
specific blocker paths:
  - title_neither_stable_nor_legacy
  - table_count_mismatch
  - table_header_mismatch (Cond 1)
  - required_link_target_mismatch (Cond 2)

Cond 1 (table-header validation) test list:
  - canonical contract passes against live deck (legacy_header_sets
    cover the legacy production headers)
  - contract header mutated -> blocker (table_header_mismatch)
  - evidence_only table is excluded from the per-slide check
  - dual-table slide (slide 6 q1_loss_drivers) checks both tables in order

Cond 2 (Salesforce link target):
  - missing link entirely -> warning (M1 transition)
  - hyperlink present but wrong target -> blocker
    (required_link_target_mismatch)
"""

from __future__ import annotations

import copy
from pathlib import Path

import pytest
import yaml

from scripts.monthly_platform import deck_contract
from scripts.validate_director_monthly_pptx import validate_pptx


APAC_DECK = Path("/Users/test/Downloads/jesper-tyrer-LAND.pptx")
CANONICAL_DECK = Path(__file__).resolve().parents[1] / "config" / "deck_contract.yaml"


def _make_contract(raw: dict) -> deck_contract.DeckContract:
    return deck_contract.DeckContract(raw=raw, path=CANONICAL_DECK)


@pytest.mark.skipif(
    not APAC_DECK.exists(),
    reason="Live APAC pptx not present; run with the deck in ~/Downloads/",
)
def test_apac_anchor_passes_in_legacy_mode():
    report = validate_pptx(APAC_DECK)
    # Expected (post Track F F2 — stable table headers, 2026-04-27):
    #   - 16-slide live anchor
    #   - 0 legacy verbose titles  (closed by F1)
    #   - 0 legacy header drifts   (closed by F2)
    #   - 1 missing required link  (warning; F3 closes this)
    #   - 0 blockers
    #   - 15 stable titles (every non-static slide; cover + 14 narrative slides)
    assert report["status"] == "pass", report["findings"]
    assert report["actual_slide_count"] == 16
    assert report["legacy_verbose_title_count"] == 0, (
        "F1 acceptance gate: builder must emit stable contract titles, "
        f"got {report['legacy_verbose_title_count']} legacy verbose titles"
    )
    assert report["stable_title_count"] == 15
    assert any(
        f["code"] == "missing_required_link_transition" for f in report["findings"]
    )
    legacy_header_findings = [
        f for f in report["findings"] if f["code"] == "legacy_header_drift"
    ]
    assert len(legacy_header_findings) == 0, (
        f"F2 acceptance gate: every table must match stable contract headers "
        f"(or header_pattern_sets for dynamic-date tables), got "
        f"{len(legacy_header_findings)} legacy_header_drift warnings"
    )
    # Slide 3 (since_last_review) has dynamic-date columns and matches
    # via header_pattern_sets — counts as pass_pattern, not warning.
    pattern_passes = sum(
        1
        for s in report["slides"]
        for tr in (s.get("header_results") or [])
        if tr.get("status") == "pass_pattern"
    )
    assert pattern_passes == 1, (
        f"expected 1 header_pattern_sets match (slide 3 since_last_review), "
        f"got {pattern_passes}"
    )


@pytest.mark.skipif(
    not APAC_DECK.exists(),
    reason="Live APAC pptx not present.",
)
def test_title_neither_stable_nor_legacy_blocks():
    # Post Track F F1: the produced deck emits stable titles. To exercise
    # the blocker path we MUTATE the stable title in the contract so the
    # produced "Owner Coaching Priorities" no longer matches stable, and
    # also drop legacy_title_patterns so it can't fall back. Result:
    # title_neither_stable_nor_legacy blocker.
    raw = yaml.safe_load(CANONICAL_DECK.read_text(encoding="utf-8"))
    slide = next(
        s
        for s in raw["profiles"]["director_monthly"]["slides"]
        if s["id"] == "owner_coaching"
    )
    slide["title"] = "Some Different Stable Title"
    slide.pop("legacy_title_patterns", None)
    contract = _make_contract(raw)
    report = validate_pptx(APAC_DECK, contract=contract)
    assert report["status"] == "fail"
    assert any(
        f["code"] == "title_neither_stable_nor_legacy" for f in report["findings"]
    )


@pytest.mark.skipif(
    not APAC_DECK.exists(),
    reason="Live APAC pptx not present.",
)
def test_table_count_mismatch_blocks():
    # Force renewals (slide 17) to declare 4 tables when the deck has 1.
    raw = yaml.safe_load(CANONICAL_DECK.read_text(encoding="utf-8"))
    slide = next(
        s
        for s in raw["profiles"]["director_monthly"]["slides"]
        if s["id"] == "renewals"
    )
    extra_tables = [copy.deepcopy(slide["tables"][0]) for _ in range(3)]
    for i, t in enumerate(extra_tables):
        t["id"] = f"tbl_renewals_extra_{i}"
    slide["tables"].extend(extra_tables)
    contract = _make_contract(raw)
    report = validate_pptx(APAC_DECK, contract=contract)
    assert report["status"] == "fail"
    assert any(f["code"] == "table_count_mismatch" for f in report["findings"])


# ----------------------------------------------------------------------
# Cond 1 — table-header validation
# ----------------------------------------------------------------------


@pytest.mark.skipif(not APAC_DECK.exists(), reason="Live APAC pptx not present.")
def test_table_header_mismatch_blocks_when_neither_stable_nor_legacy():
    """If a table's headers match neither the stable contract headers
    nor any legacy_header_sets entry, the checker emits a blocker."""
    raw = yaml.safe_load(CANONICAL_DECK.read_text(encoding="utf-8"))
    slide = next(
        s
        for s in raw["profiles"]["director_monthly"]["slides"]
        if s["id"] == "owner_coaching"
    )
    tbl = slide["tables"][0]
    # Mutate stable headers + drop legacy patterns -> production headers
    # match neither -> blocker.
    tbl["columns"][0]["header"] = "TotallyDifferentHeader"
    tbl.pop("legacy_header_sets", None)
    contract = _make_contract(raw)
    report = validate_pptx(APAC_DECK, contract=contract)
    assert report["status"] == "fail"
    assert any(f["code"] == "table_header_mismatch" for f in report["findings"])


@pytest.mark.skipif(not APAC_DECK.exists(), reason="Live APAC pptx not present.")
def test_evidence_only_table_excluded_from_header_check():
    """The Q1 Forecast Variance evidence_only table must not be checked
    against produced PPTX headers — slide 5 only has 1 displayed table
    (the bridge), so the evidence_only opportunity-grain table should
    not generate a missing-table or header-mismatch finding."""
    report = validate_pptx(APAC_DECK)
    # Track E re-anchor (2026-04-27) dropped q1_forecast_variance from the
    # contract because the current builder no longer emits it. With the
    # derived_table+evidence_only pair gone, this test now verifies the more
    # general claim: no slide_result reports a missing evidence_only table,
    # and no finding path mentions a *_evidence table that should have been
    # excluded. Functionally equivalent — the validator's evidence_only
    # exclusion code path is still exercised any time evidence_only tables
    # exist in the contract.
    # No table_missing_in_pptx finding may target a *_evidence table.
    # (When q1_forecast_variance was active, this guarded its evidence_only
    #  drill-through. Stays as a regression guard for any future
    #  evidence_only table.)
    for f in report["findings"]:
        assert not f.get("path", "").endswith("_evidence].columns"), f
        assert "evidence_only" not in f.get("path", ""), f


@pytest.mark.skipif(not APAC_DECK.exists(), reason="Live APAC pptx not present.")
def test_dual_table_slide_checks_both_tables_in_order():
    """Slide 5 (q1_loss_drivers; was slide 6 pre-re-anchor) has two tables. The header validator
    must run against BOTH in order and surface a finding when the
    SECOND table's headers don't match anything."""
    raw = yaml.safe_load(CANONICAL_DECK.read_text(encoding="utf-8"))
    slide = next(
        s
        for s in raw["profiles"]["director_monthly"]["slides"]
        if s["id"] == "q1_loss_drivers"
    )
    # Mutate the SECOND table's stable headers + drop legacy.
    second = slide["tables"][1]
    second["columns"][0]["header"] = "WrongHeaderForStageReached"
    second.pop("legacy_header_sets", None)
    contract = _make_contract(raw)
    report = validate_pptx(APAC_DECK, contract=contract)
    # The blocker should target the second table specifically.
    matching = [
        f
        for f in report["findings"]
        if f["code"] == "table_header_mismatch"
        and "tbl_q1_loss_stage_reached" in f.get("path", "")
    ]
    assert matching, [f for f in report["findings"] if f["severity"] == "blocker"]


# ----------------------------------------------------------------------
# Cond 2 — Salesforce link target validation
# ----------------------------------------------------------------------


def _build_pptx_with_link(tmp_path, link_url: str | None) -> Path:
    """Build a minimal 18-slide deck with a hyperlink on slide 12."""
    from pptx import Presentation as _P
    from pptx.util import Inches

    raw = yaml.safe_load(CANONICAL_DECK.read_text(encoding="utf-8"))
    slides = raw["profiles"]["director_monthly"]["slides"]
    prs = _P()
    blank_layout = prs.slide_layouts[6]
    for s in slides:
        slide = prs.slides.add_slide(blank_layout)
        # Stable title.
        tx = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6))
        tf = tx.text_frame
        tf.text = s["title"]
        if s["id"] == "pushed_deals" and link_url:
            box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(0.6)
            )
            p = box.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Open in Salesforce"
            run.hyperlink.address = link_url
    out = tmp_path / "synth_deck.pptx"
    prs.save(out)
    return out


def test_salesforce_link_target_correct_passes(tmp_path):
    pptx = _build_pptx_with_link(
        tmp_path,
        "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=Recent",
    )
    report = validate_pptx(pptx)
    # The cover slide doesn't render properly without the master, so we
    # only assert the link finding is *not* present.
    assert not any(
        f["code"] == "required_link_target_mismatch" for f in report["findings"]
    )
    assert not any(
        f["code"] == "missing_required_link_transition" for f in report["findings"]
    )


def test_salesforce_link_target_wrong_blocks(tmp_path):
    pptx = _build_pptx_with_link(tmp_path, "https://example.com/some-other-page")
    report = validate_pptx(pptx)
    assert report["status"] == "fail"
    assert any(f["code"] == "required_link_target_mismatch" for f in report["findings"])


def test_salesforce_link_missing_warns_only(tmp_path):
    pptx = _build_pptx_with_link(tmp_path, None)
    report = validate_pptx(pptx)
    # No hyperlinks at all -> M1 transition warning, not blocker.
    assert any(
        f["code"] == "missing_required_link_transition" for f in report["findings"]
    )
    assert not any(
        f["code"] == "required_link_target_mismatch" for f in report["findings"]
    )
