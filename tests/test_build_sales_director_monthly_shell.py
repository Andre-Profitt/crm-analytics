from __future__ import annotations

import json
from pathlib import Path
import sys

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.build_sales_director_monthly_shell import build_shell_deck


def test_build_shell_deck_creates_director_template(tmp_path: Path) -> None:
    output_path = tmp_path / "director-shell.pptx"
    result = build_shell_deck(
        director_name="Jane Doe",
        territory="APAC",
        snapshot_date="2026-04-10",
        output_path=output_path,
    )

    assert output_path.exists()
    assert result["slide_count"] == 15
    assert result["builder"] == "simcorp-native-shell"

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 15
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text.append(shape.text)
    joined = "\n".join(all_text)
    assert "Salesforce Hygiene and Activity Controls" in joined
    assert "Missing Win/Loss Reason" in joined
    assert "Overdue Close Date Open Opportunities" in joined


def test_build_shell_deck_legacy_js_builder_requires_explicit_opt_in(
    tmp_path: Path,
) -> None:
    output_path = tmp_path / "director-shell-legacy.pptx"
    result = build_shell_deck(
        director_name="Jane Doe",
        territory="APAC",
        snapshot_date="2026-04-10",
        output_path=output_path,
        allow_legacy_js_builder=True,
    )

    assert output_path.exists()
    assert result["builder"] == "js-v2"


def test_build_shell_deck_renders_forward_quarter_footnote(tmp_path: Path) -> None:
    fill_payload_path = tmp_path / "fill-payload.json"
    fill_payload_path.write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "id": "quarterly-pipeline",
                        "slots": {
                            "headline_pipeline_arr_q2": "€1.0M",
                            "q2_commit_arr": "€700K",
                            "q2_best_case_arr": "€300K",
                            "q2_omitted_arr": "€200K",
                            "quarterly_pipeline_label": "Q3",
                            "quarterly_pipeline_title": "Q3 2026",
                            "quarterly_pipeline_display_reason": "forward_quarter_fallback",
                            "quarterly_pipeline_footnote": "No Q2 2026 in-scope pipeline; showing Q3 2026 forward-quarter outlook.",
                        },
                    },
                    {
                        "id": "pipeline-coverage-intel",
                        "slots": {
                            "pipeline_coverage_statement": "Use Q3 2026 active ARR of €1.0M and weighted ARR of €0.8M as the current proxy.",
                            "weighted_pipeline_arr": "€0.8M",
                            "top_opportunities": [
                                {
                                    "opportunity": "Q3 Deal A",
                                    "arr_eur": "€700K",
                                    "stage": "3 - Commit",
                                    "next_action": "Close",
                                }
                            ],
                            "stale_arr": "€0",
                            "aging_arr": "€0",
                            "data_quality_backlog": "0",
                            "competitive_loss_watchlist": [],
                            "quarterly_pipeline_label": "Q3",
                            "quarterly_pipeline_title": "Q3 2026",
                        },
                    },
                ]
            }
        ),
        encoding="utf-8",
    )

    output_path = tmp_path / "director-shell-filled.pptx"
    result = build_shell_deck(
        director_name="Jane Doe",
        territory="APAC",
        snapshot_date="2026-04-30",
        output_path=output_path,
        fill_payload_path=fill_payload_path,
    )

    assert result["builder"] == "simcorp-native"
    assert result["slide_count"] == 15
    prs = Presentation(str(output_path))
    assert len(prs.slides) == 15
    joined = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    assert "Q3 Active Pipeline" in joined
    assert (
        "No Q2 2026 in-scope pipeline; showing Q3 2026 forward-quarter outlook."
        in joined
    )
    assert "Top Q3 opportunities" in joined
