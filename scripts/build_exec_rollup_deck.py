"""
Executive roll-up deck, one PowerPoint across all nine directors.

Summarises the monthly review into six slides so the MD sees the whole
business without the per-deal detail. Reads the director workbooks that
extract_director_live.py produced, uses the same SimCorp template as the
per-director decks.
"""

import argparse
import sys
from datetime import datetime
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))

try:
    from monthly_platform.period import sheet_names
except ModuleNotFoundError:  # pragma: no cover
    pass

# Reuse rendering helpers from the per-director builder
from build_deck_from_excel import (  # type: ignore
    DEFAULT_TEMPLATE,
    LY_END_SLIDE,
    LY_TITLE_1,
    LY_TITLE_CONTENT,
    _add_table,
    _fmt_eur,
    _meur,
    _set_ph,
    _set_ph_lines,
    _unw,
    _wtd,
)


DIRECTORS = [
    ("Jesper Tyrer", "APAC", "jesper-tyrer.xlsx"),
    ("Sarah Pittroff", "EMEA Central", "sarah-pittroff.xlsx"),
    ("Dan Peppett", "EMEA UK & Ireland", "dan-peppett.xlsx"),
    ("Christian Ebbesen", "EMEA NE", "christian-ebbesen.xlsx"),
    ("Francois Thaury", "EMEA South West", "francois-thaury.xlsx"),
    ("Mourad Essofi", "EMEA MEA", "mourad-essofi.xlsx"),
    ("Patrick Gaughan", "NA Asset Mgmt", "patrick-gaughan.xlsx"),
    ("Megan Miceli", "NA Canada", "megan-miceli.xlsx"),
    ("Adam Steinhaus", "NA Insurance", "adam-steinhaus.xlsx"),
]


def _load_director(wb_path):
    """Return the relevant sheets as lists of dicts, keyed by semantic name."""
    SN = sheet_names()
    wb = load_workbook(wb_path, data_only=True)
    # Map semantic keys to actual sheet names for the sheets we need
    _keys = [
        "pipeline_open",
        "won_lost",
        "commercial_approval",
        "renewals",
        "q1_movement",
    ]
    out = {}
    for key in _keys:
        sn = SN[key]
        if sn not in wb.sheetnames:
            out[sn] = []
            continue
        ws = wb[sn]
        rows = list(ws.iter_rows(values_only=True))
        if len(rows) < 2:
            out[sn] = []
            continue
        headers = [str(h or "").strip() for h in rows[0]]
        out[sn] = [
            {headers[i]: (v if v is not None else "") for i, v in enumerate(r)}
            for r in rows[1:]
        ]
    return out


def _quarter(s):
    s = str(s or "")
    if not s.startswith("2026"):
        return None
    try:
        m = int(s[5:7])
    except ValueError:
        return None
    return f"Q{(m - 1) // 3 + 1}"


def _is_land(r):
    return str(r.get("Type", "")).strip().lower() == "land"


def _in_q1q2(r):
    cd = str(r.get("Close Date", "") or "")[:10]
    return cd >= "2026-01-01" and cd <= "2026-06-30"


def _aggregate(workbooks_dir):
    """Roll every director up into totals + a per-director table."""
    SN = sheet_names()
    rollup = {
        "open_deals": 0,
        "open_unwtd": 0.0,
        "open_wtd": 0.0,
        "q1_won": 0,
        "q1_won_arr": 0.0,
        "q1_lost": 0,
        "q1_lost_arr": 0.0,
        "q1_slips_open": 0,
        "q1_slip_arr_risk": 0.0,
        "approved_2026": 0,
        "approved_2026_arr": 0.0,
        "conditionally_approved": 0,
        "conditionally_approved_arr": 0.0,
        "missing_approval": 0,
        "missing_approval_arr": 0.0,
    }
    per_director = []
    for director, territory, fname in DIRECTORS:
        wb_path = workbooks_dir / fname
        if not wb_path.exists():
            continue
        sheets = _load_director(wb_path)
        # Land + Q1-Q2 scope for monthly review
        land_open = [
            r for r in sheets[SN["pipeline_open"]] if _is_land(r) and _in_q1q2(r)
        ]
        won_lost_land = [
            r for r in sheets[SN["won_lost"]] if _is_land(r) and _in_q1q2(r)
        ]
        q1_won = [
            r
            for r in won_lost_land
            if "Won" in str(r.get("Stage", ""))
            and _quarter(r.get("Close Date")) == "Q1"
        ]
        q1_lost = [
            r
            for r in won_lost_land
            if "Won" not in str(r.get("Stage", ""))
            and _quarter(r.get("Close Date")) == "Q1"
        ]
        # Q1 slips still open = deals in land_open with opp name present in Q1 Movement
        q1_slip_names = {
            r.get("Opportunity")
            for r in sheets["Q1 Movement"]
            if r.get("Movement") == "Q1 Slipped"
        }
        slips_still_open = [
            r for r in land_open if r.get("Opportunity") in q1_slip_names
        ]

        approved = [
            r
            for r in sheets["Commercial Approval"]
            if str(r.get("Status", "")).strip() == "Approved 2026"
        ]
        cond = [
            r
            for r in sheets["Commercial Approval"]
            if "Conditionally" in str(r.get("Status", ""))
            or "Pending" in str(r.get("Status", ""))
        ]
        missing = [
            r
            for r in sheets["Commercial Approval"]
            if "Missing" in str(r.get("Status", ""))
        ]

        d = {
            "director": director,
            "territory": territory,
            "open_deals": len(land_open),
            "open_unwtd": sum(_unw(r) for r in land_open),
            "open_wtd": sum(_wtd(r) for r in land_open),
            "q1_won": len(q1_won),
            "q1_won_arr": sum(_unw(r) for r in q1_won),
            "q1_lost": len(q1_lost),
            "q1_lost_arr": sum(_unw(r) for r in q1_lost),
            "q1_slips_open": len(slips_still_open),
            "q1_slip_arr_risk": sum(_unw(r) for r in slips_still_open),
            "approved_2026": len(approved),
            "approved_2026_arr": sum(_unw(r) for r in approved),
            "conditionally_approved": len(cond),
            "conditionally_approved_arr": sum(_unw(r) for r in cond),
            "missing_approval": len(missing),
            "missing_approval_arr": sum(_unw(r) for r in missing),
        }
        per_director.append(d)
        for k in rollup:
            rollup[k] += d[k]
    return rollup, per_director


def slide_cover(prs, when):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_1])
    _set_ph(slide, 14, "Sales Directors Monthly")
    _set_ph(slide, 15, f"Executive Roll-up  |  {when}")
    _set_ph(slide, 16, "All nine MD-1 territories")


def slide_exec_summary(prs, rollup):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    # Data-forward headline: pipeline + Q1 outcome + exposure in one line.
    headline = (
        f"{rollup['open_deals']} open Land deals, "
        f"{_fmt_eur(rollup['open_unwtd'])} unweighted. "
        f"Q1: {rollup['q1_won']} won / {rollup['q1_lost']} lost. "
        f"{rollup['q1_slips_open']} Q1 slips still open "
        f"({_fmt_eur(rollup['q1_slip_arr_risk'])} exposed)."
    )
    _set_ph(slide, 144, headline)
    _set_ph(
        slide, 145, "Rollup across nine MD-1 sales directors, FY26 Land pipeline Q1-Q2."
    )

    lines = [
        f"Open Land pipeline: {rollup['open_deals']} deals, "
        f"{_fmt_eur(rollup['open_unwtd'])} unweighted, "
        f"{_fmt_eur(rollup['open_wtd'])} weighted.",
        "",
        f"Q1 outcome: {rollup['q1_won']} wins ({_fmt_eur(rollup['q1_won_arr'])}), "
        f"{rollup['q1_lost']} losses ({_fmt_eur(rollup['q1_lost_arr'])}).",
        f"Q1 slips still open: {rollup['q1_slips_open']} deals at risk, "
        f"{_fmt_eur(rollup['q1_slip_arr_risk'])}.",
        "",
        f"Commercial approvals: {rollup['approved_2026']} approved 2026 "
        f"({_fmt_eur(rollup['approved_2026_arr'])}), "
        f"{rollup['conditionally_approved']} conditionally approved "
        f"({_fmt_eur(rollup['conditionally_approved_arr'])}).",
        f"Missing approval at Stage 3+: {rollup['missing_approval']} deals, "
        f"{_fmt_eur(rollup['missing_approval_arr'])}.",
    ]
    _set_ph_lines(slide, 22, lines)


def slide_pipeline_by_region(prs, per_director):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    if per_director:
        sorted_dirs = sorted(per_director, key=lambda x: -x["open_unwtd"])
        leader = sorted_dirs[0]
        leader_share = (
            leader["open_unwtd"] / sum(d["open_unwtd"] for d in per_director)
            if sum(d["open_unwtd"] for d in per_director)
            else 0
        )
        _set_ph(
            slide,
            144,
            (
                f"{leader['territory']} carries {_fmt_eur(leader['open_unwtd'])} "
                f"({leader_share * 100:.0f}% of global Land book)."
            ),
        )
    else:
        _set_ph(slide, 144, "Open Land Pipeline, by Region")
    _set_ph(slide, 145, "Unweighted and weighted ARR, Q1-Q2 FY26. Sorted by open ARR.")

    rows = [
        [
            "Director",
            "Territory",
            "Deals",
            "ARR Unwtd (mEUR)",
            "ARR Wtd (mEUR)",
        ]
    ]
    for d in sorted(per_director, key=lambda x: -x["open_unwtd"]):
        rows.append(
            [
                d["director"],
                d["territory"],
                str(d["open_deals"]),
                _meur(d["open_unwtd"]),
                _meur(d["open_wtd"]),
            ]
        )
    total_deals = sum(d["open_deals"] for d in per_director)
    total_u = sum(d["open_unwtd"] for d in per_director)
    total_w = sum(d["open_wtd"] for d in per_director)
    rows.append(
        [
            "Total",
            "",
            str(total_deals),
            _meur(total_u),
            _meur(total_w),
        ]
    )
    _add_table(
        slide,
        rows,
        0.9,
        2.0,
        11.5,
        row_height=0.34,
        col_widths=[2.2, 2.4, 1.2, 2.3, 2.3],
    )


def slide_q1_retrospective(prs, per_director):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    total_won = sum(d["q1_won"] for d in per_director)
    total_lost = sum(d["q1_lost"] for d in per_director)
    total_slip_arr = sum(d["q1_slip_arr_risk"] for d in per_director)
    total_won_arr = sum(d["q1_won_arr"] for d in per_director)
    total_lost_arr = sum(d["q1_lost_arr"] for d in per_director)
    biggest_exposure = max(
        per_director,
        key=lambda d: d["q1_slip_arr_risk"] + d["q1_lost_arr"],
        default=None,
    )
    who = biggest_exposure["territory"] if biggest_exposure else "n/a"
    _set_ph(
        slide,
        144,
        (
            f"Q1 closed: {total_won} won ({_fmt_eur(total_won_arr)}) vs "
            f"{total_lost} lost ({_fmt_eur(total_lost_arr)}). "
            f"{_fmt_eur(total_slip_arr)} Q1 slip exposure still open; "
            f"{who} carries the largest share."
        ),
    )
    _set_ph(
        slide,
        145,
        "Ranked by Q1 slip + loss ARR. Review territories with highest exposure first.",
    )

    rows = [
        [
            "Director",
            "Territory",
            "Q1 Won",
            "Won ARR (mEUR)",
            "Q1 Lost",
            "Lost ARR (mEUR)",
            "Slips Open",
            "Slip ARR at Risk (mEUR)",
        ]
    ]
    for d in sorted(
        per_director, key=lambda x: -(x["q1_slip_arr_risk"] + x["q1_lost_arr"])
    ):
        rows.append(
            [
                d["director"],
                d["territory"],
                str(d["q1_won"]),
                _meur(d["q1_won_arr"]),
                str(d["q1_lost"]),
                _meur(d["q1_lost_arr"]),
                str(d["q1_slips_open"]),
                _meur(d["q1_slip_arr_risk"]),
            ]
        )
    totals = [
        "Total",
        "",
        str(sum(d["q1_won"] for d in per_director)),
        _meur(sum(d["q1_won_arr"] for d in per_director)),
        str(sum(d["q1_lost"] for d in per_director)),
        _meur(sum(d["q1_lost_arr"] for d in per_director)),
        str(sum(d["q1_slips_open"] for d in per_director)),
        _meur(sum(d["q1_slip_arr_risk"] for d in per_director)),
    ]
    rows.append(totals)
    _add_table(
        slide,
        rows,
        0.5,
        2.0,
        12.3,
        row_height=0.34,
        col_widths=[1.9, 2.0, 1.0, 1.6, 1.0, 1.6, 1.1, 2.1],
    )


def slide_commercial_approvals(prs, rollup, per_director):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    biggest_missing = max(
        per_director, key=lambda d: d["missing_approval_arr"], default=None
    )
    who_missing = (
        biggest_missing["territory"]
        if biggest_missing and biggest_missing["missing_approval"]
        else "none"
    )
    _set_ph(
        slide,
        144,
        (
            f"Approvals: {rollup['approved_2026']} approved 2026 "
            f"({_fmt_eur(rollup['approved_2026_arr'])}), "
            f"{rollup['missing_approval']} missing Stage 3+ "
            f"({_fmt_eur(rollup['missing_approval_arr'])}). "
            f"Biggest exposure: {who_missing}."
        ),
    )
    _set_ph(
        slide,
        145,
        f"{rollup['conditionally_approved']} conditionally approved deals "
        f"({_fmt_eur(rollup['conditionally_approved_arr'])}) sit in front of the committee.",
    )

    rows = [
        [
            "Director",
            "Territory",
            "Approved YTD",
            "Approved ARR (mEUR)",
            "Cond Approved",
            "Cond ARR (mEUR)",
            "Missing Stage 3+",
            "Missing ARR (mEUR)",
        ]
    ]
    for d in per_director:
        rows.append(
            [
                d["director"],
                d["territory"],
                str(d["approved_2026"]),
                _meur(d["approved_2026_arr"]),
                str(d["conditionally_approved"]),
                _meur(d["conditionally_approved_arr"]),
                str(d["missing_approval"]),
                _meur(d["missing_approval_arr"]),
            ]
        )
    rows.append(
        [
            "Total",
            "",
            str(rollup["approved_2026"]),
            _meur(rollup["approved_2026_arr"]),
            str(rollup["conditionally_approved"]),
            _meur(rollup["conditionally_approved_arr"]),
            str(rollup["missing_approval"]),
            _meur(rollup["missing_approval_arr"]),
        ]
    )
    _add_table(
        slide,
        rows,
        0.5,
        2.0,
        12.3,
        row_height=0.34,
        col_widths=[1.9, 2.0, 1.2, 1.7, 1.3, 1.7, 1.3, 1.7],
    )


def slide_action_items(prs, per_director):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    top_slip_preview = sorted(per_director, key=lambda x: -x["q1_slip_arr_risk"])[:3]
    if top_slip_preview and top_slip_preview[0]["q1_slip_arr_risk"] > 0:
        top_slip_total = sum(d["q1_slip_arr_risk"] for d in top_slip_preview)
        _set_ph(
            slide,
            144,
            (
                f"Top 3 slip-risk territories carry {_fmt_eur(top_slip_total)} "
                f"of at-risk ARR. Prioritise 1:1s here."
            ),
        )
    else:
        _set_ph(slide, 144, "Risks and Action Items")

    # Highlights: top 3 regions by slip risk, top 3 by missing approval ARR,
    # biggest conditionally approved region, biggest Q1 lost region.
    top_slip = sorted(per_director, key=lambda x: -x["q1_slip_arr_risk"])[:3]
    top_missing = sorted(per_director, key=lambda x: -x["missing_approval_arr"])[:3]
    top_cond = sorted(per_director, key=lambda x: -x["conditionally_approved_arr"])[:3]
    top_lost = sorted(per_director, key=lambda x: -x["q1_lost_arr"])[:3]

    def _line(d, amount_key, count_key=None):
        arr = d[amount_key]
        count = d[count_key] if count_key else None
        arr_s = _fmt_eur(arr)
        if count is not None:
            return f"  {d['territory']}: {count} deals, {arr_s}"
        return f"  {d['territory']}: {arr_s}"

    lines = [
        "Q1 slips still open, biggest exposures",
        *[
            _line(d, "q1_slip_arr_risk", "q1_slips_open")
            for d in top_slip
            if d["q1_slips_open"]
        ],
        "",
        "Missing commercial approval at Stage 3+",
        *[
            _line(d, "missing_approval_arr", "missing_approval")
            for d in top_missing
            if d["missing_approval"]
        ],
        "",
        "Largest conditionally approved books",
        *[
            _line(d, "conditionally_approved_arr", "conditionally_approved")
            for d in top_cond
            if d["conditionally_approved"]
        ],
        "",
        "Q1 loss concentration",
        *[_line(d, "q1_lost_arr", "q1_lost") for d in top_lost if d["q1_lost"]],
    ]
    _set_ph_lines(slide, 22, lines)


def slide_end(prs):
    try:
        prs.slides.add_slide(prs.slide_layouts[LY_END_SLIDE])
    except IndexError:
        prs.slides.add_slide(prs.slide_layouts[LY_TITLE_1])


def build_exec_deck(workbooks_dir, template_path, output_path):
    rollup, per_director = _aggregate(workbooks_dir)

    prs = Presentation(str(template_path))
    # Drop template sample slides

    pres_elem = prs.part._element
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    sld_id_lst = pres_elem.find("p:sldIdLst", nsmap)
    if sld_id_lst is not None:
        for sld_id in list(sld_id_lst):
            rel_id = sld_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            if rel_id:
                prs.part.drop_rel(rel_id)
            sld_id_lst.remove(sld_id)

    when = datetime.now().strftime("%B %Y")

    slide_cover(prs, when)
    print("  [OK] 1. Cover")
    slide_exec_summary(prs, rollup)
    print("  [OK] 2. Executive Summary")
    slide_pipeline_by_region(prs, per_director)
    print("  [OK] 3. Pipeline by Region")
    slide_q1_retrospective(prs, per_director)
    print("  [OK] 4. Q1 Retrospective")
    slide_commercial_approvals(prs, rollup, per_director)
    print("  [OK] 5. Commercial Approvals")
    slide_action_items(prs, per_director)
    print("  [OK] 6. Risks and Actions")
    slide_end(prs)
    print("  [OK] 7. End")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")

    # Sidecar JSON so tie-out and scope audit can reconcile the exec deck
    # the same way they reconcile per-director decks.
    import json as _json

    sidecar = {
        "deck": "Exec Rollup",
        "scope": "All 9 directors, FY26 Land pipeline, Q1-Q2",
        "built_at": datetime.now().isoformat(timespec="seconds"),
        "director_count": len(per_director),
        "open_land_deals": rollup["open_deals"],
        "open_land_arr": rollup["open_unwtd"],
        "open_land_arr_wtd": rollup["open_wtd"],
        "q1_land_wins": rollup["q1_won"],
        "q1_land_wins_arr": rollup["q1_won_arr"],
        "q1_land_lost": rollup["q1_lost"],
        "q1_land_lost_arr": rollup["q1_lost_arr"],
        "q1_slips_open": rollup["q1_slips_open"],
        "q1_slip_arr_risk": rollup["q1_slip_arr_risk"],
        "approved_2026": rollup["approved_2026"],
        "approved_2026_arr": rollup["approved_2026_arr"],
        "conditionally_approved": rollup["conditionally_approved"],
        "missing_stage3": rollup["missing_approval"],
        "missing_stage3_arr": rollup["missing_approval_arr"],
        "by_director": {
            d["director"]: {
                "territory": d["territory"],
                "open_land_deals": d["open_deals"],
                "open_land_arr": d["open_unwtd"],
                "q1_land_wins": d["q1_won"],
                "q1_land_lost": d["q1_lost"],
            }
            for d in per_director
        },
    }
    sidecar_path = output_path.with_suffix(".json")
    sidecar_path.write_text(_json.dumps(sidecar, indent=2) + "\n")
    print(f"Sidecar: {sidecar_path}")


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--workbooks-dir",
        type=Path,
        default=Path("output/director_live_workbooks")
        / datetime.now().strftime("%Y-%m-%d"),
    )
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
    )
    args = parser.parse_args()
    if args.output is None:
        args.output = (
            Path("output/simcorp_director_decks")
            / datetime.now().strftime("%Y-%m-%d")
            / "Exec Rollup.pptx"
        )
    build_exec_deck(args.workbooks_dir, args.template, args.output)


if __name__ == "__main__":
    main()
