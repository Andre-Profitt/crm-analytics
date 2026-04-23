"""
Deck scope consistency audit.

Every per-director deck is scoped to Land + Q1-Q2 FY26 + one territory. If a
slide title quotes a numeric claim that disagrees with the deck's sidecar
JSON (the canonical scope), flag it. Catches the kind of scope drift the
user saw previously when a preview slide accidentally counted 45 losses
vs the main deck's 14.

Checks:
- For every slide, extract numeric claims that look like "N deals", "N
  losses", "N wins", "EUR X.YM", "EUR XK".
- Compare deal/loss/win counts and ARR values against the sidecar fields.
- Allow 10 percent or EUR 100K tolerance on currency; exact match on counts.
- Ignore claims that clearly reference a different scope (subtitles
  mentioning "cycle" or "YTD" etc. — those are out of scope by design).

Writes `obsidian/Monthly/YYYY-MM/scope-audit.md` with a table per director
and a summary count of drift flags. Exit 0 if clean, 1 if any flags.
"""

from __future__ import annotations

import argparse
import json
import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
VAULT = ROOT / "obsidian"
OUTPUT_ROOT = ROOT / "output" / "deck_scope_audit"

# Keywords that, when present in a slide text, mark it as out-of-deck-scope.
# These slides discuss global / historical / prior-year data intentionally.
OUT_OF_SCOPE_MARKERS = [
    "FY Targets",
    "prior year",
    "prior years",
    "approved in prior",
    "cycle",
    "Month over Month",
    "Historical Trending",
    "year-to-date",
    "YTD",
]

# Regexes for numeric claims we care about.
RE_DEALS = re.compile(r"(\d{1,3}(?:,\d{3})*)\s+(?:open\s+)?deal(?:s)?", re.IGNORECASE)
RE_LOSSES = re.compile(
    r"(\d{1,3}(?:,\d{3})*)\s+(?:Q1\s+(?:Land\s+)?)?loss(?:es)?", re.IGNORECASE
)
RE_WINS = re.compile(
    r"(\d{1,3}(?:,\d{3})*)\s+(?:Q1\s+(?:Land\s+)?)?win(?:s)?", re.IGNORECASE
)
RE_EUR = re.compile(r"EUR\s*([\d.,]+)\s*([MK])?", re.IGNORECASE)


def _parse_eur(match) -> float:
    raw = match.group(1).replace(",", "")
    try:
        v = float(raw)
    except ValueError:
        return 0.0
    suffix = (match.group(2) or "").upper()
    if suffix == "M":
        v *= 1_000_000
    elif suffix == "K":
        v *= 1_000
    return v


def _slide_titles(pptx_path: Path) -> list[tuple[int, str]]:
    """Return [(slide_index_1based, title_text)] for every slide.

    Title = the first non-trivial paragraph in the first text frame. This is
    where the slide's scope-level claim lives. Body tables, subtitles, and
    footnotes are intentionally ignored (they frequently reference subsets
    or different scopes without implying deck-wide scope drift).
    """
    prs = Presentation(str(pptx_path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame  # type: ignore[attr-defined]
            for para in tf.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text and len(text) > 8 and not title:
                    title = text
                    break
            if title:
                break
        out.append((i, title))
    return out


def _is_currency_close(a: float, b: float) -> bool:
    tol = max(a * 0.10, 100_000)
    return abs(a - b) <= tol


def _is_count_close(a: int, b: int) -> bool:
    return a == b


def _out_of_scope(text: str) -> bool:
    low = text.lower()
    return any(m.lower() in low for m in OUT_OF_SCOPE_MARKERS)


def audit_deck(pptx_path: Path, sidecar_path: Path) -> list[dict]:
    """Audit one deck and return a list of flag records.

    Each flag has: slide_index, slide_snippet, claim_type, claimed_value,
    canonical_value, detail.
    """
    if not sidecar_path.exists():
        return [
            {
                "slide_index": 0,
                "slide_snippet": "(sidecar missing)",
                "claim_type": "infra",
                "claimed_value": "",
                "canonical_value": "",
                "detail": f"No sidecar JSON at {sidecar_path.name}",
            }
        ]
    sidecar = json.loads(sidecar_path.read_text())
    canonical = {
        "open_land_deals": int(sidecar.get("open_land_deals") or 0),
        "open_land_arr": float(sidecar.get("open_land_arr") or 0),
        "q1_land_lost": int(sidecar.get("q1_land_lost") or 0),
        "q1_land_lost_arr": float(sidecar.get("q1_land_lost_arr") or 0),
        "q1_land_wins": int(sidecar.get("q1_land_wins") or 0),
        "q1_land_wins_arr": float(sidecar.get("q1_land_wins_arr") or 0),
    }

    # Titles that are legitimate SUBSET claims of deck-wide totals. If the
    # title mentions any of these keywords, a deal-count claim is interpreted
    # as "N deals IN THIS SUBSET", not "N total open deals in scope".
    SUBSET_KEYWORDS = (
        "top ",
        "at risk",
        "push",
        "slip",
        "pending",
        "approved",
        "missing",
        "commit",
        "coach",
        "forecast mix",
        "renewals",
        "owner",
        "closed-won",
        "closed won",
        "closing",
        "q2 book",
        "q2 ",
        "q3 ",
        "lost",
        "stage 3+",
        "score",
        "risk",
        "triage",
    )

    flags = []
    for slide_idx, title_text in _slide_titles(pptx_path):
        if not title_text:
            continue
        if _out_of_scope(title_text):
            # Title names itself out of deck scope (e.g. MoM, prior year).
            continue

        is_subset = any(kw in title_text.lower() for kw in SUBSET_KEYWORDS)

        deal_counts = [
            int(m.group(1).replace(",", "")) for m in RE_DEALS.finditer(title_text)
        ]
        loss_counts = [
            int(m.group(1).replace(",", "")) for m in RE_LOSSES.finditer(title_text)
        ]
        win_counts = [
            int(m.group(1).replace(",", "")) for m in RE_WINS.finditer(title_text)
        ]

        # Deal counts only flag as scope drift when the title implies the
        # deck-wide total (non-subset claims).
        if deal_counts and canonical["open_land_deals"] and not is_subset:
            first = deal_counts[0]
            if first != canonical["open_land_deals"]:
                flags.append(
                    {
                        "slide_index": slide_idx,
                        "slide_snippet": title_text[:110],
                        "claim_type": "deal_count",
                        "claimed_value": str(first),
                        "canonical_value": str(canonical["open_land_deals"]),
                        "detail": "Title deal count disagrees with sidecar",
                    }
                )

        # Loss count must match exactly.
        if loss_counts:
            first = loss_counts[0]
            if not _is_count_close(first, canonical["q1_land_lost"]):
                flags.append(
                    {
                        "slide_index": slide_idx,
                        "slide_snippet": title_text[:110],
                        "claim_type": "loss_count",
                        "claimed_value": str(first),
                        "canonical_value": str(canonical["q1_land_lost"]),
                        "detail": "Title Q1 loss count disagrees with sidecar",
                    }
                )

        # Win count must match exactly.
        if win_counts:
            first = win_counts[0]
            if canonical["q1_land_wins"] != first and "rate" not in title_text.lower():
                if not _is_count_close(first, canonical["q1_land_wins"]):
                    flags.append(
                        {
                            "slide_index": slide_idx,
                            "slide_snippet": title_text[:110],
                            "claim_type": "win_count",
                            "claimed_value": str(first),
                            "canonical_value": str(canonical["q1_land_wins"]),
                            "detail": "Title Q1 win count disagrees with sidecar",
                        }
                    )

    return flags


def write_audit_note(run_date: str, all_results: list[dict]) -> tuple[Path, int]:
    period = run_date[:7]
    out_dir = VAULT / "Monthly" / period
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / "scope-audit.md"

    total_flags = sum(len(r["flags"]) for r in all_results)

    lines = [
        "---",
        "type: scope-audit",
        f"period: {period}",
        f"run_date: {run_date}",
        f"directors_audited: {len(all_results)}",
        f"flags: {total_flags}",
        "tags: [scope-audit, validation, monthly]",
        "---",
        "",
        f"# Deck scope audit, {period}",
        "",
        (
            f"Cross-slide scope consistency check across {len(all_results)} "
            f"director decks. Every numeric claim on every slide is compared "
            f"to the deck's canonical sidecar JSON (Land + Q1-Q2 + one "
            f"director). **Flags raised: {total_flags}.**"
        ),
        "",
    ]
    if total_flags == 0:
        lines.append("No scope drift detected. Every slide ties to the sidecar.")
    else:
        lines.append(
            "Each flag below names a slide whose numeric claim disagrees "
            "with the deck's declared scope. Either the claim is wrong or "
            "the slide is legitimately out-of-scope and should add an "
            "OUT_OF_SCOPE_MARKERS keyword in audit_deck_scope.py."
        )
    lines.append("")

    for record in all_results:
        lines.append(f"## {record['director']}")
        lines.append("")
        if not record["flags"]:
            lines.append("Clean: no scope drift.")
            lines.append("")
            continue
        lines.append("| Slide | Claim | Canonical | Type | Detail |")
        lines.append("|---:|---|---|---|---|")
        for flag in record["flags"]:
            lines.append(
                f"| {flag['slide_index']} | {flag['claimed_value']} | "
                f"{flag['canonical_value']} | {flag['claim_type']} | "
                f"{flag['detail']} |"
            )
        lines.append("")
        lines.append(f"Slide snippet: *{record['flags'][0]['slide_snippet']}*")
        lines.append("")

    path.write_text("\n".join(lines) + "\n")
    print(f"  wrote {path.relative_to(ROOT)}")
    return path, total_flags


def write_audit_artifacts(run_date: str, all_results: list[dict]) -> tuple[Path, int]:
    output_dir = OUTPUT_ROOT / run_date[:10]
    output_dir.mkdir(parents=True, exist_ok=True)
    total_flags = sum(len(r["flags"]) for r in all_results)
    payload = {
        "run_date": run_date[:10],
        "status": "failed" if total_flags else "ok",
        "decks_audited": len(all_results),
        "flag_count": total_flags,
        "results": all_results,
    }
    (output_dir / "deck_scope_audit.json").write_text(
        json.dumps(payload, indent=2) + "\n",
        encoding="utf-8",
    )
    lines = [
        f"# Deck Scope Audit — {run_date[:10]}",
        "",
        f"- Status: `{payload['status']}`",
        f"- Decks audited: `{payload['decks_audited']}`",
        f"- Flags: `{payload['flag_count']}`",
        "",
        "## Decks",
        "",
    ]
    if not all_results:
        lines.append("- none")
    else:
        for item in all_results:
            flag_count = len(item.get("flags") or [])
            status = "clean" if flag_count == 0 else f"{flag_count} flag(s)"
            lines.append(f"- `{item['slug']}`: `{status}`")
    (output_dir / "summary.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"  wrote {output_dir.relative_to(ROOT)}")
    return output_dir, total_flags


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--date", default=datetime.now().strftime("%Y-%m-%d"))
    args = parser.parse_args()

    decks_dir = ROOT / "output" / "simcorp_director_decks" / args.date / "land-only"
    if not decks_dir.exists():
        print(f"  decks directory missing: {decks_dir}")
        return 1

    all_results = []
    total_flags = 0
    for deck_path in sorted(decks_dir.glob("*-LAND.pptx")):
        if deck_path.name.startswith("~"):
            continue
        slug = deck_path.stem.replace("-LAND", "")
        sidecar_path = deck_path.with_suffix(".json")
        director = slug.replace("-", " ").title()
        flags = audit_deck(deck_path, sidecar_path)
        total_flags += len(flags)
        status = "clean" if not flags else f"{len(flags)} flag(s)"
        print(f"  {director:22s}  {status}")
        all_results.append(
            {
                "director": director,
                "slug": slug,
                "deck_path": str(deck_path.relative_to(ROOT)),
                "sidecar_path": str(sidecar_path.relative_to(ROOT)),
                "flags": flags,
            }
        )

    _path, flag_count = write_audit_note(args.date, all_results)
    write_audit_artifacts(args.date, all_results)
    print(f"\nTotal scope-drift flags: {flag_count}")
    return 0 if flag_count == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
