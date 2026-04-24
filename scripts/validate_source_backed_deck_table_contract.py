#!/usr/bin/env python3
"""Validate source-backed deck tables against the standard monthly format."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_ROOT = ROOT / "output" / "source_backed_deck_table_contract"
SCHEMA_VERSION = "monthly_platform.source_backed_deck_table_contract.v1"

HEADER_FILL = "011946"
HEADER_FONT = "FFFFFF"
BODY_FONT = "011946"
ALLOWED_BODY_FILLS = {
    "F5F8FD",
    "FFFFFF",
    "E4F7F3",
    "FFF1DF",
    "F8E8F1",
}
EXPECTED_TABLES = [
    {
        "slide": 2,
        "name": "publish_gate",
        "headers": ["Gate", "Status / Count", "What it proves"],
        "min_rows": 8,
    },
    {
        "slide": 3,
        "name": "regional_rollup",
        "headers": ["Region", "Territories", "Open ARR", "Deals", "Risk"],
        "min_rows": 1,
    },
    {
        "slide": 4,
        "name": "director_book",
        "headers": [
            "Director",
            "Territory",
            "Open ARR",
            "Deals",
            "Risk rows",
            "Tie-out",
            "Source issues",
        ],
        "min_rows": 1,
    },
    {
        "slide": 5,
        "name": "quarter_policy",
        "headers": [
            "Director",
            "Territory",
            "Display quarter",
            "Reason",
            "Current active",
            "Forward active",
        ],
        "min_rows": 1,
    },
    {
        "slide": 6,
        "name": "production_handoff",
        "headers": ["Artifact", "Reference", "Use"],
        "min_rows": 4,
    },
]


def validate_table_contract(
    *,
    deck_path: Path,
    output_path: Path | None = None,
    output_root: Path = DEFAULT_OUTPUT_ROOT,
    snapshot_date: str | None = None,
    source_run_id: str | None = None,
) -> dict[str, Any]:
    deck_path = Path(deck_path)
    findings: list[dict[str, Any]] = []
    tables_checked: list[dict[str, Any]] = []
    resolved_snapshot_date = snapshot_date or "unknown-snapshot"
    resolved_run_id = source_run_id or "unknown-run"

    if not deck_path.exists():
        findings.append(_finding("high", "deck_missing", f"Missing deck: {deck_path}"))
        return _write_result(
            payload=_result_payload(
                deck_path=deck_path,
                snapshot_date=resolved_snapshot_date,
                source_run_id=resolved_run_id,
                tables_checked=tables_checked,
                findings=findings,
            ),
            output_path=output_path,
            output_root=output_root,
        )

    try:
        presentation = Presentation(str(deck_path))
    except Exception as exc:
        findings.append(_finding("high", "deck_load_failed", str(exc)))
        return _write_result(
            payload=_result_payload(
                deck_path=deck_path,
                snapshot_date=resolved_snapshot_date,
                source_run_id=resolved_run_id,
                tables_checked=tables_checked,
                findings=findings,
            ),
            output_path=output_path,
            output_root=output_root,
        )

    table_slots = _table_slots(presentation)
    if len(table_slots) != len(EXPECTED_TABLES):
        findings.append(
            _finding(
                "high",
                "table_count_mismatch",
                f"{len(table_slots)} tables; expected {len(EXPECTED_TABLES)}",
            )
        )

    for expected in EXPECTED_TABLES:
        slot = next(
            (item for item in table_slots if item["slide"] == expected["slide"]),
            None,
        )
        if slot is None:
            findings.append(
                _finding(
                    "high",
                    "expected_table_missing",
                    f"{expected['name']} table missing on slide {expected['slide']}",
                    slide=int(expected["slide"]),
                )
            )
            continue
        table = slot["shape"].table
        checked = _validate_table(
            table=table,
            expected=expected,
            findings=findings,
        )
        tables_checked.append(checked)

    return _write_result(
        payload=_result_payload(
            deck_path=deck_path,
            snapshot_date=resolved_snapshot_date,
            source_run_id=resolved_run_id,
            tables_checked=tables_checked,
            findings=findings,
        ),
        output_path=output_path,
        output_root=output_root,
    )


def _table_slots(presentation: Presentation) -> list[dict[str, Any]]:
    slots: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                slots.append({"slide": slide_index, "shape": shape})
    return slots


def _validate_table(
    *,
    table,
    expected: dict[str, Any],
    findings: list[dict[str, Any]],
) -> dict[str, Any]:
    slide = int(expected["slide"])
    name = str(expected["name"])
    expected_headers = [str(value) for value in expected["headers"]]
    actual_headers = [table.cell(0, column).text for column in range(len(table.columns))]
    row_count = len(table.rows)
    column_count = len(table.columns)

    if actual_headers != expected_headers:
        findings.append(
            _finding(
                "high",
                "table_headers_drifted",
                f"{name}: {actual_headers} != {expected_headers}",
                slide=slide,
            )
        )
    if row_count - 1 < int(expected["min_rows"]):
        findings.append(
            _finding(
                "high",
                "table_row_count_below_minimum",
                f"{name}: {row_count - 1} data rows; expected at least {expected['min_rows']}",
                slide=slide,
            )
        )

    for column in range(column_count):
        _validate_header_cell(
            table.cell(0, column),
            table_name=name,
            slide=slide,
            column=column + 1,
            findings=findings,
        )

    for row in range(1, row_count):
        for column in range(column_count):
            _validate_body_cell(
                table.cell(row, column),
                table_name=name,
                slide=slide,
                row=row + 1,
                column=column + 1,
                findings=findings,
            )

    return {
        "name": name,
        "slide": slide,
        "headers": actual_headers,
        "row_count": row_count,
        "data_row_count": row_count - 1,
        "column_count": column_count,
    }


def _validate_header_cell(
    cell,
    *,
    table_name: str,
    slide: int,
    column: int,
    findings: list[dict[str, Any]],
) -> None:
    fill = _fill_rgb(cell)
    font = _first_run_font(cell)
    if fill != HEADER_FILL:
        findings.append(
            _finding(
                "high",
                "table_header_fill_drifted",
                f"{table_name} slide {slide} col {column}: {fill}",
                slide=slide,
            )
        )
    if font["color"] != HEADER_FONT:
        findings.append(
            _finding(
                "high",
                "table_header_font_color_drifted",
                f"{table_name} slide {slide} col {column}: {font['color']}",
                slide=slide,
            )
        )
    if font["bold"] is not True:
        findings.append(
            _finding(
                "high",
                "table_header_font_weight_drifted",
                f"{table_name} slide {slide} col {column}: bold={font['bold']}",
                slide=slide,
            )
        )


def _validate_body_cell(
    cell,
    *,
    table_name: str,
    slide: int,
    row: int,
    column: int,
    findings: list[dict[str, Any]],
) -> None:
    fill = _fill_rgb(cell)
    font = _first_run_font(cell)
    if fill not in ALLOWED_BODY_FILLS:
        findings.append(
            _finding(
                "medium",
                "table_body_fill_outside_standard_palette",
                f"{table_name} slide {slide} row {row} col {column}: {fill}",
                slide=slide,
            )
        )
    if font["color"] not in {BODY_FONT, None}:
        findings.append(
            _finding(
                "medium",
                "table_body_font_color_drifted",
                f"{table_name} slide {slide} row {row} col {column}: {font['color']}",
                slide=slide,
            )
        )


def _fill_rgb(cell) -> str | None:
    try:
        if not cell.fill.fore_color.type:
            return None
        value = cell.fill.fore_color.rgb
    except Exception:
        return None
    return str(value) if value is not None else None


def _first_run_font(cell) -> dict[str, Any]:
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            color = None
            try:
                if run.font.color and run.font.color.rgb:
                    color = str(run.font.color.rgb)
            except Exception:
                color = None
            return {"color": color, "bold": run.font.bold}
    return {"color": None, "bold": None}


def _result_payload(
    *,
    deck_path: Path,
    snapshot_date: str,
    source_run_id: str,
    tables_checked: list[dict[str, Any]],
    findings: list[dict[str, Any]],
) -> dict[str, Any]:
    high_count = sum(1 for finding in findings if finding["severity"] == "high")
    medium_count = sum(1 for finding in findings if finding["severity"] == "medium")
    return {
        "schema_version": SCHEMA_VERSION,
        "status": "blocked" if high_count or medium_count else "ok",
        "snapshot_date": snapshot_date,
        "source_run_id": source_run_id,
        "deck_path": str(deck_path),
        "checks": {
            "table_count": len(tables_checked),
            "expected_table_count": len(EXPECTED_TABLES),
            "table_contract_checked": deck_path.exists(),
            "expected_tables": [
                {
                    "name": item["name"],
                    "slide": item["slide"],
                    "headers": item["headers"],
                }
                for item in EXPECTED_TABLES
            ],
            "tables_checked": tables_checked,
        },
        "summary": {
            "finding_count": len(findings),
            "high_finding_count": high_count,
            "medium_finding_count": medium_count,
            "low_finding_count": sum(
                1 for finding in findings if finding["severity"] == "low"
            ),
            "table_count": len(tables_checked),
        },
        "findings": findings,
    }


def _write_result(
    *,
    payload: dict[str, Any],
    output_path: Path | None,
    output_root: Path,
) -> dict[str, Any]:
    if output_path is None:
        output_path = (
            Path(output_root)
            / payload["snapshot_date"]
            / payload["source_run_id"]
            / "source_backed_deck_table_contract_audit.json"
        )
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    payload["output_path"] = str(output_path)
    output_path.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")
    return payload


def _finding(
    severity: str,
    issue: str,
    evidence: str,
    *,
    slide: int | None = None,
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "severity": severity,
        "issue": issue,
        "evidence": evidence,
    }
    if slide is not None:
        payload["slide"] = slide
    return payload


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--deck-path", type=Path, required=True)
    parser.add_argument("--snapshot-date")
    parser.add_argument("--source-run-id")
    parser.add_argument("--output-path", type=Path)
    parser.add_argument("--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    payload = validate_table_contract(
        deck_path=args.deck_path,
        output_path=args.output_path,
        output_root=args.output_root,
        snapshot_date=args.snapshot_date,
        source_run_id=args.source_run_id,
    )
    print(json.dumps(payload, indent=2))
    return 0 if payload["status"] == "ok" else 1


if __name__ == "__main__":
    raise SystemExit(main())
