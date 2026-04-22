#!/usr/bin/env python3
"""Build a SimCorp director deck directly from a live Excel workbook.

This is the final stage of the ETL pipeline:
  Salesforce (live) → Excel workbook → SimCorp deck

The Excel workbook (from extract_director_live.py) is the single source of truth.
Directors can review, adjust, add comments in Excel before the deck is rendered.

Usage:
    python3 scripts/build_deck_from_excel.py \
        --workbook output/director_live_workbooks/2026-04-13/jesper-tyrer.xlsx \
        [--template ~/archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx] \
        [--output output/simcorp_director_decks/2026-04-13/jesper-tyrer.pptx]
"""

from __future__ import annotations

import argparse
import re
import sys
from collections import Counter, defaultdict
from datetime import datetime, timedelta
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ── Fiscal quarter calendar (derived from run date at CLI parse time) ──
_MONTH_LABELS = {
    1: "Jan",
    2: "Feb",
    3: "Mar",
    4: "Apr",
    5: "May",
    6: "Jun",
    7: "Jul",
    8: "Aug",
    9: "Sep",
    10: "Oct",
    11: "Nov",
    12: "Dec",
}


def _compute_quarters(run_date: datetime):
    """Derive fiscal quarter boundaries from the run date.

    Returns a dict with prior/current/forward quarter boundaries
    and the FY scope (prior Q start through forward Q end).
    """
    y = run_date.year
    m = run_date.month
    cur_q = (m - 1) // 3 + 1

    def _qbounds(year, q):
        start_m = (q - 1) * 3 + 1
        end_m = q * 3
        start = f"{year}-{start_m:02d}-01"
        if end_m == 12:
            end = f"{year}-12-31"
        else:
            end = f"{year}-{end_m + 1:02d}-01"
            from datetime import date as _d

            end = str(_d.fromisoformat(end) - timedelta(days=1))
        month_start = f"{year}-{start_m:02d}"
        month_end = f"{year}-{end_m:02d}"
        label = f"{_MONTH_LABELS[start_m]}-{_MONTH_LABELS[end_m]}"
        return {
            "q": q,
            "label": f"Q{q}",
            "start": start,
            "end": end,
            "month_start": month_start,
            "month_end": month_end,
            "range_label": label,
            "year": year,
        }

    prior = _qbounds(y, cur_q - 1) if cur_q > 1 else _qbounds(y - 1, 4)
    current = _qbounds(y, cur_q)
    forward = _qbounds(y, cur_q + 1) if cur_q < 4 else _qbounds(y + 1, 1)

    fy_start = f"{y}-01-01"
    fy_end = f"{y}-12-31"
    scope_end = forward["end"]

    return {
        "fy": y,
        "fy_start": fy_start,
        "fy_end": fy_end,
        "prior": prior,
        "current": current,
        "forward": forward,
        "scope_start": fy_start,
        "scope_end": scope_end,
    }


FQ = _compute_quarters(datetime.now())

# ── Constants ──
DARK = RGBColor(0x1A, 0x1D, 0x31)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SIMCORP_BLUE = RGBColor(0x08, 0x3E, 0xA7)
RED = RGBColor(0xEF, 0x3E, 0x4A)

LY_TITLE_1 = 0
LY_TITLE_CONTENT = 6
LY_2COL_GRAD = 10
LY_4COL_GRAD = 12
LY_END_SLIDE = 31

GRAD_COL_X = [0.9, 3.9, 6.8, 9.8]

DEFAULT_TEMPLATE = (
    Path.home()
    / "archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx"
)


# ── Excel reader ──


def read_sheet(wb, sheet_name: str) -> list[dict]:
    """Read an Excel sheet into a list of dicts (header row = keys)."""
    if sheet_name not in wb.sheetnames:
        return []
    ws = wb[sheet_name]
    rows = list(ws.iter_rows(values_only=True))
    if len(rows) < 2:
        return []
    headers = [str(h or "").strip() for h in rows[0]]
    result = []
    for row in rows[1:]:
        d = {}
        for i, h in enumerate(headers):
            val = row[i] if i < len(row) else None
            d[h] = val if val is not None else ""
        result.append(d)
    return result


def read_kpis(wb) -> dict[str, str]:
    """Read KPI values from the Summary sheet."""
    if "Summary" not in wb.sheetnames:
        return {}
    ws = wb["Summary"]
    kpis = {}
    for row in ws.iter_rows(min_row=7, max_col=2, values_only=True):
        if row[0] and row[1]:
            kpis[str(row[0]).strip()] = str(row[1]).strip()
    return kpis


def read_director_analytics(analytics_path, director_name):
    """Slice the consolidated FY26 Pipeline Review workbook for one director.

    Returns a dict with:
      risk_deals:  list of dicts from Deal Risk Scoring, filtered to director
      variance:    one row from Forecast Variance (bucket decomposition)
      top_deals:   rows from ARR Concentration where Director matches
      velocity_q1: {'dates': [...], 'series': [...]} from Pipeline Velocity Q1 block
      velocity_q2: same for Q2

    Missing tabs or missing director row yield empty values so the deck still
    renders when the analytics workbook is unavailable.
    """
    out = {
        "risk_deals": [],
        "variance": None,
        "top_deals": [],
        "velocity_q1": None,
        "velocity_q2": None,
        "director_total_open_arr": 0.0,
        "slip_owners": [],
    }
    if not analytics_path:
        return out
    p = Path(analytics_path)
    if not p.exists():
        return out
    wb = load_workbook(str(p), data_only=True, read_only=True)

    if "Deal Risk Scoring" in wb.sheetnames:
        ws = wb["Deal Risk Scoring"]
        rows = list(ws.iter_rows(values_only=True))
        for r in rows[4:]:
            if not r or r[0] is None or not isinstance(r[0], int):
                continue
            if str(r[2] or "") != director_name:
                continue
            out["risk_deals"].append(
                {
                    "rank": r[0],
                    "score": r[1],
                    "account": r[3] or "",
                    "opportunity": r[4] or "",
                    "stage": r[5] or "",
                    "owner": r[6] or "",
                    "close_date": str(r[7] or "")[:10],
                    "days_to_close": r[8],
                    "days_since_act": r[9],
                    "arr": r[10] or 0,
                    "push": r[11] or 0,
                    "reasons": r[12] or "",
                }
            )

    # Forecast Variance cells are SUMIFS formulas and have no cached value
    # until Excel/LibreOffice opens the file. Compute directly from the
    # Q1 Trend Consolidated helper columns instead — same source the
    # SUMIFS formulas reference, same answer.
    if "Q1 Trend Consolidated" in wb.sheetnames:
        qws = wb["Q1 Trend Consolidated"]
        # Find the director's territory string from the first row for them
        headers = [
            qws.cell(row=2, column=c).value for c in range(1, qws.max_column + 1)
        ]
        col = {str(h): ci for ci, h in enumerate(headers, 1) if h}
        if all(k in col for k in ("Territory", "Initial ARR", "Final ARR", "Bucket")):
            terr_ci = col["Territory"]
            init_ci = col["Initial ARR"]
            final_ci = col["Final ARR"]
            bucket_ci = col["Bucket"]
            # Director -> territory mapping is implicit via the consolidated
            # tab, but we need to filter. We don't know the territory string
            # up-front, so pull it from the first director_total read below
            # or iterate rows matching later. Simpler: iterate all rows,
            # group by territory, build buckets, then pick the territory
            # that matches the director via DIRECTORS lookup. Since DIRECTORS
            # isn't imported here, accept any territory that has data with
            # the director_name via a passed mapping if available.
            # Fallback: match by the per-director open_arr value coming from
            # Territory Scorecard, which we already read above.
            # Simpler still: accept that read_director_analytics callers pass
            # the director's own territory. Keep Territory-lookup minimal:
            # scan the Q1 Trend Consolidated rows for this director's
            # Territory by cross-referencing Territory Scorecard (director ->
            # row below). If unresolved, leave variance None.
            terr_for_director = None
            # Need the Territory string. Territory Scorecard is keyed by
            # director name; its Territory is not stored there, so instead
            # we use the Summary tab on the analytics workbook if present,
            # else infer by checking any row in Q1 Trend whose ownership
            # aligns. Simpler: iterate Forecast Variance again to pull the
            # Territory literal column 2 (string, not formula).
            if "Forecast Variance" in wb.sheetnames:
                fv = wb["Forecast Variance"]
                for rr in range(5, fv.max_row + 1):
                    if str(fv.cell(row=rr, column=1).value or "") == director_name:
                        terr_for_director = str(fv.cell(row=rr, column=2).value or "")
                        break
            if terr_for_director:
                buckets = {
                    "Won": 0.0,
                    "Lost": 0.0,
                    "Added": 0.0,
                    "RevisedUp": 0.0,
                    "RevisedDown": 0.0,
                    "Unchanged": 0.0,
                    "AlreadyClosed": 0.0,
                }
                initial = 0.0
                final = 0.0
                for r_idx in range(3, qws.max_row + 1):
                    terr = str(qws.cell(row=r_idx, column=terr_ci).value or "")
                    if terr != terr_for_director:
                        continue
                    try:
                        arr0 = float(qws.cell(row=r_idx, column=init_ci).value or 0)
                    except (TypeError, ValueError):
                        arr0 = 0.0
                    try:
                        arrN = float(qws.cell(row=r_idx, column=final_ci).value or 0)
                    except (TypeError, ValueError):
                        arrN = 0.0
                    bucket = str(qws.cell(row=r_idx, column=bucket_ci).value or "")
                    if bucket != "AlreadyClosed":
                        initial += arr0
                    if bucket not in ("Won", "Lost", "AlreadyClosed"):
                        final += arrN
                    if bucket in buckets:
                        if bucket == "Won":
                            buckets["Won"] += arr0
                        elif bucket == "Lost":
                            buckets["Lost"] += arr0
                        elif bucket == "Added":
                            buckets["Added"] += arrN
                        elif bucket == "RevisedUp":
                            buckets["RevisedUp"] += arrN - arr0
                        elif bucket == "RevisedDown":
                            buckets["RevisedDown"] += arr0 - arrN
                out["variance"] = {
                    "initial": initial,
                    "final": final,
                    "net": final - initial,
                    "won": buckets["Won"],
                    "lost": buckets["Lost"],
                    "added": buckets["Added"],
                    "up": buckets["RevisedUp"],
                    "down": buckets["RevisedDown"],
                }

    if "ARR Concentration" in wb.sheetnames:
        ws = wb["ARR Concentration"]
        rows = list(ws.iter_rows(values_only=True))
        for r in rows[4:24]:
            if not r or r[0] is None or not isinstance(r[0], int):
                continue
            if str(r[1] or "") == director_name:
                out["top_deals"].append(
                    {
                        "rank": r[0],
                        "account": r[2] or "",
                        "opportunity": r[3] or "",
                        "stage": r[4] or "",
                        "owner": r[5] or "",
                        "arr": r[6] or 0,
                        "pct": r[7] or 0,
                        "cum_pct": r[8] or 0,
                    }
                )

    if "Slip Risk by Owner" in wb.sheetnames:
        ws = wb["Slip Risk by Owner"]
        rows = list(ws.iter_rows(values_only=True))
        # Header at row 3 (index 3); data starts row 4.
        for r in rows[4:]:
            if not r or not r[0]:
                continue
            if str(r[1] or "") != director_name:
                continue
            out["slip_owners"].append(
                {
                    "owner": r[0],
                    "deals": int(r[2] or 0),
                    "arr": float(r[3] or 0),
                    "pushes": int(r[4] or 0),
                    "avg_push": float(r[5] or 0),
                    "max_push": int(r[6] or 0),
                }
            )

    if "Territory Scorecard" in wb.sheetnames:
        ws = wb["Territory Scorecard"]
        rows = list(ws.iter_rows(values_only=True))
        for r in rows[4:]:
            if not r or r[0] is None:
                continue
            if str(r[0]) == director_name:
                try:
                    out["director_total_open_arr"] = float(r[2] or 0)
                except (TypeError, ValueError):
                    pass
                break

    if "Pipeline Velocity" in wb.sheetnames:
        ws = wb["Pipeline Velocity"]
        rows = list(ws.iter_rows(values_only=True))
        i = 0
        while i < len(rows):
            r = rows[i]
            if r and isinstance(r[0], str) and r[0].endswith("ARR by Snapshot (EUR)"):
                label = "velocity_q1" if "Q1" in r[0] else "velocity_q2"
                dates_row = rows[i + 1] if i + 1 < len(rows) else None
                dates = [c for c in (dates_row[1:] if dates_row else []) if c]
                j = i + 2
                while j < len(rows) and rows[j] and rows[j][0]:
                    if str(rows[j][0]) == director_name:
                        series = [rows[j][k + 1] or 0 for k in range(len(dates))]
                        out[label] = {"dates": dates, "series": series}
                        break
                    j += 1
            i += 1

    wb.close()
    return out


def read_director_info(wb) -> tuple[str, str]:
    """Read director name and territory from Summary A1."""
    if "Summary" not in wb.sheetnames:
        return "Director", "Territory"
    ws = wb["Summary"]
    title = str(ws["A1"].value or "Director (Territory)")
    match = re.match(r"(.+?)\s*\((.+?)\)", title)
    if match:
        return match.group(1).strip(), match.group(2).strip()
    return title, ""


# ── Helpers ──


def _meur(value) -> str:
    """Format EUR value as mEUR for table cells."""
    v = float(value or 0)
    if abs(v) >= 1_000_000:
        return f"{v / 1_000_000:.1f}"
    if abs(v) >= 1_000:
        return f"{v / 1_000_000:.2f}"
    if v == 0:
        return "0"
    return f"{v / 1_000_000:.3f}"


def _fmt_eur(value) -> str:
    """Format EUR value for slide text."""
    v = float(value or 0)
    if abs(v) >= 1_000_000:
        return f"EUR {v / 1_000_000:.1f}M"
    if abs(v) >= 1_000:
        return f"EUR {v / 1_000:.0f}K"
    return f"EUR {v:,.0f}"


def _unw(r) -> float:
    """Unweighted ARR (APTS_Opportunity_ARR__c). Accepts new or legacy column name."""
    v = r.get("ARR Unweighted (EUR)")
    if v is None:
        v = r.get("Opp ARR (EUR)")
    return float(v or 0)


def _wtd(r) -> float:
    """Weighted ARR (APTS_Forecast_ARR__c). Accepts new or legacy column name."""
    v = r.get("ARR Weighted (EUR)")
    if v is None:
        v = r.get("Forecast ARR (EUR)")
    return float(v or 0)


def _ph(slide, idx):
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def _set_ph(slide, idx, text):
    ph = _ph(slide, idx)
    if ph is not None:
        ph.text = str(text) if text else ""


def _set_ph_lines(slide, idx, lines):
    ph = _ph(slide, idx)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = str(line)
        else:
            tf.add_paragraph().text = str(line)


def _add_gradient_metric(slide, col_idx, text):
    ph = _ph(slide, 61 + col_idx)
    if ph is not None:
        ph.text = str(text)
        return
    x = GRAD_COL_X[col_idx] + 0.1
    txBox = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(2.4), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = None


def _add_table(
    slide,
    rows_data,
    left,
    top,
    width,
    row_height=0.22,
    col_widths=None,
    data_font_size=Pt(10),
    highlight_rules=None,
):
    """Render a table. `highlight_rules` is an optional dict mapping
    0-based column index to a list of (predicate, hex_color) tuples. First
    matching rule wins. Predicates receive the raw cell value as a string.
    Applies only to data rows; header styling is untouched."""
    if not rows_data or not rows_data[0]:
        return None
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])

    tbl_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(max(row_height * n_rows, 0.5)),
    )
    tbl = tbl_shape.table

    from pptx.oxml.ns import qn

    # Remove all default table styling — we set everything explicitly
    tbl_pr = tbl._tbl.tblPr
    tbl_pr.set("firstRow", "1")
    tbl_pr.set("bandRow", "1")
    # Clear default style so it doesn't override our colors
    for existing in tbl_pr.findall(qn("a:tableStyleId")):
        existing.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid
    for existing in tbl_pr.findall(qn("a:tblStyle")):
        tbl_pr.remove(existing)

    if col_widths and len(col_widths) == n_cols:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)

    # Rebekka's exact colors from the screenshot:
    # Header: dark navy #1A1D31, white text, 16pt bold
    # Data odd rows: white #FFFFFF
    # Data even rows: light gray #F2F2F2
    # Borders: thin gray lines
    HEADER_BG = "1A1D31"
    EVEN_ROW_BG = "F2F2F2"
    ODD_ROW_BG = "FFFFFF"
    BORDER_COLOR = "D0D0D0"

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text) if cell_text else ""
            # Wrap long values instead of clipping. PowerPoint still
            # auto-fits row height, so nothing visually truncates.
            cell.text_frame.word_wrap = True

            # Font
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if ri == 0:
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.size = data_font_size
                        run.font.color.rgb = DARK
                    run.font.name = None

            # Cell fill
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for old in tcPr.findall(qn("a:solidFill")):
                tcPr.remove(old)

            fill_color = None
            # Highlight rules override banding colors on data rows only
            if ri > 0 and highlight_rules and ci in highlight_rules:
                for predicate, color in highlight_rules[ci]:
                    try:
                        if predicate(cell_text):
                            fill_color = color
                            break
                    except Exception:
                        continue
            if fill_color is None:
                if ri == 0:
                    fill_color = HEADER_BG
                elif ri % 2 == 0:
                    fill_color = EVEN_ROW_BG
                else:
                    fill_color = ODD_ROW_BG
            sf = tcPr.makeelement(qn("a:solidFill"), {})
            sf.append(tcPr.makeelement(qn("a:srgbClr"), {"val": fill_color}))
            tcPr.append(sf)

            # Cell borders — thin gray
            for old in tcPr.findall(qn("a:tcBdr")):
                tcPr.remove(old)
            bdr = tcPr.makeelement(qn("a:tcBdr"), {})
            for side in ("top", "bottom", "left", "right"):
                side_el = bdr.makeelement(qn(f"a:{side}"), {})
                ln = side_el.makeelement(qn("a:ln"), {"w": "6350"})
                fill = ln.makeelement(qn("a:solidFill"), {})
                fill.append(ln.makeelement(qn("a:srgbClr"), {"val": BORDER_COLOR}))
                ln.append(fill)
                side_el.append(ln)
                bdr.append(side_el)
            tcPr.append(bdr)
    return tbl


def _snapshot_to_period(date_str: str) -> str:
    """Reporting period label for the cover slide.

    Uses the snapshot month (not the prior month) because monthly reviews
    happen mid-cycle: an April 15 snapshot covers Q1 retrospective plus
    Q2-to-date, so 'April 2026' is the correct frame.
    """
    try:
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        return dt.strftime("%B %Y")
    except (ValueError, TypeError):
        return date_str or ""


# ── Analysis helpers ──


def _top_pushed_owner(pi_data: list) -> str:
    """Find the owner with the most pushed deals."""
    pushed = [r for r in pi_data if int(r.get("Push Count") or 0) > 0]
    if not pushed:
        return ""
    owners = Counter(r.get("Owner", "?") for r in pushed)
    top_owner, count = owners.most_common(1)[0]
    return f"{top_owner} owns {count} of the most-pushed deals"


def _push_tiers(pi_data: list) -> dict:
    """Classify pushed deals into Critical/Watch/Early tiers."""
    pushed = [r for r in pi_data if int(r.get("Push Count") or 0) > 0]
    critical = [r for r in pushed if int(r.get("Push Count") or 0) >= 5]
    watch = [r for r in pushed if 3 <= int(r.get("Push Count") or 0) <= 4]
    early = [r for r in pushed if 1 <= int(r.get("Push Count") or 0) <= 2]
    return {
        "critical": {
            "count": len(critical),
            "arr": sum(_wtd(r) for r in critical),
        },
        "watch": {
            "count": len(watch),
            "arr": sum(_wtd(r) for r in watch),
        },
        "early": {
            "count": len(early),
            "arr": sum(_wtd(r) for r in early),
        },
    }


def _stage_breakdown(pipeline: list) -> list[tuple[str, int, float, float]]:
    """Return (stage, count, arr_unweighted, arr_weighted) sorted by stage order."""
    stages: dict[str, dict] = {}
    for r in pipeline:
        s = str(r.get("Stage", "?"))
        if s not in stages:
            stages[s] = {"count": 0, "arr": 0.0, "wtd": 0.0}
        stages[s]["count"] += 1
        stages[s]["arr"] += _unw(r)
        stages[s]["wtd"] += _wtd(r)
    return sorted(
        [(s, d["count"], d["arr"], d["wtd"]) for s, d in stages.items()],
        key=lambda x: x[0],
    )


# ── Slide builders ──


def slide_cover(prs, director, territory, snapshot_date):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_1])
    # Director name front and centre, territory underneath, period subtle.
    _set_ph(slide, 24, f"{director}, {territory}")
    _set_ph(slide, 20, "Monthly Pipeline Review")
    _set_ph(
        slide,
        22,
        f"{_snapshot_to_period(snapshot_date)}. Land pipeline, FY{str(FQ['fy'])[-2:]} {FQ['prior']['label']}-{FQ['forward']['label']}.",
    )


def _read_prior_snapshot_from_history(director_name, current_run_date):
    """Return the prior baseline for `director_name` from the cross-run
    snapshot ledger. Picks the most recent run strictly before the current
    run_date. Returns a dict with period, open_deals, open_unwtd,
    q1_won_count, q1_won_arr, approved_2026, missing_approval — or None if
    no prior entry exists.
    """
    import json as _json

    ledger = Path(__file__).resolve().parents[1] / "obsidian" / "snapshot_history.json"
    if not ledger.exists():
        return None
    try:
        history = _json.loads(ledger.read_text())
    except (_json.JSONDecodeError, ValueError):
        return None
    snapshots = history.get("snapshots") or []
    priors = [s for s in snapshots if s.get("run_date", "") < current_run_date]
    if not priors:
        return None
    # Most recent prior (ledger is already sorted ascending).
    prior = priors[-1]
    d = prior.get("directors", {}).get(director_name)
    if not d:
        return None
    return {
        "period": prior.get("run_date") or prior.get("period") or "",
        "open_deals": d.get("open_land_deals"),
        "open_unwtd": d.get("open_land_arr_unwtd"),
        "q1_won_count": d.get("q1_won_count"),
        "q1_won_arr": d.get("q1_won_arr"),
        "approved_2026": d.get("approved_2026"),
        "missing_approval": d.get("missing_approval"),
    }


def _read_prior_snapshot(director_slug, current_period):
    """Find the most recent prior monthly snapshot for this director by
    scanning `obsidian/Monthly/YYYY-MM/<slug>.auto.md` files. Returns a dict
    of parsed headline numbers or None."""
    import re

    vault_monthly = Path(__file__).resolve().parents[1] / "obsidian" / "Monthly"
    if not vault_monthly.exists():
        return None

    prior_periods = sorted(
        [
            p.name
            for p in vault_monthly.iterdir()
            if p.is_dir() and p.name < current_period
        ],
        reverse=True,
    )
    for period in prior_periods:
        path = vault_monthly / period / f"{director_slug}.auto.md"
        if not path.exists():
            continue
        text = path.read_text()

        def _grab_int(pattern):
            m = re.search(pattern, text)
            return int(m.group(1).replace(",", "")) if m else None

        def _grab_eur(pattern):
            """Parse EUR 5.2M or EUR 123K back into a float."""
            m = re.search(pattern, text)
            if not m:
                return None
            raw = m.group(1).strip()
            raw = raw.replace("EUR", "").strip()
            try:
                if raw.endswith("M"):
                    return float(raw[:-1]) * 1_000_000
                if raw.endswith("K"):
                    return float(raw[:-1]) * 1_000
                return float(raw.replace(",", ""))
            except ValueError:
                return None

        return {
            "period": period,
            "open_deals": _grab_int(r"Open Land pipeline: (\d+) deals"),
            "open_unwtd": _grab_eur(
                r"Open Land pipeline:.*?(EUR [\d.,KM]+) unweighted"
            ),
            "q1_won_count": _grab_int(r"Q1 Land outcome: (\d+) wins"),
            "q1_won_arr": _grab_eur(r"Q1 Land outcome: \d+ wins \((EUR [\d.,KM]+)\)"),
            "approved_2026": _grab_int(r"(\d+) approved 2026"),
            "missing_approval": _grab_int(r"(\d+) missing Stage 3\+"),
        }
    return None


def slide_month_over_month(
    prs, pipeline, won_lost, approvals, director, territory, snapshot_date
):
    """Month-over-month delta on six headline KPIs. Skipped if no prior
    snapshot exists yet (e.g. first ever run)."""
    current_period = snapshot_date[:7]
    slug = director.lower().replace(" ", "-")
    # Prefer the cross-run ledger (captures every daily run); fall back to
    # prior calendar-month folder scan for back-compat with older vaults.
    prior = _read_prior_snapshot_from_history(director, snapshot_date) or (
        _read_prior_snapshot(slug, current_period)
    )
    if not prior:
        return  # First run, nothing to compare against

    # Current numbers
    cur_open_deals = len(pipeline)
    cur_open_unwtd = sum(_unw(r) for r in pipeline)
    cur_q1_won = [
        r
        for r in won_lost
        if "Won" in str(r.get("Stage", ""))
        and str(r.get("Close Date", ""))[:7] <= FQ["prior"]["month_end"]
    ]
    cur_q1_won_count = len(cur_q1_won)
    cur_q1_won_arr = sum(_unw(r) for r in cur_q1_won)
    cur_approved = sum(
        1 for r in approvals if str(r.get("Status", "")).strip() == "Approved 2026"
    )
    cur_missing = sum(1 for r in approvals if "Missing" in str(r.get("Status", "")))

    def _delta_int(cur, pri):
        if pri is None:
            return ""
        d = cur - pri
        if d == 0:
            return "unchanged"
        sign = "+" if d > 0 else ""
        return f"{sign}{d}"

    def _delta_arr(cur, pri):
        if pri is None:
            return ""
        d = cur - pri
        if abs(d) < 1000:
            return "unchanged"
        sign = "+" if d > 0 else ""
        return f"{sign}{_fmt_eur(d)}"

    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    _set_ph(
        slide,
        144,
        f"Since last review ({prior['period']}): what moved",
    )
    _set_ph(
        slide,
        145,
        (
            "Change in the six headline numbers versus the prior snapshot. "
            "Currency-converted values (ARR) may shift between runs due to "
            "ARR shifts between snapshots reflect FX movement."
        ),
    )

    rows = [
        ["Metric", f"{prior['period']}", f"{snapshot_date}", "Change"],
        [
            "Open Land deals",
            str(prior["open_deals"]) if prior["open_deals"] is not None else "-",
            str(cur_open_deals),
            _delta_int(cur_open_deals, prior["open_deals"]),
        ],
        [
            "Open Land ARR (unwtd)",
            _fmt_eur(prior["open_unwtd"]) if prior["open_unwtd"] is not None else "-",
            _fmt_eur(cur_open_unwtd),
            _delta_arr(cur_open_unwtd, prior["open_unwtd"]),
        ],
        [
            "Q1 Land wins (count)",
            str(prior["q1_won_count"]) if prior["q1_won_count"] is not None else "-",
            str(cur_q1_won_count),
            _delta_int(cur_q1_won_count, prior["q1_won_count"]),
        ],
        [
            "Q1 Land wins (ARR)",
            _fmt_eur(prior["q1_won_arr"]) if prior["q1_won_arr"] is not None else "-",
            _fmt_eur(cur_q1_won_arr),
            _delta_arr(cur_q1_won_arr, prior["q1_won_arr"]),
        ],
        [
            "Approved 2026",
            str(prior["approved_2026"]) if prior["approved_2026"] is not None else "-",
            str(cur_approved),
            _delta_int(cur_approved, prior["approved_2026"]),
        ],
        [
            "Missing Stage 3+ approvals",
            str(prior["missing_approval"])
            if prior["missing_approval"] is not None
            else "-",
            str(cur_missing),
            _delta_int(cur_missing, prior["missing_approval"]),
        ],
    ]

    def _delta_positive(v):
        s = str(v)
        return s.startswith("+") and "unchanged" not in s

    def _delta_negative(v):
        s = str(v)
        return s.startswith("-")

    _add_table(
        slide,
        rows,
        0.9,
        2.2,
        11.5,
        row_height=0.35,
        col_widths=[3.5, 2.5, 2.5, 3.0],
        data_font_size=Pt(12),
        highlight_rules={
            3: [
                (_delta_positive, "D4EDDA"),  # soft green for up
                (_delta_negative, "F8D7DA"),  # soft red for down
            ],
        },
    )


def slide_executive_summary(
    prs, pipeline, won_lost, renewals, kpis, territory, pi_data=None
):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    _set_ph(slide, 144, f"Exec. Summary | {territory}")

    # Pipeline KPIs — both ARR types
    total_opp_arr = sum(_unw(r) for r in pipeline)
    total_fc_arr = sum(_wtd(r) for r in pipeline)
    total_arr = total_opp_arr  # primary metric for backward compat
    deal_count = len(pipeline)
    stages = Counter(r.get("Stage", "") for r in pipeline)
    top_stage = stages.most_common(1)[0] if stages else ("?", 0)

    # Won/Lost
    won = [r for r in won_lost if "Won" in str(r.get("Stage", ""))]
    lost = [r for r in won_lost if "Won" not in str(r.get("Stage", ""))]
    won_arr = sum(_unw(r) for r in won)
    lost_arr = sum(_unw(r) for r in lost)
    total_decisions = len(won) + len(lost)
    win_rate = f"{len(won) / total_decisions * 100:.0f}%" if total_decisions else "n/a"

    # Forecast
    commit_deals = [r for r in pipeline if r.get("Forecast Category") == "Commit"]
    commit_opp = sum(_unw(r) for r in commit_deals)
    commit_fc = sum(_wtd(r) for r in commit_deals)
    commit_arr = commit_opp
    bc_deals = [r for r in pipeline if r.get("Forecast Category") == "Best Case"]
    bc_arr = sum(_unw(r) for r in bc_deals)

    # Renewals (support both new and legacy column names)
    renewal_acv = sum(
        float(r.get("ACV Unweighted (EUR)") or r.get("ACV (EUR)") or 0)
        for r in renewals
    )

    # Owner concentration from PI
    owner_callout = _top_pushed_owner(pi_data or [])

    # Column headers
    _set_ph(slide, 42, f"Exec. Summary | {territory}")
    _set_ph(slide, 56, f"Exec. Summary | {territory}")
    _set_ph(slide, 58, f"Exec. Summary | {territory}")
    _set_ph(slide, 60, f"Exec. Summary | {territory}")

    # Gradient metrics
    _add_gradient_metric(slide, 0, _fmt_eur(total_arr))
    _add_gradient_metric(slide, 1, f"{len(won)}W / {len(lost)}L")
    _add_gradient_metric(slide, 2, _fmt_eur(commit_arr))
    _add_gradient_metric(
        slide,
        3,
        str(len([r for r in (pi_data or []) if int(r.get("Push Count") or 0) >= 5])),
    )

    # Narrative bodies — showing both unweighted and weighted
    _set_ph(
        slide,
        22,
        f"Pipeline {_fmt_eur(total_opp_arr)} unweighted, {_fmt_eur(total_fc_arr)} weighted, across {deal_count} deals. "
        f"{top_stage[1] / max(deal_count, 1) * 100:.0f}% in {top_stage[0]}.",
    )
    _set_ph(
        slide,
        55,
        f"Win/Loss: {len(won)} won ({_fmt_eur(won_arr)}) vs {len(lost)} lost ({_fmt_eur(lost_arr)}). "
        f"Win rate: {win_rate}.",
    )
    _set_ph(
        slide,
        57,
        f"Commit (unweighted): {_fmt_eur(commit_opp)} | (weighted): {_fmt_eur(commit_fc)}. "
        f"Best case: {_fmt_eur(bc_arr)}. Renewal ACV: {_fmt_eur(renewal_acv)}.",
    )
    risk_text = (
        f"Risk: {owner_callout}, pattern review recommended."
        if owner_callout
        else "Risk: review pushed deal concentration."
    )
    _set_ph(slide, 59, risk_text)


def slide_q1_promised_vs_delivered(prs, q1_summary, territory):
    """Q1 Promised vs Delivered — what was the pipeline at Jan 1, what actually happened."""
    if not q1_summary:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    _set_ph(slide, 144, f"Q1 Promised vs Delivered | {territory}")

    promised = q1_summary.get("promised_arr", 0)
    won_arr = q1_summary.get("won_arr", 0)
    lost_arr = q1_summary.get("lost_arr", 0)
    won_count = q1_summary.get("won_count", 0)
    lost_count = q1_summary.get("lost_count", 0)
    advanced = q1_summary.get("advanced_count", 0)
    regressed = q1_summary.get("regressed_count", 0)

    _set_ph(slide, 42, f"Q1 Promised vs Delivered | {territory}")
    _set_ph(slide, 56, f"Q1 Promised vs Delivered | {territory}")
    _set_ph(slide, 58, f"Q1 Promised vs Delivered | {territory}")
    _set_ph(slide, 60, f"Q1 Promised vs Delivered | {territory}")

    _add_gradient_metric(slide, 0, _fmt_eur(promised))
    _add_gradient_metric(slide, 1, _fmt_eur(won_arr))
    _add_gradient_metric(slide, 2, _fmt_eur(lost_arr))
    _add_gradient_metric(slide, 3, str(advanced))

    _set_ph(
        slide,
        22,
        f"Q1 pipeline opened at {_fmt_eur(promised)} across {q1_summary.get('total_deals', 0)} deals. "
        f"Committed (stage 4+): {_fmt_eur(q1_summary.get('committed_arr', 0))} in {q1_summary.get('committed_count', 0)} deals.",
    )
    _set_ph(
        slide,
        55,
        f"Won {won_count} deals for {_fmt_eur(won_arr)}. "
        f"Top: {q1_summary.get('top_won', 'n/a')}.",
    )
    _set_ph(
        slide,
        57,
        f"Lost {lost_count} deals ({_fmt_eur(lost_arr)}). "
        f"Top loss: {q1_summary.get('top_lost', 'n/a')}.",
    )
    _set_ph(
        slide,
        59,
        f"{advanced} deals advanced stage. {regressed} regressed. "
        f"Net pipeline movement: {_fmt_eur(q1_summary.get('net_arr_change', 0))}.",
    )

    # Won/Lost detail table below
    top_won = q1_summary.get("top_won_deals", [])
    top_lost = q1_summary.get("top_lost_deals", [])
    if top_won or top_lost:
        rows = [["Outcome", "Opportunity", "Owner", "ARR (mEUR)"]]
        for d in top_won[:5]:
            rows.append(
                [
                    "Won",
                    d.get("name", ""),
                    d.get("owner", ""),
                    _fmt_eur(d.get("arr", 0)),
                ]
            )
        for d in top_lost[:5]:
            rows.append(
                [
                    "Lost",
                    d.get("name", ""),
                    d.get("owner", ""),
                    _fmt_eur(d.get("arr", 0)),
                ]
            )
        txBox = slide.shapes.add_textbox(
            Inches(0.9), Inches(5.4), Inches(11.0), Inches(0.2)
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Key won and lost deals in Q1"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = DARK


def _fetch_quarter_enrichment(
    session,
    instance,
    territory_soql_where,
    start_date=FQ["current"]["start"],
    end_date=FQ["current"]["end"],
):
    """Pull deal-level enrichment signals from SF for the forward-look slide.

    Returns a dict per opportunity: activity counts, AuM, competitor installed,
    NextStep age, quote count. These are signals the standard extract doesn't
    capture but that answer "is this deal actually being worked?"
    """

    q2_scope = (
        f"{territory_soql_where} AND IsClosed = false "
        f"AND CloseDate >= {start_date} AND CloseDate <= {end_date} "
        "AND Type = 'Land'"
    )
    enrichment = {}
    try:
        # Get Q2 deal IDs + basic fields
        r = session.get(
            f"{instance}/services/data/v66.0/query",
            params={
                "q": (
                    f"SELECT Id, Name, Account.Name, Account.AuM_m__c, "
                    f"StageName, CloseDate, LastActivityDate, LastModifiedDate, NextStep, "
                    f"APTS_Opportunity_ARR__c, ForecastCategoryName, PushCount "
                    f"FROM Opportunity WHERE {q2_scope} "
                    f"ORDER BY APTS_Opportunity_ARR__c DESC NULLS LAST"
                )
            },
            timeout=30,
        )
        if r.status_code != 200:
            return enrichment
        deals = r.json().get("records", [])

        for deal in deals:
            oid = deal.get("Id", "")
            from datetime import date as _date

            close = str(deal.get("CloseDate") or "")[:10]
            last_act = str(deal.get("LastActivityDate") or "")[:10]
            last_mod = str(deal.get("LastModifiedDate") or "")[:10]
            try:
                days_to_close = (
                    (_date.fromisoformat(close) - _date.today()).days if close else None
                )
            except ValueError:
                days_to_close = None
            try:
                days_since_activity = (
                    (_date.today() - _date.fromisoformat(last_act)).days
                    if last_act and last_act != ""
                    else None
                )
            except ValueError:
                days_since_activity = None

            acct = deal.get("Account") or {}
            enrichment[oid] = {
                "name": deal.get("Name", ""),
                "account": acct.get("Name", ""),
                "stage": deal.get("StageName", ""),
                "close_date": close,
                "days_to_close": days_to_close,
                "arr": float(deal.get("APTS_Opportunity_ARR__c") or 0),
                "forecast_cat": deal.get("ForecastCategoryName", ""),
                "push_count": int(deal.get("PushCount") or 0),
                "last_activity": last_act if last_act else "never",
                "days_since_activity": days_since_activity,
                "next_step": str(deal.get("NextStep") or "")[:60],
                "next_step_stale": days_since_activity is not None
                and days_since_activity > 60,
                "aum_b": acct.get("AuM_m__c"),
            }

        # Competitor check: which accounts have installed competitor products?
        opp_ids = list(enrichment.keys())
        if opp_ids:
            acct_ids = set()
            for d in deals:
                aid = (d.get("Account") or {}).get("Id") or d.get("AccountId")
                if aid:
                    acct_ids.add(aid)
            if acct_ids:
                ids_sql = ",".join(f"'{a}'" for a in acct_ids)
                r2 = session.get(
                    f"{instance}/services/data/v66.0/query",
                    params={
                        "q": (
                            f"SELECT Account__c FROM Installed_Competitor_Product__c "
                            f"WHERE Account__c IN ({ids_sql})"
                        )
                    },
                    timeout=20,
                )
                if r2.status_code == 200:
                    comp_accts = {
                        rec.get("Account__c") for rec in r2.json().get("records", [])
                    }
                    for deal in deals:
                        oid = deal.get("Id", "")
                        aid = (deal.get("Account") or {}).get("Id") or deal.get(
                            "AccountId"
                        )
                        if oid in enrichment:
                            enrichment[oid]["competitor_installed"] = aid in comp_accts
    except Exception:
        pass
    return enrichment


def slide_q2_forward_look(
    prs,
    pipeline,
    enrichment,
    territory,
    workbook_path=None,
    q_label="Q2",
    month_start=FQ["current"]["month_start"],
    month_end=FQ["current"]["month_end"],
):
    """Quarter Forward Look — per-deal readiness grid with enrichment signals.

    Shows each quarter's deals with: days to close, last activity age, NextStep
    freshness, approval status, AuM context, competitor flag. Also reads
    Forecast Category History + Stage History from the director workbook
    to surface momentum signals (who upgraded/downgraded commit this month).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])

    q2 = [
        r
        for r in pipeline
        if month_start <= str(r.get("Close Date", ""))[:7] <= month_end
    ]
    if not q2 and not enrichment:
        _set_ph(slide, 144, f"{q_label} Forward Look: {territory}")
        _set_ph(slide, 145, f"No {q_label} Land deals in scope.")
        return

    # Build the readiness assessment per deal
    deals_enriched = []
    for r in q2:
        opp_name = r.get("Opportunity") or r.get("Name") or ""
        # Match to enrichment by name (since we may not have IDs in the pipeline rows)
        enr = None
        for eid, edata in enrichment.items():
            if edata.get("name") == opp_name or edata.get("account") == r.get(
                "Account"
            ):
                enr = edata
                break
        if not enr:
            enr = {
                "days_to_close": None,
                "days_since_activity": None,
                "aum_b": None,
                "competitor_installed": False,
                "next_step_stale": False,
                "last_activity": "?",
                "next_step": "",
            }

        days_left = enr.get("days_to_close")
        days_silent = enr.get("days_since_activity")
        readiness = "✓ On track"
        if days_silent is not None and days_silent > 60:
            readiness = "🔴 Silent 60d+"
        elif days_silent is not None and days_silent > 30:
            readiness = "🟡 Quiet 30d+"
        elif (
            days_left is not None
            and days_left <= 30
            and str(r.get("Stage", "")).startswith(("1", "2", "3"))
        ):
            readiness = "🟡 Early stage, close soon"
        if enr.get("competitor_installed"):
            readiness += " | Competitor"

        deals_enriched.append(
            {
                "account": str(r.get("Account", ""))[:22],
                "opportunity": str(opp_name)[:22],
                "stage": str(r.get("Stage", "")),
                "close": str(r.get("Close Date", ""))[:10],
                "days_left": str(days_left) + "d" if days_left is not None else "?",
                "arr": _fmt_eur(_unw(r)),
                "last_activity": str(enr.get("last_activity", "?"))[:10],
                "aum": f"{enr['aum_b']:.0f}B" if enr.get("aum_b") else "-",
                "readiness": readiness,
                "next_step": str(enr.get("next_step", ""))[:35],
            }
        )

    # Headline: activity alert
    silent_count = sum(
        1
        for d in deals_enriched
        if "Silent" in d["readiness"] or "Quiet" in d["readiness"]
    )
    total_q2_arr = sum(_unw(r) for r in q2)

    if silent_count == len(deals_enriched) and len(deals_enriched) > 0:
        headline = (
            f"{q_label} {territory}: {len(q2)} deals, {_fmt_eur(total_q2_arr)}. "
            f"⚠ Zero recent activity across the entire {q_label} book."
        )
    elif silent_count > 0:
        headline = (
            f"{q_label} {territory}: {len(q2)} deals, {_fmt_eur(total_q2_arr)}. "
            f"{silent_count} of {len(q2)} deals silent or quiet."
        )
    else:
        headline = f"{q_label} {territory}: {len(q2)} deals, {_fmt_eur(total_q2_arr)}."

    _set_ph(slide, 144, headline)
    _set_ph(
        slide,
        145,
        f"{q_label} deals ranked by activity and timeline risk.",
    )

    rows = [
        [
            "Account",
            "Opportunity",
            "Stage",
            "Close",
            "Days",
            "ARR",
            "Last Activity",
            "AuM",
            "Readiness",
            "Next Step",
        ]
    ]
    for d in deals_enriched:
        rows.append(
            [
                d["account"],
                d["opportunity"],
                d["stage"],
                d["close"],
                d["days_left"],
                d["arr"],
                d["last_activity"],
                d["aum"],
                d["readiness"],
                d["next_step"],
            ]
        )
    _add_table(
        slide,
        rows,
        0.3,
        2.0,
        12.7,
        col_widths=[1.8, 1.8, 1.2, 0.8, 0.6, 1.0, 0.9, 0.6, 1.8, 2.2],
        data_font_size=Pt(8),
    )

    # Momentum signals from Forecast Category History + Stage History
    # (read from the director workbook if available)
    momentum_lines = []
    if workbook_path and Path(str(workbook_path)).exists():
        try:
            _wb = load_workbook(str(workbook_path), data_only=True, read_only=True)
            cat_rank = {
                "Omitted": 0,
                "Pipeline": 1,
                "Best Case": 2,
                "Commit": 3,
                "Closed": 4,
            }
            stg_rank = {"1": 1, "2": 2, "3": 3, "4": 4, "5": 5, "6": 6, "8": 8}
            from datetime import date as _date, timedelta as _td

            cutoff = (_date.today() - _td(days=30)).isoformat()

            if "Forecast Category History" in _wb.sheetnames:
                fch = list(_wb["Forecast Category History"].iter_rows(values_only=True))
                recent_up = [
                    r
                    for r in fch[1:]
                    if str(r[6] or "") >= cutoff
                    and cat_rank.get(str(r[5] or ""), 0)
                    > cat_rank.get(str(r[4] or ""), 0)
                ]
                recent_down = [
                    r
                    for r in fch[1:]
                    if str(r[6] or "") >= cutoff
                    and cat_rank.get(str(r[5] or ""), 0)
                    < cat_rank.get(str(r[4] or ""), 0)
                ]
                if recent_up or recent_down:
                    momentum_lines.append(
                        f"Forecast momentum (30d): {len(recent_up)} upgrades, "
                        f"{len(recent_down)} downgrades."
                    )
                    for r in recent_up[:2]:
                        momentum_lines.append(
                            f"  ↑ {str(r[1] or '')[:22]}: {r[4]} → {r[5]} ({str(r[6] or '')[:10]})"
                        )
                    for r in recent_down[:2]:
                        momentum_lines.append(
                            f"  ↓ {str(r[1] or '')[:22]}: {r[4]} → {r[5]} ({str(r[6] or '')[:10]})"
                        )

            if "Stage History" in _wb.sheetnames:
                sh = list(_wb["Stage History"].iter_rows(values_only=True))
                recent_adv = [
                    r
                    for r in sh[1:]
                    if str(r[6] or "") >= cutoff
                    and str(r[5] or "")[:1].isdigit()
                    and str(r[4] or "")[:1].isdigit()
                    and str(r[5] or "") > str(r[4] or "")
                ]
                if recent_adv:
                    momentum_lines.append(
                        f"Stage advances (30d): {len(recent_adv)} deals moved forward."
                    )
                    for r in recent_adv[:2]:
                        momentum_lines.append(
                            f"  ↑ {str(r[1] or '')[:22]}: {r[4]} → {r[5]} ({str(r[6] or '')[:10]})"
                        )
            _wb.close()
        except Exception:
            pass

    if momentum_lines:
        from pptx.util import Inches

        tb = slide.shapes.add_textbox(
            Inches(0.3), Inches(5.8), Inches(12.0), Inches(1.6)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(momentum_lines[:6]):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(7)
            para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            para.font.name = "Calibri"
            para.space_after = Pt(1)


def slide_quarter_outlook(
    prs,
    pipeline,
    won_lost,
    territory,
    q_label="Q2",
    month_start=FQ["current"]["month_start"],
    month_end=FQ["current"]["month_end"],
    month_range_label=FQ["current"]["range_label"],
):
    """Quarter Outlook — what's promised, what's delivered so far, what moved."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])

    qtr = [
        r
        for r in pipeline
        if month_start <= str(r.get("Close Date", ""))[:7] <= month_end
    ]
    qtr_opp_arr = sum(_unw(r) for r in qtr)
    qtr_fc_arr = sum(_wtd(r) for r in qtr)
    qtr_arr = qtr_opp_arr
    commit = [r for r in qtr if r.get("Forecast Category") == "Commit"]
    commit_opp = sum(_unw(r) for r in commit)
    commit_fc = sum(_wtd(r) for r in commit)
    commit_arr = commit_opp
    bc = [r for r in qtr if r.get("Forecast Category") == "Best Case"]
    bc_arr = sum(_unw(r) for r in bc)

    qtr_won = [
        r
        for r in won_lost
        if "Won" in str(r.get("Stage", ""))
        and month_start <= str(r.get("Close Date", ""))[:7] <= month_end
    ]
    qtr_lost = [
        r
        for r in won_lost
        if "Won" not in str(r.get("Stage", ""))
        and month_start <= str(r.get("Close Date", ""))[:7] <= month_end
    ]
    qtr_won_arr = sum(_unw(r) for r in qtr_won)
    qtr_lost_arr = sum(_unw(r) for r in qtr_lost)

    _set_ph(
        slide,
        144,
        f"{q_label} book {_fmt_eur(qtr_arr)} unweighted, {_fmt_eur(qtr_fc_arr)} weighted. "
        f"{len(qtr)} deals closing {month_range_label}.",
    )
    if q_label == "Q3":
        _set_ph(slide, 145, "No Q2 Land pipeline. Showing Q3 forward book.")
    _set_ph(slide, 42, f"{q_label} Outlook | {territory}")
    _set_ph(slide, 56, f"{q_label} Outlook | {territory}")
    _set_ph(slide, 58, f"{q_label} Outlook | {territory}")
    _set_ph(slide, 60, f"{q_label} Outlook | {territory}")

    _add_gradient_metric(slide, 0, _fmt_eur(qtr_arr))
    _add_gradient_metric(slide, 1, _fmt_eur(commit_arr))
    _add_gradient_metric(slide, 2, f"{len(qtr_won)}W / {len(qtr_lost)}L")
    _add_gradient_metric(slide, 3, _fmt_eur(qtr_won_arr))

    _set_ph(
        slide,
        22,
        f"{q_label} book {_fmt_eur(qtr_opp_arr)} unweighted, {_fmt_eur(qtr_fc_arr)} weighted. "
        f"{len(qtr)} deals closing {month_range_label}.",
    )
    _set_ph(
        slide,
        55,
        f"Commit (unweighted): {_fmt_eur(commit_opp)} | (weighted): {_fmt_eur(commit_fc)}. "
        f"Best case: {_fmt_eur(bc_arr)}.",
    )
    _set_ph(
        slide,
        57,
        f"{q_label} results so far: {len(qtr_won)} won ({_fmt_eur(qtr_won_arr)}) vs {len(qtr_lost)} lost ({_fmt_eur(qtr_lost_arr)}). "
        f"{'Lost exceeds won, review loss reasons.' if qtr_lost_arr > qtr_won_arr else 'Positive trajectory.'}",
    )
    _set_ph(
        slide,
        59,
        f"Won ARR: {_fmt_eur(qtr_won_arr)}. "
        f"{'Early in quarter, {:.0f}% of commit delivered.'.format(qtr_won_arr / max(commit_arr, 1) * 100) if commit_arr else 'No commit baseline.'}",
    )


def slide_q2_outlook(prs, pipeline, won_lost, territory):
    slide_quarter_outlook(
        prs,
        pipeline,
        won_lost,
        territory,
        q_label="Q2",
        month_start=FQ["current"]["month_start"],
        month_end=FQ["current"]["month_end"],
        month_range_label=FQ["current"]["range_label"],
    )


def slide_pipeline_overview(prs, pipeline, territory):
    """Pipeline stage distribution narrative — matches Rebekka slide 3."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    total_arr = sum(_unw(r) for r in pipeline)
    breakdown = _stage_breakdown(pipeline)

    total_wtd = sum(_wtd(r) for r in pipeline)
    # Find dominant stage
    if breakdown:
        top_s, top_c, top_a, _top_w = max(breakdown, key=lambda x: x[2])
        pct = top_c / max(len(pipeline), 1) * 100
        _set_ph(slide, 144, f"Pipeline Overview | {territory}")
        _set_ph(
            slide,
            145,
            f"Unweighted: {_fmt_eur(total_arr)}  |  Weighted: {_fmt_eur(total_wtd)}",
        )
        _set_ph(
            slide, 42, f"{pct:.0f}% of pipeline ({_fmt_eur(top_a)}) sits in {top_s}."
        )
    else:
        _set_ph(slide, 144, f"Pipeline Overview - {territory}")
        _set_ph(slide, 145, "No open pipeline data.")

    # Stage breakdown table
    if breakdown:
        rows = [
            [
                "Stage",
                "Deals",
                "ARR (mEUR)",
                "ARR Wtd (mEUR)",
                "% of Pipeline",
            ]
        ]
        for stage, count, arr, wtd in breakdown:
            pct = arr / max(total_arr, 1) * 100
            rows.append([stage, str(count), _meur(arr), _meur(wtd), f"{pct:.0f}%"])
        rows.append(
            [
                "Total",
                str(len(pipeline)),
                _meur(total_arr),
                _meur(total_wtd),
                "100%",
            ]
        )
        _add_table(
            slide,
            rows,
            0.9,
            2.2,
            6.5,
            row_height=0.25,
            col_widths=[1.8, 0.7, 1.3, 1.3, 1.4],
        )


def slide_forecast_accuracy(prs, won_lost, pipeline, territory):
    """Forecast Accuracy — 4 cards. Matches Rebekka slide 13."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])

    won = [r for r in won_lost if "Won" in str(r.get("Stage", ""))]
    lost = [r for r in won_lost if "Won" not in str(r.get("Stage", ""))]
    won_arr = sum(_unw(r) for r in won)
    lost_arr = sum(_unw(r) for r in lost)
    total_decisions = len(won) + len(lost)
    win_rate = f"{len(won) / total_decisions * 100:.0f}%" if total_decisions else "n/a"

    commit_deals = [r for r in pipeline if r.get("Forecast Category") == "Commit"]
    commit_arr = sum(_unw(r) for r in commit_deals)
    bc_deals = [r for r in pipeline if r.get("Forecast Category") == "Best Case"]
    bc_arr = sum(_unw(r) for r in bc_deals)
    pipe_deals = [r for r in pipeline if r.get("Forecast Category") == "Pipeline"]
    pipe_arr = sum(_unw(r) for r in pipe_deals)

    _set_ph(slide, 144, f"Forecast Accuracy | {territory}")

    _set_ph(slide, 42, f"Forecast Accuracy | {territory}")
    _set_ph(slide, 56, f"Forecast Accuracy | {territory}")
    _set_ph(slide, 58, f"Forecast Accuracy | {territory}")
    _set_ph(slide, 60, f"Forecast Accuracy | {territory}")

    _add_gradient_metric(slide, 0, _fmt_eur(won_arr))
    _add_gradient_metric(slide, 1, _fmt_eur(lost_arr))
    _add_gradient_metric(slide, 2, win_rate)
    _add_gradient_metric(slide, 3, _fmt_eur(commit_arr))

    _set_ph(
        slide,
        22,
        f"{len(won)} {'deal' if len(won) == 1 else 'deals'} closed-won this quarter. "
        f"{'Low won ARR signals early-quarter timing or delayed closes.' if won_arr < lost_arr else 'Won ARR exceeds lost. Positive trajectory.'}",
    )
    _set_ph(
        slide,
        55,
        f"{len(lost)} deals closed-lost. "
        f"{'Lost ARR significantly exceeds won ARR. Review loss reasons.' if lost_arr > won_arr else 'Lost ARR within acceptable range.'}",
    )
    _set_ph(
        slide,
        57,
        f"Win rate by deal count ({len(won)}W / {len(lost)}L). "
        f"ARR win rate: {won_arr / max(won_arr + lost_arr, 1) * 100:.0f}%.",
    )
    _set_ph(
        slide,
        59,
        f"Commit forecast for remaining open deals this quarter. "
        f"Best Case: {_fmt_eur(bc_arr)}. Pipeline: {_fmt_eur(pipe_arr)}.",
    )


def slide_pipeline_combined(prs, pipeline, territory):
    """Pipeline Overview + Stage 3+ with key deals sidebar — merged into one slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    total_opp = sum(_unw(r) for r in pipeline)
    total_fc = sum(_wtd(r) for r in pipeline)
    total_arr = total_opp
    breakdown = _stage_breakdown(pipeline)

    _set_ph(slide, 144, f"Pipeline Overview | {territory}")

    if not pipeline:
        _set_ph(
            slide,
            145,
            f"No open Land deals with {FQ['prior']['label']}-{FQ['forward']['label']} close dates.",
        )
        _set_ph(
            slide,
            42,
            "All open Land pipeline for this territory sits in Q3 or Q4.",
        )
        return

    top_s, top_c, top_a, _top_w = max(breakdown, key=lambda x: x[2])
    pct = top_c / max(len(pipeline), 1) * 100

    _set_ph(
        slide,
        145,
        f"Unweighted: {_fmt_eur(total_opp)}  |  Weighted: {_fmt_eur(total_fc)}",
    )
    if total_arr > 0:
        _set_ph(
            slide,
            42,
            f"{pct:.0f}% of pipeline ({_fmt_eur(top_a)}) sits in {top_s}.",
        )
    else:
        # Pipeline has deals but ARR field is unpopulated (typical for early
        # stage Prospecting records). Say that plainly.
        _set_ph(
            slide,
            42,
            f"{len(pipeline)} deal(s) in the book, ARR not yet populated. "
            f"{pct:.0f}% sit in {top_s}.",
        )

    # Left: stage breakdown table
    if breakdown:
        rows = [
            [
                "Stage",
                "Deals",
                "ARR (mEUR)",
                "ARR Wtd (mEUR)",
                "% of Pipeline",
            ]
        ]
        for stage, count, arr, wtd in breakdown:
            p = arr / max(total_arr, 1) * 100
            rows.append([stage, str(count), _meur(arr), _meur(wtd), f"{p:.0f}%"])
        rows.append(
            [
                "Total",
                str(len(pipeline)),
                _meur(total_arr),
                _meur(total_fc),
                "100%",
            ]
        )
        _add_table(
            slide,
            rows,
            0.9,
            2.2,
            7.3,
            row_height=0.25,
            col_widths=[1.9, 0.7, 1.4, 1.4, 1.4],
        )

    # Right: key deals sidebar (top 3 stage 3+ by ARR)
    stage3_plus = [r for r in pipeline if r.get("Stage", "") >= "3"]
    key_deals = sorted(stage3_plus, key=lambda r: _unw(r), reverse=True)[:3]
    if key_deals:
        txBox = slide.shapes.add_textbox(
            Inches(8.9), Inches(1.7), Inches(3.7), Inches(0.2)
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Key deals this quarter"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = DARK

        rows = [["Opportunity Name", "Close Date", "ARR (mEUR)"]]
        for r in key_deals:
            rows.append(
                [
                    str(r.get("Opportunity", "")),
                    str(r.get("Close Date", "")),
                    _meur(_unw(r)),
                ]
            )
        _add_table(
            slide,
            rows,
            8.7,
            2.1,
            4.6,
            row_height=0.41,
            col_widths=[1.5, 1.5, 1.5],
            data_font_size=Pt(14),
        )


def slide_pipeline_stage3_plus(prs, pipeline, territory):
    """Pipeline Volume Stage 3+ with key deals sidebar. Matches Rebekka slide 4."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])

    stage3_plus = [r for r in pipeline if r.get("Stage", "") >= "3"]
    stage3_arr = sum(_unw(r) for r in stage3_plus)

    _set_ph(slide, 144, "Pipeline Volume Stage 3+")
    _set_ph(slide, 145, "FY26 year to date")

    # Key deals sidebar (right side — top 3 by ARR)
    key_deals = sorted(stage3_plus, key=lambda r: _unw(r), reverse=True)[:3]
    if key_deals:
        # Sidebar label
        txBox = slide.shapes.add_textbox(
            Inches(8.9), Inches(1.7), Inches(3.7), Inches(0.2)
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Key deals for this quarter"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = DARK

        rows = [["Opportunity Name", "Close Date", "ARR (mEUR)"]]
        for r in key_deals:
            rows.append(
                [
                    str(r.get("Opportunity", "")),
                    str(r.get("Close Date", "")),
                    _meur(_unw(r)),
                ]
            )
        _add_table(
            slide,
            rows,
            8.7,
            2.1,
            4.6,
            row_height=0.41,
            col_widths=[1.5, 1.5, 1.5],
            data_font_size=Pt(14),
        )

    # Main content — stage breakdown narrative
    breakdown = _stage_breakdown(stage3_plus)
    if breakdown:
        lines = [
            f"Stage 3+ pipeline: {_fmt_eur(stage3_arr)} across {len(stage3_plus)} deals",
            "",
        ]
        for stage, count, arr, wtd in breakdown:
            pct = arr / max(stage3_arr, 1) * 100
            lines.append(
                f"{stage}: {count} deals, {_meur(arr)} unwtd / {_meur(wtd)} wtd mEUR ({pct:.0f}%)"
            )
        _set_ph_lines(slide, 22, lines)


def slide_top_deals(prs, pipeline):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    # Filter to ARR > 0, take top 12
    deals = [r for r in pipeline if _unw(r) > 0]
    deals.sort(key=lambda r: _unw(r), reverse=True)
    top = deals[:12]

    total_arr = sum(_unw(r) for r in top)
    if not top:
        _set_ph(slide, 144, "Key Deals")
        _set_ph(
            slide,
            145,
            f"No open Land deals with {FQ['prior']['label']}-{FQ['forward']['label']} close dates in this territory.",
        )
        return
    # Data-forward title: largest deal named in the headline + total concentration.
    biggest = top[0]
    biggest_arr = _unw(biggest)
    biggest_account = str(biggest.get("Account") or "")
    top5_share = sum(_unw(r) for r in top[:5]) / total_arr if total_arr else 0
    _set_ph(
        slide,
        144,
        (
            f"Key {len(top)} open deals = {_fmt_eur(total_arr)}. "
            f"{biggest_account} leads at {_fmt_eur(biggest_arr)}; "
            f"top 5 = {top5_share * 100:.0f}% of this book."
        ),
    )
    _set_ph(
        slide,
        145,
        "Open opportunities sorted by deal size. Focus: are these progressing or stalling?",
    )

    if top:
        rows = [
            [
                "Account",
                "Opportunity",
                "Owner",
                "Stage",
                "Close Date",
                "Age",
                "ARR (mEUR)",
                "ARR Wtd (mEUR)",
            ]
        ]
        for r in top:
            # Compute age from Created date
            created = str(r.get("Created", ""))
            age = ""
            if created:
                try:
                    age_days = (
                        datetime.now() - datetime.strptime(created[:10], "%Y-%m-%d")
                    ).days
                    age = str(age_days)
                except ValueError:
                    pass
            rows.append(
                [
                    str(r.get("Account", "")),
                    str(r.get("Opportunity", "")),
                    str(r.get("Owner", "")),
                    str(r.get("Stage", "")),
                    str(r.get("Close Date", "")),
                    age,
                    _meur(_unw(r)),
                    _meur(_wtd(r)),
                ]
            )
        _add_table(
            slide,
            rows,
            0.9,
            2.2,
            11.5,
            col_widths=[1.5, 1.5, 1.5, 1.5, 1.3, 1.0, 1.1, 1.1],
        )


def slide_deal_risk_scoring(prs, risk_deals, territory):
    """Top Q2-closing deals at risk for this director with score + reason codes.

    Filters risk_deals to Q2 (Apr-Jun) close dates so the slide is scoped to
    the current quarter, matching the monthly review's forward-looking
    commitment. Reason codes come straight from the workbook's Deal Risk
    Scoring tab; the score weights live on Parameters.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])

    # Scope to Q2+Q3 close dates
    def _in_forward(d):
        cd = str(d.get("close_date") or "")[:10]
        return cd >= FQ["current"]["start"] and cd <= FQ["forward"]["end"]

    fwd_deals = [d for d in risk_deals if _in_forward(d)]

    if not fwd_deals:
        _set_ph(
            slide,
            144,
            f"Key Deals at Risk in {FQ['current']['label']}-{FQ['forward']['label']}, no flags",
        )
        _set_ph(
            slide,
            145,
            f"No {FQ['current']['label']}-{FQ['forward']['label']} Land deals in {territory} flagged for risk this period.",
        )
        return

    top = fwd_deals[:10]
    total_arr = sum(d["arr"] for d in top)
    _set_ph(
        slide,
        144,
        f"Key {len(top)} {FQ['current']['label']}-{FQ['forward']['label']} Deals at Risk, {_fmt_eur(total_arr)} exposed",
    )
    _set_ph(
        slide,
        145,
        "Deals most likely to slip or miss this quarter, ranked by risk severity.",
    )

    rows = [
        [
            "#",
            "Score",
            "Account",
            "Opportunity",
            "Stage",
            "Close",
            "ARR (mEUR)",
            "Reasons",
        ]
    ]
    for d in top:
        rows.append(
            [
                str(d["rank"]),
                str(d["score"]),
                str(d["account"]),
                str(d["opportunity"]),
                str(d["stage"]),
                str(d["close_date"])[:10],
                _meur(d["arr"]),
                str(d["reasons"]),
            ]
        )
    _add_table(
        slide,
        rows,
        0.3,
        2.2,
        12.7,
        col_widths=[0.35, 0.55, 2.1, 2.5, 1.4, 0.9, 1.1, 3.8],
        data_font_size=Pt(9),
    )


def slide_forecast_variance(prs, variance, territory):
    """Q1 forecast variance decomposition for this director.

    Explains how the Q1 pipeline moved from initial to final snapshot by
    decomposing the change into Won / Lost / Added / Revised. Should tell
    the director whether pipeline shrinkage is win-driven (good) or
    loss-driven (bad).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    if not variance:
        _set_ph(slide, 144, "Q1 Forecast Variance")
        _set_ph(
            slide,
            145,
            f"No Q1 Historical Trending snapshot available for {territory}.",
        )
        return

    net = variance["net"]
    initial = variance["initial"]
    final = variance["final"]
    direction = "shrank" if net < 0 else "grew"
    _set_ph(
        slide,
        144,
        f"Q1 Forecast Variance, {_fmt_eur(abs(net))} {direction}",
    )

    # Decide the narrative: what drove the change?
    won = variance["won"]
    lost = variance["lost"]
    subtitle = (
        f"Q1 pipeline moved from {_fmt_eur(initial)} to {_fmt_eur(final)} "
        f"between the first and last Historical Trending snapshots. "
    )
    closed = won + lost
    if initial > 0 and closed / initial >= 0.8:
        if lost > won:
            subtitle += (
                f"{closed / initial * 100:.0f}% of Q1 pipeline closed; "
                "losses outweighed wins."
            )
        else:
            subtitle += (
                f"{closed / initial * 100:.0f}% of Q1 pipeline closed; "
                "wins outweighed losses."
            )
    elif lost > won * 2 and lost > 0:
        subtitle += "Losses outpaced wins >2x. Pipeline shrinkage is loss-driven."
    elif won > lost * 2 and won > 0:
        subtitle += "Wins outpaced losses >2x. Clean quarterly execution."
    else:
        subtitle += (
            "Wins and losses roughly balanced; remainder is deal-level "
            "revisions or new additions."
        )
    _set_ph(slide, 145, subtitle)

    rows = [
        ["Bucket", "ARR Impact"],
        ["Initial Q1 pipeline (open at first snapshot)", _fmt_eur(initial)],
        ["Closed Won (ARR removed from pipeline)", "-" + _fmt_eur(won)],
        ["Closed Lost (ARR removed from pipeline)", "-" + _fmt_eur(lost)],
        ["New deals added (after first snapshot)", "+" + _fmt_eur(variance["added"])],
        ["ARR revised up (still-open deals)", "+" + _fmt_eur(variance["up"])],
        ["ARR revised down (still-open deals)", "-" + _fmt_eur(variance["down"])],
        ["Final Q1 pipeline (still open at last snapshot)", _fmt_eur(final)],
        ["Net change (Final − Initial)", _fmt_eur(net)],
    ]
    _add_table(
        slide,
        rows,
        2.5,
        2.0,
        8.5,
        col_widths=[5.5, 3.0],
        row_height=0.34,
        data_font_size=Pt(11),
    )


STAGE_RANK_FOR_LOSS = {
    "1 - Prospecting": 1,
    "2 - Discovery": 2,
    "3 - Engagement": 3,
    "4 - Shortlisted": 4,
    "5 - Preferred": 5,
    "6 - Contracting": 6,
    "8 - Won": 8,
    "0 - Lost": -1,
    "0 - No Opportunity": -2,
}
STAGES_FOR_LOSS_RENDER = [
    "1 - Prospecting",
    "2 - Discovery",
    "3 - Engagement",
    "4 - Shortlisted",
    "5 - Preferred",
    "6 - Contracting",
]


def _stage_at_loss_from_history(dashboard_path, opportunity_names):
    """Cross-reference Dashboard's Q1 History Raw: max stage each opp reached
    before closing. Returns {opp_name: stage_str}. Empty if workbook missing
    or the tab isn't present.
    """
    if not dashboard_path or not Path(dashboard_path).exists() or not opportunity_names:
        return {}
    wb = load_workbook(str(dashboard_path), read_only=True, data_only=True)
    if "Q1 History Raw" not in wb.sheetnames:
        wb.close()
        return {}
    ws = wb["Q1 History Raw"]
    rows = list(ws.iter_rows(values_only=True))
    wb.close()
    if not rows:
        return {}
    headers = rows[0]
    opp_set = set(opportunity_names)
    best = {}
    for r in rows[1:]:
        d = dict(zip(headers, r))
        opp = d.get("Opportunity")
        if opp not in opp_set:
            continue
        if d.get("Field Changed") != "StageName":
            continue
        for fld in ("Old Value", "New Value"):
            stage = str(d.get(fld) or "")
            rank = STAGE_RANK_FOR_LOSS.get(stage, 0)
            if rank >= 1:
                cur = best.get(opp)
                if cur is None or rank > STAGE_RANK_FOR_LOSS.get(cur, 0):
                    best[opp] = stage
    return best


def slide_win_loss_diagnostic(prs, won_lost, territory, dashboard_path=None):
    """Why we lost: loss reasons + stage at loss for Q1 Land losses.

    Scope matches slide 4 (Q1 Promised vs Delivered). Left table: loss reason
    codes with count and ARR. Right table: highest stage each lost opp ever
    reached before closing (from OpportunityFieldHistory). Early-stage losses
    = qualification gap; late-stage losses = execution gap.
    """
    # Filter won_lost to Q1 Land losses (incoming list is already Land + Q1-Q2).
    losses = []
    for r in won_lost:
        stage = str(r.get("Stage", "") or "")
        if "Won" in stage:
            continue
        cd = str(r.get("Close Date", "") or "")[:10]
        if not (FQ["prior"]["start"] <= cd <= FQ["prior"]["end"]):
            continue
        losses.append(
            {
                "opportunity": r.get("Opportunity") or "",
                "arr": float(r.get("ARR Unweighted (EUR)") or 0),
                "reason": str(r.get("Reason") or "(not recorded)"),
            }
        )

    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    total_n = len(losses)
    total_arr = sum(loss["arr"] for loss in losses)

    if total_n == 0:
        _set_ph(slide, 144, f"Why We Lost Q1 | {territory}")
        _set_ph(slide, 145, "No Q1 Land losses in scope.")
        return

    # Data-forward title: includes the headline number.
    uncoded = sum(1 for loss in losses if loss["reason"] == "(not recorded)")
    headline = (
        f"{total_n} Q1 Land losses, {_fmt_eur(total_arr)}. "
        f"{uncoded} missing reason code."
        if uncoded
        else f"{total_n} Q1 Land losses, {_fmt_eur(total_arr)}."
    )
    _set_ph(slide, 144, headline)
    _set_ph(
        slide,
        145,
        "Left: loss reasons. "
        "Right: highest stage reached before loss. Early-stage = qualification gap; "
        "late-stage = execution gap.",
    )

    # Left table: loss reasons
    from collections import defaultdict

    reason_agg = defaultdict(lambda: {"n": 0, "arr": 0.0})
    for loss in losses:
        reason_agg[loss["reason"]]["n"] += 1
        reason_agg[loss["reason"]]["arr"] += loss["arr"]
    reason_tbl = [["Loss Reason", "Count", "Lost ARR"]]
    for reason, agg in sorted(reason_agg.items(), key=lambda x: -x[1]["n"]):
        reason_tbl.append([reason, str(agg["n"]), _meur(agg["arr"])])
    reason_tbl.append(["Total", str(total_n), _meur(total_arr)])
    _add_table(
        slide,
        reason_tbl,
        0.3,
        2.0,
        6.4,
        col_widths=[3.8, 1.2, 1.4],
        data_font_size=Pt(10),
    )

    # Right table: stage at loss (from OpportunityFieldHistory)
    stage_at_loss = _stage_at_loss_from_history(
        dashboard_path, [loss["opportunity"] for loss in losses]
    )
    stage_agg = defaultdict(lambda: {"n": 0, "arr": 0.0})
    unclassified_n = 0
    unclassified_arr = 0.0
    for loss in losses:
        stg = stage_at_loss.get(loss["opportunity"])
        if stg:
            stage_agg[stg]["n"] += 1
            stage_agg[stg]["arr"] += loss["arr"]
        else:
            unclassified_n += 1
            unclassified_arr += loss["arr"]
    stage_tbl = [["Stage Reached Before Loss", "Count", "Lost ARR"]]
    for stg in STAGES_FOR_LOSS_RENDER:
        if stage_agg[stg]["n"] == 0 and stage_agg[stg]["arr"] == 0:
            continue
        stage_tbl.append([stg, str(stage_agg[stg]["n"]), _meur(stage_agg[stg]["arr"])])
    if unclassified_n:
        stage_tbl.append(
            ["(no stage history)", str(unclassified_n), _meur(unclassified_arr)]
        )
    stage_tbl.append(["Total", str(total_n), _meur(total_arr)])
    _add_table(
        slide,
        stage_tbl,
        6.9,
        2.0,
        6.0,
        col_widths=[3.0, 1.2, 1.8],
        data_font_size=Pt(10),
    )


# Director-friendly translations of analyst risk-code labels. Short phrases
# so they fit in a table cell and read like a 1:1 discussion topic, not a
# rule-engine output.
REASON_CODE_LABELS = {
    "PUSH_HIGH": "Pushed 5+ times",
    "PUSH_MED": "Pushed 3-4 times",
    "OVERDUE": "Past close date",
    "CLOSE_SOON": "Closing <30d in early stage",
    "STALE": "No activity 60+ days",
    "NO_NEXT_STEP": "No next step logged",
    "LOW_FCST": "Low forecast coverage",
    "HIGH_VALUE_PUSH": "Large deal pushed",
}


def _humanize_reason(code: str) -> str:
    """Translate a risk-code key to a director-friendly phrase."""
    return REASON_CODE_LABELS.get(code.strip().upper(), code.strip())


def slide_owner_coaching(prs, slip_owners, risk_deals, territory):
    """Top 3 owners to coach this month, joined from Slip Risk + Deal Risk.

    Each owner is ranked by a composite coaching-priority score that rewards
    push count, exposed ARR, and any of their deals hitting the triage risk
    floor. For each top-N owner the slide lists reason-code hot-spots and a
    one-line coaching suggestion driven by the dominant signal. This turns
    the deck from descriptive into prescriptive for the director-level 1:1.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    if not slip_owners:
        _set_ph(slide, 144, f"Owner Coaching | {territory}")
        _set_ph(
            slide,
            145,
            "No owners in this territory crossed the push threshold this run.",
        )
        return

    # Risk-deal hot-spots per owner so we can quote reason codes per person.
    from collections import Counter, defaultdict

    owner_to_risks = defaultdict(list)
    for d in risk_deals or []:
        owner = str(d.get("owner") or "")
        if owner:
            owner_to_risks[owner].append(d)

    def _score(owner_row):
        risks = owner_to_risks.get(owner_row["owner"], [])
        triage_hits = sum(1 for r in risks if (r.get("score") or 0) >= 60)
        return owner_row["pushes"] + triage_hits * 10 + owner_row["arr"] / 100_000.0

    ranked = sorted(slip_owners, key=_score, reverse=True)[:3]
    total_pushes = sum(o["pushes"] for o in ranked)
    total_arr = sum(o["arr"] for o in ranked)
    _set_ph(
        slide,
        144,
        (
            f"{len(ranked)} owners carry {total_pushes} pushes across "
            f"{_fmt_eur(total_arr)}. Coach in this order."
        ),
    )
    _set_ph(
        slide,
        145,
        "Owners with the highest push exposure this period. Use for 1:1 conversation prep.",
    )

    tbl = [
        ["Owner", "Deals", "Open ARR", "Pushes", "Top Risk Signals", "Coaching Focus"]
    ]
    for o in ranked:
        risks = owner_to_risks.get(o["owner"], [])
        code_counter = Counter()
        for r in risks:
            for code in str(r.get("reasons") or "").split(","):
                c = code.strip()
                if c:
                    code_counter[c] += 1
        top_codes = ", ".join(
            f"{_humanize_reason(c)} (×{n})" for c, n in code_counter.most_common(3)
        )
        if not top_codes:
            top_codes = "-"

        # Derive the dominant coaching angle from signal mix.
        if any("OVERDUE" in c or "CLOSE_SOON" in c for c in code_counter):
            focus = "Close-date discipline"
        elif any("STALE" in c or "NO_NEXT_STEP" in c for c in code_counter):
            focus = "Activity cadence + next steps"
        elif any("PUSH_HIGH" in c or "HIGH_VALUE_PUSH" in c for c in code_counter):
            focus = "Commit integrity: push pattern"
        elif any("LOW_FCST" in c for c in code_counter):
            focus = "Forecast categorization accuracy"
        elif o["max_push"] >= 5:
            focus = "Commit integrity: repeated slips"
        else:
            focus = "Push review"

        tbl.append(
            [
                str(o["owner"]),
                str(o["deals"]),
                _fmt_eur(o["arr"]),
                str(o["pushes"]),
                top_codes,
                focus,
            ]
        )
    _add_table(
        slide,
        tbl,
        0.4,
        2.1,
        12.5,
        col_widths=[2.4, 0.8, 1.6, 1.0, 4.5, 2.2],
        data_font_size=Pt(10),
    )


def _inject_insight_bullets_into_summary(prs, bullets):
    """Add a compact text box with the insight bullets to the most-recently-
    added slide (which should be the Executive Summary). This replaces the
    standalone Executive Insights slide — one dense slide instead of two,
    matching Rebekka's pattern where the title carries the finding and the
    body carries supporting numbers.
    """
    from pptx.util import Inches

    slide = prs.slides[-1]
    # Place below the existing body content, left-aligned, small font
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.0), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:4]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"{i + 1}. {bullet}"
        para.font.size = Pt(8)
        para.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        para.font.name = "Calibri"
        para.space_after = Pt(2)


def slide_executive_insights(prs, bullets, director, territory):
    """Synthesized findings for this director's territory. 3-5 bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    _set_ph(slide, 144, f"Executive Insights | {territory}")
    _set_ph(
        slide,
        145,
        f"Synthesized findings for {director}. Each bullet ties to a downstream "
        "slide or to the analysis workbook.",
    )

    if not bullets:
        _set_ph_lines(
            slide,
            42,
            [
                "No standout findings this period. Pipeline composition, win rate, "
                "and slip risk all within normal range."
            ],
        )
        return

    numbered = [f"{i}.  {b}" for i, b in enumerate(bullets, 1)]
    _set_ph_lines(slide, 42, numbered)


def _compute_director_insights(
    director,
    territory,
    pipeline,
    q1_won,
    q1_lost,
    risk_deals,
    variance,
    top_deals,
    velocity_q1,
    director_total_open_arr=0.0,
):
    """Generate 3-5 synthesized findings for this director's deck.

    Produces bullets that an exec could read cold and understand the
    territory's pipeline health. Order: velocity, concentration, risk,
    Q1 outcome, variance driver.
    """
    bullets = []

    if velocity_q1 and velocity_q1.get("series") and len(velocity_q1["series"]) >= 2:
        series = velocity_q1["series"]
        dates = velocity_q1.get("dates") or []
        v0 = float(series[0] or 0)
        vN = float(series[-1] or 0)
        delta = vN - v0
        if abs(delta) > 100_000 and dates:
            direction = "shrank" if delta < 0 else "grew"
            bullets.append(
                f"Q1 pipeline {direction} from {_fmt_eur(v0)} on {dates[0]} to "
                f"{_fmt_eur(vN)} on {dates[-1]} (net {_fmt_eur(delta)}). "
                "See Q1 Forecast Variance slide."
            )

    if top_deals:
        top3 = top_deals[:3]
        top3_arr = sum(d["arr"] for d in top3)
        names = "; ".join(t["account"] for t in top3)
        # Prefer the director's full Land ARR (from Territory Scorecard) as
        # the denominator, so the ratio is meaningful even when the deck is
        # Q1-Q2 scoped. Fall back to the filtered pipeline total.
        denom = director_total_open_arr or sum(_unw(r) for r in pipeline)
        if denom and top3_arr / denom >= 0.3:
            bullets.append(
                f"Concentration risk: top 3 deals total {_fmt_eur(top3_arr)}, "
                f"{top3_arr / denom * 100:.0f}% of the territory's "
                f"{_fmt_eur(denom)} open Land book ({names}). "
                "Single-deal slippage moves the quarter."
            )
        elif top3_arr:
            bullets.append(
                f"Top 3 deals in the global concentration list for this "
                f"territory: {names}. Total {_fmt_eur(top3_arr)}."
            )

    if risk_deals:
        high_risk = [d for d in risk_deals if d["score"] >= 60]
        if high_risk:
            exposed = sum(d["arr"] for d in high_risk)
            bullets.append(
                f"{len(high_risk)} open deals score 60+ on risk "
                f"({_fmt_eur(exposed)} exposed). Triage list on Deal Risk slide."
            )

    q1w_arr = sum(_unw(r) for r in q1_won)
    q1l_arr = sum(_unw(r) for r in q1_lost)
    if q1w_arr or q1l_arr:
        total = q1w_arr + q1l_arr
        win_rate = (q1w_arr / total) if total else 0
        bullets.append(
            f"Q1 outcome: {len(q1_won)} wins ({_fmt_eur(q1w_arr)}) vs "
            f"{len(q1_lost)} losses ({_fmt_eur(q1l_arr)}); "
            f"win rate {win_rate * 100:.0f}% by ARR."
        )

    if variance:
        won = float(variance.get("won") or 0)
        lost = float(variance.get("lost") or 0)
        if lost > won * 2 and lost > 500_000:
            bullets.append(
                f"Variance driver: losses ({_fmt_eur(lost)}) outpaced wins "
                f"({_fmt_eur(won)}) >2x. Pipeline shrinkage is loss-driven, not "
                "win-driven. Review Closed Lost reasons with sales ops."
            )
        elif won > lost * 2 and won > 500_000:
            bullets.append(
                f"Variance driver: wins ({_fmt_eur(won)}) outpaced losses "
                f"({_fmt_eur(lost)}) >2x. Clean quarterly execution."
            )

    return bullets


def slide_commercial_approvals(prs, approvals, territory):
    """Commercial Approvals — matches Rebekka's slide 6 layout exactly.

    Slide A layout (Rebekka's slide 6):
      Title: "Commercial Approvals"
      Top row:
        LEFT  (0.92, 1.78): "YTD Actuals" label + 2×4 table
        RIGHT (6.85, 1.78): "FY Targets"  label + 2×4 table
      Bottom row:
        LEFT  (0.92, 3.95): "Pending Approval" + 3-col table
        RIGHT (6.85, 3.95): "Commercial Approval candidates" + 2-col table
      Footnote at y=7.10: "*The average deal size is measured at the time of..."

    Slide B (our addition): narrative summary — two columns with headline + body.
    """
    approved_2026 = [r for r in approvals if r.get("Status") == "Approved 2026"]
    approved_prior = [
        r for r in approvals if "prior" in str(r.get("Status", "")).lower()
    ]
    pending = [
        r
        for r in approvals
        if "Conditionally" in str(r.get("Status", ""))
        or "Pending" in str(r.get("Status", ""))
    ]
    missing = [r for r in approvals if "Missing" in str(r.get("Status", ""))]
    all_approved = approved_2026 + approved_prior
    candidates = pending + missing

    # ── Slide A: 4-table layout (Rebekka's slide 6) ──
    slide = prs.slides.add_slide(prs.slide_layouts[LY_2COL_GRAD])
    n_approved = len(approved_2026)
    n_pending = len(pending)
    n_missing = len(missing)
    if n_approved == 0 and n_pending == 0 and n_missing == 0:
        _set_ph(slide, 144, "Commercial Approvals: none this cycle")
    else:
        parts = []
        if n_approved:
            parts.append(f"{n_approved} approved 2026")
        if n_pending:
            parts.append(f"{n_pending} pending")
        if n_missing:
            parts.append(f"{n_missing} missing Stage 3+")
        _set_ph(
            slide,
            144,
            "Commercial Approvals: " + ", ".join(parts) + ".",
        )

    TABLE_W = 5.56  # Rebekka's measured width
    LEFT_X = 0.92
    RIGHT_X = 6.85
    LABEL_Y = 1.78
    TABLE_TOP_Y = 2.31  # Rebekka's position for top-row tables
    LABEL_BOTTOM_Y = 3.95
    TABLE_BOTTOM_Y = 4.44
    FOOTNOTE_Y = 7.10

    def _add_label(text, x, y, width=TABLE_W):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.36))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = DARK

    # ── TOP LEFT: YTD Actuals ──
    # YTD = deals APPROVED in 2026 specifically (Rebekka shows "APAC | 2 | 1.6 | 3.2")
    _add_label("YTD Actuals", LEFT_X, LABEL_Y)
    total_approved_arr = sum(_unw(r) for r in approved_2026)
    avg_deal = total_approved_arr / max(len(approved_2026), 1)
    # Rebekka's exact header text (preserved)
    ytd_rows = [
        [
            "Region",
            "# Approvals",
            "Avg Deal Size (mEUR ARR)*",
            "Actual Deal ARR Coverage (mEUR)",
        ],
        [
            territory.replace(" (LAND only)", ""),
            str(len(approved_2026)),
            _meur(avg_deal),
            _meur(total_approved_arr),
        ],
    ]
    _add_table(
        slide,
        ytd_rows,
        LEFT_X,
        TABLE_TOP_Y,
        TABLE_W,
        row_height=0.34,
        col_widths=[0.85, 1.15, 1.75, 1.81],
        data_font_size=Pt(11),
    )

    # ── TOP RIGHT: FY Targets ──
    # FY = all approved (including prior-year). Rebekka shows APAC | 6 | 2.3 | 13.8
    _add_label("FY Targets", RIGHT_X, LABEL_Y)
    total_all_arr = sum(_unw(r) for r in all_approved)
    avg_all = total_all_arr / max(len(all_approved), 1)
    target_rows = [
        [
            "Region",
            "# Approvals",
            "Avg Deal Size",
            "Actual Deal ARR Coverage (mEUR)",
        ],
        [
            territory.replace(" (LAND only)", ""),
            str(len(all_approved)),
            _meur(avg_all),
            _meur(total_all_arr),
        ],
    ]
    _add_table(
        slide,
        target_rows,
        RIGHT_X,
        TABLE_TOP_Y,
        TABLE_W,
        row_height=0.34,
        col_widths=[0.85, 1.15, 1.75, 1.81],
        data_font_size=Pt(11),
    )

    # ── BOTTOM LEFT: Conditionally Approved Deals ──
    _add_label("Pending Approval", LEFT_X, LABEL_BOTTOM_Y)
    if approved_2026:
        # "Approved subject to" column: short phrase (Rebekka: "Receiving RFP",
        # "Go / No go decision"). We pull from Next Step but keep it short.
        cond_rows = [["Opportunity Name", "Deal size (mEUR)", "Approved subject to"]]
        for r in approved_2026:
            ns = str(r.get("Next Step", "") or "").strip()
            # Keep it short — first clause only, max 40 chars
            short_ns = ns.split(".")[0].split(" - ")[0].split(" – ")[0][:40] or "-"
            cond_rows.append(
                [
                    str(r.get("Opportunity", "")),
                    _meur(_unw(r)),
                    short_ns,
                ]
            )
        _add_table(
            slide,
            cond_rows,
            LEFT_X,
            TABLE_BOTTOM_Y,
            TABLE_W,
            row_height=0.34,
            col_widths=[2.5, 1.2, 1.86],
            data_font_size=Pt(10),
        )

    # ── BOTTOM RIGHT: Commercial Approval candidates ──
    _add_label("Commercial Approval candidates", RIGHT_X, LABEL_BOTTOM_Y)
    if candidates:
        cand_rows = [["Opportunity Name", "ARR (mEUR)"]]
        for r in candidates:
            cand_rows.append(
                [
                    str(r.get("Opportunity", "")),
                    _meur(_unw(r)),
                ]
            )
        _add_table(
            slide,
            cand_rows,
            RIGHT_X,
            TABLE_BOTTOM_Y,
            TABLE_W,
            row_height=0.34,
            col_widths=[4.06, 1.50],
            data_font_size=Pt(10),
        )

    # ── Footnote (Rebekka's exact wording) ──
    footer = slide.shapes.add_textbox(
        Inches(2.31), Inches(FOOTNOTE_Y), Inches(10.10), Inches(0.20)
    )
    fp = footer.text_frame.paragraphs[0]
    fr = fp.add_run()
    fr.text = (
        "*The average deal size is measured at the time of the Commercial Approval. "
        "All ARR figures are unweighted."
    )
    fr.font.size = Pt(8)
    fr.font.italic = True
    fr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def slide_missing_approval_detail(prs, approvals):
    """Land Stage 3, Missing Commercial Approval detail table. Matches Rebekka slide 8."""
    missing = [r for r in approvals if "Missing" in str(r.get("Status", ""))]
    pending = [
        r
        for r in approvals
        if "Conditionally" in str(r.get("Status", ""))
        or "Pending" in str(r.get("Status", ""))
    ]
    action_deals = pending + missing
    if not action_deals:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    total_arr = sum(_unw(r) for r in action_deals)

    _set_ph(slide, 144, "Land Stage 3, Missing Commercial Approval")
    _set_ph(
        slide, 145, f"{len(action_deals)} candidates totaling {_fmt_eur(total_arr)}"
    )
    _set_ph(
        slide,
        42,
        "These deals are in Engagement stage without Go/No-Go approval. Action: escalate to approval committee.",
    )

    rows = [
        [
            "Account",
            "Opportunity",
            "Owner",
            "Close Date",
            "Next Step",
            "ARR (mEUR)",
        ]
    ]
    for r in action_deals:
        rows.append(
            [
                str(r.get("Account", "")),
                str(r.get("Opportunity", "")),
                str(r.get("Owner", "")),
                str(r.get("Close Date", "")),
                str(r.get("Next Step", "") or ""),
                _meur(_unw(r)),
            ]
        )
    # Highlight Close Date column (index 3): red if overdue, amber if within 30d.
    today_s = datetime.now().strftime("%Y-%m-%d")
    soon_s = (datetime.now() + timedelta(days=30)).strftime("%Y-%m-%d")

    def _overdue(v):
        s = str(v)[:10]
        return bool(s) and s < today_s

    def _soon(v):
        s = str(v)[:10]
        return bool(s) and today_s <= s <= soon_s

    _add_table(
        slide,
        rows,
        0.9,
        2.2,
        11.5,
        col_widths=[1.9, 1.9, 1.9, 1.9, 1.9, 1.9],
        highlight_rules={
            3: [
                (_overdue, "F8D7DA"),
                (_soon, "FFF3CD"),
            ],
        },
    )


def slide_renewals(prs, renewals):
    """Renewals, annual view sorted by close date with probability and a
    commentary column for the director.

    Rebekka's feedback:
      - Sort on Close Date
      - Annual view (not just current quarter) until Axioma is included
      - Add Probability column
      - Add Comments column for director input
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])

    def _acv(r):
        """Read ACV from the new column name, fall back to the legacy one."""
        v = r.get("ACV Unweighted (EUR)")
        if v is None:
            v = r.get("ACV (EUR)")
        return float(v or 0)

    # Sort annual view by close date ascending
    annual = sorted(renewals, key=lambda r: str(r.get("Close Date", "")))
    total_acv = sum(_acv(r) for r in annual)

    # Count by quarter for the subtitle
    def _qtr(r):
        s = str(r.get("Close Date", ""))[:7]
        if not s.startswith(str(FQ["fy"])):
            return None
        try:
            m = int(s[5:7])
        except ValueError:
            return None
        return f"Q{(m - 1) // 3 + 1}"

    by_q = {"Q1": 0, "Q2": 0, "Q3": 0, "Q4": 0}
    for r in annual:
        q = _qtr(r)
        if q:
            by_q[q] += 1

    if annual:
        q_breakdown = ", ".join(
            f"{by_q[q]} in {q}" for q in ("Q1", "Q2", "Q3", "Q4") if by_q[q]
        )
        _set_ph(
            slide,
            144,
            (
                f"{len(annual)} renewals due FY26, {_fmt_eur(total_acv)} ACV "
                f"({q_breakdown})."
            ),
        )
    else:
        _set_ph(slide, 144, "Renewals FY26: none this cycle")
    if annual:
        _set_ph(
            slide,
            145,
            f"Total FY26 ACV: {_fmt_eur(total_acv)} across "
            f"{len(annual)} {'renewal' if len(annual) == 1 else 'renewals'}. "
            f"{q_breakdown}.",
        )
    else:
        _set_ph(slide, 145, "No FY26 renewals in scope for this territory.")
        return

    # Single annual table, sorted by close date, with commentary column
    rows = [
        [
            "Close Date",
            "Account",
            "Opportunity",
            "Owner",
            "Stage",
            "ACV (EUR)",
            "Probability",
            "Commentary",
        ]
    ]
    for r in annual:
        rows.append(
            [
                str(r.get("Close Date", "")),
                str(r.get("Account", "")),
                str(r.get("Opportunity", "")),
                str(r.get("Owner", "")),
                str(r.get("Stage", "")),
                f"EUR {_acv(r):,.0f}",
                f"{int(r.get('Probability %') or 0)}%",
                "",  # blank, director fills in during review
            ]
        )
    _add_table(
        slide,
        rows,
        0.5,
        2.0,
        12.3,
        row_height=0.30,
        col_widths=[1.3, 1.9, 2.3, 1.5, 1.3, 1.4, 1.1, 1.5],
        data_font_size=Pt(10),
    )


def slide_churn(prs, territory=None):
    """Churn Risk slide.

    When a per-territory churn screenshot exists under
    assets/rebekka-screenshots/<territory>-churn.png (kebab-case, lowercase),
    embed it as the slide visual. Otherwise fall back to the status/action
    table explaining that Finance hasn't wired up the feed yet.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])

    territory_slug = (
        (territory or "").lower().replace(" & ", "-").replace(" ", "-")
        if territory
        else ""
    )
    screenshot_path = (
        Path(__file__).resolve().parents[1]
        / "assets"
        / "rebekka-screenshots"
        / f"{territory_slug}-churn.png"
    )

    if screenshot_path.exists():
        _set_ph(slide, 144, "Churn Risk")
        _set_ph(
            slide,
            145,
            f"{territory}, sourced from Finance. Imported from the prior "
            "monthly pack; regenerating this view live is pending the Finance "
            "data pipeline.",
        )
        # Fit within content area: slide 13.33x7.5in, content below header
        slide.shapes.add_picture(
            str(screenshot_path),
            Inches(0.7),
            Inches(1.9),
            width=Inches(12.0),
            height=Inches(5.2),
        )
        return

    _set_ph(slide, 144, "Churn Risk")
    _set_ph(
        slide, 145, "Churn data sourced from Finance. Not yet available this cycle."
    )
    _set_ph(
        slide,
        42,
        "Churn monitoring is owned by Finance and will appear here once the "
        "monthly feed is agreed. See the Business At Risk tab in the analysis "
        "workbook for an interim proxy.",
    )

    rows = [
        ["Action", "Status", "Owner", "Notes"],
        [
            "Obtain churn data from Finance",
            "Pending",
            "Sales Ops → Alex P",
            "We pull from Finance, not the other way around",
        ],
        [
            "Identify at-risk renewal accounts",
            "Pending",
            "Sales Ops + Finance",
            "Top accounts by attrition signal",
        ],
        [
            "Quantify churn exposure",
            "Pending",
            "Finance (Alex P)",
            "ACV at risk, renewal attrition rate",
        ],
        [
            "Feed churn view to leadership",
            "Blocked",
            "Sales Ops",
            "Awaiting Finance feed",
        ],
    ]
    _add_table(
        slide, rows, 0.9, 2.5, 11.5, row_height=0.41, col_widths=[3.2, 2.0, 2.3, 4.0]
    )


PI_LINKS = {
    "APAC": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Ic7kTMAR",
    "Central Europe": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Kr3YvMAJ",
    "UK & Ireland": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Kr3yjMAB",
    "Southern Europe": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Kr3sHMAR",
    "NL & Nordics": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Kr4DFMAZ",
    "Middle East & Africa": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BQA00000GXOf32AH",
    "Canada": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Kr4ErMAJ",
    "NA Asset Management": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Kr4JhMAJ",
    "Pension & Insurance": "https://simcorp.lightning.force.com/lightning/o/Opportunity/list?filterName=00BTb00000Kr4OXMAZ",
}


def slide_pi_link(prs, territory):
    """Pipeline Inspection link slide — directs users to the live PI view."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    link = PI_LINKS.get(territory, "")

    _set_ph(slide, 144, f"Pipeline Inspection | {territory}")
    _set_ph(slide, 145, "Live view in Salesforce")

    lines = [
        f"Open the Pipeline Inspection view for {territory} in Salesforce:",
        "",
        link
        if link
        else "Pipeline Inspection list view not configured for this territory.",
        "",
        "The Pipeline Inspection view shows:",
        "  • Forecast category breakdown (Commit / Best Case / Pipeline)",
        "  • Deal-level forecast ARR and close dates",
        "  • Push count and deal scoring",
        "  • Priority flag for key deals",
        "",
        "Use this view for live drill-down during the pipeline review meeting.",
    ]
    _set_ph_lines(slide, 22, lines)


def slide_pushed_deals(prs, pi_data):
    if not pi_data:
        return
    pushed = [r for r in pi_data if int(r.get("Push Count") or 0) > 0]
    pushed.sort(key=lambda r: int(r.get("Push Count") or 0), reverse=True)

    if not pushed:
        return

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    total = len(pushed)
    avg = sum(int(r.get("Push Count") or 0) for r in pushed) / max(total, 1)
    exposed = sum(_wtd(r) for r in pushed)
    critical = len([r for r in pushed if int(r.get("Push Count") or 0) >= 5])

    _set_ph(
        slide, 144, f"Pushed Deals: {total} deals | {_fmt_eur(exposed)} exposed ARR"
    )
    _set_ph(slide, 42, "Total Pushed")
    _set_ph(slide, 56, "Avg Pushes")
    _set_ph(slide, 58, "Exposed ARR")
    _set_ph(slide, 60, "Critical (5+)")
    _add_gradient_metric(slide, 0, str(total))
    _add_gradient_metric(slide, 1, f"{avg:.1f}")
    _add_gradient_metric(slide, 2, _fmt_eur(exposed))
    _add_gradient_metric(slide, 3, str(critical))
    tiers = _push_tiers(pi_data)
    owner_callout = _top_pushed_owner(pi_data)

    _set_ph(
        slide,
        22,
        f"{total} open deals pushed. {owner_callout + ', pattern review recommended.' if owner_callout else ''}",
    )
    _set_ph(
        slide,
        55,
        f"Average pushes per deal. {tiers['critical']['count']} deals pushed 5+ times "
        f"({_fmt_eur(tiers['critical']['arr'])} ARR) are the highest risk.",
    )
    _set_ph(slide, 57, f"Total ARR exposed across all {total} pushed deals.")
    _set_ph(
        slide,
        59,
        f"Critical: {tiers['critical']['count']} at 5+ pushes. "
        f"Watch: {tiers['watch']['count']} at 3-4 pushes ({_fmt_eur(tiers['watch']['arr'])}). "
        f"Early: {tiers['early']['count']} at 1-2 pushes ({_fmt_eur(tiers['early']['arr'])}).",
    )

    # Detail table
    slide2 = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    top_pushed = sorted(
        pi_data or [],
        key=lambda r: int(str(r.get("Push Count") or 0) or 0),
        reverse=True,
    )[:1]
    if top_pushed:
        top_name = str(top_pushed[0].get("Opportunity") or "")
        top_push = int(str(top_pushed[0].get("Push Count") or 0) or 0)
        _set_ph(
            slide2,
            144,
            f"Highest: {top_name} ({top_push}×). Deals pushed 5+ times warrant direct conversation with owner.",
        )
    else:
        _set_ph(slide2, 144, "Pipeline Inspection | Key Pushed Deals")
    _set_ph(slide2, 145, "Ranked by push count")

    rows = [
        [
            "Opportunity",
            "Account",
            "Stage",
            "Close",
            "Pushes",
            "ARR (mEUR)",
            "Owner",
        ]
    ]
    for r in pushed[:12]:
        rows.append(
            [
                str(r.get("Opportunity", "")),
                "",  # PI data doesn't have account
                str(r.get("Stage", "")),
                str(r.get("Close Date", "")),
                str(r.get("Push Count", "")),
                _meur(_wtd(r)),
                str(r.get("Owner", "")),
            ]
        )
    _add_table(
        slide2,
        rows,
        0.9,
        2.2,
        11.5,
        row_height=0.20,
        col_widths=[1.6, 1.6, 1.6, 1.6, 1.6, 1.6, 1.6],
    )


def slide_forecast_breakdown(prs, pi_data):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    if not pi_data:
        _set_ph(slide, 144, "Forecast Category Breakdown")
        _set_ph(slide, 22, "No pipeline inspection data available.")
        return

    fc = defaultdict(lambda: {"count": 0, "arr": 0.0})
    for r in pi_data:
        cat = str(r.get("Forecast Category", "Unknown"))
        fc[cat]["count"] += 1
        fc[cat]["arr"] += _wtd(r)

    total_arr = sum(b["arr"] for b in fc.values())
    commit_arr = fc.get("Commit", {}).get("arr", 0.0)
    commit_share = commit_arr / total_arr if total_arr else 0
    _set_ph(
        slide,
        144,
        (
            f"Forecast Breakdown: {_fmt_eur(total_arr)} across {len(pi_data)} "
            f"open deals. Commit = {_fmt_eur(commit_arr)} "
            f"({commit_share * 100:.0f}%)."
        ),
    )
    _set_ph(
        slide,
        145,
        "Commit is the floor you can bank; Best Case + Pipeline is the upside.",
    )

    rows = [["Category", "Deals", "ARR (mEUR)"]]
    for cat in ("Pipeline", "Commit", "Best Case"):
        b = fc.get(cat, {})
        if b.get("count"):
            rows.append([cat, str(b["count"]), _meur(b["arr"])])
    rows.append(["Total", str(len(pi_data)), _meur(total_arr)])

    _add_table(slide, rows, 0.9, 2.2, 5.5, row_height=0.21, col_widths=[1.8, 1.8, 1.8])


def slide_coverage_targets(prs, pipeline):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    _set_ph(slide, 144, "Pipeline Coverage & Targets")
    _set_ph(slide, 145, "Coverage ratio against assigned targets")

    # Compute what we can from the pipeline data
    total_arr = sum(_unw(r) for r in pipeline)
    weighted = sum(_unw(r) * float(r.get("Probability %") or 0) / 100 for r in pipeline)
    stale = [r for r in pipeline if int(r.get("Push Count") or 0) >= 3]
    stale_arr = sum(_unw(r) for r in stale)

    lines = [
        f"Total Open Pipeline: {_fmt_eur(total_arr)}",
        f"Weighted Pipeline (probability-adjusted): {_fmt_eur(weighted)}",
        f"Stale deals (3+ pushes): {len(stale)} deals, {_fmt_eur(stale_arr)}",
        "",
        "Quota integration pending. Coverage ratio unavailable this cycle.",
        "Coverage ratio will be available once targets are provided by Finance / Sales Ops.",
    ]
    _set_ph_lines(slide, 22, lines)


def slide_definitions(prs, snapshot_date, pipeline=None):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    _set_ph(slide, 144, "Definitions and Data Sources")

    lines = [
        "Metric conventions",
        "  ARR Unweighted: full deal value if the opportunity closes.",
        "    Aligns with the Opportunity ARR view on the Salesforce forecast page.",
        "  ARR Weighted: probability-weighted forecast value, stage and forecast",
        "    category adjusted. Aligns with the Commit, Best Case and Pipeline",
        "    columns on the forecast page.",
        "  Typical ratio: weighted is 15-35 percent of unweighted, depending on",
        "    stage mix.",
        "  Renewals: ACV in EUR, unweighted.",
        "  Omitted forecast category deals are excluded from headline pipeline.",
        "",
        "Scope",
        "  Type: Land only.",
        f"  Close date: {FQ['prior']['label']}-{FQ['forward']['label']} FY{str(FQ['fy'])[-2:]}.",
        "  Accounts with simcorp, test, or delete in the name are excluded.",
        "  Owners Sabiniewicz and Profit are excluded (test sandbox).",
        "",
        "Pipeline coverage",
    ]
    if pipeline:
        total_arr = sum(_unw(r) for r in pipeline)
        weighted = sum(
            _unw(r) * float(r.get("Probability %") or 0) / 100 for r in pipeline
        )
        stale = [r for r in pipeline if int(r.get("Push Count") or 0) >= 3]
        lines.extend(
            [
                f"  Total open pipeline: {_fmt_eur(total_arr)}",
                f"  Weighted (probability-adjusted): {_fmt_eur(weighted)}",
                f"  Stale deals (3+ pushes): {len(stale)} deals ({_fmt_eur(sum(_unw(r) for r in stale))})",
            ]
        )
    lines.extend(
        [
            "  Quota and target data not yet integrated, coverage ratio is qualified",
            "",
            "Known limitations",
            "  Finance churn overlay pending",
            f"  Data extracted: {snapshot_date}",
        ]
    )
    _set_ph_lines(slide, 22, lines)


def slide_end(prs):
    prs.slides.add_slide(prs.slide_layouts[LY_END_SLIDE])


# ── Consolidated slide functions ──


def slide_pushed_deals_with_link(prs, pi_data, territory, q1_movement=None):
    """Pushed Deals summary (4-card) + detail table with PI link in footer.

    Detail table now carries a `Last Pushed` column built from the Q1
    Movement history so the director can see when each push happened,
    not just that it happened.
    """
    if not pi_data:
        return
    pushed = [r for r in pi_data if int(r.get("Push Count") or 0) > 0]
    pushed.sort(key=lambda r: int(r.get("Push Count") or 0), reverse=True)
    if not pushed:
        return

    # Summary slide (4-card)
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    total = len(pushed)
    avg = sum(int(r.get("Push Count") or 0) for r in pushed) / max(total, 1)
    exposed = sum(_wtd(r) for r in pushed)
    tiers = _push_tiers(pi_data)
    owner_callout = _top_pushed_owner(pi_data)

    _set_ph(
        slide, 144, f"Pushed Deals: {total} deals | {_fmt_eur(exposed)} exposed ARR"
    )
    _set_ph(slide, 42, "Pushed Deals Summary")
    _set_ph(slide, 56, "Pushed Deals Summary")
    _set_ph(slide, 58, "Pushed Deals Summary")
    _set_ph(slide, 60, "Pushed Deals Summary")
    _add_gradient_metric(slide, 0, str(total))
    _add_gradient_metric(slide, 1, f"{avg:.1f}x avg")
    _add_gradient_metric(slide, 2, _fmt_eur(exposed))
    _add_gradient_metric(slide, 3, str(tiers["critical"]["count"]))
    _set_ph(
        slide,
        22,
        f"{total} open deals pushed. {owner_callout + ', pattern review recommended.' if owner_callout else ''}",
    )
    _set_ph(
        slide,
        55,
        f"Average pushes per deal. {tiers['critical']['count']} deals pushed 5+ times "
        f"({_fmt_eur(tiers['critical']['arr'])} ARR) are the highest risk.",
    )
    _set_ph(slide, 57, f"Total ARR exposed across all {total} pushed deals.")
    _set_ph(
        slide,
        59,
        f"Critical: {tiers['critical']['count']} at 5+ pushes. "
        f"Watch: {tiers['watch']['count']} at 3-4 ({_fmt_eur(tiers['watch']['arr'])}). "
        f"Early: {tiers['early']['count']} at 1-2 ({_fmt_eur(tiers['early']['arr'])}).",
    )

    # PI link footer on summary slide
    link = PI_LINKS.get(territory, "")
    if link:
        txBox = slide.shapes.add_textbox(
            Inches(0.9), Inches(6.8), Inches(11.0), Inches(0.3)
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"Open Pipeline Inspection in Salesforce: {link}"
        run.font.size = Pt(8)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x08, 0x3E, 0xA7)


def slide_q1_movement(prs, q1_movement):
    """Q1 Movement — slipped deals and post-Q1 pushes."""
    if not q1_movement:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])

    slipped = [r for r in q1_movement if r.get("Movement") == "Q1 Slipped"]
    pushed = [r for r in q1_movement if r.get("Movement") == "Post-Q1 Push"]
    slipped_arr = sum(_unw(r) for r in slipped)
    sum(_unw(r) for r in pushed)

    _set_ph(
        slide,
        144,
        f"{len(slipped)} deals slipped out of Q1 ({_fmt_eur(slipped_arr)}). {len(pushed)} pushed since.",
    )
    _set_ph(
        slide,
        145,
        "Deals that had Q1 close dates and moved, plus post-Q1 close date changes",
    )

    # Top slipped + pushed by ARR
    top = sorted(q1_movement, key=lambda r: _unw(r), reverse=True)[:12]
    if top:
        rows = [
            [
                "Account",
                "Opportunity",
                "Movement",
                "Old Close",
                "New Close",
                "Changed",
                "ARR (mEUR)",
            ]
        ]
        for r in top:
            rows.append(
                [
                    str(r.get("Account", "")),
                    str(r.get("Opportunity", "")),
                    str(r.get("Movement", "")),
                    str(r.get("Old Close", "")),
                    str(r.get("New Close", "")),
                    str(r.get("Changed On", "")),
                    _meur(_unw(r)),
                ]
            )
        _add_table(
            slide, rows, 0.9, 2.2, 11.5, col_widths=[1.6, 1.6, 1.4, 1.4, 1.4, 1.4, 1.2]
        )


def slide_forecast_combined(prs, won_lost, pipeline, pi_data, territory):
    """Forecast Accuracy (4-card) + Category Breakdown table on one slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])

    won = [r for r in won_lost if "Won" in str(r.get("Stage", ""))]
    lost = [r for r in won_lost if "Won" not in str(r.get("Stage", ""))]
    won_arr = sum(_unw(r) for r in won)
    lost_arr = sum(_unw(r) for r in lost)
    total_decisions = len(won) + len(lost)
    win_rate = f"{len(won) / total_decisions * 100:.0f}%" if total_decisions else "n/a"
    commit_arr = sum(
        _unw(r) for r in pipeline if r.get("Forecast Category") == "Commit"
    )

    _set_ph(slide, 144, f"Forecast Accuracy | {territory}")
    _set_ph(slide, 42, f"Forecast Accuracy - {territory}")
    _set_ph(slide, 56, f"Forecast Accuracy - {territory}")
    _set_ph(slide, 58, f"Forecast Accuracy - {territory}")
    _set_ph(slide, 60, f"Forecast Accuracy - {territory}")

    _add_gradient_metric(slide, 0, _fmt_eur(won_arr))
    _add_gradient_metric(slide, 1, _fmt_eur(lost_arr))
    _add_gradient_metric(slide, 2, win_rate)
    _add_gradient_metric(slide, 3, _fmt_eur(commit_arr))

    _set_ph(
        slide,
        22,
        f"{len(won)} {'deal' if len(won) == 1 else 'deals'} closed-won ({_fmt_eur(won_arr)}).",
    )
    _set_ph(slide, 55, f"{len(lost)} deals closed-lost ({_fmt_eur(lost_arr)}).")
    _set_ph(slide, 57, f"Win rate: {win_rate} ({len(won)}W / {len(lost)}L).")
    _set_ph(slide, 59, f"Commit forecast: {_fmt_eur(commit_arr)}.")

    # Forecast category breakdown table below (from PI data)
    if pi_data:
        fc = defaultdict(lambda: {"count": 0, "arr": 0.0})
        for r in pi_data:
            cat = str(r.get("Forecast Category", "Unknown"))
            fc[cat]["count"] += 1
            fc[cat]["arr"] += _wtd(r)

        slide2 = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
        total_pi_arr = sum(b["arr"] for b in fc.values())
        commit_arr = fc.get("Commit", {}).get("arr", 0.0)
        commit_share = commit_arr / total_pi_arr if total_pi_arr else 0
        _set_ph(
            slide2,
            144,
            (
                f"Forecast Breakdown: {_fmt_eur(total_pi_arr)} across {len(pi_data)} "
                f"open deals. Commit = {_fmt_eur(commit_arr)} "
                f"({commit_share * 100:.0f}%)."
            ),
        )
        _set_ph(
            slide2,
            145,
            "Commit is the floor you can bank; Best Case + Pipeline is the upside.",
        )

        rows = [["Category", "Deals", "ARR (mEUR)"]]
        for cat in ("Pipeline", "Commit", "Best Case"):
            b = fc.get(cat, {})
            if b.get("count"):
                rows.append([cat, str(b["count"]), _meur(b["arr"])])
        rows.append(["Total", str(len(pi_data)), _meur(total_pi_arr)])
        _add_table(
            slide2, rows, 0.9, 2.2, 5.5, row_height=0.21, col_widths=[1.8, 1.8, 1.8]
        )


# ── Main build ──


def build_deck(
    workbook_path: Path,
    template_path: Path,
    output_path: Path,
    land_only: bool = False,
    analytics_path: Path | None = None,
):
    wb = load_workbook(str(workbook_path), data_only=True)

    director, territory = read_director_info(wb)
    kpis = read_kpis(wb)
    analytics = read_director_analytics(analytics_path, director)

    # Read all sheets
    pipeline = read_sheet(wb, "Pipeline Open FY26")
    won_lost = read_sheet(wb, "Won Lost FY26")
    approvals = read_sheet(wb, "Commercial Approval")
    renewals = read_sheet(wb, "Renewals FY26")
    pi_data = read_sheet(wb, "Pipeline Inspection")
    q1_movement = read_sheet(wb, "Q1 Movement")

    # Monthly review scope: Land type, Q1-Q2 FY26 close dates.
    # Renewals keep their own sheet (Type=Renewal deals due in Q2).
    # PI, Pushed Deals and Q1 Movement are scoped to Land opportunities
    # that appear in our Land pipeline or Land won/lost sets, so every
    # slide tells the same story.
    if land_only:

        def _is_land(r):
            return str(r.get("Type", "")).strip().lower() == "land"

        def _has_type(rows):
            return bool(rows) and ("Type" in rows[0])

        def _in_scope(r):
            cd = str(r.get("Close Date", "") or "")[:10]
            return cd >= FQ["scope_start"] and cd <= FQ["scope_end"]

        def _not_omitted(r):
            fc = str(r.get("Forecast Category", "")).strip()
            return fc not in ("Omitted", "")

        if _has_type(pipeline):
            pipeline = [r for r in pipeline if _is_land(r)]
        if _has_type(won_lost):
            won_lost = [r for r in won_lost if _is_land(r)]

        pipeline = [r for r in pipeline if _not_omitted(r)]

        # Establish the Land universe BEFORE Q1-Q2 scoping so we can still
        # match PI / Q1 movement records that sit in Q3-Q4 by opportunity name.
        land_universe = {
            r.get("Opportunity") for r in pipeline + won_lost if r.get("Opportunity")
        }

        pipeline = [r for r in pipeline if _in_scope(r)]
        won_lost = [r for r in won_lost if _in_scope(r)]
        approvals = [r for r in approvals if _in_scope(r)]

        # Pipeline Inspection: keep deals that belong to the Land universe,
        # regardless of close date, so the coaching view stays complete.
        pi_data = [r for r in pi_data if r.get("Opportunity") in land_universe]

        # Q1 Movement events are by definition deals that once sat in Q1,
        # so quarterly filtering is not meaningful; just filter to Land.
        q1_movement = [r for r in q1_movement if r.get("Opportunity") in land_universe]

        # Renewals: keep Q2 renewals intact; they populate the Renewals slide.

        # Scope tag is intentionally NOT appended to `territory`. Scope
        # belongs on the cover slide and in the Definitions slide, not in
        # every slide title. Keep titles clean and executive.
        print(
            f"  LAND + {FQ['prior']['label']}-{FQ['forward']['label']} FILTER: pipeline={len(pipeline)}  "
            f"won/lost={len(won_lost)}  approvals={len(approvals)}  "
            f"pi={len(pi_data)}  q1mov={len(q1_movement)}  "
            f"renewals={len(renewals)}"
        )

    # Extract snapshot date from Summary A2
    ws = wb["Summary"]
    str(ws["A2"].value or "")
    snapshot_date = datetime.now().strftime("%Y-%m-%d")

    # Build Q1 Promised vs Delivered summary from won/lost data
    won = [r for r in won_lost if "Won" in str(r.get("Stage", ""))]
    lost = [r for r in won_lost if "Won" not in str(r.get("Stage", ""))]
    q1_won = [
        r for r in won if str(r.get("Close Date", ""))[:7] <= FQ["prior"]["month_end"]
    ]
    q1_lost = [
        r for r in lost if str(r.get("Close Date", ""))[:7] <= FQ["prior"]["month_end"]
    ]
    q1_won_arr = sum(_unw(r) for r in q1_won)
    q1_lost_arr = sum(_unw(r) for r in q1_lost)

    # Current pipeline is what's still open — Q1 promised is pipeline + won + lost
    current_arr = sum(_unw(r) for r in pipeline)
    q1_promised_arr = current_arr + q1_won_arr + q1_lost_arr  # rough approximation

    q1_summary = {
        "promised_arr": q1_promised_arr,
        "committed_arr": sum(
            _unw(r) for r in pipeline if str(r.get("Stage", "")) >= "4"
        ),
        "committed_count": len([r for r in pipeline if str(r.get("Stage", "")) >= "4"]),
        "total_deals": len(pipeline) + len(q1_won) + len(q1_lost),
        "won_arr": q1_won_arr,
        "won_count": len(q1_won),
        "lost_arr": q1_lost_arr,
        "lost_count": len(q1_lost),
        "net_arr_change": q1_won_arr - q1_lost_arr,
        "advanced_count": 0,  # would need historical data
        "regressed_count": 0,
        "top_won": q1_won[0].get("Opportunity", "n/a") if q1_won else "n/a",
        "top_lost": q1_lost[0].get("Opportunity", "n/a") if q1_lost else "n/a",
        "top_won_deals": [
            {
                "name": r.get("Opportunity", ""),
                "owner": r.get("Owner", ""),
                "arr": _unw(r),
            }
            for r in sorted(q1_won, key=lambda x: _unw(x), reverse=True)[:5]
        ],
        "top_lost_deals": [
            {
                "name": r.get("Opportunity", ""),
                "owner": r.get("Owner", ""),
                "arr": _unw(r),
            }
            for r in sorted(q1_lost, key=lambda x: _unw(x), reverse=True)[:5]
        ],
    }

    print(f"Director: {director} ({territory})")
    print(
        f"Pipeline: {len(pipeline)} | Won/Lost: {len(won_lost)} | Approvals: {len(approvals)} | Renewals: {len(renewals)} | PI: {len(pi_data)}"
    )
    print(
        f"Q1: {len(q1_won)} won (EUR {q1_won_arr:,.0f}) / {len(q1_lost)} lost (EUR {q1_lost_arr:,.0f})"
    )

    prs = Presentation(str(template_path))

    # Remove template sample slides
    from lxml import etree

    pres_elem = prs.part._element
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    sldIdLst = pres_elem.find("p:sldIdLst", nsmap)
    if sldIdLst is not None:
        for sldId in list(sldIdLst):
            rId = sldId.get(etree.QName(r_ns, "id"))
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except KeyError:
                    pass
        for child in list(sldIdLst):
            sldIdLst.remove(child)

    # Build slides — consolidated (12 slides)
    n = 0

    # ══════════════════════════════════════════════════════════════════════
    # DECK v3 — trimmed to ~17 slides following the Rebekka audit.
    # Merges: Insights+Summary, Q1Promised+Variance, Pushed+PI,
    #         3 approval slides → 1. Cuts: duplicate Pipeline Overview,
    #         standalone Definitions. Adds: Competitive W/L, Stage Conversion.
    # ══════════════════════════════════════════════════════════════════════

    # 1. Cover
    n += 1
    slide_cover(prs, director, territory, snapshot_date)
    print(f"  [OK] {n}. Cover")

    # 2. Executive Summary (MERGED: Insights folded into the summary slide
    #    as a subtitle block — one dense slide, Rebekka-style)
    insight_bullets = _compute_director_insights(
        director,
        territory,
        pipeline,
        q1_won,
        q1_lost,
        analytics.get("risk_deals", []),
        analytics.get("variance"),
        analytics.get("top_deals", []),
        analytics.get("velocity_q1"),
        director_total_open_arr=analytics.get("director_total_open_arr", 0.0),
    )
    n += 1
    slide_executive_summary(prs, pipeline, won_lost, renewals, kpis, territory, pi_data)
    # Inject insight bullets as a text box on the summary slide if available
    if insight_bullets:
        _inject_insight_bullets_into_summary(prs, insight_bullets)
    print(f"  [OK] {n}. Executive Summary (with insights)")

    # 3. Month over Month (skipped if no prior snapshot)
    try:
        prev_count = len(prs.slides)
        slide_month_over_month(
            prs, pipeline, won_lost, approvals, director, territory, snapshot_date
        )
        if len(prs.slides) > prev_count:
            n += 1
            print(f"  [OK] {n}. Month over Month")
    except Exception as exc:
        print(f"  [SKIP] Month over Month: {exc}")

    # 4. Q1 Retrospective (MERGED: Q1 Promised + Forecast Variance on one slide)
    n += 1
    slide_q1_promised_vs_delivered(prs, q1_summary, territory)
    print(f"  [OK] {n}. Q1 Promised vs Delivered")
    # Forecast Variance: keep as standalone — the bucket decomposition is
    # too data-rich to merge without losing the narrative. Revisit if
    # directors say they skip it.
    if analytics.get("variance"):
        n += 1
        slide_forecast_variance(prs, analytics["variance"], territory)
        print(f"  [OK] {n}. Q1 Forecast Variance")

    # 5. Why We Lost
    dashboard_wb = (
        workbook_path.resolve().parents[2]
        / "sharepoint"
        / "Dashboard and Q1 Analysis.xlsx"
    )
    n += 1
    slide_win_loss_diagnostic(prs, won_lost, territory, dashboard_path=dashboard_wb)
    print(f"  [OK] {n}. Why We Lost")

    # 6-7. Quarter Outlook + Forward Look (Q2 first, Q3 if Q2 is empty)
    _q2_deals = [
        r
        for r in pipeline
        if FQ["current"]["month_start"]
        <= str(r.get("Close Date", ""))[:7]
        <= FQ["current"]["month_end"]
    ]
    _q3_deals = [
        r
        for r in pipeline
        if FQ["forward"]["month_start"]
        <= str(r.get("Close Date", ""))[:7]
        <= FQ["forward"]["month_end"]
    ]

    _q2_arr = sum(_unw(r) for r in _q2_deals)
    _q3_arr = sum(_unw(r) for r in _q3_deals)

    _quarters_to_show = []
    if _q2_arr > 0:
        _quarters_to_show.append(
            (
                FQ["current"]["label"],
                FQ["current"]["month_start"],
                FQ["current"]["month_end"],
                FQ["current"]["range_label"],
                FQ["current"]["start"],
                FQ["current"]["end"],
            )
        )
    elif _q3_arr > 0:
        _quarters_to_show.append(
            (
                FQ["forward"]["label"],
                FQ["forward"]["month_start"],
                FQ["forward"]["month_end"],
                FQ["forward"]["range_label"],
                FQ["forward"]["start"],
                FQ["forward"]["end"],
            )
        )
    if not _quarters_to_show:
        _quarters_to_show.append(
            (
                FQ["current"]["label"],
                FQ["current"]["month_start"],
                FQ["current"]["month_end"],
                FQ["current"]["range_label"],
                FQ["current"]["start"],
                FQ["current"]["end"],
            )
        )

    for _qlabel, _mstart, _mend, _mrange, _dstart, _dend in _quarters_to_show:
        n += 1
        slide_quarter_outlook(
            prs,
            pipeline,
            won_lost,
            territory,
            q_label=_qlabel,
            month_start=_mstart,
            month_end=_mend,
            month_range_label=_mrange,
        )
        print(f"  [OK] {n}. {_qlabel} Outlook")

    # Forward Look (enriched per-deal readiness from live SF)
    try:
        import json as _json

        _terr_cfg_path = (
            workbook_path.resolve().parents[3]
            / "config"
            / "sd_monthly_territories.json"
        )
        _terr_soql = ""
        if _terr_cfg_path.exists():
            _terr_data = _json.loads(_terr_cfg_path.read_text())
            for _tname, _tcfg in _terr_data.get("territories", {}).items():
                if _tcfg.get("director") == director:
                    _terr_soql = _tcfg.get("soql_where", "")
                    break
        if _terr_soql:
            import subprocess as _sp

            _auth = _json.loads(
                _sp.check_output(
                    [
                        "sf",
                        "org",
                        "display",
                        "--target-org",
                        "apro@simcorp.com",
                        "--json",
                    ]
                )
            )["result"]
            _sf_session = __import__("requests").Session()
            _sf_session.headers.update(
                {"Authorization": f"Bearer {_auth['accessToken']}"}
            )
            for _qlabel, _mstart, _mend, _mrange, _dstart, _dend in _quarters_to_show:
                _enrich = _fetch_quarter_enrichment(
                    _sf_session,
                    _auth["instanceUrl"],
                    _terr_soql,
                    start_date=_dstart,
                    end_date=_dend,
                )
                if _enrich:
                    n += 1
                    slide_q2_forward_look(
                        prs,
                        pipeline,
                        _enrich,
                        territory,
                        workbook_path=workbook_path,
                        q_label=_qlabel,
                        month_start=_mstart,
                        month_end=_mend,
                    )
                    print(f"  [OK] {n}. {_qlabel} Forward Look ({len(_enrich)} deals)")
                else:
                    print(f"  [SKIP] {_qlabel} Forward Look: no enrichment data")
        else:
            print("  [SKIP] Forward Look: territory config not found")
    except Exception as exc:
        print(f"  [SKIP] Forward Look: {exc}")

    # CUT: Pipeline Overview (duplicate — stage data already in Exec Summary)

    # 8. Top Deals
    n += 1
    slide_top_deals(prs, pipeline)
    print(f"  [OK] {n}. Key Deals")

    # 8. Deal Risk Scoring
    if analytics.get("risk_deals"):
        n += 1
        slide_deal_risk_scoring(prs, analytics["risk_deals"], territory)
        print(f"  [OK] {n}. Deal Risk Scoring")

    # 9. Owner Coaching
    if analytics.get("slip_owners"):
        n += 1
        slide_owner_coaching(
            prs,
            analytics["slip_owners"],
            analytics.get("risk_deals") or [],
            territory,
        )
        print(f"  [OK] {n}. Owner Coaching")

    # 10. Pushed Deals (MERGED: pushed + PI top pushed on one slide)
    n += 1
    slide_pushed_deals_with_link(prs, pi_data, territory, q1_movement=q1_movement)
    print(f"  [OK] {n}. Pushed Deals & PI")

    # CUT: PI Top Pushed (merged into Pushed Deals above)

    # 11. Q1 Movement
    n += 1
    slide_q1_movement(prs, q1_movement)
    print(f"  [OK] {n}. Q1 Movement")

    # 12. Forecast Accuracy
    n += 1
    slide_forecast_combined(prs, won_lost, pipeline, pi_data, territory)
    print(f"  [OK] {n}. Forecast Accuracy & Breakdown")

    # 13. Commercial Approvals (MERGED: approvals + missing on one slide)
    n += 1
    slide_commercial_approvals(prs, approvals, territory)
    print(f"  [OK] {n}. Commercial Approvals")
    # CUT: Missing Approval Detail (merged into Commercial Approvals above)

    # 14. Renewals
    n += 1
    slide_renewals(prs, renewals)
    print(f"  [OK] {n}. Renewals")

    # CUT: Churn Risk (placeholder, no Finance feed yet)
    # CUT: Definitions (moved to handout — not presented)

    # 16. End
    n += 1
    slide_end(prs)
    print(f"  [OK] {n}. End Slide")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")

    # Sidecar: write the headline numbers that the deck was built from so
    # validate_tie_out.py can compare without regex-parsing PowerPoint text.
    # This is the deterministic contract between the builder and the validator.
    import json as _json
    from datetime import datetime as _dt

    def _q1(land_only_deals, stage_filter=None):
        pass

    land_wl_rows = [
        r for r in won_lost if str(r.get("Type", "")).strip().lower() == "land"
    ]
    q1_wins = [
        r
        for r in land_wl_rows
        if "Won" in str(r.get("Stage", ""))
        and str(r.get("Close Date", ""))[:7] <= FQ["prior"]["month_end"]
    ]
    q1_losses = [
        r
        for r in land_wl_rows
        if "Won" not in str(r.get("Stage", ""))
        and str(r.get("Close Date", ""))[:7] <= FQ["prior"]["month_end"]
    ]
    q2_renewals = [
        r
        for r in renewals
        if FQ["current"]["month_start"]
        <= str(r.get("Close Date", ""))[:7]
        <= FQ["current"]["month_end"]
    ]
    q3_renewals = [
        r
        for r in renewals
        if FQ["forward"]["month_start"]
        <= str(r.get("Close Date", ""))[:7]
        <= FQ["forward"]["month_end"]
    ]

    def _acv_renew(r):
        return float(r.get("ACV Unweighted (EUR)") or r.get("ACV (EUR)") or 0)

    sidecar = {
        "director": director,
        "territory": territory,
        "built_at": _dt.now().isoformat(timespec="seconds"),
        "land_only": True,
        "open_land_deals": len(pipeline),
        "open_land_arr": sum(_unw(r) for r in pipeline),
        "open_land_arr_wtd": sum(_wtd(r) for r in pipeline),
        "q2_open_deals": len(_q2_deals),
        "q2_open_arr": sum(_unw(r) for r in _q2_deals),
        "q3_open_deals": len(_q3_deals),
        "q3_open_arr": sum(_unw(r) for r in _q3_deals),
        "q1_land_wins": len(q1_wins),
        "q1_land_wins_arr": sum(_unw(r) for r in q1_wins),
        "q1_land_lost": len(q1_losses),
        "q1_land_lost_arr": sum(_unw(r) for r in q1_losses),
        "q2_renewals": len(q2_renewals),
        "q2_renewals_acv": sum(_acv_renew(r) for r in q2_renewals),
        "q3_renewals": len(q3_renewals),
        "q3_renewals_acv": sum(_acv_renew(r) for r in q3_renewals),
        "approved_2026": sum(
            1 for r in approvals if str(r.get("Status", "")).strip() == "Approved 2026"
        ),
        "conditionally_approved": sum(
            1 for r in approvals if "Conditionally" in str(r.get("Status", ""))
        ),
        "missing_stage3": sum(
            1 for r in approvals if "Missing" in str(r.get("Status", ""))
        ),
    }
    sidecar_path = output_path.with_suffix(".json")
    sidecar_path.write_text(_json.dumps(sidecar, indent=2))
    print(f"Sidecar: {sidecar_path}")


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--workbook", type=Path, required=True)
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--output", type=Path, default=None)
    parser.add_argument(
        "--land-only",
        action="store_true",
        help="Filter to Type=Land only (matches Jesper's APAC LAND forecast page scope)",
    )
    parser.add_argument(
        "--analytics-workbook",
        type=Path,
        default=Path("output/sharepoint/FY26 Pipeline Review, All Territories.xlsx"),
        help=(
            "Consolidated analytics workbook. When provided, the deck adds "
            "Executive Insights, Q1 Forecast Variance, and Deal Risk Scoring "
            "slides with per-director slices."
        ),
    )
    args = parser.parse_args()

    if not args.workbook.exists():
        print(f"ERROR: Workbook not found: {args.workbook}", file=sys.stderr)
        sys.exit(1)

    if args.output:
        output = args.output
    else:
        output = args.workbook.with_suffix(".pptx")

    analytics_path = (
        args.analytics_workbook if args.analytics_workbook.exists() else None
    )

    build_deck(
        args.workbook,
        args.template,
        output,
        land_only=args.land_only,
        analytics_path=analytics_path,
    )


if __name__ == "__main__":
    main()
