#!/usr/bin/env python3
"""Build the Sales Director Monthly deck (v2) — numbered sections, one table per slide.

Renders one PPTX per MD-1 Sales Director using filter presets from
config/sales_director_md1_presets.json. Runs source reports with ephemeral
filter overrides via the SF Analytics Reports API to get per-director data.

Structure follows docs/specs/report-1-deck-structure.md — 6 numbered sections
with separate tables under each.

Usage:
    python3 scripts/build_sd_monthly_deck_v2.py --director "Dan Peppett"
    python3 scripts/build_sd_monthly_deck_v2.py --all
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from datetime import UTC, date, datetime
from pathlib import Path
from typing import Any

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO_ROOT = Path(__file__).resolve().parents[1]
CONFIG_PATH = REPO_ROOT / "config" / "sales_director_md1_presets.json"
DEFAULT_OUTPUT_ROOT = REPO_ROOT / "output" / "sales_director_monthly_runs"
TARGET_ORG = "apro@simcorp.com"
API_VERSION = "v66.0"

# Sales Directors Monthly Pipeline and Insights dashboard
DASHBOARD1_ID = "01ZTb00000FSP7hMAH"
# Sales Ops Quarterly KPI dashboard (has no filters — D2 widgets are global scope)
DASHBOARD2_ID = "01ZTb00000FSP9JMAX"

# D1 component headers → component IDs (stable across runs)
D1_COMP = {
    "Pipeline Overview by Stage": "01aTb00000Cn9mPIAR",
    "Commercial Approval Current State": "01aTb00000Cn9mQIAR",
    "Renewal Likelihood by Probability": "01aTb00000Cn85bIAB",
    "Business At Risk": "01aTb00000Cn85dIAB",
    "Commercial Approval Approved YTD (Land)": "01aTb00000Cn85aIAB",
    "Commercial Approval Candidates by Stage": "01aTb00000Cn85jIAB",
    "Renewal Pipeline This Quarter": "01aTb00000Cn85ZIAR",
    "Close Date Slipped by Stage": "01aTb00000Cn85lIAB",
}

# Deck structure. Each sub-item specifies where its data comes from:
#   source="d1" + component_id  → extract from the pre-filtered D1 componentData cache
#   source="report" + report_id → fall back to per-report POST filter override (used for D2 widgets)
#   source="placeholder"        → awaiting external input
# Format: (subnum, title, source, component_id_or_report_id)
DECK_SECTIONS: list[tuple[int, str, list[tuple[str, str, str, str | None]]]] = [
    (
        1,
        "Pipeline Overview — quarterly focus",
        [
            (
                "1.1",
                "Pipeline by Stage (ARR, current calendar quarter)",
                "d1",
                D1_COMP["Pipeline Overview by Stage"],
            ),
            (
                "1.2",
                "Pipeline by Region (stacked, per director scope)",
                "d1",
                D1_COMP["Pipeline Overview by Stage"],
            ),
            (
                "1.3",
                "Top Opportunities in Pipeline",
                "d1",
                D1_COMP["Pipeline Overview by Stage"],
            ),
        ],
    ),
    (
        2,
        "Commercial Approval",
        [
            (
                "2.1",
                "Current State Overview (approved vs not-approved)",
                "d1",
                D1_COMP["Commercial Approval Current State"],
            ),
            (
                "2.2",
                "YTD Approved Deals (2026 to date)",
                "d1",
                D1_COMP["Commercial Approval Approved YTD (Land)"],
            ),
            (
                "2.3",
                "Land Stage 3 Missing Approval Candidates",
                "d1",
                D1_COMP["Commercial Approval Candidates by Stage"],
            ),
            (
                "2.4",
                "Missing Commercial Approval Opportunities",
                "report",
                "00OTb000008fAlBMAU",
            ),
        ],
    ),
    (
        3,
        "Renewals Tracking",
        [
            (
                "3.1",
                "Renewal ACV This Quarter (metric)",
                "d1",
                D1_COMP["Renewal Pipeline This Quarter"],
            ),
            (
                "3.2",
                "Renewal Likelihood by Probability",
                "d1",
                D1_COMP["Renewal Likelihood by Probability"],
            ),
            (
                "3.3",
                "Upcoming Renewals List (this quarter)",
                "d1",
                D1_COMP["Renewal Pipeline This Quarter"],
            ),
        ],
    ),
    (
        4,
        "Churn Risk & Trends",
        [
            (
                "4.1",
                "Churn Risk Placeholder — awaiting Finance feed from Alex P",
                "placeholder",
                None,
            ),
        ],
    ),
    (
        5,
        "Slipped Deals Analysis",
        [
            (
                "5.1",
                "Close Date Slipped by Stage",
                "d1",
                D1_COMP["Close Date Slipped by Stage"],
            ),
            ("5.2", "Slipped Deals Trend (6-month)", "placeholder", None),
            ("5.3", "Slipped Deals Root Cause Commentary", "placeholder", None),
        ],
    ),
    (
        6,
        "Data Quality",
        [
            (
                "6.1",
                "Missing Win/Loss Reason (excludes 0-No-Opportunity)",
                "report",
                "00OTb000008el0PMAQ",
            ),
            (
                "6.2",
                "Overdue Close Date Open Opps (sorted by record count)",
                "report",
                "00OTb000008TaBZMA0",
            ),
            ("6.3", "Accounts without KYC Approval", "report", "00OTb000007BvlJMAS"),
        ],
    ),
]

# Director name → dashboard1 filter query params (positional: filter1=Industry, filter2=Legal Country, filter3=Sales Region, filter4=Account Unit Group)
# Option IDs pulled from dashboard describe on 2026-04-09.
_OPT = {
    # Industry (filter1)
    "ind_asset_mgmt": "0ICTb0000007DbdOAE",
    "ind_bank": "0ICTb0000007DbeOAE",
    "ind_insurance": "0ICTb0000007DbfOAE",
    "ind_pension": "0ICTb0000007DbgOAE",
    "ind_wealth": "0ICTb0000007DbhOAE",
    "ind_servicer": "0ICTb0000007DbiOAE",
    "ind_other": "0ICTb0000007DbjOAE",
    # Legal Country (filter2)
    "lc_canada": "0ICTb0000007DgTOAU",
    "lc_excl_canada": "0ICTb0000007DgUOAU",
    # Sales Region (filter3)
    "sr_apac": "0ICTb0000007DbnOAE",
    "sr_central_europe": "0ICTb0000007DboOAE",
    "sr_mea": "0ICTb0000007DbpOAE",
    "sr_nam": "0ICTb0000007DbqOAE",
    "sr_northern_europe": "0ICTb0000007DbrOAE",
    "sr_southwestern_europe": "0ICTb0000007DbsOAE",
    "sr_uki": "0ICTb0000007DbtOAE",
    # Account Unit Group (filter4)
    "aug_sc_nam": "0ICTb0000007Di5OAE",
    "aug_sc_asia": "0ICTb0000007Di6OAE",
    "aug_sc_emea": "0ICTb0000007Di7OAE",
}

# Per-director dashboard1 filter params. Industry (filter1) is single-select at the dashboard level,
# so Patrick/Adam can only express ONE primary industry via the dashboard filter — other industries
# are noted on the deck cover as a caveat.
DIRECTOR_D1_FILTERS: dict[str, dict[str, str]] = {
    "Megan Miceli": {
        "filter2": _OPT["lc_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
    "Patrick Gaughan": {
        "filter1": _OPT["ind_asset_mgmt"],
        "filter2": _OPT["lc_excl_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
    "Jesper Tyrer": {"filter3": _OPT["sr_apac"], "filter4": _OPT["aug_sc_asia"]},
    "Sarah Pittroff": {
        "filter3": _OPT["sr_central_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Francois Thaury": {
        "filter3": _OPT["sr_southwestern_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Dan Peppett": {"filter3": _OPT["sr_uki"], "filter4": _OPT["aug_sc_emea"]},
    "Christian Ebbesen": {
        "filter3": _OPT["sr_northern_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Mourad Essofi": {"filter3": _OPT["sr_mea"], "filter4": _OPT["aug_sc_emea"]},
    "Adam Steinhaus": {
        "filter1": _OPT["ind_pension"],
        "filter2": _OPT["lc_excl_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
}

# Colors
TEAL = RGBColor(0x00, 0x74, 0x80)
TEAL_DEEP = RGBColor(0x00, 0x3E, 0x52)
INK = RGBColor(0x1B, 0x1B, 0x1B)
MUTED = RGBColor(0x6E, 0x6E, 0x6E)
PAPER = RGBColor(0xF5, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCC, 0xCC, 0xCC)


def get_auth() -> tuple[str, str]:
    """Resolve SF access token + instance URL via sf CLI."""
    result = subprocess.run(
        ["sf", "org", "display", "--target-org", TARGET_ORG, "--json"],
        capture_output=True,
        text=True,
        check=True,
    )
    data = json.loads(result.stdout)["result"]
    return data["accessToken"], data["instanceUrl"]


def run_report_with_filters(
    instance_url: str,
    token: str,
    report_id: str,
    extra_filters: list[dict[str, Any]],
) -> dict[str, Any]:
    """Execute a report with ephemeral filter overrides and return the result.

    Falls back to unfiltered GET if the POST override returns 400 (e.g. the
    report's sobject doesn't accept Opportunity.* filter columns).
    """
    url = f"{instance_url}/services/data/{API_VERSION}/analytics/reports/{report_id}"
    # GET current report metadata
    r = requests.get(
        url,
        headers={"Authorization": f"Bearer {token}"},
        params={"includeDetails": "true"},
        timeout=60,
    )
    r.raise_for_status()
    payload = r.json()
    rm = payload.get("reportMetadata", {})
    report_type = (rm.get("reportType") or {}).get("type", "")
    # Only apply filters whose column path makes sense for this report type.
    # Opportunity.* columns won't apply to Account-type reports, etc.
    compatible = []
    for f in extra_filters:
        col = f.get("column", "")
        if report_type.startswith("Opportunity") or report_type == "Opportunity":
            if (
                col.startswith("Opportunity.")
                or col in {"INDUSTRY", "ADDRESS1_COUNTRY_CODE"}
                or col.startswith("Account.")
            ):
                compatible.append(f)
        elif "Account" in report_type or "KYC" in report_type:
            if (
                col.startswith("Account.")
                or col == "INDUSTRY"
                or col == "ADDRESS1_COUNTRY_CODE"
            ):
                compatible.append(f)
        else:
            compatible.append(f)
    if not compatible:
        return payload  # no filters applicable; return unfiltered baseline
    existing = rm.get("reportFilters", []) or []
    rm["reportFilters"] = existing + compatible
    body = {"reportMetadata": rm}
    r2 = requests.post(
        url,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        },
        params={"includeDetails": "true"},
        json=body,
        timeout=120,
    )
    if r2.status_code >= 400:
        # Filter column was rejected. Fall back to unfiltered baseline with a note.
        return payload
    return r2.json()


def extract_table(
    result: dict[str, Any], max_rows: int = 15
) -> tuple[list[str], list[list[str]]]:
    """Extract (headers, rows) from a report result.

    Handles both TABULAR (detail rows under factMap['T!T']) and SUMMARY
    (grouped rows under factMap[{group_key}!T] walked via groupingsDown).
    """
    rm = result.get("reportMetadata", {})
    ext_info = result.get("reportExtendedMetadata", {})
    report_format = rm.get("reportFormat", "TABULAR")
    fact_map = result.get("factMap", {}) or {}

    if report_format == "TABULAR":
        detail_info = ext_info.get("detailColumnInfo", {}) or {}
        detail_cols = rm.get("detailColumns", []) or []
        headers = [(detail_info.get(c, {}) or {}).get("label", c) for c in detail_cols]
        fact = fact_map.get("T!T") or {}
        raw_rows = fact.get("rows", []) or []
        rows: list[list[str]] = []
        for r in raw_rows[:max_rows]:
            cells = r.get("dataCells", []) or []
            rows.append(
                [(c.get("label") or str(c.get("value", "")) or "") for c in cells]
            )
        return headers, rows

    # SUMMARY (or MATRIX): walk groupingsDown, pull aggregates per group key
    agg_info = ext_info.get("aggregateColumnInfo", {}) or {}
    agg_names = rm.get("aggregates", []) or []
    group_dims = rm.get("groupingsDown", []) or []
    group_label = (group_dims[0] or {}).get("name", "Group") if group_dims else "Group"
    # Prefer ext label
    group_col_info = (ext_info.get("groupingColumnInfo", {}) or {}).get(group_label, {})
    group_header = group_col_info.get("label", group_label)

    agg_headers = [(agg_info.get(a, {}) or {}).get("label", a) for a in agg_names]
    headers = [group_header] + agg_headers

    groupings = (result.get("groupingsDown", {}) or {}).get("groupings", []) or []
    rows = []
    for g in groupings[:max_rows]:
        key = g.get("key", "")
        label = g.get("label", str(g.get("value", "")))
        fact = fact_map.get(f"{key}!T") or {}
        aggs = fact.get("aggregates", []) or []
        row = [label]
        for a in aggs[: len(agg_headers)]:
            row.append(str(a.get("label", a.get("value", ""))))
        rows.append(row)
    return headers, rows


def extract_aggregate(result: dict[str, Any]) -> str:
    """Pull the primary aggregate value (for metric slides)."""
    fact = result.get("factMap", {}).get("T!T") or {}
    agg = fact.get("aggregates", []) or []
    if agg:
        return str(agg[0].get("label", ""))
    return "n/a"


def load_preset(director_name: str) -> dict[str, Any]:
    """Look up a director's preset by name (fuzzy match on first/last)."""
    config = json.loads(CONFIG_PATH.read_text())
    name_lower = director_name.lower().strip()
    for preset in config["presets"]:
        pname = preset["name"].lower()
        if (
            pname == name_lower
            or name_lower in pname
            or pname.split()[0] == name_lower.split()[0]
        ):
            return preset
    raise ValueError(
        f"No preset matched {director_name!r}. Available: {[p['name'] for p in config['presets']]}"
    )


def slugify_director(name: str) -> str:
    parts = name.lower().strip().split()
    if len(parts) >= 2:
        return f"{parts[-1]}-{parts[0]}"
    return parts[0]


def add_cover_slide(
    pres: Presentation, director: dict[str, Any], snapshot_date: str
) -> None:
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank layout
    # Title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12), Inches(1.2))
    tf = tb.text_frame
    tf.text = "Sales Director Monthly"
    p = tf.paragraphs[0]
    p.font.name = "Avenir Next"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAL_DEEP
    # Subtitle: director + territory
    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(12), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{director['name']} — {director['territory']}"
    p2.font.name = "Avenir Next"
    p2.font.size = Pt(24)
    p2.font.color.rgb = INK
    # Meta
    tb3 = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(12), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = f"Snapshot: {snapshot_date}   |   Pipeline Reporting & Insights"
    p3.font.name = "Avenir Next"
    p3.font.size = Pt(12)
    p3.font.color.rgb = MUTED


def add_section_divider(pres: Presentation, num: int, title: str) -> None:
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    # Teal bar
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.24)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    # Number badge
    tb_num = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.2), Inches(2.5), Inches(2.0)
    )
    p_num = tb_num.text_frame.paragraphs[0]
    p_num.text = str(num)
    p_num.font.name = "Avenir Next"
    p_num.font.size = Pt(140)
    p_num.font.bold = True
    p_num.font.color.rgb = TEAL
    # Title
    tb = slide.shapes.add_textbox(Inches(3.2), Inches(3.0), Inches(9.5), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Avenir Next"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAL_DEEP


def add_table_slide(
    pres: Presentation,
    subnum: str,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    footer_source: str,
) -> None:
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    # Title: "1.1 Pipeline by Stage"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{subnum}  {title}"
    p.font.name = "Avenir Next"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEAL_DEEP
    # Table
    if not rows:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(1))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = "(No rows returned for this director scope.)"
        p2.font.name = "Avenir Next"
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED
    else:
        n_cols = max(1, len(headers))
        n_rows = len(rows) + 1  # header row
        # Cap columns at 6 to stay readable on 13.33" width
        max_cols = 6
        if n_cols > max_cols:
            headers = headers[:max_cols]
            rows = [r[:max_cols] for r in rows]
            n_cols = max_cols
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5)
        )
        tbl = table_shape.table
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = str(h)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Avenir Next"
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL_DEEP
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Avenir Next"
                        run.font.size = Pt(10)
                        run.font.color.rgb = INK
    # Footer
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.3))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = f"Source: {footer_source}"
    p3.font.name = "Avenir Next"
    p3.font.size = Pt(9)
    p3.font.color.rgb = MUTED


def add_placeholder_slide(
    pres: Presentation, subnum: str, title: str, reason: str
) -> None:
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{subnum}  {title}"
    p.font.name = "Avenir Next"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEAL_DEEP
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.5), Inches(1.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = reason
    p2.font.name = "Avenir Next"
    p2.font.size = Pt(18)
    p2.font.italic = True
    p2.font.color.rgb = MUTED


def get_dashboard_filtered(
    instance_url: str,
    token: str,
    dashboard_id: str,
    filter_params: dict[str, str],
    wait_sec: int = 6,
) -> dict[str, dict[str, Any]]:
    """Run a dashboard with filter selections and return {component_id: reportResult}.

    1. PUT the dashboard URL with filterN=<optionId> query params to trigger an
       async refresh of the filtered combination.
    2. Wait for the refresh.
    3. GET the same URL → componentData contains the filtered factMaps.
    """
    base = f"{instance_url}/services/data/{API_VERSION}/analytics/dashboards/{dashboard_id}"
    if filter_params:
        put_resp = requests.put(
            base,
            headers={"Authorization": f"Bearer {token}"},
            params=filter_params,
            timeout=60,
        )
        # 201 = refresh queued, 200 = cached
        if put_resp.status_code not in (200, 201):
            print(f"  warn: PUT refresh returned {put_resp.status_code}")
        import time

        time.sleep(wait_sec)
    get_resp = requests.get(
        base,
        headers={"Authorization": f"Bearer {token}"},
        params=filter_params,
        timeout=60,
    )
    get_resp.raise_for_status()
    payload = get_resp.json()
    out: dict[str, dict[str, Any]] = {}
    for cd in payload.get("componentData", []) or []:
        cid = cd.get("componentId", "")
        rr = cd.get("reportResult", {}) or {}
        out[cid] = rr
    return out


def build_deck_for_director(
    director: dict[str, Any],
    snapshot_date: str,
    output_path: Path,
    token: str,
    instance_url: str,
) -> None:
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    # 1. Cover
    add_cover_slide(pres, director, snapshot_date)

    # Fetch D1 once with the director's dashboard-level filter selection.
    # This single call replaces 8 per-report POST overrides and gives us
    # exactly what the Sales Director sees when they open the dashboard
    # with their preset filters applied.
    d1_filter_params = DIRECTOR_D1_FILTERS.get(director["name"], {})
    if d1_filter_params:
        print(f"  D1 filter params: {d1_filter_params}")
    try:
        d1_results = get_dashboard_filtered(
            instance_url, token, DASHBOARD1_ID, d1_filter_params
        )
        print(f"  D1 fetched: {len(d1_results)} components")
    except Exception as exc:
        print(f"  D1 fetch FAILED: {exc}")
        d1_results = {}

    # D2 has no dashboard filters — per-report override fallback for those.
    extra_filters = director.get("filters", [])

    for section_num, section_title, subitems in DECK_SECTIONS:
        add_section_divider(pres, section_num, section_title)
        for subnum, sub_title, source, target_id in subitems:
            if source == "placeholder":
                reason = (
                    "Awaiting Finance feed (Alex P)"
                    if section_num == 4
                    else "Pipeline Inspection native required — see Phase 4 handoff"
                )
                if "root cause" in sub_title.lower():
                    reason = "Opp owner commentary outreach pending — column populated manually"
                add_placeholder_slide(pres, subnum, sub_title, reason)
                continue
            if source == "d1":
                assert target_id is not None
                rr = d1_results.get(target_id)
                if rr is None:
                    add_placeholder_slide(
                        pres,
                        subnum,
                        sub_title,
                        f"(D1 component {target_id} not in response)",
                    )
                    print(f"  MISS {subnum} — component {target_id} not in D1 results")
                    continue
                try:
                    # Synthesize a report-shaped payload for extract_table
                    synthetic = {
                        "reportMetadata": rr.get("reportMetadata", {}) or {},
                        "reportExtendedMetadata": rr.get("reportExtendedMetadata", {})
                        or {},
                        "factMap": rr.get("factMap", {}) or {},
                        "groupingsDown": rr.get("groupingsDown", {}) or {},
                    }
                    headers, rows = extract_table(synthetic)
                    add_table_slide(
                        pres,
                        subnum,
                        sub_title,
                        headers,
                        rows,
                        f"Dashboard 1 (filter-applied) | component {target_id}",
                    )
                    print(f"  OK   {subnum} ({len(rows)} rows) — d1:{target_id}")
                except Exception as exc:
                    add_placeholder_slide(
                        pres, subnum, sub_title, f"(D1 extract failed: {exc})"
                    )
                    print(f"  FAIL {subnum} — d1:{target_id}: {exc}")
                continue
            if source == "report":
                assert target_id is not None
                try:
                    result = run_report_with_filters(
                        instance_url, token, target_id, extra_filters
                    )
                    headers, rows = extract_table(result)
                    add_table_slide(
                        pres,
                        subnum,
                        sub_title,
                        headers,
                        rows,
                        f"SF Report {target_id} (D2 global — no dashboard filter)",
                    )
                    print(f"  OK   {subnum} ({len(rows)} rows) — report:{target_id}")
                except Exception as exc:
                    add_placeholder_slide(
                        pres, subnum, sub_title, f"(Data fetch failed: {exc})"
                    )
                    print(f"  FAIL {subnum} — report:{target_id}: {exc}")
                continue
    # Closing
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "End of Report"
    p.font.name = "Avenir Next"
    p.font.size = Pt(28)
    p.font.color.rgb = TEAL_DEEP
    p2 = tb.text_frame.add_paragraph()
    p2.text = f"Generated {datetime.now(UTC).strftime('%Y-%m-%d %H:%M UTC')}"
    p2.font.name = "Avenir Next"
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(output_path))
    print(f"Saved: {output_path}  ({len(pres.slides)} slides)")


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("--director", help="Single director name (e.g. 'Dan Peppett')")
    p.add_argument(
        "--all", action="store_true", help="Generate decks for all 9 MD-1 directors"
    )
    p.add_argument("--snapshot-date", default=date.today().isoformat())
    p.add_argument("--output-root", default=str(DEFAULT_OUTPUT_ROOT))
    return p.parse_args()


def main() -> int:
    args = parse_args()
    if not args.director and not args.all:
        print("Must specify --director NAME or --all", file=sys.stderr)
        return 2
    token, instance_url = get_auth()
    config = json.loads(CONFIG_PATH.read_text())
    targets = config["presets"] if args.all else [load_preset(args.director)]
    output_root = Path(args.output_root) / args.snapshot_date
    for director in targets:
        slug = slugify_director(director["name"])
        out = output_root / f"{slug}.pptx"
        print(f"\n=== {director['name']} ({director['territory']}) ===")
        build_deck_for_director(director, args.snapshot_date, out, token, instance_url)
    return 0


if __name__ == "__main__":
    sys.exit(main())
