#!/usr/bin/env python3
"""Build the canonical Sales Region monthly shell from the SimCorp master."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from territory_mapping import get_director_book, get_forecast_rollup_config
except ModuleNotFoundError:  # pragma: no cover - import path differs under pytest/package execution
    from scripts.territory_mapping import get_director_book, get_forecast_rollup_config


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MASTER_TEMPLATE_PATH = (
    Path.home()
    / "archive"
    / "simcorp-deck-agent-backup"
    / "reference-decks"
    / "SimCorp_PPT_Template.pptx"
)
DEFAULT_SHELL_CONTRACT_PATH = REPO_ROOT / "config" / "sales_region_monthly_shell.json"
DEFAULT_OUTPUT_ROOT = REPO_ROOT / "output" / "sales_region_monthly_shells"
DIRECTOR_ORDER = [
    "Jesper Tyrer",
    "Sarah Pittroff",
    "Francois Thaury",
    "Dan Peppett",
    "Christian Ebbesen",
    "Mourad Essofi",
    "Megan Miceli",
    "Patrick Gaughan",
    "Adam Steinhaus",
]

# SimCorp production colors
NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
RICH_BLUE = RGBColor(0x09, 0x4F, 0xB2)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
AQUA_LIGHT = RGBColor(0xCB, 0xF5, 0xF3)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
LIGHT_PANEL_2 = RGBColor(0xF5, 0xF8, 0xFD)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
MAGENTA_LIGHT = RGBColor(0xF8, 0xE8, 0xF1)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)
DIVIDER_GREY = RGBColor(0xD7, 0xE2, 0xE8)
AMBER = RGBColor(0xFB, 0x9B, 0x2A)
GREEN_OK = RGBColor(0x33, 0xD8, 0xCE)

SECTION_LABELS = {
    "executive-summary": "EXECUTIVE SUMMARY",
    "q1-review": "Q1 REVIEW",
    "quarterly-pipeline": "PIPELINE",
    "regional-book-breakdown": "BOOK BREAKDOWN",
    "pipeline-coverage-intel": "PIPELINE INTEL",
    "commercial-approval-overview": "COMMERCIAL APPROVAL",
    "missing-commercial-approvals": "APPROVAL FOLLOW-UP",
    "renewals-retention": "RENEWALS",
    "slipped-deals": "SLIPPED DEALS",
    "churn-finance": "CHURN / FINANCE",
    "appendix-notes": "APPENDIX",
}

SECTION_ACCENTS = {
    "executive-summary": AQUA,
    "q1-review": AMBER,
    "quarterly-pipeline": AQUA,
    "regional-book-breakdown": PRIMARY_BLUE,
    "pipeline-coverage-intel": AQUA,
    "commercial-approval-overview": AMBER,
    "missing-commercial-approvals": AMBER,
    "renewals-retention": MAGENTA,
    "slipped-deals": AMBER,
    "churn-finance": MAGENTA,
    "appendix-notes": PRIMARY_BLUE,
}


def load_shell_contract(path: Path = DEFAULT_SHELL_CONTRACT_PATH) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def get_layout(prs: Presentation, name: str):
    lowered = name.lower()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name or layout.name.lower() == lowered:
                return layout
    raise KeyError(f"Layout not found: {name}")


def clear_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        rel_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]


def fill_placeholder(slide, idx: int, text: str | None) -> bool:
    if text is None:
        return False
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = str(text)
            return True
    return False


def txt(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = NAVY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft Sans Serif"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def multi_paragraph(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: Iterable[str],
    *,
    size: int = 12,
    bold: bool = False,
    color: RGBColor = NAVY,
    bullet: bool = False,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.clear()
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{line}" if not bullet else f"• {line}"
        run.font.name = "Microsoft Sans Serif"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def rounded_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def dot(slide, left: float, top: float, size: float, color: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def blank_slide(prs: Presentation):
    try:
        layout = get_layout(prs, "Blank")
    except KeyError:
        layout = get_layout(prs, "Title only")
    slide = prs.slides.add_slide(layout)
    for placeholder in list(slide.placeholders):
        placeholder._element.getparent().remove(placeholder._element)
    return slide


def header(
    slide,
    *,
    eyebrow: str,
    title: str,
    narrative: str,
    accent: RGBColor = AQUA,
) -> None:
    rect(slide, 0.6, 0.42, 0.08, 0.7, accent)
    txt(slide, 0.82, 0.38, 12.0, 0.24, eyebrow, size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, 0.82, 0.62, 12.2, 0.55, title, size=28, bold=True, color=NAVY)
    txt(slide, 0.82, 1.25, 12.1, 0.4, narrative, size=13, color=GREY_TEXT)


def source_footer(slide, text: str) -> None:
    txt(slide, 0.6, 7.05, 12.1, 0.22, text, size=9, color=GREY_TEXT)


def divider(slide, top: float) -> None:
    rect(slide, 0.6, top, 12.1, 0.015, DIVIDER_GREY)


def takeaway(slide, *, top: float, bullets: list[str], accent: RGBColor = AQUA) -> None:
    divider(slide, top - 0.25)
    rounded_card(slide, 0.6, top, 12.1, 1.55, LIGHT_PANEL_2)
    rect(slide, 0.6, top, 0.1, 1.55, accent)
    txt(slide, 0.95, top + 0.14, 11.5, 0.2, "LEADERSHIP TAKEAWAYS", size=10, bold=True, color=PRIMARY_BLUE)
    multi_paragraph(slide, 0.95, top + 0.4, 11.3, 0.95, bullets, size=12, color=NAVY, bullet=True)


def kpi_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    big: str,
    label: str,
    context: str,
    panel: RGBColor = LIGHT_BLUE_PANEL,
    accent_dot: RGBColor | None = None,
) -> None:
    rounded_card(slide, left, top, width, height, panel)
    txt(slide, left + 0.22, top + 0.18, width - 0.44, 0.52, big, size=24, bold=True, color=PRIMARY_BLUE)
    txt(slide, left + 0.22, top + 0.88, width - 0.44, 0.28, label.upper(), size=10, bold=True, color=NAVY)
    txt(slide, left + 0.22, top + 1.18, width - 0.44, height - 1.35, context, size=10, color=GREY_TEXT)
    if accent_dot is not None:
        dot(slide, left + width - 0.42, top + 0.24, 0.18, accent_dot)


def hero_stat_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    overline: str,
    big: str,
    context: str,
    panel: RGBColor = LIGHT_BLUE_PANEL,
    accent: RGBColor = AQUA,
) -> None:
    rounded_card(slide, left, top, width, height, panel)
    rect(slide, left, top, 0.12, height, accent)
    txt(slide, left + 0.28, top + 0.22, width - 0.56, 0.24, overline, size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, left + 0.28, top + 0.58, width - 0.56, 0.95, big, size=34, bold=True, color=NAVY)
    txt(slide, left + 0.28, top + 1.62, width - 0.56, height - 1.9, context, size=11, color=GREY_TEXT)


def mini_stat(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    big: str,
    context: str,
    panel: RGBColor = LIGHT_PANEL_2,
    status: RGBColor | None = None,
) -> None:
    rounded_card(slide, left, top, width, height, panel)
    txt(slide, left + 0.18, top + 0.16, width - 0.36, 0.2, label.upper(), size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, left + 0.18, top + 0.42, width - 0.36, 0.42, big, size=20, bold=True, color=NAVY)
    txt(slide, left + 0.18, top + 0.84, width - 0.36, height - 1.0, context, size=9, color=GREY_TEXT)
    if status is not None:
        dot(slide, left + width - 0.28, top + 0.2, 0.14, status)


def content_panel(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    lines: list[str],
    panel: RGBColor = LIGHT_BLUE_PANEL,
    accent: RGBColor = AQUA,
) -> None:
    rounded_card(slide, left, top, width, height, panel)
    rect(slide, left, top, 0.1, height, accent)
    txt(slide, left + 0.22, top + 0.18, width - 0.44, 0.24, title.upper(), size=10, bold=True, color=PRIMARY_BLUE)
    multi_paragraph(slide, left + 0.22, top + 0.5, width - 0.44, height - 0.7, lines, size=11, color=NAVY, bullet=True)


def section_eyebrow(snapshot_date: str, number: int, section_label: str, region_name: str) -> str:
    return f"{snapshot_date}   ·   {number:02d}   ·   {region_name.upper()}   ·   {section_label}"


def region_component_books(region_name: str) -> list[str]:
    cfg = get_forecast_rollup_config(region_name)
    sales_regions = set(cfg.get("sales_regions", []))
    components: list[str] = []
    for director_name in DIRECTOR_ORDER:
        book = get_director_book(director_name)
        if book.get("sales_region") in sales_regions:
            components.append(f"{book['territory']} - {director_name}")
    return components


def region_breakdown_guidance(region_name: str) -> list[str]:
    lines = region_component_books(region_name)
    if region_name == "EMEA":
        lines.append("Forecast rule: Middle East & Africa rolls into EMEA, not APAC.")
    return lines


def add_title_slide(prs: Presentation, *, region_name: str, snapshot_date: str) -> None:
    slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
    fill_placeholder(slide, 20, "")
    fill_placeholder(slide, 22, "")
    fill_placeholder(slide, 24, "")
    txt(slide, 1.0, 1.55, 2.3, 0.3, "SIMCORP", size=11, bold=True, color=AQUA_LIGHT)
    txt(slide, 1.0, 2.0, 4.6, 0.52, "Sales Region Monthly", size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    txt(slide, 1.0, 2.55, 5.8, 0.85, region_name, size=30, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    txt(slide, 1.0, 3.25, 6.3, 0.48, f"Executive operating deck built from validated workbook facts | Snapshot {snapshot_date}", size=15, color=AQUA_LIGHT)
    txt(slide, 1.0, 5.95, 3.2, 0.28, "Monthly regional cadence", size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def add_agenda_slide(prs: Presentation, shell: dict[str, Any], *, region_name: str, snapshot_date: str) -> None:
    slide = blank_slide(prs)
    header(
        slide,
        eyebrow=f"{snapshot_date}   ·   OPERATING CADENCE   ·   {region_name.upper()}",
        title=shell.get("agenda_title", "Regional Monthly Agenda"),
        narrative=shell.get("agenda_subtitle", "Fixed operating sequence for each top-level sales region"),
        accent=PRIMARY_BLUE,
    )
    titles = [item["title"] for item in shell.get("slides", [])]
    card_w = 3.85
    card_h = 0.92
    start_x = 0.6
    start_y = 1.95
    x_gap = 0.28
    y_gap = 0.18
    for index, title in enumerate(titles):
        col = index % 3
        row = index // 3
        left = start_x + col * (card_w + x_gap)
        top = start_y + row * (card_h + y_gap)
        rounded_card(slide, left, top, card_w, card_h, LIGHT_BLUE_PANEL if index < 6 else LIGHT_PANEL_2)
        txt(slide, left + 0.18, top + 0.16, 0.4, 0.3, f"{index + 1:02d}", size=10, bold=True, color=PRIMARY_BLUE)
        txt(slide, left + 0.58, top + 0.14, card_w - 0.76, 0.48, title, size=12, bold=True, color=NAVY)
    takeaway(
        slide,
        top=6.02,
        bullets=[
            "Use the same operating sequence every month so leadership learns where to find each answer.",
            "Populate facts only from the validated regional fact pack; do not freehand rollups in PowerPoint.",
            "Treat this deck as the regional shell. Director-book detail belongs in supporting MD-1 packs.",
        ],
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Shell only — populate from validated regional fact pack before review.")


def add_executive_summary_slide(slide, slide_def: dict[str, Any]) -> None:
    cards = slide_def.get("cards", [])
    for index, card in enumerate(cards[:4]):
        kpi_card(
            slide,
            left=0.6 + index * 3.08,
            top=1.95,
            width=2.9,
            height=1.95,
            big=card["metric_hint"],
            label=card["label"],
            context=card["body"],
            panel=MAGENTA_LIGHT if index == 3 else LIGHT_BLUE_PANEL,
            accent_dot=GREEN_OK if index < 3 else MAGENTA,
        )
    takeaway(
        slide,
        top=4.35,
        bullets=[
            "Top risk: insert the single most material regional risk from the validated fact pack.",
            "Top action: insert the one leadership action that changes the quarter, not a generic summary.",
            "Keep horizons explicit: All Open, FY26, Q2. Do not collapse them into one 'pipeline' number.",
        ],
    )


def add_q1_review_slide(slide, slide_def: dict[str, Any]) -> None:
    cards = slide_def.get("cards", [])
    for index, card in enumerate(cards[:4]):
        kpi_card(
            slide,
            left=0.6 + index * 3.08,
            top=1.95,
            width=2.9,
            height=2.05,
            big=card["metric_hint"],
            label=card["label"],
            context=card["body"],
            panel=LIGHT_PANEL_2 if index == 3 else LIGHT_BLUE_PANEL,
            accent_dot=AMBER if index in (1, 2, 3) else GREEN_OK,
        )
    takeaway(
        slide,
        top=4.52,
        bullets=[
            "Show what the region delivered in Q1 before discussing Q2 ambition.",
            "Only state a promise baseline when the source is forecast-safe; otherwise qualify the ambiguity explicitly.",
            "Use slipped count and slipped ARR from the validated pack, not workbook tab noise.",
        ],
        accent=AMBER,
    )


def add_quarterly_pipeline_slide(slide, slide_def: dict[str, Any]) -> None:
    cards = slide_def.get("cards", [])
    primary = cards[0]
    secondary = cards[1:4]
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=5.8,
        height=3.15,
        overline=primary["label"].upper(),
        big=primary["metric_hint"],
        context=primary["body"],
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    for index, card in enumerate(secondary):
        mini_stat(
            slide,
            left=6.65,
            top=1.95 + index * 1.08,
            width=6.05,
            height=0.92,
            label=card["label"],
            big=card["metric_hint"],
            context=card["body"],
            panel=LIGHT_PANEL_2,
            status=AMBER if card["label"] == "Omitted" else GREEN_OK,
        )
    takeaway(
        slide,
        top=5.42,
        bullets=[
            "Lead with displayed-quarter active ARR, then show Commit and Best Case as forecast mix, not as substitute headlines.",
            "Keep Omitted visible and separate. It is an exposure signal, not active pipeline.",
            "If the current quarter is empty, disclose the forward-quarter fallback plainly in the footnote.",
        ],
    )


def add_book_breakdown_slide(slide, region_name: str) -> None:
    rounded_card(slide, 0.6, 1.95, 5.9, 4.65, LIGHT_BLUE_PANEL)
    rect(slide, 0.6, 1.95, 0.1, 4.65, PRIMARY_BLUE)
    txt(slide, 0.82, 2.12, 5.2, 0.22, "COMPONENT BOOKS", size=10, bold=True, color=PRIMARY_BLUE)
    books = region_component_books(region_name)
    for index, book in enumerate(books[:5]):
        row_top = 2.45 + index * 0.72
        rounded_card(slide, 0.92, row_top, 5.2, 0.56, RGBColor(0xF8, 0xFB, 0xFF))
        txt(slide, 1.12, row_top + 0.11, 4.8, 0.24, book, size=12, bold=True, color=NAVY)
    if region_name == "EMEA":
        txt(slide, 0.92, 6.15, 5.2, 0.26, "Forecast rule: Middle East & Africa rolls into EMEA, not APAC.", size=10, bold=True, color=GREY_TEXT)
    content_panel(
        slide,
        left=6.75,
        top=1.95,
        width=5.95,
        height=1.45,
        title="Largest book summary",
        lines=[
            "Insert the dominant MD-1 book in this region.",
            "State the reason it matters: concentration, upside, or control risk.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AQUA,
    )
    content_panel(
        slide,
        left=6.75,
        top=3.6,
        width=5.95,
        height=1.45,
        title="Weakest book summary",
        lines=[
            "Insert the book with the clearest delivery gap or hygiene burden.",
            "Keep the statement factual and book-specific.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=6.75,
        top=5.25,
        width=5.95,
        height=1.35,
        title="Rollup rule",
        lines=[
            "Preserve the Salesforce forecast hierarchy in every monthly cut.",
            "Do not reassign subregions in PowerPoint.",
        ],
        panel=LIGHT_PANEL_2,
        accent=PRIMARY_BLUE,
    )


def add_pipeline_intel_slide(slide, slide_def: dict[str, Any]) -> None:
    cards = slide_def.get("cards", [])
    for index, card in enumerate(cards[:3]):
        kpi_card(
            slide,
            left=0.6 + index * 4.08,
            top=1.95,
            width=3.85,
            height=2.45,
            big=card["metric_hint"],
            label=card["label"],
            context=card["body"],
            panel=LIGHT_BLUE_PANEL if index < 2 else LIGHT_PANEL_2,
            accent_dot=GREEN_OK if index == 0 else AMBER,
        )
    takeaway(
        slide,
        top=4.9,
        bullets=[
            "This slide is where the region earns or loses credibility: coverage, top deals, and hygiene pressure together.",
            "Keep named opportunities few and material. Use only the largest validated deals with a concrete next step.",
            "If stale ARR or data-quality backlog is the issue, say so directly instead of hiding behind generic language.",
        ],
    )


def add_commercial_approval_slide(slide, slide_def: dict[str, Any]) -> None:
    cards = slide_def.get("cards", [])
    hero = cards[0]
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=5.35,
        height=3.0,
        overline=hero["label"].upper(),
        big=hero["metric_hint"],
        context=hero["body"],
        panel=LIGHT_BLUE_PANEL,
        accent=AMBER,
    )
    for index, card in enumerate(cards[1:]):
        mini_stat(
            slide,
            left=6.2,
            top=1.95 + index * 1.02,
            width=6.5,
            height=0.88,
            label=card["label"],
            big=card["metric_hint"],
            context=card["body"],
            panel=LIGHT_PANEL_2,
            status=AMBER if "Missing" in card["label"] or "Pending" in card["label"] else GREEN_OK,
        )
    takeaway(
        slide,
        top=5.35,
        bullets=[
            "Use this slide as the control seam for stage 3+ governance across the whole region.",
            "Approved, pending, and missing should reconcile to the validated snapshot logic.",
            "If approval-rate methodology is approximate, disclose the method instead of implying precision.",
        ],
        accent=AMBER,
    )


def add_missing_approvals_slide(slide, slide_def: dict[str, Any]) -> None:
    rounded_card(slide, 0.6, 1.95, 7.2, 4.7, LIGHT_BLUE_PANEL)
    rect(slide, 0.6, 1.95, 0.1, 4.7, AMBER)
    txt(slide, 0.82, 2.12, 6.8, 0.22, "CANDIDATE LIST", size=10, bold=True, color=PRIMARY_BLUE)
    rounded_card(slide, 0.92, 2.45, 6.55, 0.42, RGBColor(0xF8, 0xFB, 0xFF))
    txt(slide, 1.12, 2.56, 1.8, 0.18, "Opportunity", size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, 3.1, 2.56, 1.0, 0.18, "ARR", size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, 4.1, 2.56, 1.35, 0.18, "Owner / Book", size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, 5.72, 2.56, 1.45, 0.18, "Next action", size=9, bold=True, color=PRIMARY_BLUE)
    for row_index in range(5):
        row_top = 3.02 + row_index * 0.64
        rounded_card(slide, 0.92, row_top, 6.55, 0.5, RGBColor(0xF8, 0xFB, 0xFF))
        txt(slide, 1.12, row_top + 0.12, 1.78, 0.18, f"Insert candidate {row_index + 1}", size=10, bold=True, color=NAVY)
        txt(slide, 3.1, row_top + 0.12, 0.92, 0.18, "EUR ARR", size=10, color=GREY_TEXT)
        txt(slide, 4.1, row_top + 0.12, 1.35, 0.18, "Owner  |  Book", size=10, color=GREY_TEXT)
        txt(slide, 5.72, row_top + 0.12, 1.45, 0.18, "Approval follow-up", size=10, color=GREY_TEXT)
    content_panel(
        slide,
        left=8.05,
        top=1.95,
        width=4.65,
        height=2.1,
        title="Formatting guardrails",
        lines=[
            "Sort by ARR descending.",
            "Show owner, book, stage, and next action in one line per deal.",
            "Do not infer approval beyond the validated snapshot.",
        ],
        panel=LIGHT_PANEL_2,
        accent=PRIMARY_BLUE,
    )
    content_panel(
        slide,
        left=8.05,
        top=4.25,
        width=4.65,
        height=2.4,
        title="Leadership action",
        lines=[
            "Route urgent candidates to the right MD-1 immediately after the readout.",
            "This slide should drive action, not just document the backlog.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AMBER,
    )


def add_renewals_slide(slide, slide_def: dict[str, Any]) -> None:
    cards = slide_def.get("cards", [])
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=4.15,
        height=3.25,
        overline=cards[0]["label"].upper(),
        big=cards[0]["metric_hint"],
        context=cards[0]["body"],
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    content_panel(
        slide,
        left=4.98,
        top=1.95,
        width=3.75,
        height=3.25,
        title=cards[1]["label"],
        lines=[cards[1]["body"], "State sparse tagging or methodology limits explicitly when needed."],
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    content_panel(
        slide,
        left=8.95,
        top=1.95,
        width=3.75,
        height=3.25,
        title=cards[2]["label"],
        lines=[cards[2]["body"], "Name the renewal owner and next checkpoint for the largest cases."],
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    takeaway(
        slide,
        top=5.55,
        bullets=[
            "Renewal amounts stay ACV, converted to EUR. Do not mix renewal ACV with new-business ARR.",
            "If the risk distribution is sparse, disclose that and focus leadership attention on named renewals.",
            "This slide should answer: what renews this quarter, what is at risk, and who owns the intervention.",
        ],
        accent=MAGENTA,
    )


def add_slipped_slide(slide, slide_def: dict[str, Any]) -> None:
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=4.2,
        height=3.2,
        overline="VALIDATED SLIPPED EXPOSURE",
        big="Count + EUR ARR",
        context="Lead with the regional slipped-deal count and slipped ARR from the validated fact pack.",
        panel=LIGHT_BLUE_PANEL,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=5.05,
        top=1.95,
        width=3.8,
        height=3.2,
        title="Watchlist",
        lines=[
            "List the largest slipped opportunities first.",
            "Include stated reason only if it exists in the source; otherwise leave commentary open.",
            "Tie the opportunity back to the responsible book owner.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=9.05,
        top=1.95,
        width=3.65,
        height=3.2,
        title="Root cause commentary",
        lines=[
            "Use owner follow-up to fill this section.",
            "If commentary is incomplete, mark that gap explicitly.",
            "Do not fabricate root cause from pipeline metadata alone.",
        ],
        panel=LIGHT_PANEL_2,
        accent=PRIMARY_BLUE,
    )
    takeaway(
        slide,
        top=5.48,
        bullets=[
            "Keep slipped ARR and slipped count factual and validated.",
            "Root-cause commentary is a controlled override from owners, not an LLM inference exercise.",
            "Use this slide to drive post-readout follow-up by book, not to hand-wave misses.",
        ],
        accent=AMBER,
    )


def add_churn_slide(slide, slide_def: dict[str, Any]) -> None:
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=12.1,
        height=2.15,
        overline="FINANCE INPUT STATUS",
        big="Finance overlay required",
        context=slide_def.get("body_guidance", ["State current status plainly."])[0],
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    mini_stat(
        slide,
        left=0.6,
        top=4.35,
        width=3.8,
        height=1.2,
        label="Source status",
        big="Known / missing",
        context="State what feed exists today and what still has to be operationalized.",
        panel=LIGHT_PANEL_2,
        status=AMBER,
    )
    mini_stat(
        slide,
        left=4.7,
        top=4.35,
        width=3.8,
        height=1.2,
        label="Finance owner",
        big="Named owner",
        context="Use the actual Finance reporting owner when known. Leave blank rather than guessing.",
        panel=LIGHT_PANEL_2,
        status=PRIMARY_BLUE,
    )
    mini_stat(
        slide,
        left=8.8,
        top=4.35,
        width=3.9,
        height=1.2,
        label="Next milestone",
        big="Operational plan",
        context="State what must happen before churn trend numbers can become publishable.",
        panel=LIGHT_PANEL_2,
        status=AMBER,
    )
    takeaway(
        slide,
        top=5.95,
        bullets=[
            "Keep this slide explicit until Finance inputs are wired into the monthly process.",
            "No placeholder churn numbers. Status and ownership are the product until the feed is real.",
        ],
        accent=MAGENTA,
    )


def add_appendix_slide(slide, slide_def: dict[str, Any]) -> None:
    content_panel(
        slide,
        left=0.6,
        top=1.95,
        width=5.95,
        height=4.8,
        title="Metric rules",
        lines=[
            "New business and expansion stay ARR, EUR converted.",
            "Renewals stay ACV, EUR converted.",
            "Omitted remains visible but excluded from active headline pipeline.",
            "Do not state Q1 promise baseline without source qualification.",
        ],
        panel=LIGHT_BLUE_PANEL,
        accent=PRIMARY_BLUE,
    )
    content_panel(
        slide,
        left=6.8,
        top=1.95,
        width=5.9,
        height=4.8,
        title="Source and lineage notes",
        lines=slide_def.get("body_guidance", []) + [
            "Use workbook snapshot, validated fact pack, and forecast hierarchy rules as source of truth.",
            "Document known data-quality limitations instead of hiding them in narration.",
        ],
        panel=LIGHT_PANEL_2,
        accent=PRIMARY_BLUE,
    )


def add_shell_slide(
    prs: Presentation,
    slide_def: dict[str, Any],
    *,
    region_name: str,
    snapshot_date: str,
    section_number: int,
) -> None:
    slide = blank_slide(prs)
    header(
        slide,
        eyebrow=section_eyebrow(
            snapshot_date,
            section_number,
            SECTION_LABELS.get(slide_def["id"], slide_def["title"].upper()),
            region_name,
        ),
        title=slide_def["title"],
        narrative=slide_def.get("subtitle", ""),
        accent=SECTION_ACCENTS.get(slide_def["id"], AQUA),
    )

    if slide_def["id"] == "executive-summary":
        add_executive_summary_slide(slide, slide_def)
    elif slide_def["id"] == "q1-review":
        add_q1_review_slide(slide, slide_def)
    elif slide_def["id"] == "quarterly-pipeline":
        add_quarterly_pipeline_slide(slide, slide_def)
    elif slide_def["id"] == "regional-book-breakdown":
        add_book_breakdown_slide(slide, region_name)
    elif slide_def["id"] == "pipeline-coverage-intel":
        add_pipeline_intel_slide(slide, slide_def)
    elif slide_def["id"] == "commercial-approval-overview":
        add_commercial_approval_slide(slide, slide_def)
    elif slide_def["id"] == "missing-commercial-approvals":
        add_missing_approvals_slide(slide, slide_def)
    elif slide_def["id"] == "renewals-retention":
        add_renewals_slide(slide, slide_def)
    elif slide_def["id"] == "slipped-deals":
        add_slipped_slide(slide, slide_def)
    elif slide_def["id"] == "churn-finance":
        add_churn_slide(slide, slide_def)
    elif slide_def["id"] == "appendix-notes":
        add_appendix_slide(slide, slide_def)
    else:  # pragma: no cover - defensive against future config drift
        raise ValueError(f"Unsupported shell slide id: {slide_def['id']}")

    source_footer(slide, "Shell only — replace guidance text with validated regional facts before leadership review.")


def build_shell_deck(
    *,
    region_name: str,
    snapshot_date: str,
    output_path: Path,
    master_template_path: Path = DEFAULT_MASTER_TEMPLATE_PATH,
    shell_contract_path: Path = DEFAULT_SHELL_CONTRACT_PATH,
) -> dict[str, Any]:
    shell = load_shell_contract(shell_contract_path)
    prs = Presentation(str(master_template_path))
    clear_slides(prs)

    add_title_slide(prs, region_name=region_name, snapshot_date=snapshot_date)
    add_agenda_slide(prs, shell, region_name=region_name, snapshot_date=snapshot_date)
    for section_number, slide_def in enumerate(shell.get("slides", []), start=1):
        add_shell_slide(
            prs,
            slide_def,
            region_name=region_name,
            snapshot_date=snapshot_date,
            section_number=section_number,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return {
        "deck_path": str(output_path),
        "slide_count": len(prs.slides),
        "region_name": region_name,
        "component_books": region_component_books(region_name),
        "master_template_path": str(master_template_path),
        "shell_contract_path": str(shell_contract_path),
        "template_version": shell.get("template_version"),
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--region-name", choices=("APAC", "EMEA", "North America"), required=True)
    parser.add_argument("--snapshot-date", required=True)
    parser.add_argument("--output-path", type=Path, required=True)
    parser.add_argument("--master-template-path", type=Path, default=DEFAULT_MASTER_TEMPLATE_PATH)
    parser.add_argument("--shell-contract-path", type=Path, default=DEFAULT_SHELL_CONTRACT_PATH)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    result = build_shell_deck(
        region_name=args.region_name,
        snapshot_date=args.snapshot_date,
        output_path=args.output_path,
        master_template_path=args.master_template_path,
        shell_contract_path=args.shell_contract_path,
    )
    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
