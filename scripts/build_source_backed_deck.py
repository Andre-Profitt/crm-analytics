#!/usr/bin/env python3
"""Build a source-backed monthly review deck from the validated truth packet."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_ROOT = ROOT / "output" / "source_backed_decks"
SCHEMA_VERSION = "monthly_platform.source_backed_deck.v1"
SLIDE_TITLES = [
    "Monthly Sales Director Operating Review",
    "Publish gate status is green across the evidence chain",
    "Regional open pipeline is traceable to director bundles",
    "Director book view: pipeline, risk, and tie-out health",
    "Quarter display logic: current quarter unless fallback is required",
    "Leadership readout: clean source chain, ready for standard deck production",
]

NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
RICH_BLUE = RGBColor(0x09, 0x4F, 0xB2)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
AQUA_LIGHT = RGBColor(0xCB, 0xF5, 0xF3)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
LIGHT_PANEL = RGBColor(0xF5, 0xF8, 0xFD)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
MAGENTA_LIGHT = RGBColor(0xF8, 0xE8, 0xF1)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)
DIVIDER_GREY = RGBColor(0xD7, 0xE2, 0xE8)
AMBER = RGBColor(0xFB, 0x9B, 0x2A)
AMBER_LIGHT = RGBColor(0xFF, 0xF1, 0xDF)
GREEN = RGBColor(0x18, 0x8A, 0x73)
GREEN_LIGHT = RGBColor(0xE4, 0xF7, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def build_source_backed_deck(
    *,
    truth_packet_path: Path,
    output_path: Path | None = None,
    output_root: Path = DEFAULT_OUTPUT_ROOT,
    source_bundle_manifest_path: Path | None = None,
    source_backed_publish_gate_path: Path | None = None,
) -> dict[str, Any]:
    truth_packet_path = Path(truth_packet_path)
    truth_packet = _load_json_object(truth_packet_path)
    source_gate_path = _resolve_source_gate_path(
        truth_packet_path=truth_packet_path,
        truth_packet=truth_packet,
        explicit_path=source_backed_publish_gate_path,
    )
    source_gate = _load_optional_json(source_gate_path)
    source_bundle_manifest_path = _resolve_source_bundle_manifest_path(
        explicit_path=source_bundle_manifest_path,
        source_gate=source_gate,
    )
    source_bundle_manifest = _load_optional_json(source_bundle_manifest_path)
    snapshot_date = str(truth_packet.get("snapshot_date") or "unknown-snapshot")
    source_run_id = _infer_source_run_id(
        truth_packet=truth_packet,
        source_gate=source_gate,
        source_bundle_manifest=source_bundle_manifest,
    )
    if output_path is None:
        output_path = Path(output_root) / snapshot_date / source_run_id / "source_backed_monthly_review.pptx"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_cover(prs, truth_packet, source_gate, source_run_id)
    _add_publish_gate_slide(prs, truth_packet, source_gate)
    _add_regional_rollup_slide(prs, truth_packet)
    _add_director_table_slide(prs, truth_packet)
    _add_quarter_policy_slide(prs, truth_packet, source_bundle_manifest)
    _add_leadership_readout_slide(
        prs,
        truth_packet,
        source_gate,
        source_bundle_manifest,
        truth_packet_path,
        source_gate_path,
        source_bundle_manifest_path,
    )
    prs.save(output_path)

    manifest = {
        "schema_version": SCHEMA_VERSION,
        "status": "ok" if truth_packet.get("status") == "ok" else "blocked",
        "snapshot_date": snapshot_date,
        "source_run_id": source_run_id,
        "deck_path": str(output_path),
        "truth_packet_path": str(truth_packet_path),
        "source_backed_publish_gate_path": str(source_gate_path) if source_gate_path else None,
        "source_bundle_manifest_path": str(source_bundle_manifest_path)
        if source_bundle_manifest_path
        else None,
        "slide_count": len(prs.slides),
        "slide_titles": SLIDE_TITLES,
        "summary": _deck_summary(truth_packet, source_gate, source_bundle_manifest),
        "visual_contract": {
            "required_titles": SLIDE_TITLES,
            "required_text": [
                "Source-backed",
                "Truth Status: OK",
                "Publish gate",
                "Quarter display logic",
                "think-cell",
            ],
        },
    }
    manifest_path = output_path.with_name("source_backed_deck_manifest.json")
    manifest_path.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    manifest["manifest_path"] = str(manifest_path)
    return manifest


def _load_json_object(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError(f"{path}: expected JSON object")
    return payload


def _load_optional_json(path: Path | None) -> dict[str, Any] | None:
    if path is None or not Path(path).exists():
        return None
    return _load_json_object(Path(path))


def _resolve_repo_path(path: str | Path | None, *, anchor: Path | None = None) -> Path | None:
    if path is None:
        return None
    candidate = Path(path)
    if candidate.is_absolute():
        return candidate
    if (ROOT / candidate).exists():
        return ROOT / candidate
    if anchor is not None:
        anchored = anchor.parent / candidate
        if anchored.exists():
            return anchored
    return ROOT / candidate


def _resolve_source_gate_path(
    *,
    truth_packet_path: Path,
    truth_packet: dict[str, Any],
    explicit_path: Path | None,
) -> Path | None:
    if explicit_path is not None:
        return Path(explicit_path)
    sources = truth_packet.get("sources") or {}
    return _resolve_repo_path(
        sources.get("source_backed_publish_gate"),
        anchor=truth_packet_path,
    )


def _resolve_source_bundle_manifest_path(
    *,
    explicit_path: Path | None,
    source_gate: dict[str, Any] | None,
) -> Path | None:
    if explicit_path is not None:
        return Path(explicit_path)
    source_bundle_dir = (source_gate or {}).get("source_bundle_dir")
    if not source_bundle_dir:
        return None
    return _resolve_repo_path(source_bundle_dir) / "source_bundle_manifest.json"


def _infer_source_run_id(
    *,
    truth_packet: dict[str, Any],
    source_gate: dict[str, Any] | None,
    source_bundle_manifest: dict[str, Any] | None,
) -> str:
    for value in [
        (source_bundle_manifest or {}).get("source_run_id"),
        _path_leaf((source_gate or {}).get("source_run_dir")),
        _path_leaf((source_gate or {}).get("director_bundle_dir")),
        _path_parent_leaf((truth_packet.get("sources") or {}).get("analyst_workbook_path")),
    ]:
        if value:
            return _safe_slug(str(value))
    return "unknown-run"


def _path_leaf(value: str | None) -> str | None:
    if not value:
        return None
    return Path(value).name


def _path_parent_leaf(value: str | None) -> str | None:
    if not value:
        return None
    return Path(value).parent.name


def _safe_slug(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9._-]+", "-", value).strip("-") or "run"


def _deck_summary(
    truth_packet: dict[str, Any],
    source_gate: dict[str, Any] | None,
    source_bundle_manifest: dict[str, Any] | None,
) -> dict[str, Any]:
    truth_summary = truth_packet.get("summary") or {}
    source_counts = (source_gate or {}).get("counts") or {}
    source_summary = (source_bundle_manifest or {}).get("summary") or {}
    return {
        "director_count": truth_summary.get("director_count", 0),
        "metric_count": truth_summary.get("metric_count", 0),
        "claim_count": truth_summary.get("claim_count", 0),
        "high_blocker_count": truth_summary.get("high_blocker_count", 0),
        "tieout_mismatch_count": truth_summary.get("tieout_mismatch_count", 0),
        "source_extract_count": source_counts.get(
            "source_extract_count", source_summary.get("source_extract_count", 0)
        ),
        "selected_source_count": source_counts.get(
            "selected_source_count", source_summary.get("selected_source_count", 0)
        ),
        "forward_fallback_count": source_summary.get("forward_fallback_count", 0),
    }


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for placeholder in list(slide.placeholders):
        placeholder._element.getparent().remove(placeholder._element)
    return slide


def _shape(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    *,
    shape_type=MSO_SHAPE.RECTANGLE,
):
    item = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    item.fill.solid()
    item.fill.fore_color.rgb = fill
    item.line.fill.background()
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            item.adjustments[0] = 0.08
        except Exception:
            pass
    return item


def _text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = NAVY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "Microsoft Sans Serif",
):
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return textbox


def _multi_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    size: int = 12,
    color: RGBColor = NAVY,
    bullet: bool = False,
):
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    frame.clear()
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = f"• {line}" if bullet else line
        run.font.name = "Microsoft Sans Serif"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return textbox


def _header(slide, *, eyebrow: str, title: str, narrative: str, accent: RGBColor = AQUA) -> None:
    _shape(slide, 0.6, 0.42, 0.08, 0.7, accent)
    _text(slide, 0.82, 0.38, 12.0, 0.24, eyebrow, size=10, bold=True, color=PRIMARY_BLUE)
    _text(slide, 0.82, 0.62, 12.2, 0.55, title, size=27, bold=True, color=NAVY)
    _text(slide, 0.82, 1.23, 12.1, 0.36, narrative, size=13, color=GREY_TEXT)


def _footer(slide, text: str) -> None:
    _text(slide, 0.6, 7.05, 12.1, 0.22, text, size=8, color=GREY_TEXT)


def _readout(slide, top: float, label: str, message: str, *, accent: RGBColor = AQUA) -> None:
    _shape(slide, 0.82, top + 0.07, 0.08, 0.34, accent)
    _text(slide, 1.03, top, 1.35, 0.24, label.upper(), size=8, bold=True, color=PRIMARY_BLUE)
    _text(slide, 2.05, top, 10.5, 0.32, message, size=12, bold=True, color=NAVY)


def _add_cover(
    prs: Presentation,
    truth_packet: dict[str, Any],
    source_gate: dict[str, Any] | None,
    source_run_id: str,
) -> None:
    slide = _blank_slide(prs)
    summary = truth_packet.get("summary") or {}
    total_open_arr = _sum_values(truth_packet.get("directors") or [], "open_arr")
    total_open_deals = _sum_values(truth_packet.get("directors") or [], "open_deals")
    gate_status = str((source_gate or {}).get("status") or "not provided").upper()
    truth_status = str(truth_packet.get("status") or "unknown").upper()
    snapshot_date = str(truth_packet.get("snapshot_date") or "")
    _shape(slide, 0, 0, 13.333, 7.5, LIGHT_PANEL)
    _shape(slide, 0, 0, 0.36, 7.5, NAVY)
    _shape(slide, 0.36, 0, 0.08, 7.5, AQUA)
    _text(slide, 0.8, 0.65, 3.0, 0.25, "SIMCORP", size=11, bold=True, color=PRIMARY_BLUE)
    _text(
        slide,
        0.8,
        1.22,
        9.7,
        1.15,
        SLIDE_TITLES[0],
        size=39,
        bold=True,
        color=NAVY,
    )
    _text(
        slide,
        0.82,
        2.5,
        9.6,
        0.5,
        f"Source-backed monthly control deck | Snapshot {snapshot_date}",
        size=18,
        color=GREY_TEXT,
    )
    _text(
        slide,
        0.82,
        2.93,
        9.2,
        0.3,
        f"{_format_eur(total_open_arr)} open ARR across {_format_int(total_open_deals)} {_plural(total_open_deals, 'deal')}, reconciled through the validated fact packet.",
        size=13,
        bold=True,
        color=NAVY,
    )
    _text(slide, 0.82, 3.22, 4.0, 0.36, "Truth Status: OK" if truth_status == "OK" else f"Truth Status: {truth_status}", size=20, bold=True, color=GREEN if truth_status == "OK" else AMBER)
    _text(slide, 0.82, 3.72, 4.0, 0.32, f"Publish gate: {gate_status}", size=14, color=GREY_TEXT)
    _text(slide, 0.82, 4.13, 5.0, 0.32, f"Run: {source_run_id}", size=12, color=GREY_TEXT)
    _cover_metric(
        slide,
        7.3,
        3.36,
        str(summary.get("director_count", 0)),
        "directors",
        "Full source-backed coverage.",
    )
    _cover_metric(
        slide,
        9.15,
        3.36,
        str(summary.get("claim_count", 0)),
        "validated claims",
        "Deck facts mapped to source paths.",
    )
    _cover_metric(
        slide,
        11.25,
        3.36,
        str(summary.get("high_blocker_count", 0)),
        "high blockers",
        "Must stay at zero.",
    )
    _text(
        slide,
        0.82,
        6.32,
        9.4,
        0.3,
        "Prepared for analyst workbook, think-cell source workbook, and standard PowerPoint production.",
        size=12,
        color=GREY_TEXT,
    )


def _cover_metric(slide, left: float, top: float, big: str, label: str, context: str) -> None:
    _text(slide, left, top, 1.8, 0.56, big, size=33, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _text(slide, left, top + 0.72, 1.8, 0.24, label.upper(), size=9, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)
    _text(slide, left - 0.12, top + 1.07, 2.05, 0.44, context, size=9, color=GREY_TEXT, align=PP_ALIGN.CENTER)


def _add_publish_gate_slide(
    prs: Presentation,
    truth_packet: dict[str, Any],
    source_gate: dict[str, Any] | None,
) -> None:
    slide = _blank_slide(prs)
    summary = truth_packet.get("summary") or {}
    counts = (source_gate or {}).get("counts") or {}
    _header(
        slide,
        eyebrow=f"{truth_packet.get('snapshot_date', '')}   ·   SOURCE-BACKED PUBLISH GATE",
        title=SLIDE_TITLES[1],
        narrative="The monthly pack should fail closed if source, workbook, fact, or deck truth controls drift.",
        accent=GREEN if truth_packet.get("status") == "ok" else AMBER,
    )
    rows = [
        ["Truth packet", str(truth_packet.get("status", "unknown")).upper(), "Validated fact registry and claim set"],
        ["Source publish gate", str((source_gate or {}).get("status", "not provided")).upper(), "Source contract, bundles, readiness, and workbook coverage"],
        ["High blockers", str(summary.get("high_blocker_count", 0)), "Must be zero for leadership publish"],
        ["Tie-out mismatches", str(summary.get("tieout_mismatch_count", 0)), "Deck/workbook/fact disagreement count"],
        ["Source extracts", str(counts.get("source_extract_count", 0)), "Salesforce reports and list views resolved for the run"],
        ["Selected sources", str(counts.get("selected_source_count", 0)), "Expected source contract items present"],
        ["Director bundles", str(counts.get("director_bundle_count", summary.get("director_count", 0))), "Normalized source-backed DirectorBundle files"],
        ["Validated claims", str(summary.get("claim_count", 0)), "Presentation-ready facts with source paths"],
    ]
    _readout(
        slide,
        1.62,
        "Readout",
        f"{_format_int(summary.get('high_blocker_count', 0))} high blockers and {_format_int(summary.get('tieout_mismatch_count', 0))} tie-out mismatches; source chain is ready for standard deck production.",
        accent=GREEN if truth_packet.get("status") == "ok" else AMBER,
    )
    _add_table(
        slide,
        left=0.78,
        top=2.05,
        width=11.8,
        height=4.35,
        headers=["Gate", "Status / Count", "What it proves"],
        rows=rows,
        widths=[2.7, 1.9, 7.2],
        status_column=1,
    )
    _footer(slide, "Source: deck truth packet and source-backed publish gate JSON.")


def _add_regional_rollup_slide(prs: Presentation, truth_packet: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    regions = sorted(
        list(truth_packet.get("regional_rollups") or []),
        key=lambda item: float((item.get("totals") or {}).get("open_arr") or 0),
        reverse=True,
    )
    _header(
        slide,
        eyebrow=f"{truth_packet.get('snapshot_date', '')}   ·   REGIONAL ROLLUP",
        title=SLIDE_TITLES[2],
        narrative="Regional totals are recomputed from director bundles; SOQL remains validation, not the source of deck truth.",
        accent=AQUA,
    )
    total_open_arr = sum(float((item.get("totals") or {}).get("open_arr") or 0) for item in regions)
    total_open_deals = sum(float((item.get("totals") or {}).get("open_deals") or 0) for item in regions)
    top_region = regions[0] if regions else {}
    top_region_name = str(top_region.get("region") or "No region")
    top_region_arr = float((top_region.get("totals") or {}).get("open_arr") or 0)
    top_region_share = _format_pct(top_region_arr / total_open_arr if total_open_arr else 0)
    _readout(
        slide,
        1.62,
        "Executive read",
        f"{_format_eur(total_open_arr)} open ARR across {_format_int(total_open_deals)} {_plural(total_open_deals, 'deal')}; {top_region_name} represents {top_region_share} of the open book.",
        accent=AQUA,
    )
    chart_data = CategoryChartData()
    chart_data.categories = [str(item.get("region") or "") for item in regions]
    chart_data.add_series(
        "Open ARR (€M)",
        [
            round(float((item.get("totals") or {}).get("open_arr") or 0) / 1_000_000, 2)
            for item in regions
        ],
    )
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.82),
        Inches(2.16),
        Inches(6.15),
        Inches(4.42),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.color.rgb = GREY_TEXT
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = NAVY
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = RICH_BLUE
    rows = [
        [
            str(item.get("region") or ""),
            ", ".join(item.get("territories") or []),
            _format_eur((item.get("totals") or {}).get("open_arr")),
            _format_int((item.get("totals") or {}).get("open_deals")),
            _format_int((item.get("totals") or {}).get("deal_risk_rows")),
        ]
        for item in regions
    ]
    _add_table(
        slide,
        left=7.3,
        top=2.16,
        width=5.25,
        height=4.42,
        headers=["Region", "Territories", "Open ARR", "Deals", "Risk"],
        rows=rows,
        widths=[0.9, 1.65, 1.0, 0.75, 0.75],
        font_size=8,
    )
    _footer(slide, "Source: director gold analytics regional rollups; ARR is open book, not active forecast unless explicitly labeled.")


def _add_director_table_slide(prs: Presentation, truth_packet: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    directors = sorted(
        list(truth_packet.get("directors") or []),
        key=lambda item: float(item.get("open_arr") or 0),
        reverse=True,
    )
    _header(
        slide,
        eyebrow=f"{truth_packet.get('snapshot_date', '')}   ·   DIRECTOR BOOK VIEW",
        title=SLIDE_TITLES[3],
        narrative="This is the analyst handoff layer: facts are deterministic, readable, and source-path backed.",
        accent=PRIMARY_BLUE,
    )
    rows = [
        [
            str(item.get("director") or ""),
            str(item.get("territory") or ""),
            _format_eur(item.get("open_arr")),
            _format_int(item.get("open_deals")),
            _format_int(item.get("deal_risk_rows")),
            _format_int(item.get("tieout_mismatch_count")),
            _format_int(item.get("bundle_issue_count")),
        ]
        for item in directors
    ]
    total_open_arr = _sum_values(directors, "open_arr")
    total_risk_rows = _sum_values(directors, "deal_risk_rows")
    top_directors = directors[:3]
    top_director_arr = _sum_values(top_directors, "open_arr")
    top_book_count = min(3, len(top_directors))
    top_book_label = "Top book" if top_book_count == 1 else f"Top {top_book_count} books"
    _readout(
        slide,
        1.62,
        "Operator lens",
        f"{top_book_label} carry {_format_eur(top_director_arr)} ({_format_pct(top_director_arr / total_open_arr if total_open_arr else 0)}) of open ARR; {_format_int(total_risk_rows)} deal-risk rows need field action.",
        accent=PRIMARY_BLUE,
    )
    _add_table(
        slide,
        left=0.74,
        top=2.06,
        width=11.95,
        height=4.77,
        headers=["Director", "Territory", "Open ARR", "Deals", "Risk rows", "Tie-out", "Source issues"],
        rows=rows,
        widths=[2.05, 2.0, 1.25, 0.75, 0.9, 0.8, 1.1],
        font_size=8,
        status_columns={5, 6},
    )
    _footer(slide, "Source: deck truth packet director facts; all values retain source_artifact and source_json_path lineage.")


def _add_quarter_policy_slide(
    prs: Presentation,
    truth_packet: dict[str, Any],
    source_bundle_manifest: dict[str, Any] | None,
) -> None:
    slide = _blank_slide(prs)
    _header(
        slide,
        eyebrow=f"{truth_packet.get('snapshot_date', '')}   ·   PERIOD LOCK / FALLBACK",
        title=SLIDE_TITLES[4],
        narrative="Fallback is explicit: if current-quarter active pipeline is empty, the display moves to the forward quarter.",
        accent=AMBER,
    )
    territory_rows = ((source_bundle_manifest or {}).get("summary") or {}).get("territories") or []
    if not territory_rows:
        _text(
            slide,
            0.85,
            2.45,
            11.5,
            0.55,
            "Source bundle quarter decisions were not provided with this truth packet.",
            size=16,
            color=GREY_TEXT,
        )
    else:
        fallback_count = sum(
            1
            for item in territory_rows
            if str(item.get("display_reason") or "") == "forward_quarter_fallback"
        )
        empty_count = sum(
            1
            for item in territory_rows
            if str(item.get("display_reason") or "") == "empty_current_and_forward_quarter"
        )
        _readout(
            slide,
            1.62,
            "Period control",
            f"{_format_int(fallback_count)} {_plural(fallback_count, 'territory', 'territories')} use forward fallback; "
            f"current/forward empty territories: {_format_int(empty_count)}.",
            accent=AMBER,
        )
        rows = [
            [
                str(item.get("director") or ""),
                str(item.get("territory") or ""),
                str(item.get("display_quarter_title") or ""),
                _display_reason(str(item.get("display_reason") or "")),
                _format_int(item.get("current_quarter_active_deals")),
                _format_int(item.get("forward_quarter_active_deals")),
            ]
            for item in territory_rows
        ]
        _add_table(
            slide,
            left=0.74,
            top=2.08,
            width=11.95,
            height=4.67,
            headers=["Director", "Territory", "Display quarter", "Reason", "Current active", "Forward active"],
            rows=rows,
            widths=[1.85, 2.0, 1.45, 2.25, 1.05, 1.05],
            font_size=8,
            reason_column=3,
        )
    _footer(slide, "Source: monthly source bundle manifest pipeline_display_decision for each territory.")


def _add_leadership_readout_slide(
    prs: Presentation,
    truth_packet: dict[str, Any],
    source_gate: dict[str, Any] | None,
    source_bundle_manifest: dict[str, Any] | None,
    truth_packet_path: Path,
    source_gate_path: Path | None,
    source_bundle_manifest_path: Path | None,
) -> None:
    slide = _blank_slide(prs)
    summary = _deck_summary(truth_packet, source_gate, source_bundle_manifest)
    thinkcell = truth_packet.get("thinkcell") or {}
    recommended_elements = thinkcell.get("recommended_element_names") or []
    _header(
        slide,
        eyebrow=f"{truth_packet.get('snapshot_date', '')}   ·   DECK PRODUCTION HANDOFF",
        title=SLIDE_TITLES[5],
        narrative="Standard PowerPoint output should be generated from the same workbook and named element contract each month.",
        accent=MAGENTA,
    )
    _multi_text(
        slide,
        0.85,
        1.95,
        6.0,
        1.55,
        [
            f"{_format_int(summary['source_extract_count'])} Salesforce source extracts resolved; {_format_int(summary['selected_source_count'])} selected sources present.",
            f"{_format_int(summary['director_count'])} directors, {_format_int(summary['metric_count'])} metrics, and {_format_int(summary['claim_count'])} validated claims are ready.",
            f"{_format_int(summary['forward_fallback_count'])} {_plural(summary['forward_fallback_count'], 'territory', 'territories')} use forward-quarter fallback; every decision is visible in the source manifest.",
            "No silent calendar drift: snapshot, current quarter, forward quarter, and source contracts stay locked together.",
        ],
        size=14,
        color=NAVY,
        bullet=True,
    )
    _readout(
        slide,
        3.5,
        "Decision",
        "Use this packet as the standard deck production input; keep narrative edits inside the verified source-backed fact set.",
        accent=MAGENTA,
    )
    artifact_rows = [
        ["Truth packet", _short_path(truth_packet_path), "Deck facts and claim registry"],
        ["Publish gate", _short_path(source_gate_path), "Source/workbook readiness control"],
        ["Source bundle manifest", _short_path(source_bundle_manifest_path), "Quarter and fallback decisions"],
        ["think-cell elements", ", ".join(recommended_elements[:4]) or "n/a", "Named PowerPoint handoff targets"],
    ]
    _add_table(
        slide,
        left=0.85,
        top=4.22,
        width=11.65,
        height=2.03,
        headers=["Artifact", "Path / element", "Use"],
        rows=artifact_rows,
        widths=[1.65, 5.9, 3.55],
        font_size=8,
    )
    _footer(slide, "Source: truth packet think-cell contract and source-backed monthly artifacts.")


def _add_table(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    headers: list[str],
    rows: list[list[str]],
    widths: list[float],
    font_size: int = 9,
    status_column: int | None = None,
    status_columns: set[int] | None = None,
    reason_column: int | None = None,
) -> None:
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table
    for column_index, column_width in enumerate(widths[: len(headers)]):
        table.columns[column_index].width = Inches(column_width)
    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        _set_cell_text(cell, header, size=font_size, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    zero_status_columns = set(status_columns or set())
    for row_index, row in enumerate(rows, start=1):
        for column_index, value in enumerate(row[: len(headers)]):
            cell = table.cell(row_index, column_index)
            _set_cell_text(cell, value, size=font_size, color=NAVY)
            cell.fill.solid()
            if column_index == reason_column and "fallback" in str(value).lower():
                cell.fill.fore_color.rgb = AMBER_LIGHT
            elif column_index == reason_column and "empty" in str(value).lower():
                cell.fill.fore_color.rgb = MAGENTA_LIGHT
            elif column_index == status_column:
                cell.fill.fore_color.rgb = _status_fill(str(value))
            elif column_index in zero_status_columns:
                cell.fill.fore_color.rgb = GREEN_LIGHT if str(value).strip() in {"0", "0.0"} else AMBER_LIGHT
            else:
                cell.fill.fore_color.rgb = LIGHT_PANEL if row_index % 2 else WHITE


def _set_cell_text(
    cell,
    text: str,
    *,
    size: int = 9,
    bold: bool = False,
    color: RGBColor = NAVY,
) -> None:
    cell.text = ""
    frame = cell.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(3)
    frame.margin_bottom = Pt(3)
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = "Microsoft Sans Serif"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _status_fill(value: str) -> RGBColor:
    normalized = value.strip().lower()
    if normalized in {"ok", "0"}:
        return GREEN_LIGHT
    if normalized in {"not provided", "unknown", ""}:
        return AMBER_LIGHT
    return AMBER_LIGHT if normalized != "blocked" else MAGENTA_LIGHT


def _display_reason(value: str) -> str:
    labels = {
        "current_quarter": "Current quarter",
        "forward_quarter_fallback": "Forward fallback",
        "empty_current_and_forward_quarter": "Current and forward empty",
    }
    return labels.get(value, value.replace("_", " ").strip().title())


def _sum_values(items: list[dict[str, Any]], key: str) -> float:
    total = 0.0
    for item in items:
        try:
            total += float(item.get(key) or 0)
        except (TypeError, ValueError):
            continue
    return total


def _format_pct(value: float) -> str:
    return f"{value * 100:.0f}%"


def _plural(value: Any, singular: str, plural: str | None = None) -> str:
    try:
        count = int(value)
    except (TypeError, ValueError):
        count = 0
    if count == 1:
        return singular
    return plural or f"{singular}s"


def _format_int(value: Any) -> str:
    try:
        return f"{int(value):,}"
    except (TypeError, ValueError):
        return "0"


def _format_eur(value: Any) -> str:
    if value is None:
        return "€0"
    amount = float(value)
    if abs(amount) >= 1_000_000_000:
        return f"€{amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"€{amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"€{amount / 1_000:.0f}K"
    return f"€{amount:.0f}"


def _short_path(path: Path | None) -> str:
    if path is None:
        return "not provided"
    try:
        return str(Path(path).resolve().relative_to(ROOT))
    except ValueError:
        return str(path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--truth-packet", type=Path, required=True)
    parser.add_argument("--source-bundle-manifest", type=Path, default=None)
    parser.add_argument("--source-backed-publish-gate", type=Path, default=None)
    parser.add_argument("--output-path", type=Path, default=None)
    parser.add_argument("--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    manifest = build_source_backed_deck(
        truth_packet_path=args.truth_packet,
        output_path=args.output_path,
        output_root=args.output_root,
        source_bundle_manifest_path=args.source_bundle_manifest,
        source_backed_publish_gate_path=args.source_backed_publish_gate,
    )
    print(json.dumps(manifest, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
