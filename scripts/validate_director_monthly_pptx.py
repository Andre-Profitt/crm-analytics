#!/usr/bin/env python3
"""Track E — produced-PPTX checker (E4).

Reads a generated Sales Director Monthly deck (.pptx) and validates
it against the active director_monthly profile in
``config/deck_contract.yaml``. Read-only: never modifies the PPTX,
never modifies the builder.

Checks:
  - exactly 18 slides (or whatever profile.expected_slide_count is)
  - per-slide table count matches the contract (evidence_only tables
    are excluded from the expectation)
  - per-slide title is either:
      a) exact match to the stable contract title, OR
      b) matches one of the slide's legacy_title_patterns regexes
    Mismatches in (a) but matches in (b) -> warning, not blocker
    (legacy-tolerant mode for current production output that still
     emits dynamic verbose titles)
  - legal_notice slide is the last slide
  - any slide with required_links -> at least one shape has a real
    hyperlink whose target matches the link's kind hint
  - cover slide has at least one shape with the contract title text

Output:
  - pptx_contract_report.json (and optional .md) summarising slide-
    by-slide title status, table-count parity, link presence, and
    legal-notice placement.

Usage:
    python scripts/validate_director_monthly_pptx.py \\
        --pptx /Users/test/Downloads/jesper-tyrer-LAND.pptx \\
        --report-out output/track_e/pptx_contract_report.json \\
        --md-out output/track_e/pptx_contract_report.md
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

REPO_ROOT = Path(__file__).resolve().parent.parent
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from pptx import Presentation  # noqa: E402

from scripts.monthly_platform import deck_contract  # noqa: E402


REPORT_SCHEMA_VERSION = "monthly_platform.pptx_contract_report.v1"


def _slide_title(slide) -> str:
    """Extract the slide's title text.

    Order of attempts:
      1. The placeholder whose ``placeholder_format.idx`` is 144 (the
         template-anchored title placeholder used by build_deck_from_excel.py).
      2. The placeholder whose ``placeholder_format.type`` is TITLE
         (PowerPoint's standard title placeholder type).
      3. The first non-empty text frame in shape order — fallback for
         slides without a title placeholder (e.g. legal/cover slides
         that draw their headline as a free textbox).
    """
    title_idx_match: str | None = None
    title_type_match: str | None = None
    first_frame: str | None = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        first_line = text.split("\n")[0].strip()

        try:
            ph_fmt = shape.placeholder_format
        except (ValueError, AttributeError):
            ph_fmt = None

        if ph_fmt is not None:
            try:
                if ph_fmt.idx == 144 and title_idx_match is None:
                    title_idx_match = first_line
            except (ValueError, AttributeError):
                pass
            try:
                # PowerPoint's PP_PLACEHOLDER.TITLE = 13 (and CTR_TITLE = 14)
                if (
                    ph_fmt.type is not None
                    and ph_fmt.type in (13, 14)
                    and title_type_match is None
                ):
                    title_type_match = first_line
            except (ValueError, AttributeError):
                pass

        if first_frame is None:
            first_frame = first_line

    return title_idx_match or title_type_match or first_frame or ""


def _slide_table_count(slide) -> int:
    return sum(1 for shape in slide.shapes if shape.has_table)


def _slide_has_hyperlink(slide) -> bool:
    """Return True if any run on the slide has a non-empty hyperlink."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                hl = getattr(run, "hyperlink", None)
                if hl is not None and getattr(hl, "address", None):
                    return True
    return False


# Salesforce list-view URL must point at the Lightning Opportunity list.
# Accepts either domain form: simcorp.lightning.force.com (Lightning) or
# simcorp.my.salesforce.com (My Domain) — both resolve to the same Lightning UI.
SALESFORCE_LIST_VIEW_RE = re.compile(
    r"^https?://[^/]*simcorp(\.lightning\.force\.com|\.my\.salesforce\.com)"
    r"(/lightning)?/o/Opportunity/list",
    re.IGNORECASE,
)


def _slide_link_addresses(slide) -> list[str]:
    """All non-empty hyperlink addresses on the slide."""
    out: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                hl = getattr(run, "hyperlink", None)
                addr = getattr(hl, "address", None) if hl is not None else None
                if addr:
                    out.append(str(addr))
    return out


def _link_satisfies_kind(address: str, kind: str | None) -> bool:
    if kind == "salesforce_list_view":
        return bool(SALESFORCE_LIST_VIEW_RE.search(address or ""))
    if kind == "salesforce_report":
        return bool(
            re.search(
                r"simcorp(\.lightning\.force\.com|\.my\.salesforce\.com)",
                address or "",
                re.IGNORECASE,
            )
        )
    if kind == "sharepoint":
        return "sharepoint.com" in (address or "").lower()
    if kind == "external_url":
        return (address or "").startswith(("http://", "https://"))
    # Unknown kind — accept any non-empty address.
    return bool(address)


def _slide_table_objects(slide) -> list:
    return [shape.table for shape in slide.shapes if shape.has_table]


def _table_header_row(table) -> list[str]:
    """Extract first-row cell text, normalised (strip + collapse newlines)."""
    return [
        " / ".join(
            part.strip() for part in table.cell(0, c).text.split("\n") if part.strip()
        )
        for c in range(len(table.columns))
    ]


def _validate_table_headers(
    *,
    slide_decl: dict[str, Any],
    actual_slide,
    findings: list[dict[str, Any]],
) -> tuple[str, list[dict[str, Any]]]:
    """Compare per-table actual headers to contract columns[].header.

    Returns (overall_status, per_table_results). overall_status is one of:
      - "pass"               every displayed table matches stable headers
      - "warning"            at least one table only matches a legacy_header_sets entry
      - "fail"               at least one table matches neither stable nor legacy
      - "n/a"                no displayed (non-evidence_only) tables on this slide

    Matches by slide order (per GPT spec).
    """
    expected_tables = [
        t for t in (slide_decl.get("tables") or []) if not t.get("evidence_only")
    ]
    actual_tables = _slide_table_objects(actual_slide)

    if not expected_tables:
        return "n/a", []

    overall = "pass"
    per_table: list[dict[str, Any]] = []
    sid = slide_decl.get("id")
    snum = slide_decl.get("slide_number")

    for idx, decl in enumerate(expected_tables):
        if idx >= len(actual_tables):
            findings.append(
                {
                    "severity": "blocker",
                    "code": "table_missing_in_pptx",
                    "path": f"slides[{snum}].tables[{decl.get('id')}]",
                    "message": (
                        f"slide {snum} ({sid}) declares table {decl.get('id')!r} "
                        f"at order {idx} but the produced deck has only "
                        f"{len(actual_tables)} table(s)"
                    ),
                }
            )
            per_table.append(
                {
                    "table_id": decl.get("id"),
                    "order": idx,
                    "status": "fail",
                    "detail": "table missing in produced deck",
                }
            )
            overall = "fail"
            continue

        stable_headers = [str(c.get("header", "")) for c in (decl.get("columns") or [])]
        actual_headers = _table_header_row(actual_tables[idx])

        if actual_headers == stable_headers:
            per_table.append(
                {
                    "table_id": decl.get("id"),
                    "order": idx,
                    "status": "pass_stable",
                    "actual_headers": actual_headers,
                }
            )
            continue

        # header_pattern_sets — canonical regex tuples for tables whose
        # headers are dynamic-by-design (e.g. snapshot-date columns). A match
        # is treated as pass_stable (the canonical contract form), NOT a
        # warning. Checked before legacy_header_sets so canonical-pattern
        # tables don't generate transition warnings.
        pattern_sets = decl.get("header_pattern_sets") or []
        matched_pattern: list[str] | None = None
        for header_set in pattern_sets:
            if len(header_set) != len(actual_headers):
                continue
            try:
                if all(re.fullmatch(p, h) for p, h in zip(header_set, actual_headers)):
                    matched_pattern = list(header_set)
                    break
            except re.error:
                continue

        if matched_pattern is not None:
            per_table.append(
                {
                    "table_id": decl.get("id"),
                    "order": idx,
                    "status": "pass_pattern",
                    "actual_headers": actual_headers,
                    "matched_pattern": matched_pattern,
                }
            )
            continue

        legacy_sets = decl.get("legacy_header_sets") or []
        matched_legacy: list[str] | None = None
        for header_set in legacy_sets:
            if len(header_set) != len(actual_headers):
                continue
            try:
                if all(re.fullmatch(p, h) for p, h in zip(header_set, actual_headers)):
                    matched_legacy = list(header_set)
                    break
            except re.error:
                # invalid regex inside legacy_header_sets — skip and let the
                # contract validator surface it on the next pass.
                continue

        if matched_legacy is not None:
            findings.append(
                {
                    "severity": "warning",
                    "code": "legacy_header_drift",
                    "path": f"slides[{snum}].tables[{decl.get('id')}].columns",
                    "message": (
                        f"slide {snum} ({sid}) table {decl.get('id')!r} headers "
                        f"match a legacy_header_sets entry rather than stable headers. "
                        f"actual={actual_headers} stable={stable_headers}"
                    ),
                }
            )
            per_table.append(
                {
                    "table_id": decl.get("id"),
                    "order": idx,
                    "status": "warning_legacy",
                    "actual_headers": actual_headers,
                    "stable_headers": stable_headers,
                    "matched_legacy": matched_legacy,
                }
            )
            if overall == "pass":
                overall = "warning"
        else:
            findings.append(
                {
                    "severity": "blocker",
                    "code": "table_header_mismatch",
                    "path": f"slides[{snum}].tables[{decl.get('id')}].columns",
                    "message": (
                        f"slide {snum} ({sid}) table {decl.get('id')!r} headers do "
                        f"not match stable columns nor any legacy_header_sets entry. "
                        f"actual={actual_headers} stable={stable_headers}"
                    ),
                }
            )
            per_table.append(
                {
                    "table_id": decl.get("id"),
                    "order": idx,
                    "status": "fail",
                    "actual_headers": actual_headers,
                    "stable_headers": stable_headers,
                }
            )
            overall = "fail"

    return overall, per_table


def validate_pptx(
    pptx_path: Path,
    *,
    contract: deck_contract.DeckContract | None = None,
) -> dict[str, Any]:
    if contract is None:
        contract = deck_contract.load()
    assert contract is not None

    profile = contract.director_monthly
    expected_count = int(profile.get("expected_slide_count", 18))
    contract_slides = profile.get("slides", [])

    findings: list[dict[str, Any]] = []
    slide_results: list[dict[str, Any]] = []

    prs = Presentation(str(pptx_path))
    actual_count = len(prs.slides)

    if actual_count != expected_count:
        findings.append(
            {
                "severity": "blocker",
                "code": "slide_count_mismatch",
                "path": "slides",
                "message": f"expected {expected_count} slides, found {actual_count}",
            }
        )

    # Build slide-by-slide expectations indexed by slide_number.
    by_number: dict[int, dict[str, Any]] = {}
    for s in contract_slides:
        by_number[int(s["slide_number"])] = s

    legacy_title_count = 0
    title_mismatch_count = 0

    for snum in range(1, min(actual_count, expected_count) + 1):
        actual_slide = prs.slides[snum - 1]
        decl = by_number.get(snum)
        if decl is None:
            findings.append(
                {
                    "severity": "blocker",
                    "code": "unexpected_slide",
                    "path": f"slides[{snum}]",
                    "message": f"PPTX has slide {snum} but contract does not declare it",
                }
            )
            slide_results.append(
                {
                    "slide_number": snum,
                    "slide_id": None,
                    "status": "fail",
                    "title": _slide_title(actual_slide),
                }
            )
            continue

        sid = decl["id"]
        stable_title = str(decl.get("title", ""))
        actual_title = _slide_title(actual_slide)
        legacy_patterns = decl.get("legacy_title_patterns") or []
        title_status = "pass"
        title_detail = ""

        if actual_title == stable_title:
            title_status = "pass_stable"
        elif decl.get("static") or sid == "legal_notice":
            # Static slides may have no programmatic title text.
            title_status = "pass_static"
        else:
            matched_legacy = next(
                (p for p in legacy_patterns if re.match(p, actual_title)), None
            )
            if matched_legacy:
                title_status = "warning_legacy"
                title_detail = f"matched legacy pattern: {matched_legacy!r}"
                legacy_title_count += 1
                findings.append(
                    {
                        "severity": "warning",
                        "code": "legacy_verbose_title",
                        "path": f"slides[{snum}].title",
                        "message": (
                            f"slide {snum} ({sid}) emits legacy verbose title "
                            f"{actual_title!r}; matches legacy_title_patterns. "
                            f"Update builder to emit stable title {stable_title!r}."
                        ),
                    }
                )
            else:
                title_status = "fail"
                title_detail = (
                    f"actual={actual_title!r} not equal to stable "
                    f"{stable_title!r} and no legacy_title_patterns match"
                )
                title_mismatch_count += 1
                findings.append(
                    {
                        "severity": "blocker",
                        "code": "title_neither_stable_nor_legacy",
                        "path": f"slides[{snum}].title",
                        "message": title_detail,
                    }
                )

        # Table count parity (excluding evidence_only).
        expected_tables = sum(
            1 for t in (decl.get("tables") or []) if not t.get("evidence_only")
        )
        actual_tables = _slide_table_count(actual_slide)
        table_status = "pass"
        table_detail = ""
        if actual_tables != expected_tables:
            table_status = "fail"
            table_detail = f"declared={expected_tables} actual={actual_tables}"
            findings.append(
                {
                    "severity": "blocker",
                    "code": "table_count_mismatch",
                    "path": f"slides[{snum}].tables",
                    "message": f"slide {snum} ({sid}) {table_detail}",
                }
            )

        # Per-table header validation (Cond 1). Skips evidence_only
        # tables. Pass if headers match the stable contract OR a
        # legacy_header_sets entry; fail (blocker) only when neither
        # matches. Skipped entirely when table_status == "fail" so we
        # don't emit cascading mismatches against missing tables.
        header_status: str = "n/a"
        header_results: list[dict[str, Any]] = []
        if expected_tables > 0 and table_status != "fail":
            header_status, header_results = _validate_table_headers(
                slide_decl=decl,
                actual_slide=actual_slide,
                findings=findings,
            )

        # Required-links presence (Cond 2 — tightened semantics).
        # For required_links.kind=salesforce_list_view, a hyperlink only
        # satisfies the requirement if the target URL is a Salesforce
        # Opportunity list view. Wrong-target = blocker. Missing link
        # entirely = warning during M1 transition.
        link_status = "n/a"
        link_detail = ""
        required_links = decl.get("required_links") or []
        if required_links:
            addresses = _slide_link_addresses(actual_slide)
            unsatisfied: list[dict[str, Any]] = []
            for link_decl in required_links:
                kind = link_decl.get("kind")
                matching = [a for a in addresses if _link_satisfies_kind(a, kind)]
                if not matching:
                    unsatisfied.append(
                        {
                            "id": link_decl.get("id"),
                            "kind": kind,
                            "label": link_decl.get("label"),
                        }
                    )

            if not unsatisfied:
                link_status = "pass"
            elif not addresses:
                # No hyperlinks at all — M1 transition warning.
                link_status = "warning"
                link_detail = (
                    f"expected {len(required_links)} required link(s), "
                    f"found no hyperlinks at all (M1 transition warning)"
                )
                findings.append(
                    {
                        "severity": "warning",
                        "code": "missing_required_link_transition",
                        "path": f"slides[{snum}].required_links",
                        "message": (
                            f"slide {snum} ({sid}) missing required hyperlink "
                            f"({len(required_links)} declared). Update builder "
                            f"to emit the Salesforce drill-through link."
                        ),
                    }
                )
            else:
                # Hyperlinks exist but none satisfy the declared kinds —
                # this is a wrong-target failure, NOT a transition gap.
                link_status = "fail"
                link_detail = (
                    f"slide has {len(addresses)} hyperlink(s) but none satisfy "
                    f"required kinds: {[u['kind'] for u in unsatisfied]}"
                )
                for u in unsatisfied:
                    findings.append(
                        {
                            "severity": "blocker",
                            "code": "required_link_target_mismatch",
                            "path": f"slides[{snum}].required_links[{u['id']}]",
                            "message": (
                                f"slide {snum} ({sid}) link {u['id']!r} "
                                f"(kind={u['kind']!r}) not satisfied by any "
                                f"hyperlink on the slide. addresses={addresses}"
                            ),
                        }
                    )

        slide_results.append(
            {
                "slide_number": snum,
                "slide_id": sid,
                "stable_title": stable_title,
                "actual_title": actual_title,
                "title_status": title_status,
                "title_detail": title_detail,
                "expected_tables": expected_tables,
                "actual_tables": actual_tables,
                "table_status": table_status,
                "table_detail": table_detail,
                "header_status": header_status,
                "header_results": header_results,
                "link_status": link_status,
                "link_detail": link_detail,
                "status": (
                    "fail"
                    if title_status == "fail"
                    or table_status == "fail"
                    or header_status == "fail"
                    or link_status == "fail"
                    else (
                        "warning"
                        if title_status == "warning_legacy"
                        or header_status == "warning"
                        or link_status == "warning"
                        else "pass"
                    )
                ),
            }
        )

    # Legal-notice placement: must be the last slide.
    legal_decl = next(
        (s for s in contract_slides if s.get("id") == "legal_notice"), None
    )
    if legal_decl and legal_decl["slide_number"] != actual_count:
        findings.append(
            {
                "severity": "warning",
                "code": "legal_notice_not_last",
                "path": "slides[legal_notice]",
                "message": (
                    f"legal_notice slide_number={legal_decl['slide_number']} "
                    f"but actual deck has {actual_count} slides"
                ),
            }
        )

    blockers = [f for f in findings if f["severity"] == "blocker"]
    warnings = [f for f in findings if f["severity"] == "warning"]

    return {
        "schema_version": REPORT_SCHEMA_VERSION,
        "pptx_path": str(pptx_path),
        "deck_contract_path": str(contract.path),
        "validated_at": datetime.now(timezone.utc).isoformat(),
        "status": "pass" if not blockers else "fail",
        "blocker_count": len(blockers),
        "warning_count": len(warnings),
        "expected_slide_count": expected_count,
        "actual_slide_count": actual_count,
        "legacy_verbose_title_count": legacy_title_count,
        "stable_title_count": sum(
            1 for r in slide_results if r["title_status"] == "pass_stable"
        ),
        "title_mismatch_count": title_mismatch_count,
        "slides": slide_results,
        "findings": findings,
    }


def render_markdown(report: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append("# PPTX contract report\n")
    lines.append(f"- pptx: `{report['pptx_path']}`")
    lines.append(f"- deck contract: `{report['deck_contract_path']}`")
    lines.append(f"- validated_at: {report['validated_at']}")
    lines.append(f"- **status: {report['status']}**")
    lines.append(
        f"- slides: {report['actual_slide_count']}/{report['expected_slide_count']}"
    )
    lines.append(
        f"- titles: stable={report['stable_title_count']} "
        f"legacy_verbose={report['legacy_verbose_title_count']} "
        f"mismatch={report['title_mismatch_count']}"
    )
    lines.append(
        f"- blockers: {report['blocker_count']} | warnings: {report['warning_count']}"
    )
    lines.append("")

    lines.append(
        "| # | Slide | Title | Tables (decl/actual) | Headers | Link | Slide status |"
    )
    lines.append("| ---: | --- | --- | --- | --- | --- | --- |")
    for r in report["slides"]:
        lines.append(
            f"| {r['slide_number']} | `{r['slide_id'] or '?'}` | {r['title_status']} | "
            f"{r['expected_tables']}/{r['actual_tables']} | {r.get('header_status', 'n/a')} | "
            f"{r['link_status']} | {r['status']} |"
        )
    lines.append("")

    if report["findings"]:
        lines.append("## Findings\n")
        for f in report["findings"]:
            lines.append(f"- **{f['severity']}** `{f['code']}` — {f['message']}")
        lines.append("")

    return "\n".join(lines) + "\n"


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Validate a produced PPTX against the deck contract."
    )
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--contract", default=None)
    parser.add_argument("--report-out", default=None)
    parser.add_argument("--md-out", default=None)
    parser.add_argument("--show-findings", action="store_true")
    args = parser.parse_args(argv)

    contract = deck_contract.load(args.contract)
    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"ERROR: pptx not found: {pptx_path}", file=sys.stderr)
        return 2

    report = validate_pptx(pptx_path, contract=contract)

    if args.report_out:
        out = Path(args.report_out)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(report, indent=2) + "\n", encoding="utf-8")
        print(f"report: {out}")
    if args.md_out:
        out = Path(args.md_out)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(render_markdown(report), encoding="utf-8")
        print(f"md: {out}")

    if args.show_findings or report["status"] == "fail":
        for f in report["findings"]:
            print(f"[{f['severity']}] {f['code']} {f.get('path', '')}: {f['message']}")

    print(
        f"pptx_contract: {report['status']} "
        f"(blockers={report['blocker_count']} warnings={report['warning_count']} "
        f"stable={report['stable_title_count']} legacy={report['legacy_verbose_title_count']} "
        f"slides={report['actual_slide_count']}/{report['expected_slide_count']})"
    )
    return 0 if report["status"] == "pass" else 1


if __name__ == "__main__":
    raise SystemExit(main())
