#!/usr/bin/env python3
"""Path B example: build a single sample title slide using the real
SC-Master Gradient_Title layout from the production SimCorp master,
filled with the same data as the canonical deck's slide 1.

Lets us compare Path A (transplanted shapes + remap) vs Path B (real
SimCorp layout with filled placeholders) before committing to a full
restructure.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

TEMPLATE = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"
RUN_DIR = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01"
)
SNAPSHOT = RUN_DIR / "report1_snapshot.json"
OUTPUT = RUN_DIR / "sales_director_monthly_path_b_sample.pptx"


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def fill(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return True
    return False


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def main():
    with open(SNAPSHOT) as f:
        snap = json.load(f)

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    cover = get_layout(prs, "SC-Master Gradient_Title")
    slide = prs.slides.add_slide(cover)

    # idx=0 title, idx=20 subheading, idx=22 date/location
    fill(slide, 0, "Sales Director Monthly Pipeline Insights")
    fill(
        slide,
        20,
        f"{snap.get('quarter_focus', 'Q1')} FY27 forward look. Pipeline coverage, commercial approvals, renewals, churn, and slipped deals.",
    )
    fill(slide, 22, f"Snapshot {snap.get('snapshot_date', '2026-04-01')}")

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
