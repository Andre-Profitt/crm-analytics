#!/usr/bin/env python3
"""Option D proof of concept.

Pull a chart step from a live CRM Analytics dashboard via Wave API, run
its SAQL query, and render the result as a NATIVE PowerPoint chart
inside a SimCorp-branded sample slide.

Uses the Pipeline & Opportunity Operations dashboard's s_region_hygiene
step ("Open vs At-Risk ARR by Region") as the example.

Run:
    python3 scripts/simcorp_crma_chart_sample.py
"""

from __future__ import annotations

import html
import json
import re
import subprocess
from pathlib import Path

import requests
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"
DASHBOARD_ID = "0FKTb0000000KwPOAU"  # Pipeline & Opportunity Operations
STEP_NAME = "s_region_hygiene"
DATASET_NAME = "Pipeline_Opportunity_Operations"
OUTPUT = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01/sample_slide_crma_chart.pptx"
)

NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)


def get_auth():
    r = subprocess.run(
        ["sf", "org", "display", "--target-org", "apro@simcorp.com", "--json"],
        capture_output=True,
        text=True,
        check=True,
    )
    d = json.loads(r.stdout[r.stdout.find("{") :])["result"]
    return d["instanceUrl"], d["accessToken"]


def get_dataset_version(inst, tok, dataset_name):
    r = requests.get(
        f"{inst}/services/data/v66.0/wave/datasets?pageSize=200",
        headers={"Authorization": f"Bearer {tok}"},
        timeout=30,
    )
    j = r.json()
    all_ds = j.get("datasets", [])
    while j.get("nextPageUrl"):
        r = requests.get(
            f"{inst}{j['nextPageUrl']}",
            headers={"Authorization": f"Bearer {tok}"},
            timeout=30,
        )
        j = r.json()
        all_ds.extend(j.get("datasets", []))
    for ds in all_ds:
        if ds.get("name") == dataset_name:
            return ds.get("id"), ds.get("currentVersionId")
    raise KeyError(dataset_name)


def get_step_query(inst, tok, dashboard_id, step_name):
    r = requests.get(
        f"{inst}/services/data/v66.0/wave/dashboards/{dashboard_id}",
        headers={"Authorization": f"Bearer {tok}"},
        timeout=30,
    )
    r.raise_for_status()
    steps = r.json().get("state", {}).get("steps", {})
    step = steps.get(step_name, {})
    q = step.get("query")
    if isinstance(q, dict):
        q = q.get("query", "")
    return html.unescape(q or "")


def rewrite_query_with_dataset(saql, dataset_id, version_id):
    """Replace load "DatasetName" with load "datasetId/versionId" and strip
    Mustache filter bindings so the query can run via /wave/query without
    a dashboard context."""
    # Swap the load target to fully qualified form
    saql = re.sub(
        r'q\s*=\s*load\s*"[^"]+"',
        f'q = load "{dataset_id}/{version_id}"',
        saql,
        count=1,
    )
    # Remove filter lines that contain {{...}} Mustache bindings
    out_lines = []
    for line in saql.splitlines():
        if "{{" in line and "}}" in line:
            continue
        out_lines.append(line)
    return "\n".join(out_lines)


def run_saql(inst, tok, saql):
    r = requests.post(
        f"{inst}/services/data/v66.0/wave/query",
        headers={"Authorization": f"Bearer {tok}", "Content-Type": "application/json"},
        json={"query": saql},
        timeout=60,
    )
    r.raise_for_status()
    return r.json().get("results", {}).get("records", [])


def txt(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def rect(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def build():
    inst, tok = get_auth()
    print(f"Auth OK ({inst})")

    ds_id, ver_id = get_dataset_version(inst, tok, DATASET_NAME)
    print(f"Dataset {DATASET_NAME}: id={ds_id}  version={ver_id}")

    saql_raw = get_step_query(inst, tok, DASHBOARD_ID, STEP_NAME)
    saql = rewrite_query_with_dataset(saql_raw, ds_id, ver_id)
    print(f"\nRunnable SAQL:\n{saql}\n")

    records = run_saql(inst, tok, saql)
    records = [
        r for r in records if r.get("SalesRegion") and r.get("WeightedOpenARR", 0) > 0
    ]
    print(f"Rows returned: {len(records)}")
    for r in records:
        print(
            f"  {r['SalesRegion']:<30} open={r.get('WeightedOpenARR', 0):>14,.0f}  "
            f"at_risk={r.get('AtRiskARR', 0):>14,.0f}"
        )

    # ------------------------------------------------------------------
    # Build a single sample slide
    # ------------------------------------------------------------------
    prs = Presentation(TEMPLATE)
    clear_slides(prs)
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # Header (vertical aqua tab + eyebrow + title + narrative)
    rect(slide, 0.6, 0.42, 0.08, 0.7, AQUA)
    txt(
        slide,
        0.82,
        0.38,
        12,
        0.28,
        "APRIL 2026   ·   LIVE CRMA CHART   ·   PIPELINE HYGIENE",
        size=10,
        bold=True,
        color=PRIMARY_BLUE,
    )
    txt(
        slide,
        0.82,
        0.62,
        12.4,
        0.58,
        "Weighted open vs at-risk ARR by region",
        size=28,
        bold=True,
        color=NAVY,
    )
    largest_region = records[0]["SalesRegion"] if records else "n/a"
    largest_open = records[0].get("WeightedOpenARR", 0) if records else 0
    narrative = (
        f"Live CRMA query. {largest_region} leads at "
        f"EUR {largest_open / 1_000_000:.1f}M weighted open."
    )
    txt(slide, 0.82, 1.25, 12.1, 0.4, narrative, size=13, color=GREY_TEXT)

    # ------------------------------------------------------------------
    # NATIVE POWERPOINT CHART (stacked horizontal bar)
    # ------------------------------------------------------------------
    chart_data = CategoryChartData()
    categories = [r["SalesRegion"] for r in records]
    open_values = [round(r.get("WeightedOpenARR", 0) or 0, 2) for r in records]
    risk_values = [round(r.get("AtRiskARR", 0) or 0, 2) for r in records]
    chart_data.categories = categories
    chart_data.add_series("Weighted Open ARR", open_values)
    chart_data.add_series("At-Risk ARR", risk_values)

    chart_left = Inches(0.6)
    chart_top = Inches(1.85)
    chart_w = Inches(12.1)
    chart_h = Inches(3.9)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        chart_left,
        chart_top,
        chart_w,
        chart_h,
        chart_data,
    )
    chart = graphic_frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Colorize series to SimCorp palette
    for idx, series in enumerate(chart.series):
        series_fmt = series.format
        fill = series_fmt.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_BLUE if idx == 0 else MAGENTA

    # Source / methodology footer
    txt(
        slide,
        0.6,
        6.0,
        12.1,
        0.22,
        f"Source: CRM Analytics dashboard 0FKTb0000000KwPOAU  ·  step '{STEP_NAME}'  "
        f"·  dataset {DATASET_NAME} (version {ver_id})",
        size=9,
        color=GREY_TEXT,
    )
    # Method line
    txt(
        slide,
        0.6,
        6.22,
        12.1,
        0.22,
        "Fetched via Wave API /wave/query  ·  rendered as a native PowerPoint chart",
        size=9,
        color=GREY_TEXT,
    )

    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")


if __name__ == "__main__":
    build()
