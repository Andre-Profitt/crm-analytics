#!/usr/bin/env python3
"""
Build a SimCorp-branded director deck using native PowerPoint template layouts.

Reads a fill-payload JSON (produced by the master builder) and renders each slide
using the SimCorp PPT template's built-in slide layouts, placeholders, and theme.

Reference format: 14 slides matching the APAC director deck structure.

Usage:
    python3 scripts/build_simcorp_director_deck.py \
        --fill-payload path/to/powerpoint-fill-payload.json \
        [--template path/to/SimCorp_PPT_Template.pptx] \
        [--pi-snapshot path/to/pi-snapshot.json] \
        [--output path/to/output.pptx]
"""

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
DARK = RGBColor(0x1A, 0x1D, 0x31)
DEAL_TABLE_FONT = Pt(10)
DEAL_TABLE_ROW_H = 0.22
SUMMARY_TABLE_FONT = Pt(14)
SUMMARY_TABLE_ROW_H = 0.41
HEADER_FONT_SIZE = Pt(8)

# Layout indices in the SimCorp template
LY_TITLE_1 = 0  # Cover
LY_TITLE_CONTENT = 6  # Title and Content
LY_2COL_GRAD = 10  # 2 x content w/ gradient line
LY_4COL_GRAD = 12  # 4 x content w/ gradient line
LY_END_SLIDE = 31  # End slide with disclaimer 1

# Gradient metric column positions (inches) for 4-col layout
GRAD_COL_X = [0.9, 3.9, 6.8, 9.8]
GRAD_COL_W = 2.6
GRAD_METRIC_Y = 2.5
GRAD_METRIC_H = 0.35

DEFAULT_TEMPLATE = (
    Path.home()
    / "archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx"
)

# Slides absorbed into executive-summary or dropped entirely
SKIP_SLIDES = {"q1-review", "quarterly-pipeline", "missing-win-loss-reason"}

# These builders receive (prs, slide_data, payload) instead of (prs, slide_data)
NEEDS_PAYLOAD = {"executive-summary"}

# Internal account filter for PI data
_INTERNAL_PATTERNS = re.compile(
    r"(simcorp|test account|internal|sandbox|demo)", re.IGNORECASE
)


def _is_internal_account(name: str) -> bool:
    """Return True if the account name looks like an internal/test account."""
    return bool(_INTERNAL_PATTERNS.search(name or ""))


def _is_fy26_pi(record: dict) -> bool:
    """Return True if the PI record's close_date is within FY26."""
    cd = record.get("close_date", "")
    if not cd:
        return True
    try:
        year = int(cd[:4])
        return year <= 2026
    except (ValueError, IndexError):
        return True


def _parse_eur(value) -> float:
    """Parse a formatted EUR value like '€8.0M' or '€337K' back to a float."""
    if isinstance(value, (int, float)):
        return float(value)
    s = str(value or "0").replace("€", "").replace("EUR ", "").replace(",", "").strip()
    if not s or s == "—" or s == "-":
        return 0.0
    multiplier = 1.0
    if s.endswith("M"):
        s = s[:-1]
        multiplier = 1_000_000
    elif s.endswith("K"):
        s = s[:-1]
        multiplier = 1_000
    elif s.endswith("B"):
        s = s[:-1]
        multiplier = 1_000_000_000
    try:
        return float(s) * multiplier
    except ValueError:
        return 0.0


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _get_slide(slides: list, slide_id: str) -> dict | None:
    """Find a slide dict by id."""
    for s in slides:
        if s.get("id") == slide_id:
            return s
    return None


def _ph(slide, idx: int):
    """Get a placeholder by index, or None if missing."""
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def _set_placeholder(slide, idx: int, text: str):
    """Set a placeholder's text if it exists on the slide."""
    ph = _ph(slide, idx)
    if ph is not None:
        ph.text = str(text) if text else ""


def _set_placeholder_paragraphs(slide, idx: int, lines: list[str]):
    """Set a placeholder with multiple paragraphs (one per line)."""
    ph = _ph(slide, idx)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = str(line)
        else:
            p = tf.add_paragraph()
            p.text = str(line)


def _add_textbox_lines(
    slide,
    lines: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: Pt = Pt(12),
):
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    if not lines:
        return box
    for idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = str(line)
        for run in paragraph.runs:
            run.font.size = font_size
            run.font.color.rgb = DARK
            run.font.name = None
    return box


def _add_table(
    slide,
    rows_data: list[list[str]],
    left: float,
    top: float,
    width: float,
    row_height: float = DEAL_TABLE_ROW_H,
    col_widths: list[float] | None = None,
    data_font_size: Pt = DEAL_TABLE_FONT,
):
    """Add a table shape to the slide. Returns the table object.

    rows_data: list of rows, each row is a list of cell strings.
               First row is treated as header.
    row_height: 0.22 for deal tables, 0.41 for summary/KPI tables.
    data_font_size: Pt(10) for deal tables, Pt(14) for summary tables.
    """
    if not rows_data or not rows_data[0]:
        return None
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    total_height = max(row_height * n_rows, 0.5)

    tbl_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(total_height),
    )
    tbl = tbl_shape.table

    if col_widths and len(col_widths) == n_cols:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text) if cell_text else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = HEADER_FONT_SIZE if ri == 0 else data_font_size
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else DARK
                    run.font.bold = ri == 0
                    run.font.name = None
            if ri == 0:
                from pptx.oxml.ns import qn

                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn("a:solidFill"), {})
                srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "083EA7"})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)
    return tbl


def formatTableValue(value, fallback: str = "—") -> str:
    """Normalize optional slot values to display-safe strings."""
    if value is None:
        return fallback
    if isinstance(value, str):
        cleaned = value.strip()
        return cleaned or fallback
    return str(value)


def _watchlist_to_rows(
    watchlist: list[dict], columns: list[tuple[str, str]]
) -> list[list[str]]:
    """Convert a list of dicts to table rows with a header row.

    columns: list of (dict_key, display_header) tuples.
    """
    header = [h for _, h in columns]
    rows = [header]
    for item in watchlist:
        row = [str(item.get(k, "")) for k, _ in columns]
        rows.append(row)
    return rows


def _truncate_text(text: str, max_len: int = 120) -> str:
    if not text:
        return ""
    if len(text) <= max_len:
        return text
    return text[: max_len - 3] + "..."


def _add_gradient_metric(slide, col_idx: int, text: str):
    """Add a large bold metric number at the gradient highlight position.

    Tries placeholder indices 61-64 first. Falls back to a textbox.
    """
    target_ph_idx = 61 + col_idx
    ph = _ph(slide, target_ph_idx)
    if ph is not None:
        ph.text = str(text)
        return

    # Fallback: add a textbox at the gradient metric position
    x = GRAD_COL_X[col_idx] + 0.1
    txBox = slide.shapes.add_textbox(
        Inches(x),
        Inches(GRAD_METRIC_Y),
        Inches(GRAD_COL_W - 0.2),
        Inches(GRAD_METRIC_H),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = None


def _fmt_eur(value, compact: bool = True) -> str:
    """Format a numeric value as EUR string for slide titles/text."""
    if isinstance(value, str):
        if value.startswith("EUR") or value.startswith("€"):
            return value
        try:
            value = float(value.replace(",", ""))
        except (ValueError, AttributeError):
            return f"EUR {value}"
    if not isinstance(value, (int, float)):
        return "EUR —"
    if compact and abs(value) >= 1_000_000:
        return f"EUR {value / 1_000_000:.1f}M"
    if compact and abs(value) >= 1_000:
        return f"EUR {value / 1_000:.0f}K"
    return f"EUR {value:,.0f}"


def _meur(value) -> str:
    """Format a value as mEUR for table cells (e.g. '2.6' for 2.6M)."""
    v = _parse_eur(value) if isinstance(value, str) else float(value or 0)
    if abs(v) >= 1_000_000:
        return f"{v / 1_000_000:.1f}"
    if abs(v) >= 1_000:
        return f"{v / 1_000_000:.2f}"
    if v == 0:
        return "0"
    return f"{v / 1_000_000:.3f}"


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _snapshot_to_period(snapshot_date: str) -> str:
    """Convert '2026-04-10' to 'March 2026' (prior month = reporting period)."""
    from datetime import datetime, timedelta

    try:
        dt = datetime.strptime(snapshot_date, "%Y-%m-%d")
        # Report on prior month
        first_of_month = dt.replace(day=1)
        prior_month = first_of_month - timedelta(days=1)
        return prior_month.strftime("%B %Y")
    except (ValueError, TypeError):
        return snapshot_date


def build_cover(prs: Presentation, payload: dict):
    """Slide 1: Cover — Layout 'Title 1' (index 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_1])
    territory = payload.get("territory", "")
    snapshot = payload.get("snapshot_date", "")
    period = _snapshot_to_period(snapshot)

    _set_placeholder(slide, 24, f"{territory} Pipeline Reporting and Insights")
    _set_placeholder(slide, 20, "Sales Director Monthly")
    _set_placeholder(slide, 22, f"{period}")


def build_agenda(prs: Presentation, payload: dict):
    """Slide 2: Agenda / operating sequence."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    _set_placeholder(slide, 144, "Review sequence")
    _set_placeholder(
        slide,
        145,
        "Contract-ordered operating storyline for the monthly director review",
    )
    titles = [
        f"{index}. {slide_data.get('title', 'Untitled slide')}"
        for index, slide_data in enumerate(payload.get("slides", []), start=1)
    ]
    midpoint = max(1, (len(titles) + 1) // 2)
    _set_placeholder_paragraphs(slide, 22, titles[:midpoint])
    _add_textbox_lines(
        slide,
        titles[midpoint:],
        left=6.6,
        top=2.0,
        width=5.0,
        height=3.8,
        font_size=Pt(12),
    )


def build_executive_summary(prs: Presentation, slide_data: dict, payload: dict):
    """Slide 2: Executive Summary — 4-column gradient layout.

    Combines exec-summary, q1-review, and quarterly-pipeline data into
    4 narrative cards with gradient metric highlights.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    exec_slots = slide_data.get("slots", {})

    # Pull absorbed slide data
    q1_slots = {}
    pipe_slots = {}
    for s in payload.get("slides", []):
        if s["id"] == "q1-review":
            q1_slots = s.get("slots", {})
        elif s["id"] == "quarterly-pipeline":
            pipe_slots = s.get("slots", {})

    territory = payload.get("territory", "")
    quarter = "Q2 2026"

    # Slide title — narrative message
    pipeline_arr = exec_slots.get("headline_pipeline_arr_all_open", "—")
    _set_placeholder(slide, 144, f"Executive Summary \u2013 {territory}, {quarter}")

    # Column headers
    _set_placeholder(slide, 42, f"Executive summary - {territory}, {quarter}")
    _set_placeholder(slide, 56, f"Executive summary - {territory}, {quarter}")
    _set_placeholder(slide, 58, f"Executive summary - {territory}, {quarter}")
    _set_placeholder(slide, 60, f"Executive summary - {territory}, {quarter}")

    # --- Column 1: Pipeline ---
    _add_gradient_metric(slide, 0, str(pipeline_arr))
    pipeline_q2 = exec_slots.get("headline_pipeline_arr_q2", "—")
    pipeline_body = (
        f"Total open pipeline stands at {pipeline_arr}. "
        f"Q2 active pipeline is {pipeline_q2}."
    )
    _set_placeholder(slide, 22, pipeline_body)

    # --- Column 2: Win/Loss + Forecast Accuracy ---
    won_count = int(q1_slots.get("q1_won_count", 0) or 0)
    lost_count = int(q1_slots.get("q1_lost_count", 0) or 0)
    total_decisions = won_count + lost_count
    win_rate = (
        f"{won_count / total_decisions * 100:.0f}%" if total_decisions > 0 else "—"
    )
    _add_gradient_metric(slide, 1, f"{won_count}W / {lost_count}L")
    won_arr = q1_slots.get("q1_won_arr", "—")
    lost_arr = q1_slots.get("q1_lost_arr", "—")
    slipped = q1_slots.get("q1_slipped_count", "0")
    slipped_arr = q1_slots.get("q1_slipped_arr", "—")
    q1_body = (
        f"Win/Loss: {won_count} won ({won_arr}) vs {lost_count} lost ({lost_arr}). "
        f"Win rate: {win_rate}. "
        f"{slipped} deals slipped ({slipped_arr})."
    )
    _set_placeholder(slide, 55, q1_body)

    # --- Column 3: Q2 Forecast ---
    commit = pipe_slots.get("q2_commit_arr", "—")
    best_case = pipe_slots.get("q2_best_case_arr", "—")
    _add_gradient_metric(slide, 2, str(commit))
    forecast_body = (
        f"Commit at {commit}, best case {best_case}. "
        f"Omitted: {pipe_slots.get('q2_omitted_arr', '—')}."
    )
    _set_placeholder(slide, 57, forecast_body)

    # --- Column 4: Renewals ---
    renewal_acv = exec_slots.get("headline_renewal_acv", "—")
    _add_gradient_metric(slide, 3, str(renewal_acv))
    risk = exec_slots.get("top_risk", "")
    action = exec_slots.get("top_action", "")
    renewal_body = f"Open renewal ACV: {renewal_acv}."
    if risk:
        renewal_body += f" Risk: {_truncate_text(risk, 100)}."
    if action:
        renewal_body += f" Action: {_truncate_text(action, 100)}."
    _set_placeholder(slide, 59, renewal_body)


def build_q1_review(prs: Presentation, slide_data: dict):
    """Slide 4: Q1 promised vs delivered."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    slots = slide_data.get("slots", {})

    won_count = formatTableValue(slots.get("q1_won_count"))
    won_arr = formatTableValue(slots.get("q1_won_arr"))
    lost_count = formatTableValue(slots.get("q1_lost_count"))
    lost_arr = formatTableValue(slots.get("q1_lost_arr"))
    slipped_count = formatTableValue(slots.get("q1_slipped_count"))
    slipped_arr = formatTableValue(slots.get("q1_slipped_arr"))
    baseline = formatTableValue(
        slots.get("q1_promise_baseline_qualification"),
        "Promise baseline still needs director-safe cache confirmation.",
    )

    _set_placeholder(
        slide,
        144,
        f"Q1 delivered {won_arr} won ARR while {slipped_arr} slipped out",
    )
    _set_placeholder(slide, 42, "Won in Q1")
    _set_placeholder(slide, 56, "Lost in Q1")
    _set_placeholder(slide, 58, "Slipped out of Q1")
    _set_placeholder(slide, 60, "Promise Baseline Qualified")

    _add_gradient_metric(slide, 0, f"{won_count} | {won_arr}")
    _add_gradient_metric(slide, 1, f"{lost_count} | {lost_arr}")
    _add_gradient_metric(slide, 2, f"{slipped_count} | {slipped_arr}")
    _add_gradient_metric(slide, 3, "Qualified")

    _set_placeholder(slide, 22, f"{won_count} deals closed won in Q1.")
    _set_placeholder(slide, 55, f"{lost_count} deals closed lost in Q1.")
    _set_placeholder(
        slide,
        57,
        f"{slipped_count} deals slipped with {slipped_arr} still requiring recovery follow-up.",
    )
    _set_placeholder(slide, 59, _truncate_text(baseline, 150))


def build_quarterly_pipeline(prs: Presentation, slide_data: dict):
    """Slide 5: displayed-quarter active pipeline and forecast mix."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    slots = slide_data.get("slots", {})

    quarter_label = formatTableValue(slots.get("quarterly_pipeline_label"), "Q2")
    quarter_title = formatTableValue(
        slots.get("quarterly_pipeline_title"), f"{quarter_label} 2026"
    )
    active = formatTableValue(slots.get("headline_pipeline_arr_q2"))
    commit = formatTableValue(slots.get("q2_commit_arr"))
    best_case = formatTableValue(slots.get("q2_best_case_arr"))
    omitted = formatTableValue(slots.get("q2_omitted_arr"))
    footnote = formatTableValue(slots.get("quarterly_pipeline_footnote"), "")
    display_reason = formatTableValue(
        slots.get("quarterly_pipeline_display_reason"), "current_quarter"
    )

    title = f"{quarter_title} active ARR is {active}"
    if commit != "—" or best_case != "—":
        title = f"{quarter_title} is light in commit at {commit} against {best_case} best case"
    _set_placeholder(slide, 144, title)

    _set_placeholder(slide, 42, f"{quarter_label} Active Pipeline")
    _set_placeholder(slide, 56, "Commit")
    _set_placeholder(slide, 58, "Best Case")
    _set_placeholder(slide, 60, "Omitted")

    _add_gradient_metric(slide, 0, active)
    _add_gradient_metric(slide, 1, commit)
    _add_gradient_metric(slide, 2, best_case)
    _add_gradient_metric(slide, 3, omitted)

    _set_placeholder(
        slide,
        22,
        f"{quarter_label} active pipeline remains separate from omitted backlog.",
    )
    _set_placeholder(slide, 55, f"{quarter_title} commit is {commit}.")
    _set_placeholder(slide, 57, f"{quarter_title} best case is {best_case}.")
    _set_placeholder(slide, 59, f"Omitted remains visible at {omitted}.")

    if footnote:
        _add_textbox_lines(
            slide,
            [footnote],
            left=0.9,
            top=5.85,
            width=11.0,
            height=0.45,
            font_size=Pt(10),
        )
    elif display_reason != "current_quarter":
        _add_textbox_lines(
            slide,
            [f"Displayed quarter reason: {display_reason}."],
            left=0.9,
            top=5.85,
            width=11.0,
            height=0.45,
            font_size=Pt(10),
        )


def build_top_deals(prs: Presentation, slide_data: dict):
    """Slide 6: pipeline coverage proxy and top-quarter opportunities."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    raw_opps = slots.get("top_opportunities", [])
    top_opps = [o for o in raw_opps if _parse_eur(o.get("arr_eur", 0)) > 0]
    count = len(top_opps)
    total_arr = sum(_parse_eur(o.get("arr_eur", 0)) for o in top_opps)
    quarter_label = formatTableValue(slots.get("quarterly_pipeline_label"), "Q2")
    weighted = formatTableValue(slots.get("weighted_pipeline_arr"))
    stale_arr = formatTableValue(slots.get("stale_arr"))
    coverage_statement = formatTableValue(slots.get("pipeline_coverage_statement"), "")

    _set_placeholder(
        slide,
        144,
        f"Coverage stays qualified with {weighted} weighted ARR and {stale_arr} stale ARR",
    )
    _set_placeholder(slide, 145, f"Top {quarter_label} opportunities")
    if coverage_statement:
        _set_placeholder(slide, 22, _truncate_text(coverage_statement, 180))

    if top_opps:
        header = [
            "Account",
            "Opportunity",
            "Owner",
            "Stage",
            "Close Date",
            "Age",
            "ARR mEUR",
        ]
        rows = [header]
        for o in top_opps[:12]:
            rows.append(
                [
                    str(o.get("account", o.get("opportunity", ""))[:30]),
                    _truncate_text(str(o.get("opportunity", "")), 30),
                    str(o.get("owner", "")),
                    str(o.get("stage", "")),
                    str(o.get("close_date", "")),
                    str(o.get("age", "")),
                    _meur(o.get("arr_eur", 0)),
                ]
            )
        _add_table(
            slide,
            rows,
            left=0.9,
            top=2.2,
            width=11.5,
            col_widths=[1.6, 1.6, 1.6, 1.6, 1.6, 1.6, 1.6],
        )
    else:
        _set_placeholder(slide, 22, "No top deal data available.")


def build_commercial_approvals(prs: Presentation, slide_data: dict):
    """Slide 4: Commercial Approvals — 2-column gradient layout.

    Left: YTD Actuals (approved deals table)
    Right: Candidates (missing approval candidates table)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_2COL_GRAD])
    slots = slide_data.get("slots", {})

    approved_count = slots.get("approved_deal_count", "0")
    approved_arr = slots.get("approved_arr", "—")

    _set_placeholder(
        slide,
        144,
        f"Commercial Approvals: {approved_count} approved YTD ({approved_arr})",
    )

    # Left panel header
    _set_placeholder(slide, 42, "YTD Actuals — Approved Deals")

    # Right panel header
    _set_placeholder(slide, 56, "Commercial Approval Candidates")

    # Left panel: approved deals table
    approved = slots.get("approved_deals_2026", [])
    if approved:
        cols = [
            ("opportunity", "Opportunity"),
            ("arr", "ARR"),
            ("stage", "Stage"),
            ("owner", "Owner"),
        ]
        rows = _watchlist_to_rows(approved[:8], cols)
        _add_table(
            slide,
            rows,
            left=0.5,
            top=2.8,
            width=5.8,
            col_widths=[2.2, 1.2, 1.2, 1.2],
        )

    # Right panel: missing approval candidates
    missing_sf = slots.get("missing_approval_candidates_sf", [])
    if missing_sf:
        cols = [
            ("opportunity", "Opportunity"),
            ("arr_eur", "ARR"),
            ("owner", "Owner"),
            ("stage", "Stage"),
        ]
        rows = _watchlist_to_rows(missing_sf[:8], cols)
        _add_table(
            slide,
            rows,
            left=6.5,
            top=2.8,
            width=5.8,
            col_widths=[2.2, 1.2, 1.2, 1.2],
        )

    # Summary KPIs below
    approval_rate = slots.get("approval_rate_stage3_plus", "")
    pending = slots.get("pending_missing_approval_count", "0")
    pending_arr = slots.get("pending_missing_approval_arr", "—")
    if approval_rate or pending:
        summary_text = (
            f"Approval rate (Stage 3+): {approval_rate}  |  "
            f"Pending/missing: {pending} deals ({pending_arr})"
        )
        _set_placeholder(slide, 22, summary_text)


def build_missing_commercial_approvals(prs: Presentation, slide_data: dict):
    """Slide 5: Missing Commercial Approval — full deal table."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    candidates = slots.get("missing_approval_candidates", [])
    count = len(candidates)

    _set_placeholder(
        slide, 144, f"Missing Commercial Approval: {count} deals without approval"
    )
    _set_placeholder(slide, 145, slide_data.get("management_question", ""))

    if candidates:
        header = [
            "Account",
            "Opportunity",
            "Owner",
            "Close Date",
            "Next Step",
            "ARR mEUR",
        ]
        rows = [header]
        for c in candidates[:10]:
            rows.append(
                [
                    _truncate_text(str(c.get("account", "")), 25),
                    _truncate_text(str(c.get("opportunity", "")), 25),
                    str(c.get("owner", "")),
                    str(c.get("close_date", "")),
                    _truncate_text(str(c.get("next_step", "")), 25),
                    _meur(c.get("arr_eur", 0)),
                ]
            )
        _add_table(
            slide,
            rows,
            left=0.9,
            top=2.2,
            width=11.5,
            col_widths=[1.9, 1.9, 1.9, 1.9, 1.9, 1.9],
        )
    else:
        _set_placeholder(
            slide, 22, "No missing commercial approval candidates identified."
        )


def build_renewals(prs: Presentation, slide_data: dict):
    """Slide 6: Renewal Pipeline — annual view, 8 columns per Rebekka."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    deal_count = slots.get("renewal_open_deal_count", "0")
    acv = slots.get("renewal_open_acv", "—")

    _set_placeholder(
        slide,
        144,
        f"Renewal Pipeline - 2026: {deal_count} open ({acv} ACV)",
    )
    _set_placeholder(slide, 145, f"Renewal ACV: {acv}")

    watchlist = slots.get("renewal_watchlist", [])
    if watchlist:
        header = [
            "Close Date",
            "Account",
            "Opportunity",
            "Owner",
            "Stage",
            "ACV mEUR",
            "Probability",
            "Comments",
        ]
        rows = [header]
        for w in watchlist[:12]:
            rows.append(
                [
                    str(w.get("close_date", "")),
                    _truncate_text(str(w.get("account", "")), 22),
                    _truncate_text(str(w.get("opportunity", "")), 22),
                    str(w.get("owner", "")),
                    str(w.get("stage", "")),
                    _meur(w.get("renewal_acv_eur", 0)),
                    str(w.get("probability", "")),
                    "",  # comment field for director to fill
                ]
            )
        _add_table(
            slide,
            rows,
            left=0.9,
            top=2.0,
            width=11.5,
            col_widths=[1.4, 1.4, 1.4, 1.4, 1.4, 1.4, 1.4, 1.4],
        )
    else:
        _set_placeholder(slide, 22, "No renewal pipeline data available.")


def build_slipped_deals(prs: Presentation, slide_data: dict):
    """Slide 10: slipped deals watchlist and recovery summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    slipped_count = formatTableValue(slots.get("q1_slipped_count"), "0")
    slipped_arr = formatTableValue(slots.get("q1_slipped_arr"))
    summary = formatTableValue(
        slots.get("slip_root_cause_summary"),
        "Owner commentary is still required for full slipped-deal root-cause depth.",
    )
    watchlist = slots.get("slipped_deal_watchlist", [])

    _set_placeholder(
        slide,
        144,
        f"{slipped_arr} slipped out of Q1 across {slipped_count} deals",
    )
    _set_placeholder(slide, 145, _truncate_text(summary, 170))

    if watchlist:
        rows = [["Opportunity", "ARR mEUR", "Stage", "Owner", "Follow-up"]]
        for row in watchlist[:8]:
            rows.append(
                [
                    _truncate_text(str(row.get("opportunity", "")), 28),
                    _meur(row.get("arr_eur", 0)),
                    str(row.get("stage", "")),
                    _truncate_text(str(row.get("owner", "")), 18),
                    _truncate_text(
                        str(row.get("next_action", "")) or "Assign recovery follow-up",
                        28,
                    ),
                ]
            )
        _add_table(
            slide,
            rows,
            left=0.9,
            top=2.5,
            width=11.2,
            col_widths=[4.1, 1.2, 1.7, 1.8, 2.4],
        )
    else:
        _set_placeholder(
            slide,
            22,
            "No slipped-deal watchlist is available in the validated payload.",
        )


def build_churn_risk(prs: Presentation, slide_data: dict):
    """Slide 7: Churn Risk — placeholder with next-steps table."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    _set_placeholder(slide, 144, "Churn Risk - Pending Finance Feed")
    _set_placeholder(
        slide,
        145,
        "Churn reporting requires Finance input on renewal attrition and at-risk accounts",
    )

    status = slots.get("finance_churn_inputs_status", "")
    owner = slots.get("finance_churn_owner", "Sales Ops + Finance")
    notes = slots.get("churn_placeholder_notes", "")

    # Direction: we pull FROM Finance and feed downstream (Alex P etc.)
    next_steps = [
        ["Action", "Status", "Owner", "Notes"],
        [
            "Obtain churn data from Finance",
            status or "Not yet provided",
            owner,
            "We pull from Finance — not the other way around",
        ],
        [
            "Identify at-risk renewal accounts",
            "Pending Finance feed",
            owner,
            "Top accounts by attrition signal",
        ],
        [
            "Quantify churn exposure",
            "Pending Finance feed",
            owner,
            "ACV at risk, renewal attrition rate",
        ],
        [
            "Feed churn view to leadership",
            "Blocked",
            "Sales Ops",
            "Deck populated once Finance feed is live",
        ],
    ]
    _add_table(
        slide,
        next_steps,
        left=0.9,
        top=2.5,
        width=11.5,
        row_height=SUMMARY_TABLE_ROW_H,
        col_widths=[3.2, 2.0, 2.3, 4.0],
        data_font_size=DEAL_TABLE_FONT,
    )
    if notes:
        _add_textbox_lines(
            slide,
            [_truncate_text(notes, 220)],
            left=0.9,
            top=4.95,
            width=11.0,
            height=0.7,
            font_size=Pt(10),
        )


def build_pi_top_pushed(prs: Presentation, pi_data: dict | None):
    """Slide 9: Pipeline Inspection — Top Pushed Deals table."""
    if not pi_data:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    summary = pi_data.get("summary", {})
    high_push = summary.get("high_push_deals", [])
    all_deals = summary.get("top_deals", [])

    # Merge and sort by push_count desc
    seen = set()
    pushed = []
    for d in high_push + all_deals:
        name = d.get("name", "")
        if name not in seen and d.get("push_count", 0) > 0:
            if not _is_internal_account(name) and _is_fy26_pi(d):
                seen.add(name)
                pushed.append(d)
    pushed.sort(key=lambda d: d.get("push_count", 0), reverse=True)

    _set_placeholder(
        slide, 144, f"Pipeline Inspection: {len(pushed)} deals with close date pushes"
    )
    _set_placeholder(slide, 145, "Deals ranked by push frequency")

    if pushed:
        rows = [
            ["Opportunity", "Account", "Stage", "Close", "Pushes", "ARR mEUR", "Owner"]
        ]
        for d in pushed[:12]:
            rows.append(
                [
                    _truncate_text(d.get("name", ""), 30),
                    _truncate_text(d.get("account", ""), 22),
                    d.get("stage", ""),
                    d.get("close_date", ""),
                    str(d.get("push_count", 0)),
                    _meur(d.get("forecast_arr", 0)),
                    d.get("owner", ""),
                ]
            )
        _add_table(
            slide,
            rows,
            left=0.9,
            top=2.2,
            width=11.5,
            row_height=0.20,
            col_widths=[1.6, 1.6, 1.6, 1.6, 1.6, 1.6, 1.6],
        )
    else:
        _set_placeholder(
            slide, 22, "No pushed deals found in pipeline inspection data."
        )


def build_forecast_accuracy(prs: Presentation, payload: dict):
    """Slide 10: Forecast Accuracy — 4 gradient cards from Q1 + pipeline data."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])

    q1_slots = {}
    pipe_slots = {}
    for s in payload.get("slides", []):
        if s["id"] == "q1-review":
            q1_slots = s.get("slots", {})
        elif s["id"] == "quarterly-pipeline":
            pipe_slots = s.get("slots", {})

    won_arr = q1_slots.get("q1_won_arr", "—")
    lost_arr = q1_slots.get("q1_lost_arr", "—")
    won_count = int(q1_slots.get("q1_won_count", 0) or 0)
    lost_count = int(q1_slots.get("q1_lost_count", 0) or 0)
    total_decisions = won_count + lost_count
    win_rate = (
        f"{won_count / total_decisions * 100:.0f}%" if total_decisions > 0 else "—"
    )
    commit = pipe_slots.get("q2_commit_arr", "—")

    _set_placeholder(
        slide,
        144,
        f"Forecast Accuracy: {win_rate} win rate | {won_arr} won this quarter",
    )

    # Column headers
    _set_placeholder(slide, 42, "Won ARR (Q1)")
    _set_placeholder(slide, 56, "Lost ARR (Q1)")
    _set_placeholder(slide, 58, "Win Rate")
    _set_placeholder(slide, 60, "Q2 Commit Forecast")

    # Gradient metrics
    _add_gradient_metric(slide, 0, str(won_arr))
    _add_gradient_metric(slide, 1, str(lost_arr))
    _add_gradient_metric(slide, 2, win_rate)
    _add_gradient_metric(slide, 3, str(commit))

    # Narrative bodies
    _set_placeholder(slide, 22, f"Won {won_count} deals for {won_arr} in Q1.")
    _set_placeholder(slide, 55, f"Lost {lost_count} deals totalling {lost_arr} in Q1.")
    _set_placeholder(
        slide,
        57,
        f"Win rate of {win_rate} across {total_decisions} closed decisions."
        if total_decisions > 0
        else "No closed decisions in the period.",
    )
    _set_placeholder(slide, 59, f"Q2 commit forecast stands at {commit}.")


def build_forecast_category_breakdown(prs: Presentation, pi_data: dict | None):
    """Slide 11: Forecast Category Breakdown — table from PI data."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])

    if not pi_data:
        _set_placeholder(slide, 144, "Forecast Category Breakdown")
        _set_placeholder(slide, 22, "No pipeline inspection data available.")
        return

    # Recompute from raw records with FY26 filter
    from collections import defaultdict

    all_records = pi_data.get("records", [])
    fy26_records = [
        r
        for r in all_records
        if _is_fy26_pi(r) and not _is_internal_account(r.get("name", ""))
    ]
    fc: dict[str, dict] = defaultdict(lambda: {"count": 0, "forecast_arr": 0.0})
    for r in fy26_records:
        cat = r.get("forecast_category", "Unknown")
        fc[cat]["count"] += 1
        fc[cat]["forecast_arr"] += r.get("forecast_arr", 0)
    open_recs = [
        r
        for r in fy26_records
        if not r.get("is_closed")
        and r.get("forecast_category") not in ("Omitted", "Closed")
    ]
    total_open_arr = sum(r.get("forecast_arr", 0) for r in open_recs)

    _set_placeholder(slide, 144, "Forecast Category Breakdown — FY26 open pipeline")
    _set_placeholder(slide, 145, f"{len(open_recs)} open FY26 deals in pipeline")

    rows = [["Category", "Deals", "ARR mEUR"]]
    for cat in ("Pipeline", "Omitted", "Commit", "Best Case", "Closed"):
        bucket = fc.get(cat, {})
        if bucket.get("count"):
            rows.append(
                [cat, str(bucket["count"]), _meur(bucket.get("forecast_arr", 0))]
            )

    rows.append(["Total Open (FY26)", str(len(open_recs)), _meur(total_open_arr)])

    if len(rows) > 2:
        _add_table(
            slide,
            rows,
            left=0.9,
            top=2.2,
            width=5.5,
            row_height=0.21,
            col_widths=[1.8, 1.8, 1.8],
        )
    else:
        _set_placeholder(slide, 22, "No forecast category data available.")


def build_pipeline_coverage_targets(prs: Presentation, slide_data: dict):
    """Slide 12: Pipeline Coverage & Targets — placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    _set_placeholder(slide, 144, "Pipeline Coverage & Targets")
    _set_placeholder(slide, 145, "What is the coverage ratio against assigned targets?")

    weighted = slots.get("weighted_pipeline_arr", "")
    stale = slots.get("stale_arr", "")

    lines = []
    if weighted:
        lines.append(f"Weighted Pipeline (probability-adjusted): {weighted}")
    if stale:
        lines.append(f"Stale 30d+ ARR: {stale}")
    lines.extend(
        [
            "",
            "Quota and target data is not yet integrated into the automated pipeline.",
            "Coverage ratio remains qualified until formal targets are provided.",
            "",
            "Next steps:",
            "  - Integrate quota/target assignments from Finance or Sales Ops",
            "  - Replace this placeholder with actual coverage ratio by territory",
        ]
    )
    _set_placeholder_paragraphs(slide, 22, lines)


def build_salesforce_hygiene(prs: Presentation, slide_data: dict):
    """Slide 11: Salesforce Hygiene — KPIs + rep table."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    no_act = slots.get("no_activity_count", "0")
    overdue = slots.get("overdue_close_count", "0")
    missing_ns = slots.get("missing_next_step_count", "0")
    total = slots.get("total_data_quality_issues", "0")

    _set_placeholder(slide, 144, f"Salesforce Hygiene: {total} data quality issues")
    _set_placeholder(slide, 145, slide_data.get("management_question", ""))

    # Summary bullets
    bullets = [
        f"No Activity: {no_act} deals",
        f"Overdue Close: {overdue} deals",
        f"Missing Next Step: {missing_ns} deals",
        f"Total Data Quality Issues: {total}",
    ]
    conc = slots.get("rep_concentration_summary", "")
    if conc:
        bullets.append(f"Concentration: {conc}")
    _set_placeholder_paragraphs(slide, 22, bullets)

    # Top issue reps table
    reps = slots.get("top_issue_reps", [])
    if reps:
        header = [
            "Rep",
            "Total Issues",
            "No Activity",
            "Overdue Close",
            "Missing Next Step",
        ]
        rows = [header]
        for r in reps[:5]:
            rows.append(
                [
                    str(r.get("rep", "")),
                    str(r.get("total_issues", "")),
                    str(r.get("no_activity", "")),
                    str(r.get("overdue_close", "")),
                    str(r.get("missing_next_step", "")),
                ]
            )
        _add_table(
            slide,
            rows,
            left=0.5,
            top=4.5,
            width=11.5,
            col_widths=[3.5, 1.8, 1.8, 2.2, 2.2],
        )


def build_missing_win_loss_reason(prs: Presentation, slide_data: dict):
    """Slide 12: win/loss reason control exceptions."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    issue_count = formatTableValue(slots.get("missing_win_loss_reason_count"), "0")
    rows = slots.get("missing_win_loss_reason_rows", [])
    rule_note = formatTableValue(
        slots.get("missing_win_loss_reason_rule_note"),
        "Treat `0 - No Opportunity` with no reason as acceptable; all other blank-reason outcomes require hygiene follow-up.",
    )

    if not rows and issue_count in {"0", "—"}:
        _set_placeholder(
            slide,
            144,
            "Win/loss reason hygiene is clean outside accepted 0 - No Opportunity cases",
        )
        _set_placeholder(slide, 145, "Outcome reason controls")
        _set_placeholder_paragraphs(
            slide,
            22,
            [
                rule_note,
                "No outcome rows currently require reason follow-up after applying the accepted 0 - No Opportunity exception.",
            ],
        )
        return

    _set_placeholder(
        slide,
        144,
        f"{issue_count} outcome rows still need decision-reason hygiene",
    )
    _set_placeholder(slide, 145, "Outcome reason controls")
    _set_placeholder(slide, 22, _truncate_text(rule_note, 180))

    table_rows = [["Opportunity", "Owner", "Stage", "Close Date", "ARR mEUR"]]
    for row in rows[:8]:
        table_rows.append(
            [
                _truncate_text(str(row.get("opportunity", "")), 26),
                _truncate_text(str(row.get("owner", "")), 18),
                str(row.get("stage", "")),
                str(row.get("close_date", "")),
                _meur(row.get("arr_eur", 0)),
            ]
        )
    _add_table(
        slide,
        table_rows,
        left=0.9,
        top=2.6,
        width=11.2,
        col_widths=[2.9, 1.8, 2.0, 2.0, 1.5],
    )


def build_overdue_close_open_opps(prs: Presentation, slide_data: dict):
    """Slide 13: overdue-close watchlist plus owner concentration."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    overdue_count = formatTableValue(slots.get("overdue_close_count"), "0")
    watchlist = slots.get("overdue_close_watchlist", [])
    owner_summary = slots.get("overdue_close_owner_summary", [])
    lead_owners = [row.get("owner", "") for row in owner_summary[:2] if row.get("owner")]
    if lead_owners:
        owner_phrase = " and ".join(lead_owners)
        title = f"{overdue_count} open opps are past close date, led by {owner_phrase}"
    else:
        title = f"{overdue_count} open opps are past close date"

    _set_placeholder(slide, 144, title)
    _set_placeholder(slide, 145, slide_data.get("management_question", ""))

    table_rows = [["Opportunity", "ARR mEUR", "Close Date", "Owner", "Follow-up"]]
    for row in watchlist[:8]:
        table_rows.append(
            [
                _truncate_text(str(row.get("opportunity", "")), 28),
                _meur(row.get("arr_eur", 0)),
                str(row.get("close_date", "")),
                _truncate_text(str(row.get("owner", "")), 18),
                _truncate_text(str(row.get("next_action", "")) or "Inspect overdue close date", 28),
            ]
        )
    _add_table(
        slide,
        table_rows,
        left=0.7,
        top=2.4,
        width=7.8,
        col_widths=[2.9, 1.1, 1.4, 1.4, 1.7],
    )

    if owner_summary:
        owner_rows = [["Owner", "Open overdue opps"]]
        for row in owner_summary[:5]:
            owner_rows.append(
                [
                    _truncate_text(str(row.get("owner", "")), 18),
                    str(row.get("record_count", "")),
                ]
            )
        _add_table(
            slide,
            owner_rows,
            left=8.9,
            top=2.4,
            width=3.0,
            row_height=0.30,
            col_widths=[1.9, 1.1],
        )


def build_definitions(prs: Presentation, slide_data: dict):
    """Slide 14: Definitions and Data Sources — appendix."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    slots = slide_data.get("slots", {})

    _set_placeholder(slide, 144, slide_data.get("title", "Appendix and Factual Notes"))
    _set_placeholder(slide, 145, "")

    lines = [
        "Metric conventions",
        "  Pipeline metrics are ARR in EUR (converted) unless stated otherwise",
        "  Renewal metrics are ACV in EUR (converted)",
        "  Omitted deals are visible but excluded from headline pipeline",
        "",
        "Data sources",
    ]
    lineage = slots.get("sources_lineage_summary", "")
    if lineage:
        lines.append(f"  {lineage}")
    lines.append(
        "  Commercial approval: Salesforce reports (approved 2026 + candidates)"
    )

    kyc = slots.get("kyc_missing_status", "")
    if kyc:
        lines.append(f"  KYC: {kyc}")

    dq = slots.get("data_quality_backlog", "")
    if dq:
        lines.append(f"  Data quality backlog: {dq} open issues")

    lines.extend(
        [
            "",
            "Known limitations",
            "  Q1 promise baseline is directional \u2014 formal targets not yet integrated",
            "  Finance churn overlay pending \u2014 churn slide is a placeholder",
            "  Quota/target data not yet available for automated coverage ratio",
        ]
    )

    _set_placeholder_paragraphs(slide, 22, lines)


def build_end_slide(prs: Presentation):
    """Slide 15: End slide with disclaimer."""
    prs.slides.add_slide(prs.slide_layouts[LY_END_SLIDE])


# ---------------------------------------------------------------------------
# Slide dispatch — maps payload slide IDs to builder functions
# ---------------------------------------------------------------------------

SLIDE_BUILDERS = {
    "executive-summary": build_executive_summary,
    "pipeline-coverage-intel": build_top_deals,
    "commercial-approval-overview": build_commercial_approvals,
    "missing-commercial-approvals": build_missing_commercial_approvals,
    "renewals-retention": build_renewals,
    "churn-finance": build_churn_risk,
    "slipped-deals": None,  # handled via PI pushed deals
    "salesforce-hygiene-activity": build_salesforce_hygiene,
    "overdue-close-open-opps": None,  # folded into hygiene
    "appendix-notes": build_definitions,
}


# ---------------------------------------------------------------------------
# Template cleanup
# ---------------------------------------------------------------------------


def _remove_all_slides(prs: Presentation):
    """Remove all pre-existing slides from the presentation (keeps layouts)."""
    from lxml import etree

    pres_elem = prs.part._element
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    sldIdLst = pres_elem.find("p:sldIdLst", nsmap)
    if sldIdLst is None:
        return

    rIds = [sldId.get(etree.QName(r_ns, "id")) for sldId in list(sldIdLst)]
    for rId in rIds:
        if rId is None:
            continue
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass

    # Clear the slide ID list element
    for child in list(sldIdLst):
        sldIdLst.remove(child)


# ---------------------------------------------------------------------------
# Main build function
# ---------------------------------------------------------------------------


def build_deck(
    payload: dict,
    template_path: Path,
    output_path: Path,
    pi_data: dict | None = None,
):
    """Build the full deck from a fill payload.

    Slide order:
      1. Cover
      2. Review sequence
      3. Executive summary
      4. Q1 review
      5. Quarterly pipeline
      6. Pipeline coverage intel
      7. Commercial approvals
      8. Missing commercial approvals
      9. Renewals
     10. Slipped deals
     11. Salesforce hygiene
     12. Missing win/loss reason
     13. Overdue close date open opportunities
     14. Churn and Finance inputs
     15. Appendix and factual notes
    """
    prs = Presentation(str(template_path))
    _remove_all_slides(prs)

    slides_list = payload.get("slides", [])

    # --- 1. Cover ---
    build_cover(prs, payload)
    print("  [OK]   cover")

    # --- 2. Agenda ---
    build_agenda(prs, payload)
    print("  [OK]   review-sequence")

    # --- 3. Executive Summary ---
    exec_data = _get_slide(slides_list, "executive-summary") or {"slots": {}}
    build_executive_summary(prs, exec_data, payload)
    print("  [OK]   executive-summary")

    # --- 4. Q1 Review ---
    q1_data = _get_slide(slides_list, "q1-review") or {"slots": {}}
    build_q1_review(prs, q1_data)
    print("  [OK]   q1-review")

    # --- 5. Quarterly Pipeline ---
    quarterly_data = _get_slide(slides_list, "quarterly-pipeline") or {"slots": {}}
    build_quarterly_pipeline(prs, quarterly_data)
    print("  [OK]   quarterly-pipeline")

    # --- 6. Pipeline Coverage Intel ---
    pci_data = _get_slide(slides_list, "pipeline-coverage-intel") or {"slots": {}}
    build_top_deals(prs, pci_data)
    print("  [OK]   pipeline-coverage-intel")

    # --- 7. Commercial Approvals ---
    cao_data = _get_slide(slides_list, "commercial-approval-overview") or {"slots": {}}
    build_commercial_approvals(prs, cao_data)
    print("  [OK]   commercial-approvals")

    # --- 8. Missing Commercial Approval ---
    mca_data = _get_slide(slides_list, "missing-commercial-approvals") or {"slots": {}}
    build_missing_commercial_approvals(prs, mca_data)
    print("  [OK]   missing-commercial-approvals")

    # --- 9. Renewal Pipeline ---
    ren_data = _get_slide(slides_list, "renewals-retention") or {"slots": {}}
    build_renewals(prs, ren_data)
    print("  [OK]   renewals")

    # --- 10. Slipped Deals ---
    slipped_data = _get_slide(slides_list, "slipped-deals") or {"slots": {}}
    build_slipped_deals(prs, slipped_data)
    print("  [OK]   slipped-deals")

    # --- 11. Salesforce Hygiene ---
    hygiene_data = _get_slide(slides_list, "salesforce-hygiene-activity") or {"slots": {}}
    build_salesforce_hygiene(prs, hygiene_data)
    print("  [OK]   salesforce-hygiene")

    # --- 12. Missing Win/Loss Reason ---
    missing_reason_data = _get_slide(slides_list, "missing-win-loss-reason") or {
        "slots": {}
    }
    build_missing_win_loss_reason(prs, missing_reason_data)
    print("  [OK]   missing-win-loss-reason")

    # --- 13. Overdue Close Open Opps ---
    overdue_data = _get_slide(slides_list, "overdue-close-open-opps") or {"slots": {}}
    build_overdue_close_open_opps(prs, overdue_data)
    print("  [OK]   overdue-close-open-opps")

    # --- 14. Churn Risk ---
    churn_data = _get_slide(slides_list, "churn-finance") or {"slots": {}}
    build_churn_risk(prs, churn_data)
    print("  [OK]   churn-risk")

    # --- 15. Definitions ---
    appendix_data = _get_slide(slides_list, "appendix-notes") or {"slots": {}}
    build_definitions(prs, appendix_data)
    print("  [OK]   appendix-notes")

    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")


def _build_pushed_summary_direct(prs: Presentation, pi_data: dict):
    """Build the Pushed Deals Summary slide directly on prs."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_4COL_GRAD])
    summary = pi_data.get("summary", {})
    high_push = summary.get("high_push_deals", [])
    all_deals = summary.get("top_deals", [])

    pushed_deals = [
        d
        for d in all_deals
        if d.get("push_count", 0) > 0
        and not _is_internal_account(d.get("name", ""))
        and _is_fy26_pi(d)
    ]
    total_pushed = len(pushed_deals)
    avg_pushes = sum(d.get("push_count", 0) for d in pushed_deals) / max(
        total_pushed, 1
    )
    exposed_arr = sum(d.get("forecast_arr", 0) for d in pushed_deals)
    critical = len([d for d in pushed_deals if d.get("push_count", 0) >= 5])

    _set_placeholder(
        slide,
        144,
        f"Pushed Deals: {total_pushed} deals pushed | {_fmt_eur(exposed_arr)} exposed",
    )

    _set_placeholder(slide, 42, "Total Pushed")
    _set_placeholder(slide, 56, "Avg Pushes / Deal")
    _set_placeholder(slide, 58, "Exposed ARR")
    _set_placeholder(slide, 60, "Critical (5+ Pushes)")

    _add_gradient_metric(slide, 0, str(total_pushed))
    _add_gradient_metric(slide, 1, f"{avg_pushes:.1f}")
    _add_gradient_metric(slide, 2, _fmt_eur(exposed_arr))
    _add_gradient_metric(slide, 3, str(critical))

    _set_placeholder(slide, 22, f"{total_pushed} deals have been pushed at least once.")
    _set_placeholder(slide, 55, f"Average of {avg_pushes:.1f} pushes per pushed deal.")
    _set_placeholder(
        slide, 57, f"{_fmt_eur(exposed_arr)} forecast ARR exposed to push risk."
    )
    _set_placeholder(
        slide,
        59,
        f"{critical} deals at critical push threshold (5+)."
        if critical
        else "No deals at critical push threshold (5+).",
    )


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------


def main():
    parser = argparse.ArgumentParser(
        description="Build SimCorp-branded director deck from fill payload"
    )
    parser.add_argument(
        "--fill-payload",
        required=True,
        help="Path to the powerpoint-fill-payload.json file",
    )
    parser.add_argument(
        "--template",
        default=str(DEFAULT_TEMPLATE),
        help="Path to SimCorp_PPT_Template.pptx",
    )
    parser.add_argument(
        "--output",
        default=None,
        help="Output .pptx path (defaults to director name slug)",
    )
    parser.add_argument(
        "--pi-snapshot",
        default=None,
        help="Path to a pipeline inspection snapshot JSON",
    )
    args = parser.parse_args()

    payload_path = Path(args.fill_payload)
    if not payload_path.exists():
        print(f"ERROR: Fill payload not found: {payload_path}", file=sys.stderr)
        sys.exit(1)

    template_path = Path(args.template)
    if not template_path.exists():
        print(f"ERROR: Template not found: {template_path}", file=sys.stderr)
        sys.exit(1)

    with open(payload_path) as f:
        payload = json.load(f)

    # Load PI snapshot if provided
    pi_data = None
    if args.pi_snapshot:
        pi_path = Path(args.pi_snapshot)
        if pi_path.exists():
            with open(pi_path) as f:
                pi_data = json.load(f)
            print(
                f"PI data: {pi_path.name} "
                f"({pi_data.get('summary', {}).get('open_pipeline_count', '?')} open deals)"
            )
        else:
            print(f"WARNING: PI snapshot not found: {pi_path}", file=sys.stderr)

    # Derive output path
    if args.output:
        output_path = Path(args.output)
    else:
        name_slug = re.sub(
            r"[^a-z0-9]+", "-", payload.get("director_name", "director").lower()
        ).strip("-")
        snapshot = payload.get("snapshot_date", "undated")
        output_path = payload_path.parent / f"{name_slug}-{snapshot}-simcorp.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"Template: {template_path}")
    print(f"Payload:  {payload_path}")
    print(
        f"Director: {payload.get('director_name', '?')} ({payload.get('territory', '?')})"
    )
    print("Building deck...\n")

    build_deck(payload, template_path, output_path, pi_data=pi_data)


if __name__ == "__main__":
    main()
