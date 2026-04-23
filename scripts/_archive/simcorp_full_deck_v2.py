#!/usr/bin/env python3
"""Full Sales Director Monthly deck at SimCorp production spec.

Pulls data from report1_snapshot.json and renders 10 slides using the v2
polish pattern (rounded cards, eyebrow kickers, aqua tabs, progress bars,
status dots, takeaway callouts, horizontal dividers).

Calendar year labels only, no fiscal year prefix.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"
RUN_DIR = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01"
)
SNAPSHOT = RUN_DIR / "report1_snapshot.json"
OUTPUT = RUN_DIR / "sales_director_monthly_simcorp_v2.pptx"

# SimCorp production color tokens
NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
LIGHT_PANEL_2 = RGBColor(0xF5, 0xF8, 0xFD)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)
DIVIDER_GREY = RGBColor(0xD7, 0xE2, 0xE8)
AMBER = RGBColor(0xFB, 0x9B, 0x2A)
GREEN_OK = RGBColor(0x33, 0xD8, 0xCE)

LINKS = [
    (
        "Sales Directors Monthly Dashboard",
        "https://simcorp.my.salesforce.com/lightning/r/Dashboard/01ZTb00000FSP7hMAH/view",
    ),
    (
        "Sales Ops Quarterly KPI Dashboard",
        "https://simcorp.my.salesforce.com/lightning/r/Dashboard/01ZTb00000FSP9JMAX/view",
    ),
    (
        "Sales Ops Data Quality and Forecast Accuracy (CRMA)",
        "https://simcorp.my.salesforce.com/analytics/dashboard/0FKTb0000000K5BOAU",
    ),
    (
        "Pipeline Coverage by Stage",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008TZc5MAG/view",
    ),
    (
        "Forecast and Closed Won",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008TZaTMAW/view",
    ),
    (
        "Land Stage 3 Missing Commercial Approval by Region",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekltMAA/view",
    ),
    (
        "Commercial Approval Candidates by Stage",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekp7MAA/view",
    ),
    (
        "Renewals by Fiscal Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008eksLMAQ/view",
    ),
    (
        "Renewal Pipeline This Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ektxMAA/view",
    ),
    (
        "Renewal ACV by Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekxBMAQ/view",
    ),
    (
        "Business At Risk",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008Ta9xMAC/view",
    ),
    (
        "Close Date Slipped by Stage",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008eknVMAQ/view",
    ),
]


def fmt_eur(amount):
    if amount is None:
        return "n/a"
    if abs(amount) >= 1_000_000_000:
        return f"EUR {amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"EUR {amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"EUR {amount / 1_000:.0f}K"
    return f"EUR {amount:.0f}"


def fmt_int(n):
    return f"{n:,}" if n is not None else "n/a"


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_rect(
    slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rounded_card(slide, left, top, width, height, fill_color):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.fill.background()
    try:
        card.adjustments[0] = 0.08
    except Exception:
        pass
    return card


def add_dot(slide, left, top, size, color):
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return dot


def add_kpi_card(
    slide,
    left,
    top,
    width,
    height,
    big,
    label,
    context,
    progress_pct=None,
    status_color=None,
):
    add_rounded_card(slide, left, top, width, height, LIGHT_BLUE_PANEL)
    # Big
    tb = slide.shapes.add_textbox(
        Inches(left + 0.22), Inches(top + 0.18), Inches(width - 0.44), Inches(0.7)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = big
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_BLUE
    # Status dot
    if status_color is not None:
        add_dot(slide, left + width - 0.44, top + 0.28, 0.22, status_color)
    # Label
    add_text(
        slide,
        left + 0.22,
        top + 0.92,
        width - 0.44,
        0.32,
        label,
        size=13,
        bold=True,
        color=NAVY,
    )
    # Context
    ctx_top = top + 1.22
    ctx_h = height - 1.3
    if progress_pct is not None:
        ctx_h -= 0.25
    add_text(
        slide,
        left + 0.22,
        ctx_top,
        width - 0.44,
        ctx_h,
        context,
        size=10,
        color=PRIMARY_BLUE,
    )
    # Progress bar
    if progress_pct is not None:
        bar_top = top + height - 0.32
        bar_left = left + 0.22
        bar_width = width - 0.44
        add_rect(slide, bar_left, bar_top, bar_width, 0.1, DIVIDER_GREY)
        fill_w = max(0.02, bar_width * min(1.0, progress_pct / 100))
        add_rect(slide, bar_left, bar_top, fill_w, 0.1, AQUA)


def slide_header(slide, eyebrow, title, narrative):
    # Vertical aqua tab
    add_rect(slide, 0.6, 0.42, 0.08, 0.7, AQUA)
    # Eyebrow
    add_text(
        slide, 0.82, 0.38, 12, 0.28, eyebrow, size=10, bold=True, color=PRIMARY_BLUE
    )
    # Title
    add_text(slide, 0.82, 0.62, 12.4, 0.58, title, size=28, bold=True, color=NAVY)
    # Narrative
    add_text(slide, 0.82, 1.25, 12.1, 0.4, narrative, size=13, color=GREY_TEXT)


def add_takeaway(slide, top, bullets):
    # Divider
    add_rect(slide, 0.6, top - 0.25, 12.1, 0.015, DIVIDER_GREY)
    # Callout panel
    callout_h = 1.9
    add_rounded_card(slide, 0.6, top, 12.1, callout_h, LIGHT_PANEL_2)
    add_rect(slide, 0.6, top, 0.1, callout_h, AQUA)
    add_text(
        slide,
        0.95,
        top + 0.16,
        11.5,
        0.28,
        "TAKEAWAY",
        size=10,
        bold=True,
        color=PRIMARY_BLUE,
    )
    add_text(
        slide, 0.95, top + 0.45, 11.5, callout_h - 0.55, bullets, size=12, color=NAVY
    )


def source_footer(slide, text):
    add_text(slide, 0.6, 7.08, 12.1, 0.22, text, size=9, color=GREY_TEXT)


# ============================================================================
# Slide builders
# ============================================================================


def build_cover(prs, snap):
    cover = prs.slides.add_slide(get_layout(prs, "SC-Master Gradient_Title"))

    def fill_ph(idx, text):
        for ph in cover.placeholders:
            if ph.placeholder_format.idx == idx:
                ph.text = text
                return

    fill_ph(0, "Sales Director Monthly Pipeline Insights")
    fill_ph(
        20,
        "Pipeline coverage, commercial approvals, renewals, churn, and slipped deals for Q1 2026.",
    )
    fill_ph(22, f"Snapshot {snap.get('snapshot_date', '2026-04-01')}")


def build_exec_summary(prs, snap):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    regions = snap["pipeline"]["deck_regions"]
    total_target = sum(r.get("target_arr", 0) or 0 for r in regions.values())
    total_call = sum(r.get("best_case_call_arr", 0) or 0 for r in regions.values())
    gap = max(0.0, total_target - total_call)
    q_window = snap.get("quarter_window", {})
    days_remaining = q_window.get("days_remaining", 0)
    coverage_pct = (total_call / total_target * 100) if total_target else 0

    slide_header(
        slide,
        eyebrow="Q1 2026   ·   EXECUTIVE SUMMARY",
        title="Pipeline coverage at a glance",
        narrative=(
            f"Combined best case call is {fmt_eur(total_call)} against a "
            f"{fmt_eur(total_target)} quarter target. Coverage {coverage_pct:.0f}%."
        ),
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 2.95
    gap_bet = 0.2
    start_x = 0.6
    cards = [
        (
            fmt_eur(total_target),
            "Quarter target",
            "Combined regional seam, ARR.",
            None,
            None,
        ),
        (
            fmt_eur(total_call),
            "Best case call",
            f"{coverage_pct:.0f}% of target coverage.",
            coverage_pct,
            None,
        ),
        (
            fmt_eur(gap),
            "Gap to target",
            "Promotion plus new pipeline required.",
            None,
            None,
        ),
        (
            f"{days_remaining}d",
            "Days remaining",
            "Left in the current quarter to close the gap.",
            None,
            AMBER if days_remaining < 35 else GREEN_OK,
        ),
    ]
    for i, (big, lbl, ctx, prog, status) in enumerate(cards):
        add_kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            big,
            lbl,
            ctx,
            progress_pct=prog,
            status_color=status,
        )

    biggest_gap_region = snap.get("biggest_gap_region", "North America")
    biggest_gap_arr = snap.get("biggest_gap_arr", 0)
    weakest_conf_region = snap.get("weakest_confidence_region", "EMEA")
    weakest_conf_pct = snap.get("weakest_confidence_pct", 0)
    bullets = (
        f"•  Biggest gap region: {biggest_gap_region} at {fmt_eur(biggest_gap_arr)}.\n"
        f"•  Weakest forecast confidence: {weakest_conf_region} at {weakest_conf_pct:.1f}%.\n"
        f"•  {days_remaining} days left in the quarter to close the gap."
    )
    add_takeaway(slide, 4.35, bullets)
    source_footer(
        slide,
        "Source: Sales Director Monthly snapshot, 2026-04-01   ·   ARR for land and expand, ACV for renewals and churn",
    )


def build_region(prs, snap, region_name, eyebrow_num, narrative_override=None):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    reg = snap["pipeline"]["deck_regions"].get(region_name, {})
    target = reg.get("target_arr", 0) or 0
    call = reg.get("best_case_call_arr", 0) or 0
    confidence = reg.get("forecast_confidence_pct", 0) or 0
    coverage_status = reg.get("coverage_status", "n/a") or "n/a"
    gap = max(0.0, target - call)
    coverage_pct = (call / target * 100) if target else 0

    narrative = narrative_override or (
        f"{region_name} best case call {fmt_eur(call)} against {fmt_eur(target)} quarter target. "
        f"Forecast confidence {confidence:.1f}%. Coverage status: {coverage_status}."
    )
    slide_header(
        slide,
        eyebrow=f"Q1 2026   ·   {eyebrow_num}   ·   REGIONAL OUTLOOK",
        title=f"{region_name} regional outlook",
        narrative=narrative,
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 2.95
    gap_bet = 0.2
    start_x = 0.6
    cards = [
        (fmt_eur(target), "Quarter target", "Executive target seam, ARR.", None, None),
        (
            fmt_eur(call),
            "Best case call",
            f"{coverage_pct:.0f}% of target coverage.",
            coverage_pct,
            None,
        ),
        (
            fmt_eur(gap),
            "Gap to target",
            "Needed from promotion or new pipeline.",
            None,
            None,
        ),
        (
            f"{confidence:.1f}%",
            "Forecast confidence",
            f"Status: {coverage_status}.",
            None,
            AMBER if confidence < 96 else GREEN_OK,
        ),
    ]
    for i, (big, lbl, ctx, prog, status) in enumerate(cards):
        add_kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            big,
            lbl,
            ctx,
            progress_pct=prog,
            status_color=status,
        )

    bullets = _region_takeaway_bullets(region_name, reg, gap)
    add_takeaway(slide, 4.35, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


def _region_takeaway_bullets(region_name, reg, gap):
    if gap <= 0:
        status_line = f"•  {region_name} is at or above target on best case call."
    else:
        status_line = f"•  {region_name} needs {fmt_eur(gap)} additional coverage to close the quarter."
    coverage = reg.get("coverage_status", "") or ""
    coverage_line = f"•  Coverage status: {coverage}."
    confidence = reg.get("forecast_confidence_pct", 0) or 0
    if confidence < 96:
        conf_line = f"•  Forecast confidence {confidence:.1f}% is below threshold. Validate promotion evidence."
    else:
        conf_line = f"•  Forecast confidence {confidence:.1f}% is healthy."
    return "\n".join([status_line, coverage_line, conf_line])


def build_commercial_approval(prs, snap):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    ca = snap.get("commercial_approval", {})
    s = ca.get("summary", {})
    approved_count = s.get("approved_count", 0) or 0
    approved_arr = s.get("approved_arr", 0) or 0
    pending_count = s.get("pending_count", 0) or 0
    pending_arr = s.get("pending_arr", 0) or 0
    stale_count = s.get("stale_count", 0) or 0
    stale_arr = s.get("stale_arr", 0) or 0
    candidate_count = ca.get("candidate_count", 0) or 0

    slide_header(
        slide,
        eyebrow="Q1 2026   ·   02   ·   COMMERCIAL GOVERNANCE",
        title="Commercial approval overview",
        narrative=(
            f"{approved_count} approved at {fmt_eur(approved_arr)}. "
            f"{pending_count} pending. {stale_count} stale."
        ),
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 2.95
    gap_bet = 0.2
    start_x = 0.6
    cards = [
        (
            fmt_int(approved_count),
            "Approved",
            f"Total value: {fmt_eur(approved_arr)}",
            None,
            GREEN_OK if approved_count > 0 else None,
        ),
        (
            fmt_int(pending_count),
            "Pending approval",
            f"Awaiting sign off: {fmt_eur(pending_arr)}",
            None,
            AMBER if pending_count > 0 else None,
        ),
        (
            fmt_int(stale_count),
            "Stale",
            f"Cleanup required: {fmt_eur(stale_arr)}",
            None,
            AMBER if stale_count > 0 else None,
        ),
        (
            fmt_int(candidate_count),
            "Land stage 3 candidates",
            "Open candidates without approval.",
            None,
            AMBER if candidate_count > 0 else GREEN_OK,
        ),
    ]
    for i, (big, lbl, ctx, prog, status) in enumerate(cards):
        add_kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            big,
            lbl,
            ctx,
            progress_pct=prog,
            status_color=status,
        )

    if candidate_count == 0 and pending_count == 0 and stale_count == 0:
        line1 = "•  Approval control is fully contained this quarter."
    else:
        line1 = f"•  {pending_count + stale_count + candidate_count} items require attention this cycle."
    bullets = "\n".join(
        [
            line1,
            "•  Drill through to the Land Stage 3 Missing Approval report for owner-level follow up.",
            "•  Live dashboard link available in the appendix.",
        ]
    )
    add_takeaway(slide, 4.35, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


def build_renewals(prs, snap):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    ren = snap.get("renewals", {}).get("summary", {})
    total = ren.get("total_renewal_pipeline_acv", 0) or 0
    deals = ren.get("total_deals", 0) or 0
    due_acv = ren.get("due_this_quarter_acv", 0) or 0
    due_ct = ren.get("due_this_quarter_count", 0) or 0
    over_acv = ren.get("overdue_carryover_acv", 0) or 0
    over_ct = ren.get("overdue_carryover_count", 0) or 0
    crit_acv = ren.get("critical_acv", 0) or 0
    crit_ct = ren.get("critical_count", 0) or 0

    slide_header(
        slide,
        eyebrow="Q1 2026   ·   03   ·   RETENTION",
        title="Renewal pipeline and risk",
        narrative=(
            f"{deals} open renewal deals at {fmt_eur(total)} ACV. "
            f"{over_ct} overdue carryover worth {fmt_eur(over_acv)}."
        ),
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 2.95
    gap_bet = 0.2
    start_x = 0.6
    cards = [
        (
            fmt_eur(total),
            "Total open renewal ACV",
            f"{deals} deals. Methodology: ACV.",
            None,
            None,
        ),
        (
            fmt_eur(due_acv),
            "Due this quarter",
            f"{due_ct} deals closing in the current quarter.",
            None,
            AMBER if due_ct > 0 else GREEN_OK,
        ),
        (
            fmt_eur(over_acv),
            "Overdue carryover",
            f"{over_ct} deals carried over from prior quarters.",
            None,
            AMBER if over_ct > 0 else GREEN_OK,
        ),
        (
            fmt_eur(crit_acv),
            "Critical risk",
            f"{crit_ct} flagged as critical escalation.",
            None,
            AMBER if crit_ct > 0 else GREEN_OK,
        ),
    ]
    for i, (big, lbl, ctx, prog, status) in enumerate(cards):
        add_kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            big,
            lbl,
            ctx,
            progress_pct=prog,
            status_color=status,
        )

    bullets = (
        f"•  All {deals} open renewals carry some level of risk this cycle.\n"
        f"•  {over_ct} overdue carryover deals are the largest concentration of escalation pressure.\n"
        "•  Value methodology: Renewal ACV per SimCorp standard."
    )
    add_takeaway(slide, 4.35, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


def build_churn(prs, snap):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    trend = snap.get("churn", {}).get("trend", [])
    last_three = (
        trend[-3:]
        if len(trend) >= 3
        else trend
        + [{"quarter_label": "n/a", "churned_acv": 0, "churned_deals": 0}]
        * (3 - len(trend))
    )
    finance_status = snap.get("churn", {}).get("finance_feed_status", "pending")

    slide_header(
        slide,
        eyebrow="Q1 2026   ·   04   ·   CHURN",
        title="Churn trend, last three quarters",
        narrative=(
            f"Finance feed status: {finance_status}. "
            "CRM-side churn shown below. Finance overlay is pending this cycle."
        ),
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 3.95
    gap_bet = 0.22
    start_x = 0.6
    for i, q in enumerate(last_three):
        qlabel = q.get("quarter_label", "n/a")
        churned = q.get("churned_acv", 0) or 0
        deals = q.get("churned_deals", 0) or 0
        prev_churned = last_three[i - 1].get("churned_acv", 0) if i > 0 else 0
        delta = churned - prev_churned
        if i == 0:
            delta_ctx = "Baseline for the 3 quarter window."
        else:
            delta_ctx = (
                f"Change vs prior quarter: {'down' if delta < 0 else 'up'} "
                f"{fmt_eur(abs(delta))}"
            )
        status = None
        if i == len(last_three) - 1:
            status = GREEN_OK if delta < 0 else AMBER
        add_kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            fmt_eur(churned),
            f"{qlabel}  ·  {deals} deals",
            delta_ctx,
            status_color=status,
        )

    bullets = (
        f"•  Latest quarter ({last_three[-1].get('quarter_label')}) churned ACV: {fmt_eur(last_three[-1].get('churned_acv', 0))}.\n"
        "•  CRM signal only. Full churn narrative requires the Finance overlay when available.\n"
        "•  Value methodology: Renewal ACV."
    )
    add_takeaway(slide, 4.35, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


def build_slipped(prs, snap):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    slip = snap.get("slipped_deals", {}).get("summary_by_region", {})
    order = ["EMEA", "North America", "APAC"]
    biggest_region = max(
        order, key=lambda r: slip.get(r, {}).get("slipped_arr", 0) or 0, default="n/a"
    )
    biggest_arr = slip.get(biggest_region, {}).get("slipped_arr", 0) or 0

    slide_header(
        slide,
        eyebrow="Q1 2026   ·   05   ·   SLIPPED DEALS",
        title="Slipped exposure by region",
        narrative=(
            f"Largest slipped pool: {biggest_region} at {fmt_eur(biggest_arr)}. "
            "Root cause commentary pending from opportunity owners."
        ),
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 3.95
    gap_bet = 0.22
    start_x = 0.6
    for i, region in enumerate(order):
        data = slip.get(region, {})
        arr = data.get("slipped_arr", 0) or 0
        count = data.get("slipped_opp_count", 0) or 0
        pushes = data.get("avg_push_count", 0) or 0
        status = AMBER if region == biggest_region and count > 0 else None
        add_kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            fmt_eur(arr),
            f"{region}  ·  {count} deals",
            f"Average push count: {pushes:.1f}",
            status_color=status,
        )

    bullets = (
        f"•  {biggest_region} carries the largest slipped exposure at {fmt_eur(biggest_arr)}.\n"
        "•  Close Date Slipped by Stage is the live drilldown in the appendix.\n"
        "•  Owner commentary for root cause is still pending this cycle."
    )
    add_takeaway(slide, 4.35, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


def build_appendix(prs):
    slide = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    add_rect(slide, 0.6, 0.42, 0.08, 0.7, AQUA)
    add_text(
        slide, 0.82, 0.38, 12, 0.28, "APPENDIX", size=10, bold=True, color=PRIMARY_BLUE
    )
    add_text(
        slide,
        0.82,
        0.62,
        12,
        0.58,
        "Live Salesforce links",
        size=28,
        bold=True,
        color=NAVY,
    )
    add_text(
        slide,
        0.82,
        1.25,
        12.1,
        0.4,
        "Click any link to open the live dashboard or report.",
        size=13,
        color=GREY_TEXT,
    )

    body_tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.9), Inches(12.1), Inches(5)
    )
    body_tf = body_tb.text_frame
    body_tf.word_wrap = True
    for i, (label, url) in enumerate(LINKS):
        para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        run = para.add_run()
        run.text = label
        run.font.name = "Microsoft Sans Serif"
        run.font.size = Pt(13)
        run.font.color.rgb = PRIMARY_BLUE
        run.font.underline = True
        run.hyperlink.address = url
        para.space_after = Pt(8)


def build():
    with open(SNAPSHOT) as f:
        snap = json.load(f)
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    build_cover(prs, snap)
    build_exec_summary(prs, snap)
    build_region(prs, snap, "EMEA", "01")
    build_region(prs, snap, "North America", "01")
    build_region(prs, snap, "APAC", "01")
    build_commercial_approval(prs, snap)
    build_renewals(prs, snap)
    build_churn(prs, snap)
    build_slipped(prs, snap)
    build_appendix(prs)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
