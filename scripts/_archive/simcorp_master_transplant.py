#!/usr/bin/env python3
"""Transplant slide content from a working pptx onto the SimCorp template's
Blank layout so existing shapes inherit the SimCorp master (brand background,
theme colors, fonts) while the data stays exactly where it was.

Approach:
- Open the source deck (the canonical 812KB pptxgenjs deck).
- Open the SimCorp template, clear its sample slides.
- For each source slide:
  1. Add a new slide using the SimCorp "Blank" layout (idx 24, 0 placeholders,
     inherits SimCorp master with brand chrome).
  2. Drop any inherited placeholders.
  3. Re-add pictures via image blob (avoids dangling rId references).
  4. Deep-copy all other shape XML elements onto the new slide's spTree.
- Save the result.

Run:
    python3 scripts/simcorp_master_transplant.py [source.pptx] [output.pptx]
"""

from __future__ import annotations

import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

TEMPLATE_PATH = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"

DEFAULT_SOURCE = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01/sales_director_monthly_pipeline_insights_2026-04-01.pptx"
)
DEFAULT_OUTPUT = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01/sales_director_monthly_pipeline_insights_2026-04-01_simcorp.pptx"
)


def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    raise KeyError("Blank layout not found in template")


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def transplant_slide(src_slide, dst_prs, blank_layout):
    """Create a new slide in dst_prs and transplant shapes from src_slide."""
    new_slide = dst_prs.slides.add_slide(blank_layout)

    # Drop any inherited placeholders so the canvas is fully clean
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    pic_skipped = 0
    other_count = 0
    skipped = 0

    for shape in src_slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Skip pictures: the canonical deck embedded the SimCorp logo
                # on every slide via pptxgenjs, but the SimCorp template
                # master already provides branding chrome. Re-adding these
                # pictures would create duplicate logos.
                pic_skipped += 1
                continue
            # Deep copy the XML element for shapes, text, tables, etc.
            new_el = deepcopy(shape._element)
            new_slide.shapes._spTree.append(new_el)
            other_count += 1
        except Exception as e:
            skipped += 1
            print(f"    SKIP shape {shape.shape_type}: {e}")

    return pic_skipped, other_count, skipped


def main():
    src_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_SOURCE
    dst_path = Path(sys.argv[2]) if len(sys.argv) > 2 else DEFAULT_OUTPUT

    print(f"Source:   {src_path}")
    print(f"Template: {TEMPLATE_PATH}")
    print(f"Output:   {dst_path}")
    print()

    src = Presentation(str(src_path))
    dst = Presentation(TEMPLATE_PATH)

    print(f"Source slides:    {len(src.slides)}")
    print(f"Template slides:  {len(dst.slides)} (will be cleared)")

    clear_slides(dst)
    blank = get_blank_layout(dst)
    print(f"Using layout:     {blank.name!r}")
    print()

    for i, src_slide in enumerate(src.slides, 1):
        pic_skipped, other, skipped = transplant_slide(src_slide, dst, blank)
        print(
            f"  Slide {i}: transplanted {other} shapes, dropped {pic_skipped} pictures (logos)"
            f"{f' (skipped {skipped})' if skipped else ''}"
        )

    dst.save(str(dst_path))
    print()
    print(f"Saved {len(dst.slides)} slides to:")
    print(f"  {dst_path}")


if __name__ == "__main__":
    main()
