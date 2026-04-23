"""Preview deck: 3 analyses not yet in the per-director deck.

Scope matches the main deck's Q1 retrospective (Land only, Q1 2026 close
dates, one director) so numbers reconcile to slide 4 / 5. Stage Conversion
stays global because the per-territory sample is too small for the advance-
rate denominators to be meaningful, and the slide is explicitly labelled.

Output: output/simcorp_director_decks/<date>/preview/apac-missing-analyses.pptx
"""

from __future__ import annotations

import sys
from datetime import datetime
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt

# Re-use the main builder's helpers so styling is identical.
sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_deck_from_excel import (  # type: ignore
    DEFAULT_TEMPLATE,
    LY_TITLE_1,
    LY_TITLE_CONTENT,
    LY_END_SLIDE,
    _set_ph,
    _add_table,
    _fmt_eur,
    _meur,
)


ROOT = Path(__file__).resolve().parents[1]
DASH = ROOT / "output/sharepoint/Dashboard and Q1 Analysis.xlsx"
FY26 = ROOT / "output/sharepoint/FY26 Pipeline Review, All Territories.xlsx"
DIRECTOR_WB = ROOT / "output/director_live_workbooks/2026-04-16/jesper-tyrer.xlsx"

TERRITORY = "APAC"  # Matches the Region column on Dashboard workbook tabs
DIRECTOR_NAME = "Jesper Tyrer"  # Matches the Director column on Slip Risk

# Stage ordering for "highest stage reached" lookups. Must stay in sync with
# STAGE_ORDER in build_sharepoint_analysis.py.
STAGE_RANK = {
    "1 - Prospecting": 1,
    "2 - Discovery": 2,
    "3 - Engagement": 3,
    "4 - Shortlisted": 4,
    "5 - Preferred": 5,
    "6 - Contracting": 6,
    "8 - Won": 8,
    "0 - Lost": -1,
    "0 - No Opportunity": -2,
}
STAGES_IN_ORDER = [
    "1 - Prospecting",
    "2 - Discovery",
    "3 - Engagement",
    "4 - Shortlisted",
    "5 - Preferred",
    "6 - Contracting",
]


def _read_rows(wb_path, sheet_name, header_row=3):
    wb = load_workbook(wb_path, data_only=True, read_only=True)
    ws = wb[sheet_name]
    rows = list(ws.iter_rows(values_only=True))
    wb.close()
    headers = list(rows[header_row - 1])
    data = []
    for r in rows[header_row:]:
        if not r or all(v is None for v in r):
            continue
        d = {h: r[i] if i < len(r) else None for i, h in enumerate(headers) if h}
        data.append(d)
    return data, headers


def _cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_1])
    for ph in slide.placeholders:
        try:
            idx = ph.placeholder_format.idx
        except AttributeError:
            continue
        if idx == 0:
            ph.text = "APAC, Proposed Additional Slides"
        elif idx == 1:
            ph.text = (
                f"Preview of three analyses currently in the workbook but not "
                f"in the deck. {datetime.now().strftime('%d %B %Y')}."
            )


def _q1_land_losses_for_director(director_wb_path):
    """Pull Q1 2026 Land losses from a director workbook, matching the
    scope of slide 4 (Q1 Promised vs Delivered). Returns list of dicts:
    account, opportunity, owner, close_date, arr, reason, competitor.
    """
    wb = load_workbook(str(director_wb_path), read_only=True, data_only=True)
    ws = wb["Won Lost FY26"]
    rows = list(ws.iter_rows(values_only=True))
    wb.close()
    headers = rows[0]
    losses = []
    for r in rows[1:]:
        d = dict(zip(headers, r))
        if str(d.get("Type", "")).strip().lower() != "land":
            continue
        if "Won" in str(d.get("Stage", "")):
            continue
        cd = str(d.get("Close Date", "") or "")[:10]
        if not (cd >= "2026-01-01" and cd <= "2026-03-31"):
            continue
        losses.append(
            {
                "account": d.get("Account", ""),
                "opportunity": d.get("Opportunity", ""),
                "owner": d.get("Owner", ""),
                "close_date": cd,
                "arr": float(d.get("ARR Unweighted (EUR)") or 0),
                "reason": str(d.get("Reason") or "(not recorded)"),
                "competitor": str(d.get("Lost To Competitor") or ""),
            }
        )
    return losses


def _highest_stage_per_opp(opp_names):
    """Cross-ref Dashboard's Q1 History Raw to find the highest stage each
    opp ever reached (before going to Lost). Returns {opp_name: stage_str}.
    """
    if not opp_names:
        return {}
    wb = load_workbook(DASH, read_only=True, data_only=True)
    ws = wb["Q1 History Raw"]
    rows = list(ws.iter_rows(values_only=True))
    wb.close()
    headers = rows[0]
    opp_names_set = set(opp_names)
    best = {}
    for r in rows[1:]:
        d = dict(zip(headers, r))
        if d.get("Opportunity") not in opp_names_set:
            continue
        if d.get("Field Changed") != "StageName":
            continue
        for fld in ("Old Value", "New Value"):
            stage = str(d.get(fld) or "")
            rank = STAGE_RANK.get(stage, 0)
            if rank >= 1:  # only real positive stages count as "reached"
                cur = best.get(d.get("Opportunity"))
                if cur is None or rank > STAGE_RANK.get(cur, 0):
                    best[d.get("Opportunity")] = stage
    return best


def slide_why_we_lost(prs):
    """Loss Reasons + Stage at Loss, Q1 2026 Land losses for this director.

    Scope matches main deck slide 4. Left table: loss reasons with count and
    ARR. Right table: stage reached before losing (from history), with count
    and ARR sum per stage.
    """
    losses = _q1_land_losses_for_director(DIRECTOR_WB)
    opp_to_loss = {loss["opportunity"]: loss for loss in losses}
    highest = _highest_stage_per_opp(list(opp_to_loss.keys()))

    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    total_n = len(losses)
    total_arr = sum(loss["arr"] for loss in losses)
    _set_ph(
        slide,
        144,
        f"Q1 Land Losses, {DIRECTOR_NAME} ({TERRITORY}): {total_n} deals, "
        f"{_fmt_eur(total_arr)}",
    )
    _set_ph(
        slide,
        145,
        "Reconciles to slide 4 (Q1 Promised vs Delivered). Left: loss reasons "
        "with count and ARR. Right: stage reached before losing, from "
        "OpportunityFieldHistory. Diagnostic: early-stage losses = "
        "qualification gap; late-stage = execution gap.",
    )

    # Left table: Loss Reasons (Q1 Land only)
    from collections import defaultdict

    reason_agg = defaultdict(lambda: {"n": 0, "arr": 0.0})
    for loss in losses:
        reason_agg[loss["reason"]]["n"] += 1
        reason_agg[loss["reason"]]["arr"] += loss["arr"]
    reason_tbl = [["Loss Reason", "Count", "Lost ARR"]]
    for reason, agg in sorted(reason_agg.items(), key=lambda x: -x[1]["n"]):
        reason_tbl.append([reason, str(agg["n"]), _meur(agg["arr"])])
    reason_tbl.append(["Total", str(total_n), _meur(total_arr)])
    _add_table(
        slide,
        reason_tbl,
        0.3,
        2.0,
        6.4,
        col_widths=[3.8, 1.2, 1.4],
        data_font_size=Pt(10),
    )

    # Right table: Stage at Loss (highest stage reached per lost opp)
    stage_agg = defaultdict(lambda: {"n": 0, "arr": 0.0})
    unclassified_n = 0
    unclassified_arr = 0.0
    for loss in losses:
        stage = highest.get(loss["opportunity"])
        if stage:
            stage_agg[stage]["n"] += 1
            stage_agg[stage]["arr"] += loss["arr"]
        else:
            unclassified_n += 1
            unclassified_arr += loss["arr"]
    stage_tbl = [["Stage Reached Before Loss", "Count", "Lost ARR"]]
    for stg in STAGES_IN_ORDER:
        if stage_agg[stg]["n"] == 0 and stage_agg[stg]["arr"] == 0:
            continue
        stage_tbl.append([stg, str(stage_agg[stg]["n"]), _meur(stage_agg[stg]["arr"])])
    if unclassified_n:
        stage_tbl.append(
            ["(no stage history)", str(unclassified_n), _meur(unclassified_arr)]
        )
    stage_tbl.append(["Total", str(total_n), _meur(total_arr)])
    _add_table(
        slide,
        stage_tbl,
        6.9,
        2.0,
        6.0,
        col_widths=[3.0, 1.2, 1.8],
        data_font_size=Pt(10),
    )


def slide_slip_risk_owners(prs):
    """Slip Risk by Owner, filtered to this director's team."""
    rows, _ = _read_rows(FY26, "Slip Risk by Owner", header_row=4)
    mine = [r for r in rows if str(r.get("Director", "")) == DIRECTOR_NAME]
    mine.sort(key=lambda r: -int(r.get("Total Pushes") or 0))

    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    if not mine:
        _set_ph(slide, 144, f"Slip Risk Owners, {TERRITORY}")
        _set_ph(
            slide,
            145,
            "No owners in this territory breached the push threshold this run.",
        )
        return

    total_arr = sum(float(r.get("Open ARR (EUR)") or 0) for r in mine)
    total_pushes = sum(int(r.get("Total Pushes") or 0) for r in mine)
    _set_ph(
        slide,
        144,
        f"Slip Risk Owners, {TERRITORY}: {len(mine)} owners, "
        f"{total_pushes} total pushes across {_fmt_eur(total_arr)}",
    )
    _set_ph(
        slide,
        145,
        "Ranked by total push count across open Land deals. High-push owners "
        "carry the most slip exposure. Coaching conversation: push reasons, "
        "close-date discipline, stage-age.",
    )

    tbl = [
        ["Owner", "# Deals", "Open ARR", "Total Pushes", "Avg/Deal", "Max on 1 Deal"]
    ]
    for r in mine:
        tbl.append(
            [
                str(r.get("Owner", "")),
                str(int(r.get("# Deals") or 0)),
                _fmt_eur(r.get("Open ARR (EUR)") or 0),
                str(int(r.get("Total Pushes") or 0)),
                f"{float(r.get('Avg Push/Deal') or 0):.1f}",
                str(int(r.get("Max Pushes on 1 Deal") or 0)),
            ]
        )
    _add_table(
        slide,
        tbl,
        0.7,
        2.1,
        11.9,
        col_widths=[3.3, 1.3, 2.3, 1.8, 1.4, 1.8],
        data_font_size=Pt(11),
    )


def slide_stage_conversion(prs):
    """Stage Conversion funnel, global (source tab is global)."""
    rows, _ = _read_rows(DASH, "Stage Conversion")
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_CONTENT])
    _set_ph(slide, 144, "Stage Conversion Funnel, FY26 Year-to-Date (Global)")
    _set_ph(
        slide,
        145,
        "Deals that entered each stage vs. those that advanced to a later "
        "stage or reached Won. Low advance rate = bottleneck stage. Source "
        "is OpportunityFieldHistory; not territory-filtered because the "
        "sample per territory is too small to be reliable.",
    )
    tbl = [
        [
            "Stage",
            "Deals Entered",
            "Advanced",
            "Advance Rate",
            "Reached Won",
            "Win Rate",
        ]
    ]
    for r in rows:
        tbl.append(
            [
                str(r.get("Stage", "")),
                str(r.get("Deals Entered") or ""),
                str(r.get("Advanced to Next+") or ""),
                str(r.get("Advance Rate") or ""),
                str(r.get("Reached Won") or ""),
                str(r.get("Win Rate") or ""),
            ]
        )
    _add_table(
        slide,
        tbl,
        0.7,
        2.1,
        11.9,
        col_widths=[2.8, 1.8, 1.8, 2.0, 1.8, 1.7],
        data_font_size=Pt(11),
    )


def _end(prs):
    prs.slides.add_slide(prs.slide_layouts[LY_END_SLIDE])


def main():
    # Strip template sample slides exactly like the main builder does.
    prs = Presentation(str(DEFAULT_TEMPLATE))
    from lxml import etree

    pres_elem = prs.part._element
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    sldIdLst = pres_elem.find("p:sldIdLst", nsmap)
    if sldIdLst is not None:
        for sldId in list(sldIdLst):
            rId = sldId.get(etree.QName(r_ns, "id"))
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except KeyError:
                    pass
        for child in list(sldIdLst):
            sldIdLst.remove(child)

    _cover(prs)
    slide_why_we_lost(prs)
    slide_slip_risk_owners(prs)
    slide_stage_conversion(prs)
    _end(prs)

    today = datetime.now().strftime("%Y-%m-%d")
    out_dir = ROOT / "output/simcorp_director_decks" / today / "preview"
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "apac-missing-analyses.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
