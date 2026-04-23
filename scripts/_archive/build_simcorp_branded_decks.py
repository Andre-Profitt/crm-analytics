#!/usr/bin/env python3
"""Rebuild Report 1 and Report 2 PowerPoint decks using the official SimCorp
template (34 master layouts, GS2024 font scheme, 240306 Simcorp colors).

Template source: archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx
No em-dashes anywhere. ACV used for renewals (SimCorp methodology).

Outputs:
  output/sales_director_monthly_runs/.../sales_director_monthly_simcorp_branded.pptx
  output/sales_ops_quarterly_deck_2026-03-31/sales_ops_quarterly_simcorp_branded.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

TEMPLATE_PATH = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx"
REPO_ROOT = Path("/Users/test/crm-analytics")


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout not found: {name}")


def fill(slide, idx, text):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            return True
    return False


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def fmt_eur(amount):
    if amount is None:
        return "n/a"
    if abs(amount) >= 1_000_000_000:
        return f"EUR {amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"EUR {amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"EUR {amount / 1_000:.0f}K"
    return f"EUR {amount:.0f}"


def fmt_pct(pct):
    if pct is None:
        return "n/a"
    return f"{pct:.1f}%"


# ============================================================================
# REPORT 1: Sales Directors Monthly
# ============================================================================


def build_report1(snapshot_path, output_path):
    with open(snapshot_path) as f:
        snap = json.load(f)

    prs = Presentation(TEMPLATE_PATH)
    clear_slides(prs)

    quarter_focus = snap.get("quarter_focus", "Q1")
    snapshot_date = snap.get("snapshot_date", "2026-04-01")
    q_window = snap.get("quarter_window", {})
    days_remaining = q_window.get("days_remaining", 0)

    # Compute pipeline headline metrics for the exec summary
    regions = snap.get("pipeline", {}).get("deck_regions", {})
    total_target = sum(r.get("target_arr", 0) for r in regions.values())
    total_projected = sum(
        r.get("actual_arr", 0) + r.get("weighted_open_arr", 0) for r in regions.values()
    )
    gap = total_target - total_projected
    gap_pct = (gap / total_target * 100) if total_target else 0
    biggest_gap_region = snap.get("biggest_gap_region", "North America")
    biggest_gap_arr = snap.get("biggest_gap_arr", 0)
    weakest_conf_region = snap.get("weakest_confidence_region", "EMEA")
    weakest_conf_pct = snap.get("weakest_confidence_pct", 0)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
    fill(slide, 20, "Sales Directors Monthly Pipeline and Insights")
    fill(slide, 22, f"{quarter_focus} FY27 Forward Look")
    fill(slide, 24, f"Snapshot {snapshot_date}")

    # --- Slide 2: Executive Summary (answer-first) ---
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    fill(slide, 42, "Executive summary")
    # Big stat headers
    fill(slide, 61, fmt_eur(total_target))
    fill(slide, 62, fmt_eur(total_projected))
    fill(slide, 63, fmt_eur(gap))
    fill(slide, 64, f"{days_remaining}d")
    # Titles under each
    fill(slide, 42, f"Total target, {quarter_focus} FY27")
    fill(slide, 56, "Projected call")
    fill(slide, 58, "Gap to target")
    fill(slide, 60, "Days remaining")
    # Content bullets
    fill(
        slide,
        22,
        f"Biggest gap region: {biggest_gap_region} at {fmt_eur(biggest_gap_arr)}.",
    )
    fill(
        slide,
        55,
        f"Projected call is {(total_projected / total_target * 100 if total_target else 0):.0f}% of target across all regions.",
    )
    fill(slide, 57, f"Gap is {gap_pct:.0f}% of the total quarterly target.")
    fill(
        slide,
        59,
        f"Weakest forecast confidence: {weakest_conf_region} at {fmt_pct(weakest_conf_pct)}.",
    )

    # --- Slide 3: Agenda ---
    slide = prs.slides.add_slide(get_layout(prs, "Agenda 1"))
    agenda_idxs = [40, 69, 71, 74, 76, 101]
    agenda_items = [
        "Pipeline coverage by region",
        "Commercial approval overview",
        "Renewals tracking and risk",
        "Churn trends and outlook",
        "Slipped deals analysis",
        "Actions and follow ups",
    ]
    for idx, text in zip(agenda_idxs, agenda_items):
        fill(slide, idx, text)

    # --- Slide 3: Divider 01 Pipeline ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "01")
    fill(slide, 20, "Pipeline coverage by region")

    # --- Slide 4: Pipeline by region (3 x content with gradient line) ---
    slide = prs.slides.add_slide(get_layout(prs, "3 x content w/ gradient line"))
    fill(slide, 42, f"Pipeline coverage, {quarter_focus} FY27")
    regions = snap.get("pipeline", {}).get("deck_regions", {})
    region_order = ["APAC", "EMEA", "North America"]
    title_idxs = [42, 56, 58]
    content_idxs = [22, 55, 57]
    header_idxs = [61, 62, 63]
    for i, reg_name in enumerate(region_order):
        reg = regions.get(reg_name, {})
        target = reg.get("target_arr", 0)
        actual = reg.get("actual_arr", 0) + reg.get("weighted_open_arr", 0)
        coverage = reg.get("coverage_status", "n/a")
        confidence = reg.get("forecast_confidence_pct", 0)
        # Big stat header = coverage status
        fill(slide, header_idxs[i], coverage)
        # Title = region name
        fill(slide, title_idxs[i], reg_name)
        # Content = target/actual/confidence bullets
        content = (
            f"Target: {fmt_eur(target)}\n"
            f"Projected: {fmt_eur(actual)}\n"
            f"Forecast confidence: {fmt_pct(confidence)}"
        )
        fill(slide, content_idxs[i], content)

    # --- Slide 5: Divider 02 Commercial Approval ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "02")
    fill(slide, 20, "Commercial approval overview")

    # --- Slide 6: Commercial approval summary ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 42, "Commercial approval status")
    ca = snap.get("commercial_approval", {})
    candidate_count = ca.get("candidate_count", 0)
    approved_count = len(ca.get("approved_deals", []))
    control_exc = ca.get("control_exceptions", [])
    exc_count = len(control_exc) if isinstance(control_exc, list) else 0
    content = (
        f"Approved in quarter: {approved_count}\n"
        f"Candidates awaiting approval: {candidate_count}\n"
        f"Control exceptions flagged: {exc_count}\n\n"
        "Link: live candidate list on the Salesforce dashboard"
    )
    fill(slide, 22, content)

    # --- Slide 7: Divider 03 Renewals ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "03")
    fill(slide, 20, "Renewals tracking and risk")

    # --- Slide 8: Renewals summary ---
    slide = prs.slides.add_slide(get_layout(prs, "2 x content w/ gradient line"))
    fill(slide, 42, "Renewal pipeline")
    ren = snap.get("renewals", {})
    ren_sum = ren.get("summary", {}) if isinstance(ren.get("summary"), dict) else {}
    total_open_acv = (
        ren_sum.get("total_open_acv")
        or snap.get("total_open_renewal_pipeline_acv")
        or 0
    )
    critical_acv = ren_sum.get("critical_acv") or snap.get("critical_renewal_acv") or 0
    # Big stats
    fill(slide, 61, fmt_eur(total_open_acv))
    fill(slide, 62, fmt_eur(critical_acv))
    # Titles
    fill(slide, 42, "Open renewal ACV")
    fill(slide, 56, "Critical renewal ACV")
    # Content
    fill(
        slide,
        22,
        "Total ACV across all open renewals this quarter.\nValue methodology: Renewal ACV per SimCorp standard.",
    )
    fill(
        slide,
        55,
        "ACV of renewals flagged as critical risk (churn candidates, escalation required).",
    )

    # --- Slide 9: Divider 04 Churn ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "04")
    fill(slide, 20, "Churn trends and outlook")

    # --- Slide 10: Churn summary ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 42, "Churn overview")
    churn = snap.get("churn", {})
    finance_status = churn.get("finance_feed_status", "pending")
    note = (
        f"Finance feed status: {finance_status}\n\n"
        "Churn risk signals from CRM are tracked in the live Salesforce dashboard "
        "under Business At Risk and Customer Account Health.\n\n"
        "Full churn narrative requires the Finance overlay, which is pending this cycle."
    )
    fill(slide, 22, note)

    # --- Slide 11: Divider 05 Slipped ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "05")
    fill(slide, 20, "Slipped deals analysis")

    # --- Slide 12: Slipped deals summary ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 42, "Slipped deals, live view")
    biggest_region = snap.get("biggest_slipped_region") or "North America"
    biggest_arr = snap.get("biggest_slipped_arr") or 0
    content = (
        f"Biggest slipped region: {biggest_region}\n"
        f"Biggest slipped ARR: {fmt_eur(biggest_arr)}\n\n"
        "Live drilldown on the Salesforce dashboard:\n"
        "Close Date Slipped by Stage, Close Date Slipped Aging (5Y, Ae, KC)\n\n"
        "Root cause commentary is collected from opportunity owners via the "
        "owner commentary flow and attached to the dashboard drill views."
    )
    fill(slide, 22, content)

    # --- Slide 13: Actions and follow ups ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 42, "Actions and follow ups")
    publish_blockers = snap.get("external_inputs", {})
    blockers = []
    if publish_blockers.get("finance_churn", {}).get("status") != "complete":
        blockers.append("Finance churn feed pending from Finance owner")
    if publish_blockers.get("slipped_commentary", {}).get("status") != "complete":
        blockers.append("Slipped deal owner commentary pending")
    blockers_text = "\n".join(f"- {b}" for b in blockers) or "All inputs received"
    content = (
        f"Publish blockers this cycle:\n{blockers_text}\n\n"
        f"Methodology lock: renewals and churn in Renewal ACV, land and expand in ARR.\n\n"
        f"Next review: end of {quarter_focus} FY27."
    )
    fill(slide, 22, content)

    # --- Slide 14: Closing ---
    slide = prs.slides.add_slide(get_layout(prs, "End slide with disclaimer 1"))
    fill(slide, 28, "Thank you\n\nwww.simcorp.com")

    prs.save(str(output_path))
    print(f"Report 1 saved: {output_path}")
    print(f"  slides: {len(prs.slides)}")


# ============================================================================
# REPORT 2: Sales Ops Quarterly
# ============================================================================


def build_report2(summary_path, output_path):
    with open(summary_path) as f:
        summary = json.load(f)

    prs = Presentation(TEMPLATE_PATH)
    clear_slides(prs)

    snapshot_display = summary.get("snapshot_display", "April 1, 2026")
    metrics = summary.get("key_metrics", {})

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
    fill(slide, 20, "Sales Ops Quarterly Review")
    fill(
        slide,
        22,
        "Data quality, process compliance, forecast accuracy, pipeline hygiene",
    )
    fill(slide, 24, f"Snapshot {snapshot_display}")

    # --- Slide 2: Executive Summary (answer-first KPI strip) ---
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    # Big stat headers with the 4 KPI scores
    fill(slide, 61, f"{metrics.get('data_completeness_score', 0):.1f}")
    fill(slide, 62, f"{metrics.get('process_compliance_rate', 0):.1f}")
    fill(slide, 63, f"{metrics.get('forecast_accuracy', 0):.1f}")
    fill(slide, 64, f"{metrics.get('pipeline_hygiene_rate', 0):.1f}")
    # Titles
    fill(slide, 42, "Data completeness")
    fill(slide, 56, "Process compliance")
    fill(slide, 58, "Forecast accuracy")
    fill(slide, 60, "Pipeline hygiene")
    # Content: answer-first framing of the score
    top_area = metrics.get("top_exception_queue_area", "Commercial Process")
    top_count = metrics.get("top_exception_queue_count", 0)
    fill(
        slide,
        22,
        f"Strongest KPI. {metrics.get('data_quality_queue_count', 0)} exception items in the queue.",
    )
    fill(
        slide,
        55,
        f"Weakest control area: {top_area} with {top_count} flagged items.",
    )
    fill(
        slide,
        57,
        f"Forecast risk queue: {metrics.get('forecast_risk_queue_count', 0)} opportunities.",
    )
    fill(
        slide,
        59,
        f"Stale pipeline queue: {metrics.get('stale_pipeline_queue_count', 0)} opportunities.",
    )

    # --- Slide 3: Agenda ---
    slide = prs.slides.add_slide(get_layout(prs, "Agenda 1"))
    agenda_idxs = [40, 69, 71, 74, 76, 101]
    agenda_items = [
        "CRM data quality",
        "Process compliance",
        "Forecast accuracy",
        "Pipeline hygiene",
        "Exception queues",
        "Actions and follow ups",
    ]
    for idx, text in zip(agenda_idxs, agenda_items):
        fill(slide, idx, text)

    # --- Slide 3: Headline KPIs (4 x content w/ gradient line) ---
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    fill(slide, 42, "Headline KPIs this quarter")
    # Big stats
    fill(slide, 61, f"{metrics.get('data_completeness_score', 0):.1f}")
    fill(slide, 62, f"{metrics.get('process_compliance_rate', 0):.1f}")
    fill(slide, 63, f"{metrics.get('forecast_accuracy', 0):.1f}")
    fill(slide, 64, f"{metrics.get('pipeline_hygiene_rate', 0):.1f}")
    # Titles
    fill(slide, 42, "Data completeness")
    fill(slide, 56, "Process compliance")
    fill(slide, 58, "Forecast accuracy")
    fill(slide, 60, "Pipeline hygiene")
    # Content (bullet under each)
    fill(slide, 22, "Opportunity field completeness across required attributes")
    fill(slide, 55, "Stage transition and approval control adherence")
    fill(slide, 57, "Forecast vs closed won over rolling trailing windows")
    fill(slide, 59, "Aging, stage progression, stalled pipeline signals")

    # --- Slide 4: Divider 01 Data Quality ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "01")
    fill(slide, 20, "CRM data quality")

    # --- Slide 5: Data quality detail ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(
        slide,
        42,
        f"Data completeness score: {metrics.get('data_completeness_score', 0):.1f}",
    )
    dq_q = metrics.get("data_quality_queue_count", 0)
    content = (
        f"Exception queue: {dq_q} opportunities flagged\n\n"
        "Live drilldown on the Salesforce dashboard:\n"
        "Missing Quote Type, Missing Amount on Open Opps, Missing Decision Reason, "
        "Won Loss Info Missing CFQ, KYC Not Completed"
    )
    fill(slide, 22, content)

    # --- Slide 6: Divider 02 Process Compliance ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "02")
    fill(slide, 20, "Process compliance")

    # --- Slide 7: Process compliance detail ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(
        slide,
        42,
        f"Process compliance rate: {metrics.get('process_compliance_rate', 0):.1f}",
    )
    cp_q = metrics.get("commercial_process_queue_count", 0)
    content = (
        f"Top exception queue: {metrics.get('top_exception_queue_area', 'Commercial Process')}\n"
        f"Exception count: {cp_q} opportunities\n\n"
        "Live drilldown on the Salesforce dashboard:\n"
        "Overdue Opportunities, Overdue Close Date Open Opps, No Activity 30 Plus Days, "
        "Probability Mismatch by Stage"
    )
    fill(slide, 22, content)

    # --- Slide 8: Divider 03 Forecast Accuracy ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "03")
    fill(slide, 20, "Forecast accuracy")

    # --- Slide 9: Forecast accuracy detail ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 42, f"Forecast accuracy: {metrics.get('forecast_accuracy', 0):.1f}")
    fr_q = metrics.get("forecast_risk_queue_count", 0)
    content = (
        f"Forecast risk exception queue: {fr_q} opportunities\n\n"
        "Live drilldown on the Salesforce dashboard:\n"
        "Forecast Accuracy, Low Probability In Quarter, Forecast and Closed Won"
    )
    fill(slide, 22, content)

    # --- Slide 10: Divider 04 Pipeline Hygiene ---
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "04")
    fill(slide, 20, "Pipeline hygiene")

    # --- Slide 11: Pipeline hygiene detail ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(
        slide,
        42,
        f"Pipeline hygiene rate: {metrics.get('pipeline_hygiene_rate', 0):.1f}",
    )
    sp_q = metrics.get("stale_pipeline_queue_count", 0)
    content = (
        f"Stale pipeline queue: {sp_q} opportunities\n\n"
        "Live drilldown on the Salesforce dashboard:\n"
        "Aging Pipeline 365 Plus Days, Stale Opportunities, High Value Stale Deals"
    )
    fill(slide, 22, content)

    # --- Slide 12: Exception queue rollup ---
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 42, "Exception queue rollup")
    top_area = metrics.get("top_exception_queue_area", "Commercial Process")
    top_count = metrics.get("top_exception_queue_count", 0)
    content = (
        f"Top exception area: {top_area}\n"
        f"Total flagged items: {top_count}\n\n"
        f"Data quality queue: {metrics.get('data_quality_queue_count', 0)}\n"
        f"Commercial process queue: {metrics.get('commercial_process_queue_count', 0)}\n"
        f"Forecast risk queue: {metrics.get('forecast_risk_queue_count', 0)}\n"
        f"Stale pipeline queue: {metrics.get('stale_pipeline_queue_count', 0)}"
    )
    fill(slide, 22, content)

    # --- Slide 13: Closing ---
    slide = prs.slides.add_slide(get_layout(prs, "End slide with disclaimer 1"))
    fill(slide, 28, "Thank you\n\nwww.simcorp.com")

    prs.save(str(output_path))
    print(f"Report 2 saved: {output_path}")
    print(f"  slides: {len(prs.slides)}")


def main():
    r1_snapshot = (
        REPO_ROOT
        / "output/sales_director_monthly_runs/2026-04-06T18-31-12Z_2026-04-01/report1_snapshot.json"
    )
    r1_output = (
        REPO_ROOT
        / "output/sales_director_monthly_runs/2026-04-06T18-31-12Z_2026-04-01/sales_director_monthly_simcorp_branded.pptx"
    )
    build_report1(r1_snapshot, r1_output)

    r2_summary = (
        REPO_ROOT
        / "output/sales_ops_quarterly_deck_2026-03-31/sales_ops_quarterly_review_2026-04-01.summary.json"
    )
    r2_output = (
        REPO_ROOT
        / "output/sales_ops_quarterly_deck_2026-03-31/sales_ops_quarterly_simcorp_branded.pptx"
    )
    build_report2(r2_summary, r2_output)


if __name__ == "__main__":
    main()
