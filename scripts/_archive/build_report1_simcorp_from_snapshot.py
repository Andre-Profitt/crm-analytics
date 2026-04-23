#!/usr/bin/env python3
"""Build Report 1 (Sales Director Monthly) as a SimCorp-branded deck by
pouring snapshot data into the SimCorp template's pre-styled layouts.

Approach:
- The brand styling lives entirely in the SimCorp template's slide masters.
- We just pick the right layout per slide and fill placeholders with real
  data from report1_snapshot.json.
- 8 content slides + closing, mapped 1:1 to the pipeline / approval /
  renewals / churn / slipped sections.
- No em-dashes anywhere. Renewals expressed in ACV per SimCorp methodology.

Run:
    python3 scripts/build_report1_simcorp_from_snapshot.py [snapshot_path] [output_path]
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation

TEMPLATE_PATH = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx"
DEFAULT_SNAPSHOT = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01/report1_snapshot.json"
)
DEFAULT_OUTPUT = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01/sales_director_monthly_simcorp.pptx"
)


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout not found: {name}")


def fill(slide, idx, text):
    if text is None:
        return False
    text = str(text).replace("\u2014", "-").replace("\u2013", "-")
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            return True
    return False


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def fmt_eur(amount):
    if amount is None:
        return "n/a"
    if abs(amount) >= 1_000_000_000:
        return f"EUR {amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"EUR {amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"EUR {amount / 1_000:.0f}K"
    return f"EUR {amount:.0f}"


def fmt_int(n):
    return f"{n:,}" if n is not None else "n/a"


def fmt_pct(p):
    return f"{p:.1f}%" if p is not None else "n/a"


def build(snapshot_path: Path, output_path: Path):
    with open(snapshot_path) as f:
        snap = json.load(f)

    prs = Presentation(TEMPLATE_PATH)
    clear_slides(prs)

    snapshot_date = snap.get("snapshot_date", "2026-04-01")
    quarter_focus = snap.get("quarter_focus", "Q1")
    q_window = snap.get("quarter_window", {})
    days_remaining = q_window.get("days_remaining", 0)

    regions = snap.get("pipeline", {}).get("deck_regions", {})
    total_target = sum(r.get("target_arr", 0) or 0 for r in regions.values())
    total_call = sum((r.get("best_case_call_arr") or 0) for r in regions.values())
    gap = total_target - total_call
    coverage_pct = (total_call / total_target * 100) if total_target else 0

    ren = snap.get("renewals", {}).get("summary", {})
    ca = snap.get("commercial_approval", {}).get("summary", {})
    slipped_by_region = snap.get("slipped_deals", {}).get("summary_by_region", {})
    churn_trend = snap.get("churn", {}).get("trend", [])

    # ================================================================
    # SLIDE 1: Title
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
    fill(slide, 20, "Sales Director Monthly Pipeline and Insights")
    fill(
        slide,
        22,
        f"{quarter_focus} FY27 forward look. Pipeline, approvals, renewals, churn, slipped deals.",
    )
    fill(slide, 24, f"Snapshot {snapshot_date}")

    # ================================================================
    # SLIDE 2: Executive summary (4 KPI cards)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    # Big stat headers
    fill(slide, 61, fmt_eur(total_target))
    fill(slide, 62, fmt_eur(total_call))
    fill(slide, 63, fmt_eur(gap))
    fill(slide, 64, f"{days_remaining}d")
    # Card titles
    fill(slide, 42, "Quarter target")
    fill(slide, 56, "Best case call")
    fill(slide, 58, "Gap to target")
    fill(slide, 60, "Days remaining")
    # Card context
    fill(slide, 22, f"Total ARR target across all 3 regions for {quarter_focus} FY27.")
    fill(
        slide,
        55,
        f"Combined regional best case call. Coverage at {coverage_pct:.0f}% of target.",
    )
    fill(
        slide,
        57,
        f"Below target by {(abs(gap) / total_target * 100 if total_target else 0):.0f}%.",
    )
    fill(slide, 59, "Days left in quarter to close the gap.")

    # ================================================================
    # SLIDE 3: Pipeline coverage by region (3 cards)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "3 x content w/ gradient line"))
    region_order = ["EMEA", "North America", "APAC"]
    title_idxs = [42, 56, 58]
    content_idxs = [22, 55, 57]
    header_idxs = [61, 62, 63]
    for i, reg_name in enumerate(region_order):
        reg = regions.get(reg_name, {})
        target = reg.get("target_arr", 0) or 0
        call = reg.get("best_case_call_arr", 0) or 0
        confidence = reg.get("forecast_confidence_pct", 0) or 0
        coverage = reg.get("coverage_status", "n/a") or "n/a"
        # Big stat = call / target ratio as a coverage percent
        ratio_pct = (call / target * 100) if target else 0
        fill(slide, header_idxs[i], f"{ratio_pct:.0f}%")
        fill(slide, title_idxs[i], reg_name)
        fill(
            slide,
            content_idxs[i],
            f"Target: {fmt_eur(target)}\n"
            f"Best case call: {fmt_eur(call)}\n"
            f"Forecast confidence: {fmt_pct(confidence)}\n"
            f"Status: {coverage}",
        )

    # ================================================================
    # SLIDE 4: Commercial approval (3 cards)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "3 x content w/ gradient line"))
    fill(slide, 61, fmt_int(ca.get("approved_count", 0)))
    fill(slide, 62, fmt_int(ca.get("pending_count", 0)))
    fill(slide, 63, fmt_int(ca.get("stale_count", 0)))
    fill(slide, 42, "Approved")
    fill(slide, 56, "Pending")
    fill(slide, 58, "Stale")
    fill(
        slide,
        22,
        f"Approved value: {fmt_eur(ca.get('approved_arr', 0))}\nLand stage 3 deals approved this quarter.",
    )
    fill(
        slide,
        55,
        f"Pending value: {fmt_eur(ca.get('pending_arr', 0))}\nAwaiting commercial approval.",
    )
    fill(
        slide,
        57,
        f"Stale value: {fmt_eur(ca.get('stale_arr', 0))}\nRequires cleanup or escalation.",
    )

    # ================================================================
    # SLIDE 5: Renewals (4 cards)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    fill(slide, 61, fmt_eur(ren.get("total_renewal_pipeline_acv", 0)))
    fill(slide, 62, fmt_eur(ren.get("due_this_quarter_acv", 0)))
    fill(slide, 63, fmt_eur(ren.get("overdue_carryover_acv", 0)))
    fill(slide, 64, fmt_eur(ren.get("critical_acv", 0)))
    fill(slide, 42, "Total renewal pipeline")
    fill(slide, 56, "Due this quarter")
    fill(slide, 58, "Overdue carryover")
    fill(slide, 60, "Critical risk")
    fill(
        slide,
        22,
        f"{fmt_int(ren.get('total_deals', 0))} open renewal deals. Methodology: ACV per SimCorp standard.",
    )
    fill(
        slide,
        55,
        f"{fmt_int(ren.get('due_this_quarter_count', 0))} deals closing this quarter.",
    )
    fill(
        slide,
        57,
        f"{fmt_int(ren.get('overdue_carryover_count', 0))} deals carried over from prior quarters.",
    )
    fill(
        slide,
        59,
        f"{fmt_int(ren.get('critical_count', 0))} flagged as critical risk and escalation required.",
    )

    # ================================================================
    # SLIDE 6: Churn trend (3 cards = last 3 quarters)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "3 x content w/ gradient line"))
    last_three = churn_trend[-3:] if len(churn_trend) >= 3 else churn_trend
    while len(last_three) < 3:
        last_three.insert(
            0, {"quarter_label": "n/a", "churned_acv": 0, "churned_deals": 0}
        )
    for i, q in enumerate(last_three):
        fill(slide, header_idxs[i], fmt_eur(q.get("churned_acv", 0)))
        fill(slide, title_idxs[i], q.get("quarter_label", "n/a"))
        fill(
            slide,
            content_idxs[i],
            f"{fmt_int(q.get('churned_deals', 0))} closed lost renewal deals this quarter.\nValue methodology: Renewal ACV.",
        )

    # ================================================================
    # SLIDE 7: Slipped deals by region (3 cards)
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "3 x content w/ gradient line"))
    for i, reg_name in enumerate(region_order):
        reg_slip = slipped_by_region.get(reg_name, {})
        slipped_arr = reg_slip.get("slipped_arr", 0) or 0
        slipped_count = reg_slip.get("slipped_opp_count", 0) or 0
        avg_push = reg_slip.get("avg_push_count", 0) or 0
        fill(slide, header_idxs[i], fmt_eur(slipped_arr))
        fill(slide, title_idxs[i], reg_name)
        fill(
            slide,
            content_idxs[i],
            f"{fmt_int(slipped_count)} slipped opportunities.\nAverage push count: {avg_push:.1f}\nLive drilldown on the Salesforce dashboard.",
        )

    # ================================================================
    # SLIDE 8: Closing
    # ================================================================
    slide = prs.slides.add_slide(get_layout(prs, "End slide with disclaimer 1"))
    fill(slide, 28, "Thank you\n\nwww.simcorp.com")

    prs.save(str(output_path))
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


def main():
    snapshot_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_SNAPSHOT
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else DEFAULT_OUTPUT
    build(snapshot_path, output_path)


if __name__ == "__main__":
    main()
