#!/usr/bin/env python3
"""Build a single sample slide ("EMEA Regional Outlook") at production
SimCorp spec, using the Commercial Update template's 'Title Only' layout
as the base and manually adding content at production sizing.

Design decisions:
- 28pt title, bold, #011946 (SimCorp dark navy)
- Aqua accent bar #6FCCDD under the title
- 4 KPI cards with 32pt big numbers, 14pt labels, 10pt context
- Light blue card backgrounds (#E6EEFE)
- Breathing room: ~0.6 in margins, whitespace between cards
- One narrative line above the KPI strip
- One insight callout below the KPI strip
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"
RUN_DIR = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01"
)
SNAPSHOT = RUN_DIR / "report1_snapshot.json"
OUTPUT = RUN_DIR / "sample_slide_emea_v1.pptx"

# SimCorp production color tokens
NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)


def fmt_eur(amount):
    if amount is None:
        return "n/a"
    if abs(amount) >= 1_000_000_000:
        return f"EUR {amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"EUR {amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"EUR {amount / 1_000:.0f}K"
    return f"EUR {amount:.0f}"


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def set_font(run, name="Microsoft Sans Serif", size=14, bold=False, color=NAVY):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)
    return tb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    return shape


def add_kpi_card(slide, left, top, width, height, big_number, label, context):
    """Production-spec KPI card: light blue panel, big stat + label + context."""
    # Background panel
    add_rect(slide, left, top, width, height, LIGHT_BLUE_PANEL)

    # Big number (32pt, primary blue)
    num_tb = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.7)
    )
    num_tf = num_tb.text_frame
    num_tf.margin_left = num_tf.margin_right = Pt(0)
    num_tf.margin_top = num_tf.margin_bottom = Pt(0)
    p = num_tf.paragraphs[0]
    r = p.add_run()
    r.text = big_number
    set_font(r, size=32, bold=True, color=PRIMARY_BLUE)

    # Label (14pt, dark navy)
    lbl_tb = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.85), Inches(width - 0.4), Inches(0.35)
    )
    lbl_tf = lbl_tb.text_frame
    lbl_tf.margin_left = lbl_tf.margin_right = Pt(0)
    lbl_tf.margin_top = lbl_tf.margin_bottom = Pt(0)
    lbl_tf.word_wrap = True
    p = lbl_tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    set_font(r, size=14, bold=True, color=NAVY)

    # Context (10pt, secondary)
    ctx_tb = slide.shapes.add_textbox(
        Inches(left + 0.2),
        Inches(top + 1.20),
        Inches(width - 0.4),
        Inches(height - 1.25),
    )
    ctx_tf = ctx_tb.text_frame
    ctx_tf.margin_left = ctx_tf.margin_right = Pt(0)
    ctx_tf.margin_top = ctx_tf.margin_bottom = Pt(0)
    ctx_tf.word_wrap = True
    p = ctx_tf.paragraphs[0]
    r = p.add_run()
    r.text = context
    set_font(r, size=10, bold=False, color=PRIMARY_BLUE)


def build():
    with open(SNAPSHOT) as f:
        snap = json.load(f)

    emea = snap["pipeline"]["deck_regions"]["EMEA"]
    target = emea["target_arr"]
    best_case_call = emea["best_case_call_arr"]
    confidence = emea["forecast_confidence_pct"]
    coverage = emea["coverage_status"]
    gap = target - best_case_call

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # Start from Title Only layout (production base with just the title bar)
    layout = get_layout(prs, "Title Only")
    slide = prs.slides.add_slide(layout)

    # Remove the layout's inherited placeholders so we have a clean canvas
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # Title (28pt bold navy)
    add_text(
        slide,
        left=0.6,
        top=0.4,
        width=12.1,
        height=0.6,
        text="EMEA Regional Outlook",
        size=28,
        bold=True,
        color=NAVY,
    )

    # Aqua accent bar under title
    add_rect(slide, left=0.6, top=1.05, width=1.2, height=0.08, fill_color=AQUA)

    # Narrative line (14pt primary blue)
    add_text(
        slide,
        left=0.6,
        top=1.3,
        width=12.1,
        height=0.5,
        text=(
            f"EMEA is promotable. Best case call {fmt_eur(best_case_call)} "
            f"against {fmt_eur(target)} quarter target. "
            f"Confidence {confidence:.1f}%."
        ),
        size=14,
        bold=False,
        color=PRIMARY_BLUE,
    )

    # 4 KPI cards in a strip at y=2.1
    card_top = 2.2
    card_h = 1.85
    card_w = 2.95
    gap_between = 0.18
    start_x = 0.6
    cards = [
        (
            fmt_eur(target),
            "Quarter target",
            "Executive target seam (ARR)",
        ),
        (
            fmt_eur(best_case_call),
            "Best case call",
            f"{(best_case_call / target * 100):.0f}% of target coverage",
        ),
        (
            fmt_eur(gap),
            "Gap to target",
            "Needed from promotion or new pipeline",
        ),
        (
            f"{confidence:.1f}%",
            "Forecast confidence",
            f"Coverage status: {coverage}",
        ),
    ]
    for i, (num, lbl, ctx) in enumerate(cards):
        add_kpi_card(
            slide,
            left=start_x + i * (card_w + gap_between),
            top=card_top,
            width=card_w,
            height=card_h,
            big_number=num,
            label=lbl,
            context=ctx,
        )

    # Insight callout at bottom (small rectangle with aqua fill + narrative)
    callout_top = 4.5
    callout_h = 1.8
    add_rect(
        slide,
        left=0.6,
        top=callout_top,
        width=12.1,
        height=callout_h,
        fill_color=LIGHT_BLUE_PANEL,
    )
    # Accent aqua bar on the left edge of the callout
    add_rect(
        slide, left=0.6, top=callout_top, width=0.12, height=callout_h, fill_color=AQUA
    )

    # Callout title
    add_text(
        slide,
        left=0.95,
        top=callout_top + 0.12,
        width=11.6,
        height=0.35,
        text="What to watch this week",
        size=14,
        bold=True,
        color=NAVY,
    )
    # Callout body
    bullets_text = (
        "•  EMEA best case call is 85% of quarter target; forecast confidence at 94.5% is the weakest of the 3 regions.\n"
        "•  CACEIS Netherlands (owner: Michiel van der Berg) and Groupama AM are the two largest promotion candidates.\n"
        "•  29 days remaining in Q1 FY27 to close the promotion-dependent gap."
    )
    add_text(
        slide,
        left=0.95,
        top=callout_top + 0.48,
        width=11.6,
        height=1.2,
        text=bullets_text,
        size=12,
        bold=False,
        color=NAVY,
    )

    # Small "Source" footer
    add_text(
        slide,
        left=0.6,
        top=7.02,
        width=12.1,
        height=0.22,
        text="Source: Sales Director Monthly snapshot, 2026-04-01",
        size=9,
        bold=False,
        color=GREY_TEXT,
    )

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
