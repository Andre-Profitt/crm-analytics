#!/usr/bin/env python3
"""Build SimCorp-branded Sales Director Monthly decks using the official template.

Pulls live data from D1 (Sales Directors Monthly, filter-cascaded per director)
and D2 (Sales Ops Quarterly, per-report fallback). One PPTX per MD-1 director.

Usage:
    .venv_slides/bin/python scripts/build_simcorp_sd_monthly.py --director "Dan Peppett"
    .venv_slides/bin/python scripts/build_simcorp_sd_monthly.py --all
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
import time
from datetime import UTC, date, datetime
from pathlib import Path
from typing import Any

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO_ROOT = Path(__file__).resolve().parents[1]
CONFIG_PATH = REPO_ROOT / "config" / "sales_director_md1_presets.json"
DEFAULT_OUTPUT_ROOT = REPO_ROOT / "output" / "sales_director_monthly_runs"
TEMPLATE_PATH = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx"
TARGET_ORG = "apro@simcorp.com"
API_VERSION = "v66.0"

DASHBOARD1_ID = "01ZTb00000FSP7hMAH"

# D1 component IDs
D1_COMP = {
    "pipeline": "01aTb00000Cn9mPIAR",
    "approval_state": "01aTb00000Cn9mQIAR",
    "renewal_likelihood": "01aTb00000Cn85bIAB",
    "business_at_risk": "01aTb00000Cn85dIAB",
    "approval_ytd": "01aTb00000Cn85aIAB",
    "approval_candidates": "01aTb00000Cn85jIAB",
    "renewal_pipeline": "01aTb00000Cn85ZIAR",
    "slipped": "01aTb00000Cn85lIAB",
}

# Dashboard filter option IDs
_OPT = {
    "ind_asset_mgmt": "0ICTb0000007DbdOAE",
    "ind_pension": "0ICTb0000007DbgOAE",
    "lc_canada": "0ICTb0000007DgTOAU",
    "lc_excl_canada": "0ICTb0000007DgUOAU",
    "sr_apac": "0ICTb0000007DbnOAE",
    "sr_central_europe": "0ICTb0000007DboOAE",
    "sr_mea": "0ICTb0000007DbpOAE",
    "sr_nam": "0ICTb0000007DbqOAE",
    "sr_northern_europe": "0ICTb0000007DbrOAE",
    "sr_southwestern_europe": "0ICTb0000007DbsOAE",
    "sr_uki": "0ICTb0000007DbtOAE",
    "aug_sc_nam": "0ICTb0000007Di5OAE",
    "aug_sc_asia": "0ICTb0000007Di6OAE",
    "aug_sc_emea": "0ICTb0000007Di7OAE",
}

DIRECTOR_D1_FILTERS: dict[str, dict[str, str]] = {
    "Jesper Tyrer": {"filter3": _OPT["sr_apac"], "filter4": _OPT["aug_sc_asia"]},
    "Sarah Pittroff": {
        "filter3": _OPT["sr_central_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Francois Thaury": {
        "filter3": _OPT["sr_southwestern_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Dan Peppett": {"filter3": _OPT["sr_uki"], "filter4": _OPT["aug_sc_emea"]},
    "Christian Ebbesen": {
        "filter3": _OPT["sr_northern_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Mourad Essofi": {"filter3": _OPT["sr_mea"], "filter4": _OPT["aug_sc_emea"]},
    "Megan Miceli": {
        "filter2": _OPT["lc_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
    "Patrick Gaughan": {
        "filter1": _OPT["ind_asset_mgmt"],
        "filter2": _OPT["lc_excl_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
    "Adam Steinhaus": {
        "filter1": _OPT["ind_pension"],
        "filter2": _OPT["lc_excl_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
}

DECK_SECTIONS: list[tuple[int, str, list[tuple[str, str, str, str | None]]]] = [
    (
        1,
        "Pipeline coverage by region",
        [
            ("1.1", "Pipeline by Stage", "d1", D1_COMP["pipeline"]),
            ("1.2", "Top Opportunities in Pipeline", "d1", D1_COMP["pipeline"]),
        ],
    ),
    (
        2,
        "Commercial approval overview",
        [
            (
                "2.1",
                "Commercial Approval - Current State",
                "d1",
                D1_COMP["approval_state"],
            ),
            ("2.2", "Approved Deals YTD (Land)", "d1", D1_COMP["approval_ytd"]),
            (
                "2.3",
                "Land Stage 3 - Missing Approval Candidates",
                "d1",
                D1_COMP["approval_candidates"],
            ),
            (
                "2.4",
                "Missing Commercial Approval - Detail",
                "report",
                "00OTb000008fAlBMAU",
            ),
        ],
    ),
    (
        3,
        "Renewals tracking and risk",
        [
            (
                "3.1",
                "Renewal Likelihood by Probability",
                "d1",
                D1_COMP["renewal_likelihood"],
            ),
            (
                "3.2",
                "Upcoming Renewals This Quarter",
                "d1",
                D1_COMP["renewal_pipeline"],
            ),
            ("3.3", "Business At Risk", "d1", D1_COMP["business_at_risk"]),
        ],
    ),
    (
        4,
        "Churn trends and outlook",
        [
            ("4.1", "Churn Risk - Awaiting Finance Feed", "placeholder", None),
        ],
    ),
    (
        5,
        "Slipped deals analysis",
        [
            ("5.1", "Close Date Slipped - Aging Pipeline", "d1", D1_COMP["slipped"]),
            ("5.2", "Slipped Deals Trend", "placeholder", None),
            ("5.3", "Root Cause Commentary", "placeholder", None),
        ],
    ),
    (
        6,
        "Data quality and compliance",
        [
            ("6.1", "Missing Win/Loss Reason", "report", "00OTb000008el0PMAQ"),
            ("6.2", "Overdue Close Date Open Opps", "report", "00OTb000008TaBZMA0"),
            ("6.3", "Accounts without KYC Approval", "report", "00OTb000007BvlJMAS"),
        ],
    ),
]


def get_auth() -> tuple[str, str]:
    result = subprocess.run(
        ["sf", "org", "display", "--target-org", TARGET_ORG, "--json"],
        capture_output=True,
        text=True,
        check=True,
    )
    data = json.loads(result.stdout)["result"]
    return data["accessToken"], data["instanceUrl"]


def get_dashboard_filtered(
    instance_url: str,
    token: str,
    dashboard_id: str,
    filter_params: dict[str, str],
    wait_sec: int = 6,
) -> dict[str, dict[str, Any]]:
    base = f"{instance_url}/services/data/{API_VERSION}/analytics/dashboards/{dashboard_id}"
    if filter_params:
        requests.put(
            base,
            headers={"Authorization": f"Bearer {token}"},
            params=filter_params,
            timeout=60,
        )
        time.sleep(wait_sec)
    r = requests.get(
        base,
        headers={"Authorization": f"Bearer {token}"},
        params=filter_params,
        timeout=60,
    )
    r.raise_for_status()
    payload = r.json()
    return {
        cd["componentId"]: (cd.get("reportResult") or {})
        for cd in payload.get("componentData", [])
    }


def run_report(
    instance_url: str, token: str, report_id: str, extra_filters: list[dict[str, Any]]
) -> dict[str, Any]:
    url = f"{instance_url}/services/data/{API_VERSION}/analytics/reports/{report_id}"
    r = requests.get(
        url,
        headers={"Authorization": f"Bearer {token}"},
        params={"includeDetails": "true"},
        timeout=60,
    )
    r.raise_for_status()
    payload = r.json()
    rm = payload.get("reportMetadata", {})
    report_type = (rm.get("reportType") or {}).get("type", "")
    compatible = [
        f
        for f in extra_filters
        if f.get("column", "").startswith("Account.")
        or f.get("column", "").startswith("Opportunity.")
        or f.get("column") in {"INDUSTRY", "ADDRESS1_COUNTRY_CODE"}
    ]
    if not compatible:
        return payload
    rm["reportFilters"] = (rm.get("reportFilters") or []) + compatible
    r2 = requests.post(
        url,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        },
        params={"includeDetails": "true"},
        json={"reportMetadata": rm},
        timeout=120,
    )
    return r2.json() if r2.status_code == 200 else payload


def extract_table(
    result: dict[str, Any], max_rows: int = 12
) -> tuple[list[str], list[list[str]]]:
    rm = result.get("reportMetadata", {})
    ext = result.get("reportExtendedMetadata", {})
    fmt = rm.get("reportFormat", "TABULAR")
    fm = result.get("factMap", {}) or {}
    if fmt == "TABULAR":
        detail_info = ext.get("detailColumnInfo", {}) or {}
        cols = rm.get("detailColumns", []) or []
        headers = [(detail_info.get(c, {}) or {}).get("label", c) for c in cols]
        rows_raw = (fm.get("T!T") or {}).get("rows", []) or []
        rows = [
            [
                (c.get("label") or str(c.get("value", "")))
                for c in r.get("dataCells", [])
            ]
            for r in rows_raw[:max_rows]
        ]
        return headers, rows
    agg_info = ext.get("aggregateColumnInfo", {}) or {}
    agg_names = rm.get("aggregates", []) or []
    group_dims = rm.get("groupingsDown", []) or []
    group_name = group_dims[0].get("name", "Group") if group_dims else "Group"
    group_col = (ext.get("groupingColumnInfo", {}) or {}).get(group_name, {})
    headers = [group_col.get("label", group_name)] + [
        (agg_info.get(a, {}) or {}).get("label", a) for a in agg_names
    ]
    groupings = (result.get("groupingsDown", {}) or {}).get("groupings", []) or []
    rows = []
    for g in groupings[:max_rows]:
        row = [g.get("label", str(g.get("value", "")))]
        aggs = (fm.get(f"{g['key']}!T") or {}).get("aggregates", []) or []
        for a in aggs[: len(agg_names)]:
            row.append(str(a.get("label", a.get("value", ""))))
        rows.append(row)
    return headers, rows


def fmt_eur(amount):
    if amount is None:
        return "n/a"
    if abs(amount) >= 1_000_000:
        return f"EUR {amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"EUR {amount / 1_000:.0f}K"
    return f"EUR {amount:.0f}"


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout not found: {name}")


def fill(slide, idx, text):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = str(text)
            return True
    return False


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def add_table_to_slide(slide, headers, rows):
    """Add a table into the content area of a Title and Content slide."""
    if not rows:
        fill(slide, 22, "(No data for this director scope)")
        return
    n_cols = min(len(headers), 6)
    n_rows = len(rows) + 1
    # Position in the content area
    left = Inches(0.9)
    top = Inches(2.2)
    width = Inches(11.5)
    height = Inches(4.5)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    INK = RGBColor(0x1B, 0x1B, 0x1B)
    TEAL_DEEP = RGBColor(0x00, 0x3E, 0x52)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    for i, h in enumerate(headers[:n_cols]):
        cell = tbl.cell(0, i)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL_DEEP
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row[:n_cols]):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = INK


def load_preset(director_name: str) -> dict[str, Any]:
    config = json.loads(CONFIG_PATH.read_text())
    name_lower = director_name.lower().strip()
    for preset in config["presets"]:
        pname = preset["name"].lower()
        if (
            pname == name_lower
            or name_lower in pname
            or pname.split()[0] == name_lower.split()[0]
        ):
            return preset
    raise ValueError(f"No preset matched: {director_name}")


def slugify(name: str) -> str:
    parts = name.lower().strip().split()
    return f"{parts[-1]}-{parts[0]}" if len(parts) >= 2 else parts[0]


def build_deck(
    director: dict[str, Any],
    snapshot_date: str,
    output_path: Path,
    token: str,
    instance_url: str,
) -> None:
    prs = Presentation(TEMPLATE_PATH)
    clear_slides(prs)

    name = director["name"]
    territory = director["territory"]

    # --- Cover ---
    slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
    fill(slide, 20, "Sales Director Monthly - Pipeline and Insights")
    fill(slide, 22, f"Snapshot {snapshot_date}")
    fill(slide, 24, f"{name}\n{territory}")

    # --- Agenda ---
    slide = prs.slides.add_slide(get_layout(prs, "Agenda 1"))
    agenda_idxs = [40, 69, 71, 74, 76, 101]
    agenda_items = [s[1] for s in DECK_SECTIONS]
    for idx, text in zip(agenda_idxs, agenda_items):
        fill(slide, idx, text)

    # --- Fetch D1 filtered ---
    d1_filters = DIRECTOR_D1_FILTERS.get(name, {})
    print(f"  Fetching D1 with filters: {d1_filters}")
    try:
        d1 = get_dashboard_filtered(instance_url, token, DASHBOARD1_ID, d1_filters)
        print(f"  D1: {len(d1)} components")
    except Exception as e:
        print(f"  D1 FAILED: {e}")
        d1 = {}

    extra_filters = director.get("filters", [])

    # --- Sections ---
    for section_num, section_title, subitems in DECK_SECTIONS:
        # Divider
        slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
        fill(slide, 24, f"0{section_num}" if section_num < 10 else str(section_num))
        fill(slide, 20, section_title)

        for subnum, sub_title, source, target_id in subitems:
            if source == "placeholder":
                slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
                fill(slide, 144, sub_title)
                fill(slide, 145, f"Section {subnum}")
                reason = (
                    "Awaiting Finance feed from Alex P"
                    if section_num == 4
                    else "Pipeline Inspection in progress"
                    if "Trend" in sub_title
                    else "Owner commentary outreach pending"
                )
                fill(slide, 22, reason)
                fill(slide, 42, "Placeholder")
                continue

            # Fetch data
            headers, rows = [], []
            source_label = ""
            if source == "d1" and target_id:
                rr = d1.get(target_id)
                if rr:
                    synthetic = {
                        "reportMetadata": rr.get("reportMetadata", {}),
                        "reportExtendedMetadata": rr.get("reportExtendedMetadata", {}),
                        "factMap": rr.get("factMap", {}),
                        "groupingsDown": rr.get("groupingsDown", {}),
                    }
                    headers, rows = extract_table(synthetic)
                    source_label = "Dashboard (filtered)"
                else:
                    source_label = "Dashboard component not found"
                print(f"  {subnum} d1:{target_id} -> {len(rows)} rows")
            elif source == "report" and target_id:
                try:
                    result = run_report(instance_url, token, target_id, extra_filters)
                    headers, rows = extract_table(result)
                    source_label = "SF Report (global)"
                except Exception as e:
                    source_label = f"Report fetch failed: {e}"
                print(f"  {subnum} report:{target_id} -> {len(rows)} rows")

            slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
            fill(slide, 144, sub_title)
            fill(slide, 145, f"Section {subnum} - {source_label}")
            fill(slide, 42, f"{name} - {territory}")
            add_table_to_slide(slide, headers, rows)

    # --- End slide ---
    slide = prs.slides.add_slide(get_layout(prs, "End slide with disclaimer 1"))
    fill(
        slide,
        28,
        f"Generated {datetime.now(UTC).strftime('%Y-%m-%d %H:%M UTC')}\nData sourced from Salesforce dashboards with director-level filter presets applied.\nConfidential - SimCorp internal use only.",
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"  Saved: {output_path} ({len(prs.slides)} slides)")


def main() -> int:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("--director", help="Single director name")
    p.add_argument("--all", action="store_true", help="All 9 directors")
    p.add_argument("--snapshot-date", default=date.today().isoformat())
    p.add_argument("--output-root", default=str(DEFAULT_OUTPUT_ROOT))
    args = p.parse_args()
    if not args.director and not args.all:
        print("Must specify --director or --all", file=sys.stderr)
        return 2
    token, instance_url = get_auth()
    config = json.loads(CONFIG_PATH.read_text())
    targets = config["presets"] if args.all else [load_preset(args.director)]
    out_root = Path(args.output_root) / args.snapshot_date
    for director in targets:
        slug = slugify(director["name"])
        out = out_root / f"{slug}-simcorp.pptx"
        print(f"\n=== {director['name']} ({director['territory']}) ===")
        build_deck(director, args.snapshot_date, out, token, instance_url)
    return 0


if __name__ == "__main__":
    sys.exit(main())
