#!/usr/bin/env python3
"""Convert existing pptxgenjs-generated decks to SimCorp-branded versions
by extracting real slide content + recreating each slide using the SimCorp
template's 34 master layouts.

Strategy:
1. Read each source slide, extract every shape's text in reading order
   (top-to-bottom, left-to-right) along with font size so we can infer
   headlines vs labels vs values.
2. Pair small labels with adjacent larger values to form KPI cards.
3. Map each source slide to an appropriate SimCorp layout based on card count.
4. Render the new slide using fill_placeholder_by_idx with the extracted
   content.

Design rules:
- No em-dashes in output text (strip them).
- Use Renewal ACV (not ARR) phrasing per SimCorp methodology.
- Preserve every KPI card value from the original.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

TEMPLATE_PATH = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx"
REPO_ROOT = Path("/Users/test/crm-analytics")

EMDASH_RE = re.compile(r"[\u2013\u2014\u2015]")


def strip_emdash(text: str) -> str:
    if not text:
        return text
    return EMDASH_RE.sub("-", text)


@dataclass
class ShapeText:
    left: float  # inches
    top: float  # inches
    width: float
    height: float
    text: str
    max_font_size: float  # pt


@dataclass
class Card:
    label: str
    value: str
    context: str = ""


@dataclass
class SlideContent:
    title: str = ""
    subtitle: str = ""
    narrative: str = ""
    cards: list[Card] = field(default_factory=list)
    raw_lines: list[str] = field(default_factory=list)


# ----------------------------------------------------------------------------
# Extraction
# ----------------------------------------------------------------------------


def extract_shapes(slide) -> list[ShapeText]:
    shapes = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        txt = strip_emdash(txt)
        # Determine max font size across runs
        max_sz = 0
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    pt = r.font.size.pt
                    if pt > max_sz:
                        max_sz = pt
        try:
            left = sh.left.inches if sh.left is not None else 0
            top = sh.top.inches if sh.top is not None else 0
            width = sh.width.inches if sh.width is not None else 0
            height = sh.height.inches if sh.height is not None else 0
        except Exception:
            left = top = width = height = 0
        shapes.append(
            ShapeText(
                left=left,
                top=top,
                width=width,
                height=height,
                text=txt,
                max_font_size=max_sz or 12,
            )
        )
    return shapes


def extract_slide_content(slide) -> SlideContent:
    shapes = extract_shapes(slide)
    # Sort top-to-bottom, left-to-right
    shapes.sort(key=lambda s: (round(s.top, 1), round(s.left, 1)))
    content = SlideContent()

    # Largest font at top = title
    if shapes:
        title_shape = max(shapes, key=lambda s: s.max_font_size)
        content.title = title_shape.text
        # Subtitle = next largest distinct text near top
        top_shapes = [s for s in shapes if s.top <= (title_shape.top + 1.0)]
        for s in top_shapes:
            if s is title_shape:
                continue
            if s.max_font_size >= 14 and s.text != title_shape.text:
                if not content.subtitle or len(s.text) > len(content.subtitle):
                    content.subtitle = s.text
                    break

    # Build KPI cards: pair "label" shapes (small text, <=12pt) with nearby
    # "value" shapes (large text, >=18pt)
    labels = [s for s in shapes if 8 <= s.max_font_size <= 13 and len(s.text) < 80]
    values = [s for s in shapes if s.max_font_size >= 18 and len(s.text) < 50]

    for lbl in labels:
        # Find closest value below or right of label within reasonable distance
        best = None
        best_dist = 999
        for v in values:
            if v is lbl:
                continue
            dy = v.top - lbl.top
            dx = abs(v.left - lbl.left)
            # Prefer value directly below label (card format)
            if -0.2 <= dy <= 1.5 and dx <= 2.0:
                dist = dy * 2 + dx
                if dist < best_dist:
                    best = v
                    best_dist = dist
        if best:
            content.cards.append(Card(label=lbl.text, value=best.text))

    # Deduplicate cards (label+value)
    seen = set()
    uniq_cards = []
    for c in content.cards:
        key = (c.label, c.value)
        if key not in seen:
            seen.add(key)
            uniq_cards.append(c)
    content.cards = uniq_cards[:8]  # cap at 8 cards per slide

    # Narrative = longest text block not already used as title or card
    used_texts = {content.title, content.subtitle}
    for c in content.cards:
        used_texts.add(c.label)
        used_texts.add(c.value)
    narrative_candidates = [
        s.text for s in shapes if s.text not in used_texts and len(s.text) > 60
    ]
    if narrative_candidates:
        content.narrative = max(narrative_candidates, key=len)

    # All text lines for debugging
    content.raw_lines = [s.text for s in shapes]

    return content


# ----------------------------------------------------------------------------
# Rendering (SimCorp template)
# ----------------------------------------------------------------------------


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout not found: {name}")


def fill(slide, idx, text):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = strip_emdash(text or "")
            return True
    return False


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def render_slide(prs, sc: SlideContent, slide_index: int, deck_title: str):
    """Render a SlideContent into the SimCorp template using an appropriate layout."""
    num_cards = len(sc.cards)

    # Slide 1 always gets Title 1
    if slide_index == 0:
        slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
        fill(slide, 20, sc.title or deck_title)
        fill(slide, 22, sc.subtitle or sc.narrative)
        # Use the biggest card value as the date stamp
        date_val = ""
        for c in sc.cards:
            if "date" in c.label.lower() or "snapshot" in c.label.lower():
                date_val = c.value
                break
        fill(slide, 24, date_val)
        return

    # Choose layout based on card count
    if num_cards >= 5:
        layout_name = "5 x content w/ gradient line"
        title_idxs = [42, 56, 58, 60, 65]
        content_idxs = [22, 55, 57, 59, 66]
        header_idxs = [61, 62, 63, 64, 67]
        slot_count = 5
    elif num_cards >= 4:
        layout_name = "4 x content w/ gradient line"
        title_idxs = [42, 56, 58, 60]
        content_idxs = [22, 55, 57, 59]
        header_idxs = [61, 62, 63, 64]
        slot_count = 4
    elif num_cards == 3:
        layout_name = "3 x content w/ gradient line"
        title_idxs = [42, 56, 58]
        content_idxs = [22, 55, 57]
        header_idxs = [61, 62, 63]
        slot_count = 3
    elif num_cards == 2:
        layout_name = "2 x content w/ gradient line"
        title_idxs = [42, 56]
        content_idxs = [22, 55]
        header_idxs = [61, 62]
        slot_count = 2
    else:
        layout_name = "Title and Content"
        slide = prs.slides.add_slide(get_layout(prs, layout_name))
        fill(slide, 42, sc.title or f"Slide {slide_index + 1}")
        body = sc.narrative
        if sc.cards:
            for c in sc.cards:
                body += f"\n{c.label}: {c.value}"
        elif sc.raw_lines:
            # Fall back: include all non-title lines
            body = "\n".join(
                ln for ln in sc.raw_lines if ln != sc.title and len(ln) < 200
            )[:2000]
        fill(slide, 22, body or sc.subtitle or "")
        return

    slide = prs.slides.add_slide(get_layout(prs, layout_name))
    # Title
    # Note: in "N x content w/ gradient line" layouts, idx 42 is the FIRST title (col 1).
    # There's no separate "slide title" placeholder. So we put the slide title in the
    # header of column 1 via the big stat placeholder and use idx 42 for the first card.
    # Actually a simpler approach: write the slide title via the footer area idx 144
    # if available, or render the slide as a 4 x content without a page title.
    #
    # Pragmatic choice: first card's gradient header (61) stays as the stat, and we
    # write the slide title to 144 (footer). To also get a visible title, use 42
    # for the first card title and put the slide title concept in the narrative below.
    #
    # To keep things simple: make the first card's title be "{slide_title}: {card1.label}"
    # so the slide title appears in the top-left. Not ideal but gets content in.
    for i, card in enumerate(sc.cards[:slot_count]):
        if i < len(header_idxs):
            fill(slide, header_idxs[i], card.value)
        if i < len(title_idxs):
            label = card.label
            if i == 0 and sc.title:
                label = f"{sc.title}: {card.label}"
            fill(slide, title_idxs[i], label)
        if i < len(content_idxs):
            fill(slide, content_idxs[i], card.context or "")

    # Narrative goes to footer
    if sc.narrative:
        fill(slide, 144, sc.narrative[:250])


# ----------------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------------


def convert_deck(src_path: str, dst_path: str, deck_title: str):
    src = Presentation(src_path)
    print(f"\n=== Converting {Path(src_path).name} ===")
    print(f"  source slides: {len(src.slides)}")

    # Extract content from each source slide
    slide_contents = []
    for i, slide in enumerate(src.slides):
        sc = extract_slide_content(slide)
        print(
            f"  slide {i + 1}: title={sc.title[:50]!r} cards={len(sc.cards)} narrative_len={len(sc.narrative)}"
        )
        slide_contents.append(sc)

    # Render into SimCorp template
    prs = Presentation(TEMPLATE_PATH)
    clear_slides(prs)

    for i, sc in enumerate(slide_contents):
        render_slide(prs, sc, i, deck_title)

    # Append closing slide
    closing = prs.slides.add_slide(get_layout(prs, "End slide with disclaimer 1"))
    fill(closing, 28, "Thank you\n\nwww.simcorp.com")

    prs.save(dst_path)
    print(f"  output slides: {len(prs.slides)}")
    print(f"  saved: {dst_path}")


def main():
    convert_deck(
        str(
            REPO_ROOT
            / "output/sales_director_monthly_runs/2026-04-06T18-31-12Z_2026-04-01/sales_director_monthly_pipeline_insights_2026-04-01.pptx"
        ),
        str(
            REPO_ROOT
            / "output/sales_director_monthly_runs/2026-04-06T18-31-12Z_2026-04-01/sales_director_monthly_simcorp_v2.pptx"
        ),
        "Sales Director Monthly Pipeline and Insights",
    )
    convert_deck(
        str(
            REPO_ROOT
            / "output/sales_ops_quarterly_deck_2026-03-31/sales_ops_quarterly_review_2026-04-01.pptx"
        ),
        str(
            REPO_ROOT
            / "output/sales_ops_quarterly_deck_2026-03-31/sales_ops_quarterly_simcorp_v2.pptx"
        ),
        "Sales Ops Quarterly Review",
    )


if __name__ == "__main__":
    main()
