#!/usr/bin/env python3
"""Sales Director Monthly deck v3 — matches the original brief exactly
and introduces layout variety so the deck doesn't feel like 9 identical
4-card strips.

Slide structure (10 slides):
 1. Cover (SC-Master Gradient_Title)
 2. EMEA pipeline outlook — 4 KPI card strip with progress bar
 3. North America pipeline outlook — two-column: hero gap stat + mini grid
 4. APAC pipeline outlook — horizontal bars comparing target / call / projected
 5. Commercial approval global overview — 3 big hero stats + summary
 6. Commercial approval candidates by region — region comparison with counts
 7. Renewal pipeline and risk — KPI cards with magenta accent
 8. Churn trend, last three quarters — horizontal trend bars
 9. Slipped deals by region — horizontal comparison
10. Appendix — clickable Salesforce links

No fiscal year labels. No "days remaining". Calendar year only.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"
RUN_DIR = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01"
)
SNAPSHOT = RUN_DIR / "report1_snapshot.json"
OUTPUT = RUN_DIR / "sales_director_monthly_simcorp_v3.pptx"

# SimCorp production colors
NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
RICH_BLUE = RGBColor(0x09, 0x4F, 0xB2)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
AQUA_LIGHT = RGBColor(0xCB, 0xF5, 0xF3)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
LIGHT_PANEL_2 = RGBColor(0xF5, 0xF8, 0xFD)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
MAGENTA_LIGHT = RGBColor(0xF8, 0xE8, 0xF1)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)
DIVIDER_GREY = RGBColor(0xD7, 0xE2, 0xE8)
AMBER = RGBColor(0xFB, 0x9B, 0x2A)
GREEN_OK = RGBColor(0x33, 0xD8, 0xCE)

LINKS = [
    (
        "Sales Directors Monthly Dashboard",
        "https://simcorp.my.salesforce.com/lightning/r/Dashboard/01ZTb00000FSP7hMAH/view",
    ),
    (
        "Sales Ops Quarterly KPI Dashboard",
        "https://simcorp.my.salesforce.com/lightning/r/Dashboard/01ZTb00000FSP9JMAX/view",
    ),
    (
        "Sales Ops Data Quality and Forecast Accuracy (CRMA)",
        "https://simcorp.my.salesforce.com/analytics/dashboard/0FKTb0000000K5BOAU",
    ),
    (
        "Pipeline Coverage by Stage",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008TZc5MAG/view",
    ),
    (
        "Forecast and Closed Won",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008TZaTMAW/view",
    ),
    (
        "Land Stage 3 Missing Commercial Approval by Region",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekltMAA/view",
    ),
    (
        "Commercial Approval Candidates by Stage",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekp7MAA/view",
    ),
    (
        "Renewals by Fiscal Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008eksLMAQ/view",
    ),
    (
        "Renewal Pipeline This Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ektxMAA/view",
    ),
    (
        "Renewal ACV by Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekxBMAQ/view",
    ),
    (
        "Business At Risk",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008Ta9xMAC/view",
    ),
    (
        "Close Date Slipped by Stage",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008eknVMAQ/view",
    ),
]


# ============================================================================
# Helpers
# ============================================================================


def fmt_eur(amount):
    if amount is None:
        return "n/a"
    if abs(amount) >= 1_000_000_000:
        return f"EUR {amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"EUR {amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"EUR {amount / 1_000:.0f}K"
    return f"EUR {amount:.0f}"


def fmt_int(n):
    return f"{n:,}" if n is not None else "n/a"


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def txt(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def rect(slide, left, top, width, height, fill, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(
        shape, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def rounded_card(slide, left, top, width, height, fill):
    c = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    try:
        c.adjustments[0] = 0.08
    except Exception:
        pass
    return c


def dot(slide, left, top, size, color):
    d = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()
    return d


def blank_slide(prs):
    s = prs.slides.add_slide(get_layout(prs, "Title Only"))
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s


def header(slide, eyebrow, title, narrative, tab_color=AQUA):
    rect(slide, 0.6, 0.42, 0.08, 0.7, tab_color)
    txt(slide, 0.82, 0.38, 12, 0.28, eyebrow, size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, 0.82, 0.62, 12.4, 0.58, title, size=28, bold=True, color=NAVY)
    txt(slide, 0.82, 1.25, 12.1, 0.4, narrative, size=13, color=GREY_TEXT)


def takeaway(slide, top, bullets, accent=AQUA):
    rect(slide, 0.6, top - 0.25, 12.1, 0.015, DIVIDER_GREY)
    h = 1.9
    rounded_card(slide, 0.6, top, 12.1, h, LIGHT_PANEL_2)
    rect(slide, 0.6, top, 0.1, h, accent)
    txt(
        slide,
        0.95,
        top + 0.16,
        11.5,
        0.28,
        "TAKEAWAY",
        size=10,
        bold=True,
        color=PRIMARY_BLUE,
    )
    txt(slide, 0.95, top + 0.45, 11.5, h - 0.55, bullets, size=12, color=NAVY)


def source_footer(slide, text):
    txt(slide, 0.6, 7.08, 12.1, 0.22, text, size=9, color=GREY_TEXT)


def kpi_card(
    slide,
    left,
    top,
    width,
    height,
    big,
    label,
    context,
    progress_pct=None,
    status_color=None,
    panel=LIGHT_BLUE_PANEL,
):
    rounded_card(slide, left, top, width, height, panel)
    # Big stat
    tb = slide.shapes.add_textbox(
        Inches(left + 0.22), Inches(top + 0.18), Inches(width - 0.44), Inches(0.7)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = big
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_BLUE
    if status_color is not None:
        dot(slide, left + width - 0.44, top + 0.28, 0.22, status_color)
    txt(
        slide,
        left + 0.22,
        top + 0.92,
        width - 0.44,
        0.32,
        label,
        size=13,
        bold=True,
        color=NAVY,
    )
    ctx_top = top + 1.22
    ctx_h = height - 1.3
    if progress_pct is not None:
        ctx_h -= 0.25
    txt(
        slide,
        left + 0.22,
        ctx_top,
        width - 0.44,
        ctx_h,
        context,
        size=10,
        color=PRIMARY_BLUE,
    )
    if progress_pct is not None:
        bar_top = top + height - 0.32
        bar_left = left + 0.22
        bar_width = width - 0.44
        rect(slide, bar_left, bar_top, bar_width, 0.1, DIVIDER_GREY)
        fill_w = max(0.02, bar_width * min(1.0, progress_pct / 100))
        rect(slide, bar_left, bar_top, fill_w, 0.1, AQUA)


# ============================================================================
# Slide 1: Cover
# ============================================================================


def build_cover(prs, snap):
    cover = prs.slides.add_slide(get_layout(prs, "SC-Master Gradient_Title"))

    def fill_ph(idx, text):
        for ph in cover.placeholders:
            if ph.placeholder_format.idx == idx:
                ph.text = text
                return

    fill_ph(0, "Sales Director Monthly Pipeline Insights")
    fill_ph(
        20,
        "Pipeline coverage by region, commercial approvals, renewals, churn risk, slipped deals. Forward looking view for Sales Directors.",
    )
    fill_ph(
        22,
        f"April 2026 monthly view  ·  Snapshot {snap.get('snapshot_date', '2026-04-01')}",
    )


# ============================================================================
# Slide 2: EMEA — 4 KPI card strip with progress bar
# ============================================================================


def build_region_card_strip(prs, snap, region_name, section_no):
    slide = blank_slide(prs)
    reg = snap["pipeline"]["deck_regions"].get(region_name, {})
    target = reg.get("target_arr", 0) or 0
    call = reg.get("best_case_call_arr", 0) or 0
    confidence = reg.get("forecast_confidence_pct", 0) or 0
    coverage = reg.get("coverage_status", "n/a") or "n/a"
    gap = max(0.0, target - call)
    coverage_pct = (call / target * 100) if target else 0

    header(
        slide,
        eyebrow=f"APRIL 2026   ·   {section_no}   ·   PIPELINE OUTLOOK",
        title=f"{region_name} pipeline outlook",
        narrative=f"Best case call {fmt_eur(call)} against {fmt_eur(target)} target. Confidence {confidence:.1f}%.",
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 2.95
    gap_bet = 0.2
    start_x = 0.6
    cards = [
        (fmt_eur(target), "Quarter target", "Executive target seam, ARR.", None, None),
        (
            fmt_eur(call),
            "Best case call",
            f"{coverage_pct:.0f}% coverage.",
            coverage_pct,
            None,
        ),
        (
            fmt_eur(gap),
            "Gap to target",
            "Needed from promotion or new pipeline.",
            None,
            None,
        ),
        (
            f"{confidence:.1f}%",
            "Forecast confidence",
            f"Status: {coverage}.",
            None,
            AMBER if confidence < 96 else GREEN_OK,
        ),
    ]
    for i, (big, lbl, ctx, prog, status) in enumerate(cards):
        kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            big,
            lbl,
            ctx,
            progress_pct=prog,
            status_color=status,
        )

    gap_msg = (
        f"•  {region_name} needs {fmt_eur(gap)} additional coverage to close."
        if gap > 0
        else f"•  {region_name} is at or above target on best case call."
    )
    conf_msg = (
        f"•  Forecast confidence {confidence:.1f}% is below the 96% threshold."
        if confidence < 96
        else f"•  Forecast confidence {confidence:.1f}% is healthy."
    )
    bullets = f"{gap_msg}\n{conf_msg}\n•  Coverage status: {coverage}."
    takeaway(slide, 4.35, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


# ============================================================================
# Slide 3: NA — two-column hero + mini grid (different rhythm)
# ============================================================================


def build_region_hero_column(prs, snap, region_name, section_no):
    slide = blank_slide(prs)
    reg = snap["pipeline"]["deck_regions"].get(region_name, {})
    target = reg.get("target_arr", 0) or 0
    call = reg.get("best_case_call_arr", 0) or 0
    confidence = reg.get("forecast_confidence_pct", 0) or 0
    coverage = reg.get("coverage_status", "n/a") or "n/a"
    gap = max(0.0, target - call)
    coverage_pct = (call / target * 100) if target else 0

    header(
        slide,
        eyebrow=f"APRIL 2026   ·   {section_no}   ·   PIPELINE OUTLOOK",
        title=f"{region_name} pipeline outlook",
        narrative=f"Largest gap of the 3 regions. {coverage_pct:.0f}% of target coverage, status: {coverage}.",
    )

    # LEFT COLUMN: Hero gap stat (the headline for NA)
    hero_left = 0.6
    hero_top = 1.9
    hero_w = 5.9
    hero_h = 3.9
    rounded_card(slide, hero_left, hero_top, hero_w, hero_h, LIGHT_BLUE_PANEL)
    # Label above the hero number
    txt(
        slide,
        hero_left + 0.3,
        hero_top + 0.35,
        hero_w - 0.6,
        0.3,
        "GAP TO QUARTER TARGET",
        size=11,
        bold=True,
        color=PRIMARY_BLUE,
    )
    # Big hero number
    tb = slide.shapes.add_textbox(
        Inches(hero_left + 0.3),
        Inches(hero_top + 0.7),
        Inches(hero_w - 0.6),
        Inches(1.6),
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = fmt_eur(gap)
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(64)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_BLUE
    # Context under hero number
    txt(
        slide,
        hero_left + 0.3,
        hero_top + 2.4,
        hero_w - 0.6,
        1.3,
        "This is the promotion plus new pipeline gap the region must close "
        "to reach target. It is the largest gap of the 3 regions this cycle.",
        size=12,
        color=NAVY,
    )

    # RIGHT COLUMN: 2x2 mini stat grid
    grid_left = 6.7
    grid_top = 1.9
    mini_w = 3.0
    mini_h = 1.9
    mini_gap = 0.1
    # Top-left: Target
    rounded_card(slide, grid_left, grid_top, mini_w, mini_h, LIGHT_PANEL_2)
    txt(
        slide,
        grid_left + 0.22,
        grid_top + 0.2,
        mini_w - 0.44,
        0.3,
        "QUARTER TARGET",
        size=9,
        bold=True,
        color=PRIMARY_BLUE,
    )
    txt(
        slide,
        grid_left + 0.22,
        grid_top + 0.55,
        mini_w - 0.44,
        0.6,
        fmt_eur(target),
        size=22,
        bold=True,
        color=NAVY,
    )
    txt(
        slide,
        grid_left + 0.22,
        grid_top + 1.3,
        mini_w - 0.44,
        0.4,
        "Executive target seam, ARR.",
        size=10,
        color=GREY_TEXT,
    )
    # Top-right: Best case call
    rounded_card(
        slide, grid_left + mini_w + mini_gap, grid_top, mini_w, mini_h, LIGHT_PANEL_2
    )
    txt(
        slide,
        grid_left + mini_w + mini_gap + 0.22,
        grid_top + 0.2,
        mini_w - 0.44,
        0.3,
        "BEST CASE CALL",
        size=9,
        bold=True,
        color=PRIMARY_BLUE,
    )
    txt(
        slide,
        grid_left + mini_w + mini_gap + 0.22,
        grid_top + 0.55,
        mini_w - 0.44,
        0.6,
        fmt_eur(call),
        size=22,
        bold=True,
        color=NAVY,
    )
    txt(
        slide,
        grid_left + mini_w + mini_gap + 0.22,
        grid_top + 1.3,
        mini_w - 0.44,
        0.4,
        f"{coverage_pct:.0f}% of target.",
        size=10,
        color=GREY_TEXT,
    )
    # Bottom-left: Confidence
    rounded_card(
        slide, grid_left, grid_top + mini_h + mini_gap, mini_w, mini_h, LIGHT_PANEL_2
    )
    txt(
        slide,
        grid_left + 0.22,
        grid_top + mini_h + mini_gap + 0.2,
        mini_w - 0.44,
        0.3,
        "FORECAST CONFIDENCE",
        size=9,
        bold=True,
        color=PRIMARY_BLUE,
    )
    txt(
        slide,
        grid_left + 0.22,
        grid_top + mini_h + mini_gap + 0.55,
        mini_w - 0.44,
        0.6,
        f"{confidence:.1f}%",
        size=22,
        bold=True,
        color=NAVY,
    )
    if confidence < 96:
        dot(
            slide,
            grid_left + mini_w - 0.55,
            grid_top + mini_h + mini_gap + 0.28,
            0.22,
            AMBER,
        )
    # Bottom-right: Status
    rounded_card(
        slide,
        grid_left + mini_w + mini_gap,
        grid_top + mini_h + mini_gap,
        mini_w,
        mini_h,
        LIGHT_PANEL_2,
    )
    txt(
        slide,
        grid_left + mini_w + mini_gap + 0.22,
        grid_top + mini_h + mini_gap + 0.2,
        mini_w - 0.44,
        0.3,
        "COVERAGE STATUS",
        size=9,
        bold=True,
        color=PRIMARY_BLUE,
    )
    txt(
        slide,
        grid_left + mini_w + mini_gap + 0.22,
        grid_top + mini_h + mini_gap + 0.55,
        mini_w - 0.44,
        0.6,
        coverage,
        size=18,
        bold=True,
        color=NAVY,
    )

    bullets = (
        f"•  {region_name} is the largest gap region at {fmt_eur(gap)}.\n"
        f"•  Coverage at {coverage_pct:.0f}% of target. New pipeline promotion is the main lever.\n"
        f"•  Forecast confidence {confidence:.1f}%."
    )
    takeaway(slide, 6.05, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


# ============================================================================
# Slide 4: APAC — horizontal bars comparing target / call / actual
# ============================================================================


def build_region_horizontal_bars(prs, snap, region_name, section_no):
    slide = blank_slide(prs)
    reg = snap["pipeline"]["deck_regions"].get(region_name, {})
    target = reg.get("target_arr", 0) or 0
    call = reg.get("best_case_call_arr", 0) or 0
    actual = (reg.get("actual_arr", 0) or 0) + (reg.get("weighted_open_arr", 0) or 0)
    confidence = reg.get("forecast_confidence_pct", 0) or 0
    coverage = reg.get("coverage_status", "n/a") or "n/a"

    header(
        slide,
        eyebrow=f"APRIL 2026   ·   {section_no}   ·   PIPELINE OUTLOOK",
        title=f"{region_name} pipeline outlook",
        narrative=f"{coverage}. Forecast confidence {confidence:.1f}%.",
    )

    # Bar chart area
    chart_left = 0.6
    chart_top = 1.95
    chart_w = 12.1
    chart_h = 3.4

    rounded_card(slide, chart_left, chart_top, chart_w, chart_h, LIGHT_BLUE_PANEL)
    txt(
        slide,
        chart_left + 0.35,
        chart_top + 0.25,
        chart_w - 0.7,
        0.3,
        "TARGET VS CALL VS PROJECTED (ARR)",
        size=10,
        bold=True,
        color=PRIMARY_BLUE,
    )

    # Bars: label on left, bar in middle, value on right
    max_val = max(target, call, actual, 1)
    bars = [
        ("Quarter target", target, RICH_BLUE),
        ("Best case call", call, AQUA),
        ("Projected (actual + weighted open)", actual, PRIMARY_BLUE),
    ]
    bar_left = chart_left + 3.0
    bar_area_w = 6.8
    bar_top_start = chart_top + 0.85
    row_h = 0.75
    bar_h = 0.38

    for i, (label, value, color) in enumerate(bars):
        row_y = bar_top_start + i * row_h
        # Label
        txt(
            slide,
            chart_left + 0.35,
            row_y + 0.04,
            2.55,
            0.32,
            label,
            size=12,
            bold=True,
            color=NAVY,
        )
        # Track (grey)
        rect(slide, bar_left, row_y, bar_area_w, bar_h, DIVIDER_GREY)
        # Fill
        fill_w = max(0.05, bar_area_w * (value / max_val))
        rect(slide, bar_left, row_y, fill_w, bar_h, color)
        # Value on the right
        txt(
            slide,
            chart_left + chart_w - 1.85,
            row_y + 0.04,
            1.65,
            0.32,
            fmt_eur(value),
            size=13,
            bold=True,
            color=NAVY,
            align=PP_ALIGN.RIGHT,
        )

    bullets = (
        f"•  {region_name} projected {fmt_eur(actual)} closes close to the {fmt_eur(target)} target.\n"
        f"•  Best case call of {fmt_eur(call)} suggests {(call / target * 100 if target else 0):.0f}% coverage.\n"
        f"•  Forecast confidence {confidence:.1f}%. Status: {coverage}."
    )
    takeaway(slide, 5.75, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


# ============================================================================
# Slide 5: Commercial approval — global overview (big hero stats)
# ============================================================================


def build_commercial_approval_global(prs, snap):
    slide = blank_slide(prs)
    ca = snap.get("commercial_approval", {})
    s = ca.get("summary", {})
    approved_count = s.get("approved_count", 0) or 0
    approved_arr = s.get("approved_arr", 0) or 0
    pending_count = s.get("pending_count", 0) or 0
    pending_arr = s.get("pending_arr", 0) or 0
    stale_count = s.get("stale_count", 0) or 0
    stale_arr = s.get("stale_arr", 0) or 0
    candidate_count = ca.get("candidate_count", 0) or 0

    header(
        slide,
        eyebrow="APRIL 2026   ·   02   ·   COMMERCIAL GOVERNANCE",
        title="Commercial approval overview",
        narrative=f"{approved_count} approved at {fmt_eur(approved_arr)}. {pending_count} pending, {stale_count} stale.",
    )

    # Single large hero card on top showing approved count
    hero_top = 1.9
    hero_h = 2.3
    rounded_card(slide, 0.6, hero_top, 12.1, hero_h, LIGHT_BLUE_PANEL)
    # Left: big number
    tb = slide.shapes.add_textbox(
        Inches(0.95), Inches(hero_top + 0.25), Inches(4.5), Inches(1.6)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = fmt_int(approved_count)
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(96)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_BLUE
    # Right: label + context
    txt(
        slide,
        5.8,
        hero_top + 0.45,
        6.5,
        0.35,
        "APPROVED THIS CYCLE",
        size=11,
        bold=True,
        color=PRIMARY_BLUE,
    )
    txt(
        slide,
        5.8,
        hero_top + 0.85,
        6.5,
        0.55,
        f"Total value: {fmt_eur(approved_arr)}",
        size=22,
        bold=True,
        color=NAVY,
    )
    txt(
        slide,
        5.8,
        hero_top + 1.5,
        6.5,
        0.4,
        "Land stage 3 deals with commercial approval granted in the current cycle.",
        size=12,
        color=GREY_TEXT,
    )
    if approved_count > 0:
        dot(slide, 12.35, hero_top + 0.5, 0.25, GREEN_OK)

    # Secondary row: pending / stale / candidates as 3 smaller cards
    sub_top = 4.4
    sub_h = 1.1
    sub_w = 3.95
    sub_gap = 0.22
    sub_start_x = 0.6

    def mini_stat(left, top, width, height, label, number, value_text, status=None):
        rounded_card(slide, left, top, width, height, LIGHT_PANEL_2)
        txt(
            slide,
            left + 0.22,
            top + 0.15,
            width - 0.44,
            0.25,
            label,
            size=9,
            bold=True,
            color=PRIMARY_BLUE,
        )
        txt(
            slide,
            left + 0.22,
            top + 0.4,
            width - 0.44,
            0.4,
            number,
            size=22,
            bold=True,
            color=NAVY,
        )
        txt(
            slide,
            left + 0.22,
            top + 0.78,
            width - 0.44,
            0.3,
            value_text,
            size=10,
            color=GREY_TEXT,
        )
        if status:
            dot(slide, left + width - 0.42, top + 0.22, 0.2, status)

    mini_stat(
        sub_start_x + 0 * (sub_w + sub_gap),
        sub_top,
        sub_w,
        sub_h,
        "PENDING APPROVAL",
        fmt_int(pending_count),
        f"Awaiting sign off: {fmt_eur(pending_arr)}",
        status=AMBER if pending_count else None,
    )
    mini_stat(
        sub_start_x + 1 * (sub_w + sub_gap),
        sub_top,
        sub_w,
        sub_h,
        "STALE",
        fmt_int(stale_count),
        f"Cleanup required: {fmt_eur(stale_arr)}",
        status=AMBER if stale_count else None,
    )
    mini_stat(
        sub_start_x + 2 * (sub_w + sub_gap),
        sub_top,
        sub_w,
        sub_h,
        "LAND STAGE 3 CANDIDATES",
        fmt_int(candidate_count),
        "Open candidates without approval.",
        status=AMBER if candidate_count else GREEN_OK,
    )

    if pending_count == 0 and stale_count == 0 and candidate_count == 0:
        bullet1 = "•  Approval control is fully contained this cycle."
    else:
        bullet1 = f"•  {pending_count + stale_count + candidate_count} items require attention this cycle."
    bullets = (
        f"{bullet1}\n"
        f"•  {approved_count} approved deals worth {fmt_eur(approved_arr)} passed the commercial approval gate.\n"
        "•  Regional breakdown of outstanding candidates is on the next slide."
    )
    takeaway(slide, 5.9, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


# ============================================================================
# Slide 6: Commercial approval — candidates by region (table-ish list)
# ============================================================================


def build_commercial_approval_regional(prs, snap):
    slide = blank_slide(prs)
    ca = snap.get("commercial_approval", {})
    by_region = ca.get("candidates_by_region", {}) or {}
    order = ["EMEA", "North America", "APAC"]

    def reg_info(r):
        v = by_region.get(r)
        if isinstance(v, list):
            return (
                len(v),
                sum((d.get("arr", 0) or 0) for d in v if isinstance(d, dict)),
                v,
            )
        if isinstance(v, dict):
            return (
                v.get("count", 0) or v.get("candidate_count", 0) or 0,
                v.get("arr", 0) or v.get("total_arr", 0) or 0,
                v.get("deals") or v.get("candidates") or [],
            )
        return 0, 0, []

    total_count = sum(reg_info(r)[0] for r in order)
    total_arr = sum(reg_info(r)[1] for r in order)

    header(
        slide,
        eyebrow="APRIL 2026   ·   02   ·   COMMERCIAL GOVERNANCE",
        title="Land Stage 3 candidates without commercial approval",
        narrative=(
            f"{total_count} open candidates across 3 regions."
            if total_count > 0
            else "No open Land Stage 3 candidates this cycle. Approval control is contained."
        ),
    )

    # Table-ish region rows
    table_top = 2.0
    row_h = 1.2
    for i, region in enumerate(order):
        row_y = table_top + i * (row_h + 0.15)
        count, arr, _ = reg_info(region)
        rounded_card(slide, 0.6, row_y, 12.1, row_h, LIGHT_BLUE_PANEL)
        # Left accent stripe (magenta if has candidates, aqua if clean)
        rect(slide, 0.6, row_y, 0.15, row_h, MAGENTA if count > 0 else AQUA)
        # Region label
        txt(
            slide,
            0.95,
            row_y + 0.2,
            5,
            0.35,
            region.upper(),
            size=11,
            bold=True,
            color=PRIMARY_BLUE,
        )
        txt(
            slide,
            0.95,
            row_y + 0.5,
            5,
            0.5,
            f"{fmt_int(count)} candidates",
            size=22,
            bold=True,
            color=NAVY,
        )
        # Value
        txt(
            slide,
            6.0,
            row_y + 0.2,
            4,
            0.35,
            "EXPOSURE",
            size=11,
            bold=True,
            color=PRIMARY_BLUE,
        )
        txt(
            slide,
            6.0,
            row_y + 0.5,
            4,
            0.5,
            fmt_eur(arr) if arr else "—",
            size=22,
            bold=True,
            color=NAVY,
        )
        # Status
        txt(
            slide,
            10.3,
            row_y + 0.2,
            2.3,
            0.35,
            "STATUS",
            size=11,
            bold=True,
            color=PRIMARY_BLUE,
        )
        status_text = "CONTAINED" if count == 0 else "FOLLOW UP"
        status_color = GREEN_OK if count == 0 else AMBER
        txt(
            slide,
            10.3,
            row_y + 0.5,
            2.3,
            0.5,
            status_text,
            size=16,
            bold=True,
            color=NAVY,
        )
        dot(slide, 12.25, row_y + 0.55, 0.22, status_color)

    if total_count == 0:
        bullets = (
            "•  No open Land Stage 3 candidates in any region this cycle.\n"
            "•  Approval control is fully contained.\n"
            "•  Live drilldown link is in the appendix."
        )
    else:
        bullets = (
            f"•  {total_count} candidates across 3 regions represent {fmt_eur(total_arr)} in open exposure.\n"
            "•  Follow up by owner via the live Land Stage 3 Missing Approval report.\n"
            "•  Link in the appendix."
        )
    takeaway(slide, 5.95, bullets)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


# ============================================================================
# Slide 7: Renewals — 4 KPI cards with magenta accent
# ============================================================================


def build_renewals(prs, snap):
    slide = blank_slide(prs)
    ren = snap.get("renewals", {}).get("summary", {})
    total = ren.get("total_renewal_pipeline_acv", 0) or 0
    deals = ren.get("total_deals", 0) or 0
    due_acv = ren.get("due_this_quarter_acv", 0) or 0
    due_ct = ren.get("due_this_quarter_count", 0) or 0
    over_acv = ren.get("overdue_carryover_acv", 0) or 0
    over_ct = ren.get("overdue_carryover_count", 0) or 0
    crit_acv = ren.get("critical_acv", 0) or 0
    crit_ct = ren.get("critical_count", 0) or 0

    header(
        slide,
        eyebrow="APRIL 2026   ·   03   ·   RENEWALS TRACKING",
        title="Renewal pipeline and risk",
        narrative=f"{deals} open renewal deals at {fmt_eur(total)} ACV. All flagged as critical this cycle.",
        tab_color=MAGENTA,
    )

    card_top = 1.9
    card_h = 1.95
    card_w = 2.95
    gap_bet = 0.2
    start_x = 0.6
    cards = [
        (
            fmt_eur(total),
            "Total open renewal ACV",
            f"{deals} deals. Methodology: ACV.",
            None,
            None,
        ),
        (
            fmt_eur(due_acv),
            "Due this quarter",
            f"{due_ct} deals closing this quarter.",
            None,
            AMBER if due_ct > 0 else None,
        ),
        (
            fmt_eur(over_acv),
            "Overdue carryover",
            f"{over_ct} deals carried from prior quarters.",
            None,
            AMBER if over_ct > 0 else None,
        ),
        (
            fmt_eur(crit_acv),
            "Critical risk",
            f"{crit_ct} flagged as critical escalation.",
            None,
            AMBER if crit_ct > 0 else None,
        ),
    ]
    for i, (big, lbl, ctx, prog, status) in enumerate(cards):
        kpi_card(
            slide,
            start_x + i * (card_w + gap_bet),
            card_top,
            card_w,
            card_h,
            big,
            lbl,
            ctx,
            progress_pct=prog,
            status_color=status,
            panel=MAGENTA_LIGHT,
        )

    bullets = (
        f"•  All {deals} open renewals carry some level of risk this cycle.\n"
        f"•  {over_ct} overdue carryover deals represent the largest escalation pressure.\n"
        "•  Value methodology: Renewal ACV per SimCorp standard."
    )
    takeaway(slide, 4.35, bullets, accent=MAGENTA)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


# ============================================================================
# Slide 8: Churn trend — horizontal bars across 3 quarters
# ============================================================================


def build_churn(prs, snap):
    slide = blank_slide(prs)
    trend = snap.get("churn", {}).get("trend", [])
    last_three = trend[-3:] if len(trend) >= 3 else trend
    while len(last_three) < 3:
        last_three.insert(
            0, {"quarter_label": "n/a", "churned_acv": 0, "churned_deals": 0}
        )
    finance_status = snap.get("churn", {}).get("finance_feed_status", "pending")

    header(
        slide,
        eyebrow="APRIL 2026   ·   04   ·   CHURN RISK AND TRENDS",
        title="Churn trend, last three quarters",
        narrative=f"Finance feed status: {finance_status}. CRM observed churn only.",
        tab_color=MAGENTA,
    )

    chart_left = 0.6
    chart_top = 1.95
    chart_w = 12.1
    chart_h = 3.4
    rounded_card(slide, chart_left, chart_top, chart_w, chart_h, MAGENTA_LIGHT)
    txt(
        slide,
        chart_left + 0.35,
        chart_top + 0.25,
        chart_w - 0.7,
        0.3,
        "CHURNED ACV BY QUARTER",
        size=10,
        bold=True,
        color=MAGENTA,
    )

    max_val = max((q.get("churned_acv", 0) or 0) for q in last_three) or 1
    bar_left = chart_left + 2.0
    bar_area_w = 7.8
    bar_top_start = chart_top + 0.9
    row_h = 0.75
    bar_h = 0.38

    for i, q in enumerate(last_three):
        qlabel = q.get("quarter_label", "n/a")
        churned = q.get("churned_acv", 0) or 0
        deals = q.get("churned_deals", 0) or 0
        row_y = bar_top_start + i * row_h
        txt(
            slide,
            chart_left + 0.35,
            row_y + 0.04,
            1.55,
            0.32,
            qlabel,
            size=13,
            bold=True,
            color=NAVY,
        )
        rect(slide, bar_left, row_y, bar_area_w, bar_h, DIVIDER_GREY)
        fill_w = max(0.05, bar_area_w * (churned / max_val))
        rect(slide, bar_left, row_y, fill_w, bar_h, MAGENTA)
        txt(
            slide,
            chart_left + chart_w - 2.8,
            row_y + 0.04,
            2.6,
            0.32,
            f"{fmt_eur(churned)}  ·  {deals} deals",
            size=12,
            bold=True,
            color=NAVY,
            align=PP_ALIGN.RIGHT,
        )

    latest = last_three[-1]
    earlier = last_three[0]
    delta = (latest.get("churned_acv", 0) or 0) - (earlier.get("churned_acv", 0) or 0)
    direction = "down" if delta < 0 else "up"
    bullets = (
        f"•  Latest quarter ({latest.get('quarter_label')}) churned ACV: {fmt_eur(latest.get('churned_acv', 0))} across {latest.get('churned_deals', 0)} deals.\n"
        f"•  Trend is {direction} {fmt_eur(abs(delta))} over the 3 quarter window.\n"
        "•  Full churn narrative requires the Finance overlay which is still pending. CRM signal only for now."
    )
    takeaway(slide, 5.75, bullets, accent=MAGENTA)
    source_footer(
        slide,
        "Source: Sales Director Monthly snapshot, 2026-04-01   ·   Value methodology: Renewal ACV",
    )


# ============================================================================
# Slide 9: Slipped deals — horizontal bar comparison
# ============================================================================


def build_slipped(prs, snap):
    slide = blank_slide(prs)
    slip = snap.get("slipped_deals", {}).get("summary_by_region", {})
    order = ["EMEA", "North America", "APAC"]

    biggest_region = max(
        order, key=lambda r: slip.get(r, {}).get("slipped_arr", 0) or 0, default="n/a"
    )
    biggest_arr = slip.get(biggest_region, {}).get("slipped_arr", 0) or 0

    header(
        slide,
        eyebrow="APRIL 2026   ·   05   ·   SLIPPED DEALS ANALYSIS",
        title="Slipped exposure by region",
        narrative=f"Largest slipped pool: {biggest_region} at {fmt_eur(biggest_arr)}. Root cause commentary pending from opportunity owners.",
        tab_color=AMBER,
    )

    chart_left = 0.6
    chart_top = 1.95
    chart_w = 12.1
    chart_h = 3.4
    rounded_card(slide, chart_left, chart_top, chart_w, chart_h, LIGHT_BLUE_PANEL)
    txt(
        slide,
        chart_left + 0.35,
        chart_top + 0.25,
        chart_w - 0.7,
        0.3,
        "SLIPPED ARR BY REGION",
        size=10,
        bold=True,
        color=PRIMARY_BLUE,
    )

    max_val = max((slip.get(r, {}).get("slipped_arr", 0) or 0) for r in order) or 1
    bar_left = chart_left + 2.4
    bar_area_w = 7.4
    bar_top_start = chart_top + 0.9
    row_h = 0.75
    bar_h = 0.38

    for i, region in enumerate(order):
        data = slip.get(region, {})
        arr = data.get("slipped_arr", 0) or 0
        count = data.get("slipped_opp_count", 0) or 0
        pushes = data.get("avg_push_count", 0) or 0
        row_y = bar_top_start + i * row_h
        # Region label
        txt(
            slide,
            chart_left + 0.35,
            row_y + 0.04,
            1.95,
            0.32,
            region,
            size=13,
            bold=True,
            color=NAVY,
        )
        # Track
        rect(slide, bar_left, row_y, bar_area_w, bar_h, DIVIDER_GREY)
        # Fill (amber for biggest, primary for others)
        color = AMBER if region == biggest_region else PRIMARY_BLUE
        fill_w = max(0.05, bar_area_w * (arr / max_val))
        rect(slide, bar_left, row_y, fill_w, bar_h, color)
        # Right side values
        txt(
            slide,
            chart_left + chart_w - 2.8,
            row_y + 0.04,
            2.6,
            0.32,
            f"{fmt_eur(arr)}  ·  {count} deals",
            size=12,
            bold=True,
            color=NAVY,
            align=PP_ALIGN.RIGHT,
        )
        # Push count below
        txt(
            slide,
            chart_left + chart_w - 2.8,
            row_y + 0.36,
            2.6,
            0.25,
            f"avg push {pushes:.1f}",
            size=9,
            color=GREY_TEXT,
            align=PP_ALIGN.RIGHT,
        )

    bullets = (
        f"•  {biggest_region} carries the largest slipped exposure at {fmt_eur(biggest_arr)}.\n"
        "•  Root cause commentary is still pending from opportunity owners this cycle.\n"
        "•  Live drilldown: Close Date Slipped by Stage (link in appendix)."
    )
    takeaway(slide, 5.75, bullets, accent=AMBER)
    source_footer(slide, "Source: Sales Director Monthly snapshot, 2026-04-01")


# ============================================================================
# Slide 10: Appendix — Salesforce links
# ============================================================================


def build_appendix(prs):
    slide = blank_slide(prs)
    rect(slide, 0.6, 0.42, 0.08, 0.7, AQUA)
    txt(slide, 0.82, 0.38, 12, 0.28, "APPENDIX", size=10, bold=True, color=PRIMARY_BLUE)
    txt(
        slide,
        0.82,
        0.62,
        12,
        0.58,
        "Live Salesforce links",
        size=28,
        bold=True,
        color=NAVY,
    )
    txt(
        slide,
        0.82,
        1.25,
        12.1,
        0.4,
        "Click any link to open the live Salesforce dashboard or report.",
        size=13,
        color=GREY_TEXT,
    )

    body_tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.9), Inches(12.1), Inches(5)
    )
    body_tf = body_tb.text_frame
    body_tf.word_wrap = True
    for i, (label, url) in enumerate(LINKS):
        para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        run = para.add_run()
        run.text = label
        run.font.name = "Microsoft Sans Serif"
        run.font.size = Pt(13)
        run.font.color.rgb = PRIMARY_BLUE
        run.font.underline = True
        run.hyperlink.address = url
        para.space_after = Pt(8)


# ============================================================================
# Main
# ============================================================================


def build():
    with open(SNAPSHOT) as f:
        snap = json.load(f)
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    build_cover(prs, snap)
    build_region_card_strip(prs, snap, "EMEA", "01")  # 4 cards
    build_region_hero_column(prs, snap, "North America", "01")  # hero column
    build_region_horizontal_bars(prs, snap, "APAC", "01")  # horizontal bars
    build_commercial_approval_global(prs, snap)  # hero stat + mini row
    build_commercial_approval_regional(prs, snap)  # table rows
    build_renewals(prs, snap)  # 4 cards in magenta
    build_churn(prs, snap)  # horizontal bar trend
    build_slipped(prs, snap)  # horizontal bar comparison
    build_appendix(prs)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
