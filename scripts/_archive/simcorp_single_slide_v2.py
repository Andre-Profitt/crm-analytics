#!/usr/bin/env python3
"""Sample slide v2 with polish moves:
- Eyebrow text above title (small caps aqua kicker)
- Vertical aqua tab next to title (not horizontal bar under it)
- Rounded corner cards
- Horizontal progress bar inside "Best case call" card
- Colored status dot next to "Forecast confidence" card
- Horizontal divider line between KPI strip and insight callout
- "Takeaway" label in small caps on insight callout
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"
RUN_DIR = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01"
)
SNAPSHOT = RUN_DIR / "report1_snapshot.json"
OUTPUT = RUN_DIR / "sample_slide_emea_v2.pptx"

NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
LIGHT_PANEL_2 = RGBColor(0xF5, 0xF8, 0xFD)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)
DIVIDER_GREY = RGBColor(0xD7, 0xE2, 0xE8)
AMBER = RGBColor(0xFB, 0x9B, 0x2A)  # warning status
GREEN_OK = RGBColor(0x33, 0xD8, 0xCE)  # success


def fmt_eur(amount):
    if amount is None:
        return "n/a"
    if abs(amount) >= 1_000_000_000:
        return f"EUR {amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"EUR {amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"EUR {amount / 1_000:.0f}K"
    return f"EUR {amount:.0f}"


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
    letter_spacing=None,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_rect(
    slide,
    left,
    top,
    width,
    height,
    fill_color,
    shape_type=MSO_SHAPE.RECTANGLE,
    line_color=None,
):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    return shape


def add_rounded_card(slide, left, top, width, height, fill_color):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.fill.background()
    # Set adjustment for corner radius (0.06 = small, 0.1 = medium)
    try:
        card.adjustments[0] = 0.08
    except Exception:
        pass
    return card


def add_kpi_card(
    slide,
    left,
    top,
    width,
    height,
    big_number,
    label,
    context,
    progress_pct=None,
    status_color=None,
):
    """Rounded KPI card with big stat, label, context, optional progress bar or status dot."""
    add_rounded_card(slide, left, top, width, height, LIGHT_BLUE_PANEL)

    # Big number
    num_tb = slide.shapes.add_textbox(
        Inches(left + 0.22),
        Inches(top + 0.18),
        Inches(width - 0.44),
        Inches(0.7),
    )
    num_tf = num_tb.text_frame
    num_tf.margin_left = num_tf.margin_right = Pt(0)
    num_tf.margin_top = num_tf.margin_bottom = Pt(0)
    p = num_tf.paragraphs[0]
    r = p.add_run()
    r.text = big_number
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_BLUE

    # Optional status dot on the right of the big number
    if status_color is not None:
        dot_size = 0.22
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + width - 0.44),
            Inches(top + 0.28),
            Inches(dot_size),
            Inches(dot_size),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = status_color
        dot.line.fill.background()

    # Label
    lbl_tb = slide.shapes.add_textbox(
        Inches(left + 0.22),
        Inches(top + 0.92),
        Inches(width - 0.44),
        Inches(0.32),
    )
    lbl_tf = lbl_tb.text_frame
    lbl_tf.margin_left = lbl_tf.margin_right = Pt(0)
    lbl_tf.margin_top = lbl_tf.margin_bottom = Pt(0)
    lbl_tf.word_wrap = True
    p = lbl_tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = NAVY

    # Context
    ctx_top = top + 1.22
    ctx_h = height - 1.30
    if progress_pct is not None:
        ctx_h -= 0.25  # leave room for progress bar
    ctx_tb = slide.shapes.add_textbox(
        Inches(left + 0.22),
        Inches(ctx_top),
        Inches(width - 0.44),
        Inches(ctx_h),
    )
    ctx_tf = ctx_tb.text_frame
    ctx_tf.margin_left = ctx_tf.margin_right = Pt(0)
    ctx_tf.margin_top = ctx_tf.margin_bottom = Pt(0)
    ctx_tf.word_wrap = True
    p = ctx_tf.paragraphs[0]
    r = p.add_run()
    r.text = context
    r.font.name = "Microsoft Sans Serif"
    r.font.size = Pt(10)
    r.font.bold = False
    r.font.color.rgb = PRIMARY_BLUE

    # Optional progress bar at the bottom
    if progress_pct is not None:
        bar_top = top + height - 0.32
        bar_left = left + 0.22
        bar_width = width - 0.44
        bar_height = 0.10
        # Track
        add_rect(slide, bar_left, bar_top, bar_width, bar_height, DIVIDER_GREY)
        # Fill
        fill_w = max(0.02, bar_width * min(1.0, progress_pct / 100))
        add_rect(slide, bar_left, bar_top, fill_w, bar_height, AQUA)


def build():
    with open(SNAPSHOT) as f:
        snap = json.load(f)
    emea = snap["pipeline"]["deck_regions"]["EMEA"]
    target = emea["target_arr"]
    best_case_call = emea["best_case_call_arr"]
    confidence = emea["forecast_confidence_pct"]
    coverage = emea["coverage_status"]
    gap = max(0.0, target - best_case_call)
    coverage_pct = (best_case_call / target * 100) if target else 0

    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    layout = get_layout(prs, "Title Only")
    slide = prs.slides.add_slide(layout)

    # Clear inherited placeholders
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # --- Vertical aqua tab next to title ---
    add_rect(slide, left=0.6, top=0.42, width=0.08, height=0.7, fill_color=AQUA)

    # --- Eyebrow (small uppercase kicker) ---
    add_text(
        slide,
        left=0.82,
        top=0.38,
        width=12,
        height=0.28,
        text="Q1 FY27   ·   PIPELINE COVERAGE",
        size=10,
        bold=True,
        color=PRIMARY_BLUE,
    )

    # --- Title ---
    add_text(
        slide,
        left=0.82,
        top=0.62,
        width=12,
        height=0.58,
        text="EMEA Regional Outlook",
        size=28,
        bold=True,
        color=NAVY,
    )

    # --- Narrative line ---
    add_text(
        slide,
        left=0.82,
        top=1.25,
        width=12.1,
        height=0.4,
        text=(
            f"EMEA is promotable. Best case call {fmt_eur(best_case_call)} "
            f"against {fmt_eur(target)} quarter target. Confidence "
            f"{confidence:.1f}%."
        ),
        size=13,
        bold=False,
        color=GREY_TEXT,
    )

    # --- KPI cards ---
    card_top = 1.9
    card_h = 1.95
    card_w = 2.95
    gap_between = 0.2
    start_x = 0.6
    cards = [
        {
            "num": fmt_eur(target),
            "label": "Quarter target",
            "ctx": "Executive target seam, ARR.",
            "progress": None,
            "status": None,
        },
        {
            "num": fmt_eur(best_case_call),
            "label": "Best case call",
            "ctx": f"{coverage_pct:.0f}% of target coverage.",
            "progress": coverage_pct,
            "status": None,
        },
        {
            "num": fmt_eur(gap),
            "label": "Gap to target",
            "ctx": "Needed from promotion or new pipeline.",
            "progress": None,
            "status": None,
        },
        {
            "num": f"{confidence:.1f}%",
            "label": "Forecast confidence",
            "ctx": f"Status: {coverage}.",
            "progress": None,
            "status": AMBER if confidence < 96 else GREEN_OK,
        },
    ]
    for i, c in enumerate(cards):
        add_kpi_card(
            slide,
            left=start_x + i * (card_w + gap_between),
            top=card_top,
            width=card_w,
            height=card_h,
            big_number=c["num"],
            label=c["label"],
            context=c["ctx"],
            progress_pct=c["progress"],
            status_color=c["status"],
        )

    # --- Horizontal divider between KPI strip and callout ---
    add_rect(
        slide,
        left=0.6,
        top=4.1,
        width=12.1,
        height=0.015,
        fill_color=DIVIDER_GREY,
    )

    # --- Insight callout block ---
    callout_top = 4.35
    callout_h = 2.0
    add_rounded_card(
        slide,
        left=0.6,
        top=callout_top,
        width=12.1,
        height=callout_h,
        fill_color=LIGHT_PANEL_2,
    )
    # Left aqua accent strip
    add_rect(
        slide,
        left=0.6,
        top=callout_top,
        width=0.1,
        height=callout_h,
        fill_color=AQUA,
    )

    # "Takeaway" eyebrow
    add_text(
        slide,
        left=0.95,
        top=callout_top + 0.16,
        width=11.5,
        height=0.28,
        text="TAKEAWAY",
        size=10,
        bold=True,
        color=PRIMARY_BLUE,
    )
    # Callout body (3 concise bullets)
    bullets = (
        "•  Gap is promotion dependent. No new pipeline is strictly required this quarter.\n"
        "•  CACEIS Netherlands and Groupama AM are the largest promotion candidates in EMEA.\n"
        "•  29 days remaining in Q1 FY27. Forecast confidence trails APAC and North America at 94.5%."
    )
    add_text(
        slide,
        left=0.95,
        top=callout_top + 0.45,
        width=11.5,
        height=callout_h - 0.55,
        text=bullets,
        size=12,
        bold=False,
        color=NAVY,
    )

    # --- Source footer ---
    add_text(
        slide,
        left=0.6,
        top=7.08,
        width=12.1,
        height=0.22,
        text="Source: Sales Director Monthly snapshot, 2026-04-01   ·   Value methodology: Renewals and churn in Renewal ACV, land and expand in ARR",
        size=9,
        bold=False,
        color=GREY_TEXT,
    )

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
