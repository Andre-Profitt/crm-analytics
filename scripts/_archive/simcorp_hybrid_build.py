#!/usr/bin/env python3
"""Hybrid build: slide 1 uses the real SC-Master Gradient_Title layout,
slides 2-N are transplanted from the canonical deck and color/font-remapped
via the Path A pattern.

Output:
    sales_director_monthly_simcorp_hybrid.pptx
"""

from __future__ import annotations

import json
import re
import shutil
import zipfile
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

TEMPLATE = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/Commercial Update - Dec 2025.pptx"
RUN_DIR = Path(
    "/Users/test/crm-analytics/output/sales_director_monthly_runs/2026-04-06T20-00-11Z_2026-04-01"
)
CANONICAL = RUN_DIR / "sales_director_monthly_pipeline_insights_2026-04-01.pptx"
SNAPSHOT = RUN_DIR / "report1_snapshot.json"
OUTPUT = RUN_DIR / "sales_director_monthly_simcorp_hybrid.pptx"

COLOR_MAP = {
    "0A6C74": "0E3788",
    "0A4D57": "011946",
    "B45A43": "9D2E7B",
    "A7852C": "9D2E7B",
    "EEF3F5": "E6EEFE",
    "FBF8F2": "FFFFFF",
    "F7FAFB": "FFFFFF",
    "DCEEF0": "6FCCDD",
    "F7E3DD": "E6EEFE",
    "F4E8BF": "E6EEFE",
    "E3F0E7": "E6EEFE",
    "CDE7EA": "6FCCDD",
    "123040": "011946",
    "5C7482": "0E3788",
    "0F2430": "011946",
}
FONT_MAP = {"Avenir Next": "Microsoft Sans Serif"}


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def fill_ph(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return True
    return False


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def transplant(src_slide, dst_prs, blank_layout):
    new_slide = dst_prs.slides.add_slide(blank_layout)
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)
    for shape in src_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue  # drop duplicate logos
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)
    return new_slide


def build():
    with open(SNAPSHOT) as f:
        snap = json.load(f)

    src = Presentation(str(CANONICAL))
    dst = Presentation(TEMPLATE)
    clear_slides(dst)

    cover_layout = get_layout(dst, "SC-Master Gradient_Title")
    blank_layout = get_layout(dst, "Blank")

    # --- Slide 1: real SimCorp cover layout ---
    cover = dst.slides.add_slide(cover_layout)
    fill_ph(cover, 0, "Sales Director Monthly Pipeline Insights")
    fill_ph(
        cover,
        20,
        f"{snap.get('quarter_focus', 'Q1')} FY27 forward look. Pipeline coverage, commercial approvals, renewals, churn, and slipped deals.",
    )
    fill_ph(cover, 22, f"Snapshot {snap.get('snapshot_date', '2026-04-01')}")

    # --- Slides 2..N: transplant from canonical slides 2..N onto Blank ---
    src_slides = list(src.slides)
    for src_slide in src_slides[1:]:
        transplant(src_slide, dst, blank_layout)

    # Save intermediate
    stage_path = OUTPUT.with_name(OUTPUT.stem + "_stage.pptx")
    dst.save(str(stage_path))

    # --- Apply color + font remap to slides 2..N (skip slide 1 since it
    #     uses real SimCorp layout and placeholders that inherit theme colors).
    apply_remap(stage_path, OUTPUT, skip_first_slide=True)
    stage_path.unlink()


def apply_remap(src_path: Path, dst_path: Path, skip_first_slide: bool = False):
    shutil.copy2(src_path, dst_path)
    tmp_path = dst_path.with_suffix(".tmp.pptx")

    with zipfile.ZipFile(dst_path, "r") as zin:
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.startswith(
                    "ppt/slides/slide"
                ) and item.filename.endswith(".xml"):
                    # Extract slide number from slideNN.xml
                    m = re.search(r"slide(\d+)\.xml", item.filename)
                    slide_num = int(m.group(1)) if m else -1
                    if skip_first_slide and slide_num == 1:
                        pass  # don't remap cover slide
                    else:
                        text = data.decode("utf-8")
                        for old, new in COLOR_MAP.items():
                            text = re.sub(
                                rf'val="({old})"',
                                f'val="{new}"',
                                text,
                                flags=re.IGNORECASE,
                            )
                        for old, new in FONT_MAP.items():
                            text = re.sub(
                                rf'typeface="{re.escape(old)}"',
                                f'typeface="{new}"',
                                text,
                                flags=re.IGNORECASE,
                            )
                        data = text.encode("utf-8")
                zout.writestr(item, data)

    tmp_path.replace(dst_path)


if __name__ == "__main__":
    build()
    print(f"Saved: {OUTPUT}")
