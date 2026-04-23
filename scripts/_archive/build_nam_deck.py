#!/usr/bin/env python3
"""Build Sales Director Monthly deck — insight-driven, SimCorp branded.

Pulls live data from D1 (dashboard-filtered per director) + Pipeline Inspection SOQL.
Generates a polished, narrative deck with headline insights and supporting detail.

Usage:
    python3 scripts/build_nam_deck.py --director "Dan Peppett"
    python3 scripts/build_nam_deck.py --all
"""

from __future__ import annotations

import json
import subprocess
import time
from datetime import date
from pathlib import Path
from typing import Any

import requests
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO_ROOT = Path(__file__).resolve().parents[1]
CONFIG_PATH = REPO_ROOT / "config" / "sales_director_md1_presets.json"
TEMPLATE_PATH = "/Users/test/archive/simcorp-deck-agent-backup/reference-decks/SimCorp_PPT_Template.pptx"
DEFAULT_OUTPUT_ROOT = REPO_ROOT / "output" / "sales_director_monthly_runs"
TARGET_ORG = "apro@simcorp.com"
API = "v66.0"
DASHBOARD1_ID = "01ZTb00000FSP7hMAH"

# Dashboard filter option IDs (stable across runs)
_OPT: dict[str, str] = {
    # Industry (filter1)
    "ind_asset_mgmt": "0ICTb0000007DbdOAE",
    "ind_bank": "0ICTb0000007DbeOAE",
    "ind_insurance": "0ICTb0000007DbfOAE",
    "ind_pension": "0ICTb0000007DbgOAE",
    "ind_wealth": "0ICTb0000007DbhOAE",
    "ind_servicer": "0ICTb0000007DbiOAE",
    "ind_other": "0ICTb0000007DbjOAE",
    # Legal Country (filter2)
    "lc_canada": "0ICTb0000007DgTOAU",
    "lc_excl_canada": "0ICTb0000007DgUOAU",
    # Sales Region (filter3)
    "sr_apac": "0ICTb0000007DbnOAE",
    "sr_central_europe": "0ICTb0000007DboOAE",
    "sr_mea": "0ICTb0000007DbpOAE",
    "sr_nam": "0ICTb0000007DbqOAE",
    "sr_northern_europe": "0ICTb0000007DbrOAE",
    "sr_southwestern_europe": "0ICTb0000007DbsOAE",
    "sr_uki": "0ICTb0000007DbtOAE",
    # Account Unit Group (filter4)
    "aug_sc_nam": "0ICTb0000007Di5OAE",
    "aug_sc_asia": "0ICTb0000007Di6OAE",
    "aug_sc_emea": "0ICTb0000007Di7OAE",
}

# Per-director dashboard filter params
DIRECTOR_D1_FILTERS: dict[str, dict[str, str]] = {
    "Megan Miceli": {
        "filter2": _OPT["lc_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
    "Patrick Gaughan": {
        "filter1": [
            _OPT["ind_asset_mgmt"],
            _OPT["ind_bank"],
            _OPT["ind_wealth"],
            _OPT["ind_servicer"],
            _OPT["ind_other"],
        ],
        "filter2": _OPT["lc_excl_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
    "Jesper Tyrer": {"filter3": _OPT["sr_apac"], "filter4": _OPT["aug_sc_asia"]},
    "Sarah Pittroff": {
        "filter3": _OPT["sr_central_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Francois Thaury": {
        "filter3": _OPT["sr_southwestern_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Dan Peppett": {"filter3": _OPT["sr_uki"], "filter4": _OPT["aug_sc_emea"]},
    "Christian Ebbesen": {
        "filter3": _OPT["sr_northern_europe"],
        "filter4": _OPT["aug_sc_emea"],
    },
    "Mourad Essofi": {"filter3": _OPT["sr_mea"], "filter4": _OPT["aug_sc_emea"]},
    "Adam Steinhaus": {
        "filter1": [_OPT["ind_pension"], _OPT["ind_insurance"]],
        "filter2": _OPT["lc_excl_canada"],
        "filter3": _OPT["sr_nam"],
        "filter4": _OPT["aug_sc_nam"],
    },
}

TEAL_DEEP = RGBColor(0x00, 0x3E, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x1B, 0x1B)


def auth():
    r = subprocess.run(
        ["sf", "org", "display", "--target-org", TARGET_ORG, "--json"],
        capture_output=True,
        text=True,
        check=True,
    )
    d = json.loads(r.stdout)["result"]
    return d["accessToken"], d["instanceUrl"]


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout not found: {name}")


def fill(slide, idx, text):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = str(text)
            return True
    return False


def clear_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def fmt_eur(v):
    if not v:
        return "EUR 0"
    if abs(v) >= 1_000_000:
        return f"EUR {v / 1_000_000:.1f}M"
    if abs(v) >= 1_000:
        return f"EUR {v / 1_000:.0f}K"
    return f"EUR {v:,.0f}"


def _tight_cell(cell, font_size, bold=False, color=INK, bg=None):
    """Format a cell with minimal margins."""
    from pptx.oxml.ns import qn

    tcPr = cell._tc.get_or_add_tcPr()
    margin = 5000  # ~0.003in — minimal
    tcPr.set(qn("a:marL"), str(margin * 2))
    tcPr.set(qn("a:marR"), str(margin * 2))
    tcPr.set(qn("a:marT"), str(margin))
    tcPr.set(qn("a:marB"), str(margin))
    for p in cell.text_frame.paragraphs:
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_table(slide, headers, rows, left=0.9, top=2.2, width=11.5, height=None):
    if not rows:
        return
    n_cols = min(len(headers), 7)
    n_rows = len(rows) + 1
    row_h = 0.20
    h = row_h * n_rows + 0.05
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(h)
    )
    tbl = tbl_shape.table
    for i, h_text in enumerate(headers[:n_cols]):
        cell = tbl.cell(0, i)
        cell.text = str(h_text)
        _tight_cell(cell, 11, bold=True, color=WHITE, bg=TEAL_DEEP)
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row[:n_cols]):
            cell = tbl.cell(ri, ci)
            cell.text = str(val or "")
            _tight_cell(cell, 10, color=INK)


def _plotly_chart(fig, path, width=1800, height=700):
    """Write a plotly figure to PNG and return the path."""
    fig.update_layout(
        font_family="Avenir Next, Helvetica, Arial",
        font_color="#003E52",
        plot_bgcolor="white",
        paper_bgcolor="white",
        margin=dict(l=20, r=20, t=40, b=20),
    )
    fig.write_image(path, width=width, height=height, scale=2)
    return path


_chart_counter = [0]


def _chart_path():
    _chart_counter[0] += 1
    return f"/tmp/plotly_chart_{_chart_counter[0]}.png"


# SimCorp palette
_TEAL = "#003E52"
_TEAL_MID = "#007480"
_TEAL_LIGHT = "#8FBFC4"
_CORAL = "#E86C6C"
_GOLD = "#D4A843"
_SLATE = "#5A7A8A"
_GREEN = "#2ECC71"


def add_bar_chart(
    slide,
    categories,
    values,
    series_name="ARR",
    left=0.9,
    top=2.2,
    width=11.5,
    height=4.0,
):
    """Pipeline treemap — each stage sized proportionally by ARR, always readable."""
    import plotly.express as px

    labels = [
        f"{c}<br>EUR {v / 1e6:.1f}M" if v >= 1e6 else f"{c}<br>EUR {v / 1e3:.0f}K"
        for c, v in zip(categories, values)
    ]
    colors = [
        _TEAL,
        _TEAL_MID,
        _TEAL_LIGHT,
        _GOLD,
        _SLATE,
        _CORAL,
        "#C45B5B",
        "#7AA3B0",
    ][: len(values)]
    fig = px.treemap(
        names=labels,
        parents=[""] * len(labels),
        values=values,
        color_discrete_sequence=colors,
    )
    fig.update_traces(
        textinfo="label",
        textfont=dict(size=16, color="white"),
        marker=dict(cornerradius=5),
    )
    fig.update_layout(margin=dict(l=5, r=5, t=5, b=5))
    path = _chart_path()
    _plotly_chart(fig, path, width=int(width * 150), height=int(height * 160))
    slide.shapes.add_picture(
        path, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def add_pie_chart(slide, categories, values, left=0.9, top=2.2, width=5.0, height=4.0):
    """Modern donut with center number."""
    import plotly.graph_objects as go

    total = sum(values)
    fig = go.Figure(
        go.Pie(
            labels=categories,
            values=values,
            hole=0.55,
            textinfo="label+value",
            textfont=dict(size=14),
            marker=dict(
                colors=[_TEAL_MID, _CORAL, _GOLD, _SLATE][: len(values)],
                line=dict(color="white", width=2),
            ),
            pull=[0.03] * len(values),
        )
    )
    fig.add_annotation(
        text=f"<b>{total}</b><br>Total",
        showarrow=False,
        font=dict(size=22, color=_TEAL),
    )
    fig.update_layout(
        showlegend=True, legend=dict(orientation="h", y=-0.1, font=dict(size=12))
    )
    path = _chart_path()
    _plotly_chart(fig, path, width=int(width * 150), height=int(height * 160))
    slide.shapes.add_picture(
        path, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def add_forecast_chart(
    slide, categories, values, left=6.8, top=2.2, width=5.8, height=4.0
):
    """Waterfall chart for forecast categories."""
    import plotly.graph_objects as go

    color_map = {
        "Commit": _TEAL,
        "Best Case": _TEAL_MID,
        "Pipeline": _GOLD,
        "Omitted": _SLATE,
        "Won": _GREEN,
        "Lost": _CORAL,
    }
    colors = [color_map.get(c, _TEAL_MID) for c in categories]
    fig = go.Figure(
        go.Bar(
            x=categories,
            y=values,
            marker_color=colors,
            text=[
                f"EUR {v / 1e6:.1f}M"
                if v >= 1e6
                else f"EUR {v / 1e3:.0f}K"
                if v >= 1e3
                else f"EUR {v:,.0f}"
                for v in values
            ],
            textposition="outside",
            textfont=dict(size=11, color=_TEAL),
        )
    )
    fig.update_layout(
        yaxis=dict(visible=False, showgrid=False),
        xaxis=dict(tickangle=-25, tickfont=dict(size=11)),
        bargap=0.3,
    )
    path = _chart_path()
    _plotly_chart(fig, path, width=int(width * 150), height=int(height * 160))
    slide.shapes.add_picture(
        path, Inches(left), Inches(top), Inches(width), Inches(height)
    )


import argparse  # noqa: E402 — used in parse_args / main
import sys  # noqa: E402

_ = CategoryChartData, XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: F841


# ---------------------------------------------------------------------------
# Director helpers
# ---------------------------------------------------------------------------


def load_preset(director_name: str) -> dict[str, Any]:
    """Look up a director's preset by name (fuzzy match on first/last)."""

    config = json.loads(CONFIG_PATH.read_text())
    name_lower = director_name.lower().strip()
    for preset in config["presets"]:
        pname = preset["name"].lower()
        if (
            pname == name_lower
            or name_lower in pname
            or pname.split()[0] == name_lower.split()[0]
        ):
            return preset
    raise ValueError(
        f"No preset matched {director_name!r}. "
        f"Available: {[p['name'] for p in config['presets']]}"
    )


def slugify_director(name: str) -> str:
    parts = name.lower().strip().split()
    if len(parts) >= 2:
        return f"{parts[-1]}-{parts[0]}"
    return parts[0]


def _soql_where(director: dict[str, Any]) -> str:
    """Build SOQL WHERE fragment from a director's preset filters.

    Maps Reports API column names to SOQL fields:
      Account.Region__c        -> Account.Region__c
      INDUSTRY                 -> Account.Industry
      ADDRESS1_COUNTRY_CODE    -> Account.BillingCountryCode
      Opportunity.Account_Unit_Group__c -> Account_Unit_Group__c (on Opportunity)
    """
    clauses: list[str] = []
    for f in director.get("filters", []):
        col = f.get("column", "")
        op = f.get("operator", "equals")
        val = f.get("value", "")
        if col == "Account.Region__c":
            if op == "equals":
                clauses.append(f"Account.Region__c = '{val}'")
            elif op == "notEqual":
                clauses.append(f"Account.Region__c != '{val}'")
        elif col == "INDUSTRY":
            if op == "equals":
                if "," in val:
                    in_list = ", ".join(f"'{v.strip()}'" for v in val.split(","))
                    clauses.append(f"Account.Industry IN ({in_list})")
                else:
                    clauses.append(f"Account.Industry = '{val}'")
        elif col == "ADDRESS1_COUNTRY_CODE":
            if op == "equals":
                clauses.append(f"Account.BillingCountryCode = '{val}'")
            elif op == "notEqual":
                clauses.append(f"Account.BillingCountryCode != '{val}'")
        elif col == "Opportunity.Account_Unit_Group__c":
            if op == "equals":
                clauses.append(f"Account_Unit_Group__c = '{val}'")
    return " AND ".join(clauses) if clauses else "Account.Region__c != null"


# ---------------------------------------------------------------------------
# Core deck builder (parameterized per director)
# ---------------------------------------------------------------------------


def build_deck_for_director(
    director: dict[str, Any], output_path: Path, token: str, base_url: str
) -> None:
    """Build a full SD Monthly deck for one director."""
    headers = {"Authorization": f"Bearer {token}"}
    territory = director["territory"]
    d1_filters = DIRECTOR_D1_FILTERS.get(director["name"], {})
    region_where = _soql_where(director)

    # --- Fetch D1 filtered for this director ---
    print(f"Fetching D1 (filtered for {director['name']})...")
    requests.put(
        f"{base_url}/services/data/{API}/analytics/dashboards/{DASHBOARD1_ID}",
        headers=headers,
        params=d1_filters,
        timeout=60,
    )
    time.sleep(6)
    d1 = requests.get(
        f"{base_url}/services/data/{API}/analytics/dashboards/{DASHBOARD1_ID}",
        headers=headers,
        params=d1_filters,
        timeout=60,
    ).json()

    # Map components
    comp = {}
    comp_ids = {
        "01aTb00000Cn9mPIAR": "pipeline",
        "01aTb00000Cn9mQIAR": "approval_state",
        "01aTb00000Cn85bIAB": "renewal_likelihood",
        "01aTb00000Cn85dIAB": "business_at_risk",
        "01aTb00000Cn85aIAB": "approval_ytd",
        "01aTb00000Cn85jIAB": "approval_candidates",
        "01aTb00000Cn85ZIAR": "renewal_pipeline",
        "01aTb00000Cn85lIAB": "slipped",
    }
    for cd in d1.get("componentData", []):
        cid = cd["componentId"]
        if cid in comp_ids:
            comp[comp_ids[cid]] = cd.get("reportResult", {})

    # --- Extract pipeline metrics ---
    p_fm = comp.get("pipeline", {}).get("factMap", {})
    p_gd = (comp.get("pipeline", {}).get("groupingsDown") or {}).get("groupings", [])
    p_grand = p_fm.get("T!T", {}).get("aggregates", [])
    total_pipeline = p_grand[0]["value"] if p_grand else 0
    stages = []
    for g in p_gd:
        aggs = p_fm.get(f"{g['key']}!T", {}).get("aggregates", [])
        arr = aggs[0]["value"] if aggs else 0
        stages.append((g["label"], arr))
    biggest_stage = max(stages, key=lambda x: x[1]) if stages else ("?", 0)
    biggest_pct = (biggest_stage[1] / total_pipeline * 100) if total_pipeline else 0

    # --- Commercial approval ---
    a_fm = comp.get("approval_state", {}).get("factMap", {})
    a_gd = (comp.get("approval_state", {}).get("groupingsDown") or {}).get(
        "groupings", []
    )
    approved = 0
    pending = 0
    for g in a_gd:
        aggs = a_fm.get(f"{g['key']}!T", {}).get("aggregates", [])
        count = aggs[0]["value"] if aggs else 0
        if g["label"] == "true":
            approved = int(count)
        else:
            pending = int(count)
    total_deals = approved + pending
    approval_rate = (approved / total_deals * 100) if total_deals else 0

    # Approved YTD
    ay_fm = comp.get("approval_ytd", {}).get("factMap", {})
    ay_grand = ay_fm.get("T!T", {}).get("aggregates", [])
    ytd_arr = ay_grand[0]["value"] if ay_grand else 0
    ytd_count = int(ay_grand[1]["value"]) if len(ay_grand) > 1 else 0

    # Candidates missing approval
    c_fm = comp.get("approval_candidates", {}).get("factMap", {})
    c_grand = c_fm.get("T!T", {}).get("aggregates", [])
    cand_arr = c_grand[0]["value"] if c_grand else 0
    cand_count = int(c_grand[1]["value"]) if len(c_grand) > 1 else 0
    cand_rows = c_fm.get("T!T", {}).get("rows", [])
    cand_detail = []
    for row in cand_rows[:10]:
        cells = row.get("dataCells", [])
        cand_detail.append([c.get("label", "") for c in cells[:6]])

    # --- Renewals ---
    r_fm = comp.get("renewal_pipeline", {}).get("factMap", {})
    r_grand = r_fm.get("T!T", {}).get("aggregates", [])
    renewal_acv = r_grand[0]["value"] if r_grand else 0
    renewal_rows = r_fm.get("T!T", {}).get("rows", [])

    # --- Slipped (aging >90d from D1) ---
    s_fm = comp.get("slipped", {}).get("factMap", {})
    s_grand = s_fm.get("T!T", {}).get("aggregates", [])
    aging_count = int(s_grand[2]["value"]) if len(s_grand) > 2 else 0
    aging_rows = s_fm.get("T!T", {}).get("rows", [])
    aging_detail = []
    for row in aging_rows[:10]:
        cells = row.get("dataCells", [])
        aging_detail.append([c.get("label", "") for c in cells[:7]])

    # --- Pipeline Inspection SOQL ---
    print(f"Fetching Pipeline Inspection ({territory}, PushCount>0)...")
    pi_soql = f"SELECT Name, Account.Name, StageName, CloseDate, PushCount, APTS_Opportunity_ARR__c, ForecastCategoryName, Owner.Name FROM Opportunity WHERE {region_where} AND IsClosed = false AND PushCount > 0 ORDER BY PushCount DESC LIMIT 15"
    pi_records = (
        requests.get(
            f"{base_url}/services/data/{API}/query/",
            headers=headers,
            params={"q": pi_soql},
            timeout=60,
        )
        .json()
        .get("records", [])
    )

    # Push summary stats
    push_summary = requests.get(
        f"{base_url}/services/data/{API}/query/",
        headers=headers,
        params={
            "q": f"SELECT COUNT(Id) total, SUM(APTS_Opportunity_ARR__c) arr, AVG(PushCount) avg_push FROM Opportunity WHERE IsClosed = false AND PushCount > 0 AND {region_where}"
        },
        timeout=60,
    ).json()["records"][0]
    push_total = push_summary.get("total", 0)
    push_arr = push_summary.get("arr", 0) or 0
    push_avg = push_summary.get("avg_push", 0) or 0

    # Push buckets from full SOQL (not just top 15)
    bucket_recs = (
        requests.get(
            f"{base_url}/services/data/{API}/query/",
            headers=headers,
            params={
                "q": f"SELECT PushCount, COUNT(Id) ct, SUM(APTS_Opportunity_ARR__c) arr FROM Opportunity WHERE IsClosed = false AND PushCount > 0 AND {region_where} GROUP BY PushCount ORDER BY PushCount DESC"
            },
            timeout=60,
        )
        .json()
        .get("records", [])
    )
    heavy_ct = sum(r["ct"] for r in bucket_recs if r["PushCount"] >= 5)
    heavy_arr = sum((r["arr"] or 0) for r in bucket_recs if r["PushCount"] >= 5)
    mid_ct = sum(r["ct"] for r in bucket_recs if 3 <= r["PushCount"] < 5)
    mid_arr = sum((r["arr"] or 0) for r in bucket_recs if 3 <= r["PushCount"] < 5)
    light_ct = sum(r["ct"] for r in bucket_recs if 1 <= r["PushCount"] < 3)
    light_arr = sum((r["arr"] or 0) for r in bucket_recs if 1 <= r["PushCount"] < 3)

    top_pushers = pi_records[:3]
    owner_counts: dict[str, int] = {}
    for r in pi_records:
        owner = (r.get("Owner") or {}).get("Name", "?")
        owner_counts[owner] = owner_counts.get(owner, 0) + 1
    top_owner = (
        max(owner_counts.items(), key=lambda x: x[1]) if owner_counts else ("?", 0)
    )

    # --- Forecast accuracy: Won / Lost / Forecast categories this quarter ---
    print(f"Fetching forecast accuracy ({territory}, this quarter)...")
    won_data = requests.get(
        f"{base_url}/services/data/{API}/query/",
        headers=headers,
        params={
            "q": f"SELECT SUM(APTS_Opportunity_ARR__c) arr, COUNT(Id) ct FROM Opportunity WHERE StageName = '8 - Won' AND CloseDate = THIS_QUARTER AND {region_where}"
        },
        timeout=60,
    ).json()["records"][0]
    lost_data = requests.get(
        f"{base_url}/services/data/{API}/query/",
        headers=headers,
        params={
            "q": f"SELECT SUM(APTS_Opportunity_ARR__c) arr, COUNT(Id) ct FROM Opportunity WHERE StageName = '0 - Lost' AND CloseDate = THIS_QUARTER AND {region_where}"
        },
        timeout=60,
    ).json()["records"][0]
    fc_data = (
        requests.get(
            f"{base_url}/services/data/{API}/query/",
            headers=headers,
            params={
                "q": f"SELECT ForecastCategoryName, SUM(APTS_Opportunity_ARR__c) arr, COUNT(Id) ct FROM Opportunity WHERE IsClosed = false AND CloseDate = THIS_QUARTER AND {region_where} GROUP BY ForecastCategoryName"
            },
            timeout=60,
        )
        .json()
        .get("records", [])
    )

    won_arr = won_data.get("arr", 0) or 0
    won_ct = won_data.get("ct", 0) or 0
    lost_arr = lost_data.get("arr", 0) or 0
    lost_ct = lost_data.get("ct", 0) or 0
    win_rate_ct = (won_ct / (won_ct + lost_ct) * 100) if (won_ct + lost_ct) > 0 else 0
    win_rate_arr = (
        (won_arr / (won_arr + lost_arr) * 100) if (won_arr + lost_arr) > 0 else 0
    )
    commit_arr = sum(
        (r.get("arr") or 0)
        for r in fc_data
        if r.get("ForecastCategoryName") == "Commit"
    )
    best_case_arr = sum(
        (r.get("arr") or 0)
        for r in fc_data
        if r.get("ForecastCategoryName") == "Best Case"
    )
    pipeline_arr = sum(
        (r.get("arr") or 0)
        for r in fc_data
        if r.get("ForecastCategoryName") == "Pipeline"
    )
    omitted_arr = sum(
        (r.get("arr") or 0)
        for r in fc_data
        if r.get("ForecastCategoryName") == "Omitted"
    )

    # =====================================================================
    # BUILD DECK
    # =====================================================================
    print("Building deck...")
    prs = Presentation(TEMPLATE_PATH)
    clear_slides(prs)
    today = date.today().isoformat()

    # --- COVER ---
    slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
    fill(slide, 20, "Sales Director Monthly")
    fill(slide, 22, f"Snapshot {today}")
    fill(slide, 24, f"{territory}\nPipeline Reporting and Insights")

    # --- EXEC SUMMARY ---
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    fill(slide, 42, f"Executive summary - {territory}, Q2 2026")
    fill(slide, 61, fmt_eur(total_pipeline))
    fill(slide, 62, f"{won_ct}W / {lost_ct}L")
    fill(slide, 63, fmt_eur(commit_arr))
    fill(slide, 64, str(heavy_ct))
    fill(
        slide,
        22,
        f"Pipeline: {fmt_eur(total_pipeline)} open. {biggest_pct:.0f}% in {biggest_stage[0]} - conversion to later stages is key. {cand_count} deals ({fmt_eur(cand_arr)}) missing commercial approval.",
    )
    fill(
        slide,
        55,
        f"Win/Loss: {won_ct} won ({fmt_eur(won_arr)}) vs {lost_ct} lost ({fmt_eur(lost_arr)}) this quarter. ARR win rate {win_rate_arr:.1f}% - lost deals significantly outweigh wins.",
    )
    fill(
        slide,
        57,
        f"Commit forecast: {fmt_eur(commit_arr)} committed, {fmt_eur(best_case_arr)} best case, {fmt_eur(pipeline_arr)} in pipeline. {'No renewals due.' if renewal_acv == 0 else fmt_eur(renewal_acv) + ' renewal ACV due.'}",
    )
    fill(
        slide,
        59,
        f"Risk: {heavy_ct} deals pushed 5+ times ({fmt_eur(heavy_arr)} ARR). {top_owner[0]} owns {top_owner[1]} of the most-pushed deals.",
    )

    # --- AGENDA ---
    slide = prs.slides.add_slide(get_layout(prs, "Agenda 1"))
    items = {
        40: "Pipeline coverage by stage",
        69: "Commercial approval overview",
        71: "Renewals tracking",
        74: "Churn risk and trends",
        76: "Slipped deals and Pipeline Inspection",
        101: "Actions and next steps",
    }
    for idx, text in items.items():
        fill(slide, idx, text)

    # ---- 01 PIPELINE ----
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "01")
    fill(slide, 20, "Pipeline coverage by stage")

    # Pipeline chart slide
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 144, f"Pipeline Overview - {territory}, Q2 2026")
    fill(slide, 145, f"Total open pipeline: {fmt_eur(total_pipeline)}")
    insight = f"{biggest_pct:.0f}% of pipeline ({fmt_eur(biggest_stage[1])}) sits in {biggest_stage[0]}."
    if biggest_pct > 60:
        insight += " Pipeline is concentrated - conversion to later stages is the key focus area."
    fill(slide, 42, insight)
    chart_stages = [s for s in stages if s[1] > 0]
    add_bar_chart(
        slide, [s[0] for s in chart_stages], [s[1] for s in chart_stages], "ARR"
    )

    # Top deals
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 144, "Top Deals in Pipeline")
    fill(slide, 145, f"Largest opportunities by ARR - {territory}")
    top_deal_headers = [
        "Account",
        "Opportunity",
        "Owner",
        "Stage",
        "Close Date",
        "Age",
        "ARR",
    ]
    fill(
        slide,
        42,
        "Open opportunities sorted by deal size. Focus: are these progressing or stalling?",
    )
    add_table(slide, top_deal_headers, aging_detail[:10])

    # ---- 02 COMMERCIAL APPROVAL ----
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "02")
    fill(slide, 20, "Commercial approval overview")

    # Overview
    slide = prs.slides.add_slide(get_layout(prs, "2 x content w/ gradient line"))
    fill(slide, 42, f"Commercial Approval - {territory}")
    fill(slide, 61, str(approved))
    fill(slide, 62, str(pending))
    fill(
        slide,
        22,
        f"Deals with commercial approval. {ytd_count} approved YTD totaling {fmt_eur(ytd_arr)}.",
    )
    fill(
        slide,
        55,
        f"Deals pending approval. {cand_count} at Land Stage 3 worth {fmt_eur(cand_arr)} need immediate attention.",
    )
    fill(slide, 56, "Pending approval")
    add_pie_chart(
        slide,
        ["Approved", "Pending"],
        [approved, pending],
        left=8.5,
        top=3.2,
        width=4.0,
        height=3.0,
    )

    # Candidates detail
    if cand_detail:
        slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
        fill(slide, 144, "Land Stage 3 - Missing Commercial Approval")
        fill(slide, 145, f"{cand_count} candidates totaling {fmt_eur(cand_arr)}")
        fill(
            slide,
            42,
            "These deals are in Engagement stage without Go/No-Go approval. Action: escalate to approval committee.",
        )
        cand_headers = [
            "Account",
            "Opportunity",
            "Owner",
            "Close Date",
            "Next Step",
            "ARR",
        ]
        add_table(slide, cand_headers, cand_detail)

    # ---- 03 RENEWALS ----
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "03")
    fill(slide, 20, "Renewals tracking")

    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 144, f"Renewal Pipeline - {territory}, Q2 2026")
    fill(slide, 145, f"Renewal ACV this quarter: {fmt_eur(renewal_acv)}")
    if renewal_acv == 0:
        fill(
            slide,
            42,
            f"No open renewals due this quarter for {territory}. Clean quarter - focus remains on new business pipeline conversion.",
        )
        fill(slide, 22, "No renewals to display.")
    else:
        fill(
            slide,
            42,
            f"Renewals worth {fmt_eur(renewal_acv)} due this quarter. Review probability and timing below.",
        )
        r_detail = []
        for row in renewal_rows[:10]:
            cells = row.get("dataCells", [])
            r_detail.append([c.get("label", "") for c in cells[:6]])
        add_table(
            slide,
            ["Account", "Opportunity", "Owner", "Stage", "Close Date", "ACV"],
            r_detail,
        )

    # ---- 04 CHURN ----
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "04")
    fill(slide, 20, "Churn risk and trends")

    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 144, "Churn Risk - Pending Finance Feed")
    fill(slide, 145, "Status: Awaiting data from Finance (Alex P)")
    fill(
        slide,
        42,
        "Churn reporting requires Finance input on renewal attrition and at-risk accounts. Outreach to Alex P is in progress to establish a recurring feed for this section.",
    )
    fill(
        slide,
        22,
        "This section will be populated once the Finance data pipeline is established.",
    )

    # ---- 05 MOVED OUT / PUSHED DEALS ----
    slide = prs.slides.add_slide(get_layout(prs, "Divider 1"))
    fill(slide, 24, "05")
    fill(slide, 20, "Moved out and slipped deals")

    # Summary slide
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    fill(slide, 42, f"Pushed Deals Summary - {territory}")
    fill(slide, 61, str(push_total))
    fill(slide, 62, f"{push_avg:.1f}x avg")
    fill(slide, 63, fmt_eur(push_arr))
    fill(slide, 64, str(heavy_ct))
    fill(
        slide,
        22,
        f"{push_total} open deals pushed. {top_owner[0]} owns {top_owner[1]} of the top pushed deals - pattern review recommended.",
    )
    fill(
        slide,
        55,
        f"Average pushes per deal. {heavy_ct} deals pushed 5+ times ({fmt_eur(heavy_arr)} ARR) are the highest risk.",
    )
    fill(slide, 57, f"Total ARR exposed across all {push_total} pushed deals.")
    fill(
        slide,
        59,
        f"Critical: {heavy_ct} at 5+ pushes. Watch: {mid_ct} at 3-4 pushes ({fmt_eur(mid_arr)}). Early: {light_ct} at 1-2 pushes ({fmt_eur(light_arr)}).",
    )

    # Detail table
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 144, "Pipeline Inspection - Top Pushed Deals")
    fill(slide, 145, f"Top {min(len(pi_records), 12)} by push count")
    top_push_name = top_pushers[0]["Name"][:35] if top_pushers else "n/a"
    top_push_ct = top_pushers[0].get("PushCount", 0) if top_pushers else 0
    fill(
        slide,
        42,
        f"Highest: {top_push_name} ({top_push_ct}x). Deals pushed 5+ times warrant direct conversation with owner.",
    )
    pi_table = []
    for r in pi_records[:12]:
        pi_table.append(
            [
                (r.get("Name") or "")[:35],
                ((r.get("Account") or {}).get("Name") or "")[:22],
                r.get("StageName", "")[:18],
                (r.get("CloseDate") or "")[:10],
                str(r.get("PushCount", 0)),
                fmt_eur(r.get("APTS_Opportunity_ARR__c") or 0),
                ((r.get("Owner") or {}).get("Name") or "")[:18],
            ]
        )
    add_table(
        slide,
        ["Opportunity", "Account", "Stage", "Close", "Pushes", "ARR", "Owner"],
        pi_table,
    )

    # ---- 05b FORECAST ACCURACY ----
    slide = prs.slides.add_slide(get_layout(prs, "4 x content w/ gradient line"))
    fill(slide, 42, f"Forecast Accuracy - {territory}, Q2 2026")
    fill(slide, 61, fmt_eur(won_arr))
    fill(slide, 62, fmt_eur(lost_arr))
    fill(slide, 63, f"{win_rate_ct:.0f}%")
    fill(slide, 64, fmt_eur(commit_arr))
    fill(
        slide,
        22,
        f"{won_ct} deals closed-won this quarter. Low won ARR signals early-quarter timing or delayed closes.",
    )
    fill(
        slide,
        55,
        f"{lost_ct} deals closed-lost. Lost ARR significantly exceeds won ARR - review loss reasons.",
    )
    fill(
        slide,
        57,
        f"Win rate by deal count ({won_ct}W / {lost_ct}L). ARR win rate: {win_rate_arr:.1f}%.",
    )
    fill(
        slide,
        59,
        f"Commit forecast for remaining open deals this quarter. Best Case: {fmt_eur(best_case_arr)}. Pipeline: {fmt_eur(pipeline_arr)}.",
    )

    # Forecast breakdown table
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 144, "Forecast Category Breakdown - Open Pipeline This Quarter")
    fill(slide, 145, f"{territory}, Q2 2026")
    fill(
        slide,
        42,
        f"Total open: {fmt_eur(commit_arr + best_case_arr + pipeline_arr + omitted_arr)}. Commit ({fmt_eur(commit_arr)}) is the floor; Best Case + Pipeline is the upside.",
    )
    fc_table = []
    for r in fc_data:
        fc_table.append(
            [
                r.get("ForecastCategoryName", "?"),
                str(r.get("ct", 0)),
                fmt_eur(r.get("arr") or 0),
            ]
        )
    fc_table.append(["Won (closed)", str(won_ct), fmt_eur(won_arr)])
    fc_table.append(["Lost (closed)", str(lost_ct), fmt_eur(lost_arr)])
    add_table(
        slide, ["Category", "Deals", "ARR"], fc_table, left=0.9, top=2.2, width=5.5
    )
    fc_cats = [r.get("ForecastCategoryName", "?") for r in fc_data] + ["Won", "Lost"]
    fc_vals = [(r.get("arr") or 0) for r in fc_data] + [won_arr, lost_arr]
    add_forecast_chart(
        slide, fc_cats, fc_vals, left=6.8, top=2.2, width=5.8, height=4.0
    )

    PI_VIEWS = {
        "Jesper Tyrer": (
            None,
            "00BTb00000Ic7kTMAR",
        ),  # APAC ARR CFQ Forecast (pre-existing)
        "Sarah Pittroff": ("/tmp/pi_ce.png", "00BTb00000Kr3YvMAJ"),
        "Francois Thaury": ("/tmp/pi_swe.png", "00BTb00000Kr3sHMAR"),
        "Dan Peppett": ("/tmp/pi_uki.png", "00BTb00000Kr3yjMAB"),
        "Christian Ebbesen": ("/tmp/pi_nl_nordics.png", "00BTb00000Kr4DFMAZ"),
        "Mourad Essofi": (None, None),  # no MEA PI view yet
        "Megan Miceli": ("/tmp/pi_canada.png", "00BTb00000Kr4ErMAJ"),
        "Patrick Gaughan": ("/tmp/pi_nam.png", "00BTb00000Kr4JhMAJ"),
        "Adam Steinhaus": ("/tmp/pi_pi.png", "00BTb00000Kr4OXMAZ"),
    }
    pi_img_path, pi_view_id = PI_VIEWS.get(director["name"], (None, None))
    slide = prs.slides.add_slide(get_layout(prs, "Title and Content"))
    fill(slide, 144, f"Pipeline Inspection - {territory}")
    if pi_img_path and Path(pi_img_path).exists() and pi_view_id:
        fill(
            slide,
            145,
            f"Source: Salesforce Pipeline Inspection - BoB ARR CFQ Forecast {territory}",
        )
        slide.shapes.add_picture(
            pi_img_path, Inches(0.3), Inches(1.6), Inches(12.7), Inches(5.2)
        )
        from pptx.oxml.ns import qn

        PI_URL = f"https://simcorp.lightning.force.com/lightning/o/Opportunity/pipelineInspection?filterName={pi_view_id}"
        btn = slide.shapes.add_textbox(
            Inches(3.5), Inches(6.9), Inches(6.3), Inches(0.4)
        )
        run = btn.text_frame.paragraphs[0].add_run()
        run.text = "OPEN PIPELINE INSPECTION (LIVE)"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = TEAL_DEEP
        run.font.underline = True
        hlinkClick = run._r.makeelement(qn("a:hlinkClick"), {})
        run._r.get_or_add_rPr().append(hlinkClick)
        rel = slide.part.relate_to(
            PI_URL,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
        hlinkClick.set(qn("r:id"), rel)
    else:
        fill(slide, 145, "Regional PI list view in progress")
        fill(
            slide,
            42,
            f"Pipeline Inspection for {territory} will be available once the regional list view is configured.",
        )

    # --- END (branded, no text) ---
    slide = prs.slides.add_slide(get_layout(prs, "End slide with disclaimer 1"))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Saved: {output_path} ({len(prs.slides)} slides)")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("--director", help="Single director name (e.g. 'Dan Peppett')")
    p.add_argument(
        "--all", action="store_true", help="Generate decks for all 9 MD-1 directors"
    )
    p.add_argument("--snapshot-date", default=date.today().isoformat())
    p.add_argument("--output-root", default=str(DEFAULT_OUTPUT_ROOT))
    return p.parse_args()


def main() -> int:
    args = parse_args()
    if not args.director and not args.all:
        print("Must specify --director NAME or --all", file=sys.stderr)
        return 2

    token, base_url = auth()
    config = json.loads(CONFIG_PATH.read_text())
    targets = config["presets"] if args.all else [load_preset(args.director)]
    output_root = Path(args.output_root) / args.snapshot_date

    for director in targets:
        out = (
            output_root
            / f"Sales Director Monthly - {director['name']} ({director['territory']}).pptx"
        )
        print(f"\n=== {director['name']} ({director['territory']}) ===")
        build_deck_for_director(director, out, token, base_url)

    return 0


if __name__ == "__main__":
    sys.exit(main())
