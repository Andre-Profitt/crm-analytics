#!/usr/bin/env python3
"""Final polish pass for the hybrid SimCorp deck.

Does three things:

1. Remove small decorative corner bars (top-right at ~10.70, 0.42 and
   bottom-left at ~0.55, 6.92). Leaves data-chart bars alone.
2. Untruncate text ending in '...' or '…' by matching the prefix against
   string values extracted from report1_snapshot.json. Only replaces when
   there is exactly one unique match.
3. Append an Appendix slide listing the live Salesforce dashboards and
   backing reports as clickable hyperlinks.

Run:
    python3 scripts/simcorp_final_polish.py <input.pptx> <snapshot.json> <output.pptx>
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ----------------------------------------------------------------------------
# Geometry thresholds (inches) for the corner accent bars we target.
# ----------------------------------------------------------------------------
TR_BAR = {  # top-right accent
    "l_min": 10.5,
    "l_max": 10.9,
    "t_min": 0.3,
    "t_max": 0.55,
    "w_max": 1.2,
    "h_max": 0.12,
}
BL_BAR = {  # bottom-left accent
    "l_min": 0.4,
    "l_max": 0.7,
    "t_min": 6.80,
    "t_max": 7.05,
    "w_max": 1.2,
    "h_max": 0.12,
}

ELLIPSIS_CHARS = ("\u2026", "...")

LINKS = [
    (
        "Sales Directors Monthly Dashboard",
        "https://simcorp.my.salesforce.com/lightning/r/Dashboard/01ZTb00000FSP7hMAH/view",
    ),
    (
        "Sales Ops Quarterly KPI Dashboard",
        "https://simcorp.my.salesforce.com/lightning/r/Dashboard/01ZTb00000FSP9JMAX/view",
    ),
    (
        "Sales Ops Data Quality and Forecast Accuracy (CRMA)",
        "https://simcorp.my.salesforce.com/analytics/dashboard/0FKTb0000000K5BOAU",
    ),
    (
        "Pipeline Coverage by Stage",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008TZc5MAG/view",
    ),
    (
        "Forecast and Closed Won",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008TZaTMAW/view",
    ),
    (
        "Land Stage 3 Missing Commercial Approval (by Region)",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekltMAA/view",
    ),
    (
        "Commercial Approval Candidates (by Stage)",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekp7MAA/view",
    ),
    (
        "Renewals by Fiscal Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008eksLMAQ/view",
    ),
    (
        "Renewal Pipeline This Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ektxMAA/view",
    ),
    (
        "Renewal ACV by Quarter",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008ekxBMAQ/view",
    ),
    (
        "Business At Risk",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008Ta9xMAC/view",
    ),
    (
        "Close Date Slipped (by Stage)",
        "https://simcorp.my.salesforce.com/lightning/r/Report/00OTb000008eknVMAQ/view",
    ),
]


def collect_strings(value, bucket: set):
    """Walk JSON recursively, collecting distinct string values >= 4 chars."""
    if isinstance(value, str):
        s = value.strip()
        if len(s) >= 4:
            bucket.add(s)
    elif isinstance(value, dict):
        for v in value.values():
            collect_strings(v, bucket)
    elif isinstance(value, list):
        for v in value:
            collect_strings(v, bucket)


def load_snapshot_strings(snapshot_path: Path) -> set[str]:
    with open(snapshot_path) as f:
        data = json.load(f)
    bucket: set[str] = set()
    collect_strings(data, bucket)
    return bucket


def is_corner_bar(shape) -> bool:
    try:
        l = shape.left.inches if shape.left else 0
        t = shape.top.inches if shape.top else 0
        w = shape.width.inches if shape.width else 0
        h = shape.height.inches if shape.height else 0
    except Exception:
        return False
    # Must be a filled shape
    try:
        if shape.fill.type != 1:
            return False
        rgb = str(shape.fill.fore_color.rgb)
        if rgb.upper() == "FFFFFF":
            return False
    except Exception:
        return False

    def fits(spec):
        return (
            spec["l_min"] <= l <= spec["l_max"]
            and spec["t_min"] <= t <= spec["t_max"]
            and w <= spec["w_max"]
            and h <= spec["h_max"]
        )

    return fits(TR_BAR) or fits(BL_BAR)


def remove_corner_bars(slide) -> int:
    to_remove = [sh for sh in slide.shapes if is_corner_bar(sh)]
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)
    return len(to_remove)


def untruncate_run_text(text: str, candidates: set[str]) -> str | None:
    """Return the full value for a truncated text run, or None if no unique match."""
    stripped = text
    trimmed = False
    for e in ELLIPSIS_CHARS:
        if stripped.endswith(e):
            stripped = stripped[: -len(e)]
            trimmed = True
            break
    if not trimmed:
        return None
    prefix = stripped.rstrip(" -|·:,")
    if len(prefix) < 4:
        return None
    matches = [c for c in candidates if c.startswith(prefix)]
    # Drop any match that is effectively the same as the prefix
    matches = [m for m in matches if len(m) > len(prefix)]
    if len(matches) == 1:
        full = matches[0]
        # Preserve any leading-trailing context from original run
        return text[: -len(text) + len(prefix)] + full if False else full
    return None


def untruncate_slide(slide, candidates: set[str]) -> int:
    replaced = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                full = untruncate_run_text(run.text, candidates)
                if full:
                    run.text = full
                    replaced += 1
    return replaced


def add_appendix_slide(prs):
    # Use Blank layout (strip any inherited shapes already handled)
    blank = None
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == "Blank":
                blank = layout
                break
        if blank is not None:
            break
    if blank is None:
        raise RuntimeError("No 'Blank' layout found in template")
    slide = prs.slides.add_slide(blank)
    # Drop any lingering layout placeholders
    for sh in list(slide.shapes):
        if sh.is_placeholder:
            sh._element.getparent().remove(sh._element)

    # Title
    title_tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.5), Inches(12), Inches(0.8)
    )
    title_tf = title_tb.text_frame
    title_tf.text = "Appendix: Live Salesforce links"
    p = title_tf.paragraphs[0]
    p.runs[0].font.name = "Microsoft Sans Serif"
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x01, 0x19, 0x46)

    # Subtitle
    sub_tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(0.4))
    sub_tf = sub_tb.text_frame
    sub_tf.text = "Click any link to open the live Salesforce dashboard or report."
    sub_tf.paragraphs[0].runs[0].font.name = "Microsoft Sans Serif"
    sub_tf.paragraphs[0].runs[0].font.size = Pt(12)
    sub_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x0E, 0x37, 0x88)

    # Links block
    body_tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(12), Inches(5))
    body_tf = body_tb.text_frame
    body_tf.word_wrap = True
    for i, (label, url) in enumerate(LINKS):
        para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        run = para.add_run()
        run.text = label
        run.font.name = "Microsoft Sans Serif"
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0x0E, 0x37, 0x88)
        run.font.underline = True
        run.hyperlink.address = url
        # Small gap between entries
        para.space_after = Pt(6)


def main():
    if len(sys.argv) < 4:
        print(
            "Usage: simcorp_final_polish.py <input.pptx> <snapshot.json> <output.pptx>"
        )
        sys.exit(1)

    src = Path(sys.argv[1])
    snapshot = Path(sys.argv[2])
    dst = Path(sys.argv[3])

    print(f"Source:   {src}")
    print(f"Snapshot: {snapshot}")
    print(f"Output:   {dst}")

    candidates = load_snapshot_strings(snapshot)
    print(f"Snapshot string candidates loaded: {len(candidates)}")

    prs = Presentation(str(src))

    total_bars = 0
    total_untruncs = 0
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            continue  # leave cover alone
        bars = remove_corner_bars(slide)
        unt = untruncate_slide(slide, candidates)
        total_bars += bars
        total_untruncs += unt
        print(f"  Slide {i}: removed {bars} corner bars, untruncated {unt} text runs")

    add_appendix_slide(prs)
    print(f"  Added appendix slide with {len(LINKS)} links")

    prs.save(str(dst))
    print()
    print(f"Total bars removed:   {total_bars}")
    print(f"Total runs untruncated: {total_untruncs}")
    print(f"Saved: {dst}")


if __name__ == "__main__":
    main()
