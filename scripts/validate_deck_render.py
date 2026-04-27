#!/usr/bin/env python3
"""Track F / F5 — render / overflow gates.

Geometry-level checks against a produced .pptx. Detects layout-time
issues that the PPTX-contract checker (E4) doesn't:

  - title placeholder (ph 144) sits inside the expected title region
    (top edge within the slide's header band)
  - tables fit inside the slide bounds (no off-slide overflow)
  - non-static slides have a footer/source-note text frame in the
    bottom band of the slide
  - the legal-notice slide carries the SimCorp disclaimer text

Read-only. Never modifies the .pptx, never modifies the builder.

What this validator deliberately does NOT do:

  - Rendered text-overflow detection. python-pptx exposes shape
    bounds, not rendered text dimensions. A reliable text-overflow
    check requires rasterising the slide (libreoffice / pdf+ocr).
    F6 will use that pipeline for visual regression. F5 is geometry-
    only.

Slide dimensions: 13.33" wide × 7.5" tall (default 16:9). Regions
defined in EMU (914400 per inch) for direct comparison with python-
pptx Length values.

Usage:
    python3 scripts/validate_deck_render.py \\
        --pptx ~/Downloads/jesper-tyrer-LAND.pptx \\
        --report-out output/track_f/render_report.json
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

REPO_ROOT = Path(__file__).resolve().parent.parent
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

from scripts.monthly_platform import deck_contract  # noqa: E402


REPORT_SCHEMA_VERSION = "monthly_platform.deck_render_report.v1"

# Slide regions (EMU = 914400 per inch). 13.33" × 7.5" 16:9 default.
SLIDE_WIDTH_EMU = Inches(13.333)
SLIDE_HEIGHT_EMU = Inches(7.5)
TITLE_REGION_BOTTOM_EMU = Inches(1.5)  # title top must be <= this
FOOTER_REGION_TOP_EMU = Inches(6.5)  # footer top must be >= this
LEGAL_DISCLAIMER_TOKENS = ["SimCorp", "Confidential", "disclaimer"]


@dataclass
class RenderFinding:
    severity: str
    code: str
    path: str
    message: str

    def as_dict(self) -> dict[str, Any]:
        return {
            "severity": self.severity,
            "code": self.code,
            "path": self.path,
            "message": self.message,
        }


@dataclass
class SlideRenderResult:
    slide_number: int
    slide_id: str | None
    title_status: str  # pass | fail | n/a (static)
    title_top_in: float | None
    table_status: str  # pass | fail | n/a (no tables)
    footer_status: str  # pass | fail | n/a (static)
    overflowing_tables: list[dict[str, Any]] = field(default_factory=list)
    overall: str = "pass"

    def as_dict(self) -> dict[str, Any]:
        return {
            "slide_number": self.slide_number,
            "slide_id": self.slide_id,
            "title_status": self.title_status,
            "title_top_in": self.title_top_in,
            "table_status": self.table_status,
            "overflowing_tables": self.overflowing_tables,
            "footer_status": self.footer_status,
            "overall": self.overall,
        }


def _shape_title_ph(slide):
    """Find the slide's title-bearing shape.

    Order of attempts:
      1. placeholder with idx == 144 (build_deck_from_excel.py's title slot
         on most layouts)
      2. placeholder whose type is TITLE (13) or CTR_TITLE (14)
      3. on layouts like 'Title 1' where neither (1) nor (2) matches,
         fall back to the first non-empty text frame whose top is already
         inside the title region. Covers the cover slide (ph 20 BODY).
      4. None — caller treats this as missing-title.
    """
    idx_144 = None
    title_type = None
    in_region_first = None
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            ph = shape.placeholder_format
        except (ValueError, AttributeError):
            ph = None
        if ph is not None:
            try:
                if ph.idx == 144 and idx_144 is None:
                    idx_144 = shape
            except (ValueError, AttributeError):
                pass
            try:
                if ph.type is not None and ph.type in (13, 14) and title_type is None:
                    title_type = shape
            except (ValueError, AttributeError):
                pass
        # Fallback: first non-empty text frame whose top sits in the
        # title region. Captures ph-20 BODY title on the 'Title 1' layout
        # (cover slide) and similar layouts where the title placeholder
        # type isn't TITLE.
        top = shape.top if shape.top is not None else 0
        if in_region_first is None and top <= TITLE_REGION_BOTTOM_EMU:
            in_region_first = shape
    return idx_144 or title_type or in_region_first


def _slide_text_frames_with_top(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            top = shape.top if shape.top is not None else 0
            out.append((shape, top))
    return out


def validate_render(
    pptx_path: Path,
    *,
    contract: deck_contract.DeckContract | None = None,
) -> dict[str, Any]:
    if contract is None:
        contract = deck_contract.load()
    assert contract is not None

    profile = contract.director_monthly
    contract_slides_by_number: dict[int, dict[str, Any]] = {
        int(s["slide_number"]): s for s in profile.get("slides", [])
    }

    findings: list[RenderFinding] = []
    slide_results: list[SlideRenderResult] = []

    prs = Presentation(str(pptx_path))

    for snum, slide in enumerate(prs.slides, start=1):
        decl = contract_slides_by_number.get(snum)
        sid = decl["id"] if decl else None
        is_static = bool(decl and (decl.get("static") or sid in ("legal_notice",)))

        # Cover slide's 'Title 1' layout is center-cover by design — title
        # sits ~mid-slide, not in the top header band. Exempt from the
        # title-region check while still requiring the slide to have title-
        # bearing text content (the PPTX-contract checker covers that).
        is_cover = sid == "cover"

        slide_overall = "pass"

        # ---- title region ----
        title_status = "n/a"
        title_top_in: float | None = None
        if not is_static and not is_cover:
            ph = _shape_title_ph(slide)
            if ph is None:
                title_status = "fail"
                slide_overall = "fail"
                findings.append(
                    RenderFinding(
                        severity="blocker",
                        code="title_placeholder_missing",
                        path=f"slides[{snum}].title",
                        message=(
                            f"slide {snum} ({sid or '?'}) has no placeholder 144 "
                            f"with a non-empty title text"
                        ),
                    )
                )
            else:
                top = ph.top if ph.top is not None else 0
                title_top_in = round(top / 914400, 3)
                if top > TITLE_REGION_BOTTOM_EMU:
                    title_status = "fail"
                    slide_overall = "fail"
                    findings.append(
                        RenderFinding(
                            severity="blocker",
                            code="title_drift_outside_region",
                            path=f"slides[{snum}].title",
                            message=(
                                f"slide {snum} ({sid or '?'}) title top "
                                f'{title_top_in}" exceeds title-region bottom '
                                f'{TITLE_REGION_BOTTOM_EMU / 914400}"'
                            ),
                        )
                    )
                else:
                    title_status = "pass"

        # ---- table off-slide overflow ----
        overflowing: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            left = shape.left or 0
            top = shape.top or 0
            width = shape.width or 0
            height = shape.height or 0
            right = left + width
            bottom = top + height
            if (
                right > SLIDE_WIDTH_EMU
                or bottom > SLIDE_HEIGHT_EMU
                or left < 0
                or top < 0
            ):
                overflowing.append(
                    {
                        "left_in": round(left / 914400, 3),
                        "top_in": round(top / 914400, 3),
                        "right_in": round(right / 914400, 3),
                        "bottom_in": round(bottom / 914400, 3),
                    }
                )
        if overflowing:
            findings.append(
                RenderFinding(
                    severity="blocker",
                    code="table_off_slide",
                    path=f"slides[{snum}].tables",
                    message=(
                        f"slide {snum} ({sid or '?'}) has {len(overflowing)} "
                        f'table(s) extending past slide bounds (13.333" × 7.5"): '
                        f"{overflowing}"
                    ),
                )
            )
            table_status = "fail"
            slide_overall = "fail"
        elif any(shape.has_table for shape in slide.shapes):
            table_status = "pass"
        else:
            table_status = "n/a"

        # ---- footer / source note presence ----
        footer_status = "n/a"
        if not is_static:
            footer_present = any(
                top >= FOOTER_REGION_TOP_EMU
                for _, top in _slide_text_frames_with_top(slide)
            )
            footer_status = "pass" if footer_present else "fail"
            if not footer_present:
                # Track F render-gate convergence (post-footer-emission):
                # this is now a blocker. Builder emits a source-note footer
                # on every non-static slide via _apply_source_note_footer.
                slide_overall = "fail"
                findings.append(
                    RenderFinding(
                        severity="blocker",
                        code="footer_missing",
                        path=f"slides[{snum}].footer",
                        message=(
                            f"slide {snum} ({sid or '?'}) has no text frame "
                            f"in the footer band (top >= "
                            f'{FOOTER_REGION_TOP_EMU / 914400}")'
                        ),
                    )
                )

        # ---- legal notice disclaimer text ----
        # Check both slide shapes and the slide's layout shapes — the
        # SimCorp 'End slide with disclaimer 1' template layout carries the
        # disclaimer on layout-level placeholders that python-pptx does NOT
        # surface via slide.shapes when the slide just inherits the layout.
        # Falling back to layout text covers the production path where the
        # disclaimer is rendered from the template, not authored on each slide.
        if sid == "legal_notice":
            slide_text = " ".join(
                shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
            )
            layout_text = " ".join(
                shape.text_frame.text
                for shape in slide.slide_layout.shapes
                if shape.has_text_frame
            )
            full_text = f"{slide_text} {layout_text}"
            if not any(token in full_text for token in LEGAL_DISCLAIMER_TOKENS):
                slide_overall = "fail"
                findings.append(
                    RenderFinding(
                        severity="blocker",
                        code="legal_disclaimer_missing",
                        path=f"slides[{snum}].legal_notice",
                        message=(
                            f"slide {snum} (legal_notice) has no text in slide or "
                            f"layout containing any of {LEGAL_DISCLAIMER_TOKENS}"
                        ),
                    )
                )

        slide_results.append(
            SlideRenderResult(
                slide_number=snum,
                slide_id=sid,
                title_status=title_status,
                title_top_in=title_top_in,
                table_status=table_status,
                overflowing_tables=overflowing,
                footer_status=footer_status,
                overall=slide_overall,
            )
        )

    blockers = [f for f in findings if f.severity == "blocker"]
    warnings = [f for f in findings if f.severity == "warning"]

    return {
        "schema_version": REPORT_SCHEMA_VERSION,
        "validated_at": datetime.now(timezone.utc).isoformat(),
        "pptx_path": str(pptx_path),
        "deck_contract_path": str(contract.path),
        "status": "pass" if not blockers else "fail",
        "blocker_count": len(blockers),
        "warning_count": len(warnings),
        "slide_count": len(prs.slides),
        "slides": [s.as_dict() for s in slide_results],
        "findings": [f.as_dict() for f in findings],
    }


def render_markdown(report: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append("# Deck render report\n")
    lines.append(f"- pptx: `{report['pptx_path']}`")
    lines.append(f"- deck contract: `{report['deck_contract_path']}`")
    lines.append(f"- validated_at: {report['validated_at']}")
    lines.append(f"- **status: {report['status']}**")
    lines.append(
        f"- blockers: {report['blocker_count']} | warnings: {report['warning_count']}"
    )
    lines.append("")
    lines.append("| # | Slide | Title | Title top | Tables | Footer | Overall |")
    lines.append("| ---: | --- | --- | ---: | --- | --- | --- |")
    for s in report["slides"]:
        lines.append(
            f"| {s['slide_number']} | `{s['slide_id'] or '?'}` | "
            f"{s['title_status']} | {s['title_top_in'] or '—'} | "
            f"{s['table_status']} | {s['footer_status']} | {s['overall']} |"
        )
    lines.append("")
    if report["findings"]:
        lines.append("## Findings\n")
        for f in report["findings"]:
            lines.append(f"- **{f['severity']}** `{f['code']}` — {f['message']}")
    return "\n".join(lines) + "\n"


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Validate produced PPTX geometry (Track F F5 render gates)."
    )
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--contract", default=None)
    parser.add_argument("--report-out", default=None)
    parser.add_argument("--md-out", default=None)
    parser.add_argument("--show-findings", action="store_true")
    args = parser.parse_args(argv)

    contract = deck_contract.load(args.contract)
    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"ERROR: pptx not found: {pptx_path}", file=sys.stderr)
        return 2

    report = validate_render(pptx_path, contract=contract)

    if args.report_out:
        out = Path(args.report_out)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(report, indent=2) + "\n", encoding="utf-8")
        print(f"report: {out}")
    if args.md_out:
        out = Path(args.md_out)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(render_markdown(report), encoding="utf-8")
        print(f"md: {out}")

    if args.show_findings or report["status"] == "fail":
        for f in report["findings"]:
            print(f"[{f['severity']}] {f['code']} {f.get('path', '')}: {f['message']}")

    print(
        f"deck_render: {report['status']} "
        f"(blockers={report['blocker_count']} warnings={report['warning_count']} "
        f"slides={report['slide_count']})"
    )
    return 0 if report["status"] == "pass" else 1


if __name__ == "__main__":
    raise SystemExit(main())
