#!/usr/bin/env python3
"""Build a gold-example populated regional deck from a validated regional snapshot."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

try:
    from build_sales_region_monthly_shell import (
        AMBER,
        AQUA,
        GREEN_OK,
        GREY_TEXT,
        LIGHT_BLUE_PANEL,
        LIGHT_PANEL_2,
        MAGENTA,
        MAGENTA_LIGHT,
        NAVY,
        PRIMARY_BLUE,
        build_shell_deck,  # noqa: F401
        blank_slide,
        clear_slides,
        content_panel,
        fill_placeholder,
        get_layout,
        header,
        hero_stat_card,
        kpi_card,
        mini_stat,
        multi_paragraph,
        rect,
        rounded_card,
        source_footer,
        takeaway,
        txt,
    )
except ModuleNotFoundError:  # pragma: no cover
    from scripts.build_sales_region_monthly_shell import (
        AMBER,
        AQUA,
        GREEN_OK,
        GREY_TEXT,
        LIGHT_BLUE_PANEL,
        LIGHT_PANEL_2,
        MAGENTA,
        MAGENTA_LIGHT,
        NAVY,
        PRIMARY_BLUE,
        build_shell_deck,  # noqa: F401
        blank_slide,
        clear_slides,
        content_panel,
        fill_placeholder,
        get_layout,
        header,
        hero_stat_card,
        kpi_card,
        mini_stat,
        multi_paragraph,
        rect,
        rounded_card,
        source_footer,
        takeaway,
        txt,
    )

from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MASTER_TEMPLATE_PATH = (
    Path.home()
    / "archive"
    / "simcorp-deck-agent-backup"
    / "reference-decks"
    / "SimCorp_PPT_Template.pptx"
)
DEFAULT_OUTPUT_ROOT = REPO_ROOT / "output" / "sales_region_gold_decks"
DEFAULT_SHELL_CONTRACT_PATH = REPO_ROOT / "config" / "sales_region_monthly_shell.json"


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def fmt_eur(amount: float | int | None) -> str:
    if amount is None:
        return "n/a"
    amount = float(amount)
    if abs(amount) >= 1_000_000_000:
        return f"€{amount / 1_000_000_000:.2f}B"
    if abs(amount) >= 1_000_000:
        return f"€{amount / 1_000_000:.1f}M"
    if abs(amount) >= 1_000:
        return f"€{amount / 1_000:.0f}K"
    return f"€{amount:.0f}"


def fmt_pct(value: float) -> str:
    return f"{value:.1f}%"


def pct(numerator: float, denominator: float) -> float:
    if not denominator:
        return 0.0
    return numerator / denominator * 100.0


def parse_percent(text: str) -> float:
    return float(text.replace("%", "").strip())


def parse_eur_compact(text: str) -> float:
    stripped = text.strip().replace("€", "").replace(",", "")
    if stripped.endswith("M"):
        return float(stripped[:-1]) * 1_000_000
    if stripped.endswith("K"):
        return float(stripped[:-1]) * 1_000
    return float(stripped)


def get_scorecard_metric(snapshot: dict[str, Any], section_key: str, metric_key: str):
    return snapshot["scorecard"]["sections"][section_key]["metrics"][metric_key]


def load_component_book_metrics(snapshot: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for component in snapshot["component_books"]:
        ds = load_json(Path(component["snapshot_path"]))
        pipeline_metrics = ds["scorecard"]["sections"]["pipeline-health"]["metrics"]
        compliance_metrics = ds["scorecard"]["sections"]["process-compliance"]["metrics"]
        rows.append(
            {
                "territory": component["territory"],
                "director_name": component["director_name"],
                "q2_active": parse_eur_compact(
                    pipeline_metrics["Pipeline ARR — Q2 2026 Close Dates Only (excl. Omitted)"]
                ),
                "fy26_active": parse_eur_compact(
                    pipeline_metrics["Pipeline ARR — FY26 Close Dates Only (excl. Omitted)"]
                ),
                "approval_rate": parse_percent(compliance_metrics["Approval Rate (stage 3+)"]),
                "missing_approval": int(compliance_metrics["Missing Approval (Land, stage 3+)"]),
            }
        )
    rows.sort(key=lambda row: row["q2_active"], reverse=True)
    return rows


def q2_renewals(snapshot: dict[str, Any]) -> list[dict[str, Any]]:
    explicit = (snapshot.get("renewals") or {}).get("q2_open_renewals") or []
    if explicit:
        return list(explicit)
    rows = []
    for deal in snapshot["renewals"]["open_renewals"]:
        close_date = deal.get("Close Date")
        if close_date and "2026-04-01" <= close_date <= "2026-06-30":
            rows.append(deal)
    rows.sort(key=lambda row: row.get("Renewal ACV (€ converted)") or 0, reverse=True)
    return rows


def add_cover(prs: Presentation, *, region_name: str, snapshot_date: str) -> None:
    slide = prs.slides.add_slide(get_layout(prs, "Title 1"))
    fill_placeholder(slide, 20, "")
    fill_placeholder(slide, 22, "")
    fill_placeholder(slide, 24, "")
    txt(slide, 1.0, 1.4, 2.8, 0.25, "SIMCORP", size=11, bold=True, color=AQUA)
    txt(slide, 1.0, 1.95, 5.2, 0.6, "Sales Region Monthly", size=26, bold=True, color=AQUA)
    txt(slide, 1.0, 2.55, 6.8, 0.7, region_name, size=34, bold=True, color=NAVY)
    txt(
        slide,
        1.0,
        3.25,
        7.0,
        0.5,
        f"Validated operating deck | Snapshot {snapshot_date}",
        size=16,
        color=GREY_TEXT,
    )
    txt(
        slide,
        1.0,
        5.95,
        4.4,
        0.26,
        "Quarterly pipeline, approvals, renewals, slipped deals, and operating controls",
        size=12,
        color=GREY_TEXT,
    )


def add_agenda(prs: Presentation, shell: dict[str, Any], *, region_name: str, snapshot_date: str) -> None:
    slide = blank_slide(prs)
    header(
        slide,
        eyebrow=f"{snapshot_date}   ·   {region_name.upper()}   ·   LEADERSHIP CADENCE",
        title="Leadership agenda",
        narrative="Fixed sequence for the regional operating review.",
        accent=PRIMARY_BLUE,
    )
    titles = [item["title"] for item in shell["slides"]]
    for index, title in enumerate(titles):
        col = index % 3
        row = index // 3
        left = 0.6 + col * 4.08
        top = 1.95 + row * 0.95
        rounded_card(slide, left, top, 3.85, 0.78, LIGHT_BLUE_PANEL if index < 6 else LIGHT_PANEL_2)
        txt(slide, left + 0.16, top + 0.14, 0.35, 0.2, f"{index + 1:02d}", size=10, bold=True, color=PRIMARY_BLUE)
        txt(slide, left + 0.56, top + 0.13, 3.05, 0.36, title, size=11, bold=True, color=NAVY)
    takeaway(
        slide,
        top=5.95,
        bullets=[
            "Lead with what the region needs leadership to decide this month.",
            "Use the same order every month so the operating rhythm becomes predictable.",
            "All numbers in this deck come from the validated regional snapshot, not ad hoc PowerPoint edits.",
        ],
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Validated regional snapshot and fact pack.")


def add_exec_summary(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    all_open = parse_eur_compact(get_scorecard_metric(snapshot, "pipeline-health", "Pipeline ARR — All Open (any close date)"))
    fy26 = parse_eur_compact(get_scorecard_metric(snapshot, "pipeline-health", "Pipeline ARR — FY26 Close Dates Only (excl. Omitted)"))
    q2_active = parse_eur_compact(get_scorecard_metric(snapshot, "pipeline-health", "Pipeline ARR — Q2 2026 Close Dates Only (excl. Omitted)"))
    weighted = parse_eur_compact(get_scorecard_metric(snapshot, "pipeline-health", "Weighted Pipeline (probability-adj)"))
    new_pipeline = parse_eur_compact(get_scorecard_metric(snapshot, "pipeline-health", "New Pipeline This Quarter (excl. Omitted)"))
    stale = parse_eur_compact(get_scorecard_metric(snapshot, "risk", "Stale 30d+ (ARR)"))
    aging = parse_eur_compact(get_scorecard_metric(snapshot, "risk", "Aging 365+ (ARR)"))
    approval_rate = parse_percent(get_scorecard_metric(snapshot, "process-compliance", "Approval Rate (stage 3+)"))
    missing = int(get_scorecard_metric(snapshot, "process-compliance", "Missing Approval (Land, stage 3+)"))
    renewal_open = snapshot["renewals"]["summary_metrics"]["open_acv"]

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   EXECUTIVE SUMMARY",
        title=f"{snapshot['region_name']} enters Q2 with {fmt_eur(q2_active)} active ARR and zero new pipeline creation",
        narrative="Headline position, control pressure, and month-specific actions.",
        accent=AQUA,
    )
    cards = [
        (fmt_eur(all_open), "All-open pipeline", f"{snapshot['scorecard']['sections']['pipeline-health']['metrics']['Deal Count']} deals across the full open book."),
        (fmt_eur(fy26), "FY26 active pipeline", f"Weighted open sits at {fmt_eur(weighted)} across FY26 close dates."),
        (fmt_eur(q2_active), "Q2 active pipeline", f"New pipeline this quarter is {fmt_eur(new_pipeline)}."),
        (fmt_eur(renewal_open), "Open renewal book", f"Renewals stay ACV. Total open renewal ACV is {fmt_eur(renewal_open)}."),
    ]
    for index, (big, label, context) in enumerate(cards):
        kpi_card(
            slide,
            left=0.6 + index * 3.08,
            top=1.95,
            width=2.9,
            height=1.95,
            big=big,
            label=label,
            context=context,
            panel=MAGENTA_LIGHT if index == 3 else LIGHT_BLUE_PANEL,
            accent_dot=GREEN_OK if index < 3 else MAGENTA,
        )
    takeaway(
        slide,
        top=4.35,
        bullets=[
            f"Q2 active ARR is only {fmt_pct(pct(q2_active, all_open))} of the total open book, so the near-term position is much thinner than the headline pipeline suggests.",
            f"Commercial approval control is weak at {fmt_pct(approval_rate)} with {missing} missing-approval candidates still open.",
            f"Hygiene pressure remains high with {fmt_eur(stale)} stale ARR and {fmt_eur(aging)} aged 365+ ARR.",
        ],
    )
    source_footer(slide, "Regional snapshot scorecard metrics, EUR converted.")


def add_q1_review(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    actuals = snapshot["q1_review"]["actuals"]
    won_arr = actuals["won_arr"]
    lost_arr = actuals["lost_arr"]
    slipped_arr = actuals["slipped_arr"]
    count_win_rate = pct(actuals["won_count"], actuals["won_count"] + actuals["lost_count"])
    arr_win_rate = pct(won_arr, won_arr + lost_arr)
    baseline = snapshot["q1_review"]["promise_baseline"]

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   Q1 REVIEW",
        title=f"Q1 delivered {fmt_eur(won_arr)} ARR while {fmt_eur(slipped_arr)} slipped and {fmt_eur(lost_arr)} was lost",
        narrative="Delivered outcomes first; commitment baseline remains qualified.",
        accent=AMBER,
    )
    cards = [
        (f"{actuals['won_count']} / {fmt_eur(won_arr)}", "Won in Q1", "Delivered closes in quarter."),
        (f"{actuals['lost_count']} / {fmt_eur(lost_arr)}", "Lost in Q1", "Lost volume and value in quarter."),
        (f"{actuals['slipped_count']} / {fmt_eur(slipped_arr)}", "Slipped out of Q1", "Opportunities moved out of quarter."),
        ("Qualified", "Promise baseline", "Use with caution: pipeline-inspection population, not a clean commitment baseline."),
    ]
    for index, (big, label, context) in enumerate(cards):
        kpi_card(
            slide,
            left=0.6 + index * 3.08,
            top=1.95,
            width=2.9,
            height=2.0,
            big=big,
            label=label,
            context=context,
            panel=LIGHT_PANEL_2 if index == 3 else LIGHT_BLUE_PANEL,
            accent_dot=AMBER if index > 0 else GREEN_OK,
        )
    content_panel(
        slide,
        left=0.6,
        top=4.4,
        width=5.9,
        height=1.75,
        title="Forecast-safe baseline note",
        lines=[
            f"Closed population in the Q1 baseline was {fmt_eur(baseline[0]['ARR (€ converted)'])}.",
            f"Commit population was {fmt_eur(baseline[1]['ARR (€ converted)'])}; Pipeline was only {fmt_eur(baseline[2]['ARR (€ converted)'])}.",
            "Use these as context, not as an unqualified promise statement.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=6.75,
        top=4.4,
        width=5.95,
        height=1.75,
        title="Readout message",
        lines=[
            f"Count win rate was {fmt_pct(count_win_rate)}; ARR win rate was only {fmt_pct(arr_win_rate)}.",
            "The region needs tighter qualification and stronger control on what is allowed into quarter.",
        ],
        panel=LIGHT_PANEL_2,
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Q1 actuals from validated regional rollup; baseline remains qualified.")


def add_quarterly_pipeline(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    cov = snapshot["q2_outlook"]["coverage"]
    q2_active = cov["Active Pipeline ARR (Q2 close, excl. Omitted)"]
    commit = cov["Commit ARR (Q2 close)"]
    best_case = cov["Best Case ARR (Q2 close)"]
    pipeline = cov["Pipeline ARR (Q2 close)"]
    omitted = cov["Omitted ARR (Q2 close)"]
    best_case_pct = pct(best_case, q2_active)
    commit_pct = pct(commit, q2_active)

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   PIPELINE",
        title=f"Q2 outlook leans on Best Case: only {fmt_eur(commit)} of {fmt_eur(q2_active)} active ARR sits in Commit",
        narrative="Quarterly forecast mix, active book, and omitted separation.",
        accent=AQUA,
    )
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=5.6,
        height=3.1,
        overline="Q2 ACTIVE PIPELINE",
        big=fmt_eur(q2_active),
        context="Current-quarter active ARR, excluding Omitted. This is the number that should anchor the regional quarter discussion.",
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    mini_data = [
        ("Commit", fmt_eur(commit), f"{fmt_pct(commit_pct)} of active Q2 ARR"),
        ("Best Case", fmt_eur(best_case), f"{fmt_pct(best_case_pct)} of active Q2 ARR"),
        ("Pipeline", fmt_eur(pipeline), "Open but not yet in Best Case or Commit"),
        ("Omitted", fmt_eur(omitted), "Visible exposure, excluded from active headline"),
    ]
    for index, (label, big, context) in enumerate(mini_data):
        mini_stat(
            slide,
            left=6.45,
            top=1.95 + index * 0.8,
            width=6.25,
            height=0.68,
            label=label,
            big=big,
            context=context,
            panel=LIGHT_PANEL_2,
            status=AMBER if label == "Omitted" else GREEN_OK,
        )
    takeaway(
        slide,
        top=5.28,
        bullets=[
            f"Best Case represents {fmt_pct(best_case_pct)} of the active Q2 book, so the quarter depends more on deal conversion than on committed closes.",
            f"Commit is only {fmt_pct(commit_pct)} of Q2 active ARR, which is too light for a comfortable quarter.",
            f"Omitted still represents {fmt_eur(omitted)} of Q2 close-date exposure and should stay visible as a control signal.",
        ],
    )
    source_footer(slide, "Regional Q2 outlook, ARR converted to EUR.")


def add_book_breakdown(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    rows = load_component_book_metrics(snapshot)
    total_q2 = sum(row["q2_active"] for row in rows)
    largest = rows[0]
    weakest = min(rows, key=lambda row: row["approval_rate"])

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   BOOK BREAKDOWN",
        title=f"{largest['territory']} carries {fmt_pct(pct(largest['q2_active'], total_q2))} of regional Q2 active ARR",
        narrative="Book-level contribution, control quality, and where leadership pressure sits.",
        accent=PRIMARY_BLUE,
    )
    rounded_card(slide, 0.6, 1.95, 8.2, 4.85, LIGHT_BLUE_PANEL)
    rect(slide, 0.6, 1.95, 0.1, 4.85, PRIMARY_BLUE)
    txt(slide, 0.82, 2.12, 2.6, 0.2, "BOOK", size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, 4.05, 2.12, 1.2, 0.2, "Q2 ARR", size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, 5.35, 2.12, 1.2, 0.2, "FY26 ARR", size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, 6.65, 2.12, 1.15, 0.2, "APPROVAL", size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, 7.85, 2.12, 0.7, 0.2, "MISS", size=10, bold=True, color=PRIMARY_BLUE)
    for index, row in enumerate(rows):
        row_top = 2.45 + index * 0.78
        rounded_card(slide, 0.92, row_top, 7.55, 0.58, LIGHT_PANEL_2)
        txt(slide, 1.08, row_top + 0.12, 2.7, 0.2, row["territory"], size=11, bold=True, color=NAVY)
        txt(slide, 3.35, row_top + 0.12, 0.6, 0.2, row["director_name"].split()[0], size=10, color=GREY_TEXT)
        txt(slide, 4.05, row_top + 0.12, 1.0, 0.2, fmt_eur(row["q2_active"]), size=11, bold=True, color=NAVY)
        txt(slide, 5.35, row_top + 0.12, 1.0, 0.2, fmt_eur(row["fy26_active"]), size=11, bold=True, color=NAVY)
        txt(slide, 6.65, row_top + 0.12, 0.9, 0.2, fmt_pct(row["approval_rate"]), size=11, bold=True, color=NAVY)
        txt(slide, 7.95, row_top + 0.12, 0.3, 0.2, str(row["missing_approval"]), size=11, bold=True, color=NAVY)
    content_panel(
        slide,
        left=9.05,
        top=1.95,
        width=3.65,
        height=1.45,
        title="Largest book",
        lines=[
            f"{largest['territory']} contributes {fmt_eur(largest['q2_active'])} of Q2 active ARR.",
            "That book should carry disproportionate leadership attention this month.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AQUA,
    )
    content_panel(
        slide,
        left=9.05,
        top=3.62,
        width=3.65,
        height=1.45,
        title="Weakest control seam",
        lines=[
            f"{weakest['territory']} has the lowest approval rate at {fmt_pct(weakest['approval_rate'])}.",
            "Use this as the first place to tighten control discipline.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=9.05,
        top=5.3,
        width=3.65,
        height=1.5,
        title="Hierarchy rule",
        lines=[
            "Middle East & Africa belongs under EMEA in the forecast hierarchy.",
            "Do not reassign subregions in the deck layer.",
        ],
        panel=LIGHT_PANEL_2,
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Director-book snapshots rolled into forecast-region hierarchy.")


def add_pipeline_intel(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    top_opps = ((snapshot.get("pipeline_detail") or {}).get("q2_active_opportunities") or [])[:3]
    if not top_opps:
        top_opps = snapshot["pipeline_detail"]["top_opportunities"][:3]
    stale = parse_eur_compact(get_scorecard_metric(snapshot, "risk", "Stale 30d+ (ARR)"))
    aging = parse_eur_compact(get_scorecard_metric(snapshot, "risk", "Aging 365+ (ARR)"))
    backlog = snapshot["data_quality"]["total"]["Total Issues"]
    new_pipeline = parse_eur_compact(get_scorecard_metric(snapshot, "pipeline-health", "New Pipeline This Quarter (excl. Omitted)"))
    no_activity = snapshot["data_quality"]["total"]["No Activity"]

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   PIPELINE INTEL",
        title="Coverage is concentrated in a few large deals while hygiene pressure remains high",
        narrative="Material opportunities, control pressure, and competitive/loss signals.",
        accent=AQUA,
    )
    content_panel(
        slide,
        left=0.6,
        top=1.95,
        width=3.85,
        height=2.65,
        title="Coverage position",
        lines=[
            f"Q2 active ARR is {fmt_eur(snapshot['q2_outlook']['coverage']['Active Pipeline ARR (Q2 close, excl. Omitted)'])}.",
            f"Weighted pipeline is {fmt_eur(parse_eur_compact(get_scorecard_metric(snapshot, 'pipeline-health', 'Weighted Pipeline (probability-adj)')))}.",
            f"New pipeline this quarter is {fmt_eur(new_pipeline)}.",
        ],
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    content_panel(
        slide,
        left=4.73,
        top=1.95,
        width=3.95,
        height=2.65,
        title="Largest validated deals",
        lines=[
            f"{opp['Opportunity']} — {fmt_eur(opp['ARR (€ converted)'])}"
            for opp in top_opps
        ],
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    content_panel(
        slide,
        left=8.95,
        top=1.95,
        width=3.75,
        height=2.65,
        title="Risk pressure",
        lines=[
            f"Stale 30d+ ARR: {fmt_eur(stale)}",
            f"Aging 365+ ARR: {fmt_eur(aging)}",
            f"Data-quality backlog: {backlog} issue-points, including {no_activity} no-activity flags.",
        ],
        panel=LIGHT_PANEL_2,
        accent=AMBER,
    )
    takeaway(
        slide,
        top=5.05,
        bullets=[
            "The region is not short of headline pipeline; it is short of clean, near-term, well-controlled pipeline.",
            "Top-opportunity concentration means execution on a handful of deals will swing the quarter.",
            "Hygiene backlog is large enough that it should be treated as a commercial risk, not just an ops backlog.",
        ],
    )
    source_footer(slide, "Regional snapshot top opportunities and data-quality totals.")


def add_approval_overview(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    approval_rate = parse_percent(get_scorecard_metric(snapshot, "process-compliance", "Approval Rate (stage 3+)"))
    missing = int(get_scorecard_metric(snapshot, "process-compliance", "Missing Approval (Land, stage 3+)"))
    approved_summary = next(row for row in snapshot["commercial_approval"]["summary"] if row["Category"] == "Approved")
    top_candidate = snapshot["commercial_approval"]["missing_candidates"][0]

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   COMMERCIAL APPROVAL",
        title=f"Approval control is weak: {missing} stage 3+ candidates still need action",
        narrative="Regional stage 3+ governance with explicit method note.",
        accent=AMBER,
    )
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=5.4,
        height=3.1,
        overline="APPROVAL RATE",
        big=fmt_pct(approval_rate),
        context="Simple average of component-book approval-rate metrics from validated director scorecards.",
        panel=LIGHT_BLUE_PANEL,
        accent=AMBER,
    )
    mini_stat(
        slide,
        left=6.25,
        top=1.95,
        width=6.45,
        height=0.9,
        label="Approved deals",
        big=f"{approved_summary['Deal Count']} / {fmt_eur(approved_summary['ARR (€ converted)'])}",
        context="Approved deal count and ARR across the region.",
        panel=LIGHT_PANEL_2,
        status=GREEN_OK,
    )
    mini_stat(
        slide,
        left=6.25,
        top=2.95,
        width=6.45,
        height=0.9,
        label="Missing candidates",
        big=str(missing),
        context="Land stage 3+ candidates without clean approval coverage.",
        panel=LIGHT_PANEL_2,
        status=AMBER,
    )
    mini_stat(
        slide,
        left=6.25,
        top=3.95,
        width=6.45,
        height=0.9,
        label="Largest gap",
        big=fmt_eur(top_candidate["ARR (€ converted)"]),
        context=f"{top_candidate['Opportunity']} ({top_candidate['Owner']})",
        panel=LIGHT_PANEL_2,
        status=AMBER,
    )
    takeaway(
        slide,
        top=5.42,
        bullets=[
            "The control issue is not approved volume; it is the open stage 3+ backlog still missing approval discipline.",
            "Use the next slide as the operating action list, not as appendix detail.",
            "Keep the approval-rate method explicit because it is an average of component-book rates, not a direct denominator rollup.",
        ],
        accent=AMBER,
    )
    source_footer(slide, "Commercial approval rollup and missing-candidate list, EUR converted.")


def add_missing_approvals(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    candidates = snapshot["commercial_approval"]["missing_candidates"][:5]

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   APPROVAL FOLLOW-UP",
        title="Union and BBVA dominate the missing-approval backlog",
        narrative="Largest missing-approval candidates ranked by ARR exposure.",
        accent=AMBER,
    )
    rounded_card(slide, 0.6, 1.95, 12.1, 3.9, LIGHT_BLUE_PANEL)
    rect(slide, 0.6, 1.95, 0.1, 3.9, AMBER)
    rounded_card(slide, 0.92, 2.25, 11.45, 0.42, LIGHT_PANEL_2)
    headers = [("Opportunity", 1.08, 3.0), ("ARR", 4.15, 0.9), ("Owner", 5.15, 1.35), ("Stage", 6.6, 1.15), ("Next action", 7.85, 4.2)]
    for label, left, width in headers:
        txt(slide, left, 2.36, width, 0.16, label, size=9, bold=True, color=PRIMARY_BLUE)
    for index, candidate in enumerate(candidates):
        row_top = 2.82 + index * 0.58
        rounded_card(slide, 0.92, row_top, 11.45, 0.46, LIGHT_PANEL_2)
        txt(slide, 1.08, row_top + 0.11, 3.0, 0.18, candidate["Opportunity"], size=10, bold=True, color=NAVY)
        txt(slide, 4.15, row_top + 0.11, 0.9, 0.18, fmt_eur(candidate["ARR (€ converted)"]), size=10, color=NAVY)
        txt(slide, 5.15, row_top + 0.11, 1.35, 0.18, candidate["Owner"], size=10, color=GREY_TEXT)
        txt(slide, 6.6, row_top + 0.11, 1.15, 0.18, candidate["Stage"], size=10, color=GREY_TEXT)
        next_step = candidate.get("Next Step") or "Owner follow-up required"
        txt(slide, 7.85, row_top + 0.11, 4.2, 0.18, next_step[:120], size=9, color=GREY_TEXT)
    takeaway(
        slide,
        top=6.1,
        bullets=[
            f"Top two candidates alone represent {fmt_eur(candidates[0]['ARR (€ converted)'] + candidates[1]['ARR (€ converted)'])} of exposure.",
            "Use book-owner follow-up immediately after the regional review; this slide should drive action.",
        ],
        accent=AMBER,
    )
    source_footer(slide, "Top 5 missing-approval candidates from validated regional snapshot.")


def add_renewals(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    q2_deals = q2_renewals(snapshot)
    q2_acv = sum((deal.get("Renewal ACV (€ converted)") or 0) for deal in q2_deals)
    total_open = snapshot["renewals"]["summary_metrics"]["open_acv"]
    open_count = snapshot["renewals"]["summary_metrics"]["open_deal_count"]
    q2_count = snapshot["renewals"]["summary_metrics"].get("q2_open_deal_count", len(q2_deals))
    q2_acv = snapshot["renewals"]["summary_metrics"].get("q2_open_acv", q2_acv)
    top_open = sorted(snapshot["renewals"]["open_renewals"], key=lambda row: row.get("Renewal ACV (€ converted)") or 0, reverse=True)[:3]

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   RENEWALS",
        title=f"Q2 renewals are only {fmt_eur(q2_acv)} ACV across {q2_count} deals, but the total open renewal book is {fmt_eur(total_open)} and tagging is sparse",
        narrative="Renewal ACV, quarter focus, and named watchlist.",
        accent=MAGENTA,
    )
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=4.0,
        height=3.2,
        overline="Q2 RENEWALS",
        big=f"{q2_count} deals / {fmt_eur(q2_acv)}",
        context="Current-quarter renewal book only. Renewals remain ACV, EUR converted.",
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    mini_stat(
        slide,
        left=4.85,
        top=1.95,
        width=3.8,
        height=1.05,
        label="Total open renewals",
        big=f"{open_count} / {fmt_eur(total_open)}",
        context="Full open renewal book across all future close dates.",
        panel=MAGENTA_LIGHT,
        status=PRIMARY_BLUE,
    )
    mini_stat(
        slide,
        left=8.9,
        top=1.95,
        width=3.8,
        height=1.05,
        label="Risk tagging",
        big="Unspecified",
        context="All 22 open renewals currently sit in the unspecified bucket.",
        panel=MAGENTA_LIGHT,
        status=AMBER,
    )
    content_panel(
        slide,
        left=4.85,
        top=3.25,
        width=7.85,
        height=1.9,
        title="Named watchlist",
        lines=[
            f"{deal['Opportunity']} — {fmt_eur(deal.get('Renewal ACV (€ converted)') or 0)} ({deal['Owner']})"
            for deal in top_open
        ],
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    takeaway(
        slide,
        top=5.52,
        bullets=[
            "The quarter-specific renewal burden is manageable, but the total open book is much larger than the Q2 subset.",
            "Risk tagging is too sparse for a clean leadership read, so named-renewal ownership matters more than the bucket labels today.",
        ],
        accent=MAGENTA,
    )
    source_footer(slide, "Renewal ACV in EUR converted; q2 subset derived from close date.")


def add_slipped(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    pushed = snapshot["q1_review"]["pushed_deals"][:5]
    slipped_arr = snapshot["q1_review"]["actuals"]["slipped_arr"]
    slipped_count = snapshot["q1_review"]["actuals"]["slipped_count"]

    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   SLIPPED DEALS",
        title=f"AXA, GIB UK, and Longbow anchor {fmt_eur(slipped_arr)} of slipped Q1 exposure",
        narrative="Largest slipped deals first; root-cause commentary still needs owner input.",
        accent=AMBER,
    )
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=3.2,
        height=3.3,
        overline="VALIDATED SLIPPED EXPOSURE",
        big=f"{slipped_count} / {fmt_eur(slipped_arr)}",
        context="Regional slipped count and ARR from the validated fact pack.",
        panel=LIGHT_BLUE_PANEL,
        accent=AMBER,
    )
    rounded_card(slide, 4.05, 1.95, 8.65, 3.3, LIGHT_PANEL_2)
    txt(slide, 4.28, 2.13, 2.9, 0.18, "OPPORTUNITY", size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, 7.35, 2.13, 0.9, 0.18, "ARR", size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, 8.35, 2.13, 0.95, 0.18, "OLD -> NEW", size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, 9.48, 2.13, 3.0, 0.18, "OWNER / STAGE", size=9, bold=True, color=PRIMARY_BLUE)
    for index, deal in enumerate(pushed):
        row_top = 2.45 + index * 0.54
        rounded_card(slide, 4.28, row_top, 8.12, 0.42, LIGHT_BLUE_PANEL)
        txt(slide, 4.42, row_top + 0.11, 2.8, 0.16, deal["Opportunity"], size=9, bold=True, color=NAVY)
        txt(slide, 7.35, row_top + 0.11, 0.9, 0.16, fmt_eur(deal["ARR (€ converted)"]), size=9, color=NAVY)
        txt(slide, 8.35, row_top + 0.11, 0.95, 0.16, f"{deal['Old Close']} -> {deal['New Close']}", size=8, color=GREY_TEXT)
        txt(slide, 9.48, row_top + 0.11, 2.8, 0.16, f"{deal['Owner']} | {deal['Stage']}", size=8, color=GREY_TEXT)
    takeaway(
        slide,
        top=5.55,
        bullets=[
            "The slipped-deal list is factual; root-cause commentary is still dependent on owner follow-up.",
            "Use this slide to assign follow-up by book and by opportunity, not to infer causes from metadata alone.",
        ],
        accent=AMBER,
    )
    source_footer(slide, "Top slipped deals from validated regional pushed-deals rollup.")


def add_churn(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   CHURN / FINANCE",
        title="Churn narrative is not publishable yet without a Finance-owned overlay",
        narrative="Keep this as an explicit operating gap until Finance input is wired into the monthly cadence.",
        accent=MAGENTA,
    )
    hero_stat_card(
        slide,
        left=0.6,
        top=1.95,
        width=12.1,
        height=2.0,
        overline="FINANCE OVERLAY REQUIRED",
        big="Churn story blocked",
        context="CRM signals alone are not enough for a credible churn trend. The monthly deck needs a Finance-owned input to make this slide publishable.",
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    mini_stat(
        slide,
        left=0.6,
        top=4.35,
        width=3.85,
        height=1.18,
        label="Current status",
        big="Missing",
        context="No Finance-backed recurring churn feed in the regional build yet.",
        panel=LIGHT_PANEL_2,
        status=AMBER,
    )
    mini_stat(
        slide,
        left=4.72,
        top=4.35,
        width=3.85,
        height=1.18,
        label="Required owner",
        big="Finance + Sales Ops",
        context="Confirm the reporting owner and bring the feed into the monthly cadence.",
        panel=LIGHT_PANEL_2,
        status=PRIMARY_BLUE,
    )
    mini_stat(
        slide,
        left=8.85,
        top=4.35,
        width=3.85,
        height=1.18,
        label="Next step",
        big="Operationalize",
        context="Lock source, owner, and refresh rhythm before adding trend numbers here.",
        panel=LIGHT_PANEL_2,
        status=AMBER,
    )
    takeaway(
        slide,
        top=5.95,
        bullets=[
            "Keep the gap visible rather than filling it with speculative churn commentary.",
            "The process fix is part of the monthly operating model, not just a one-off deck edit.",
        ],
        accent=MAGENTA,
    )
    source_footer(slide, "Placeholder kept explicit until Finance feed is operationalized.")


def add_appendix(prs: Presentation, snapshot: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    header(
        slide,
        eyebrow=f"{snapshot['snapshot_date']}   ·   {snapshot['region_name'].upper()}   ·   APPENDIX",
        title="Metric and lineage notes",
        narrative="Definitions, guardrails, and known limitations for this regional deck.",
        accent=PRIMARY_BLUE,
    )
    content_panel(
        slide,
        left=0.6,
        top=1.95,
        width=5.9,
        height=4.6,
        title="Metric rules",
        lines=[
            "New business and expansion stay ARR, EUR converted.",
            "Renewals stay ACV, EUR converted.",
            "Omitted remains visible but excluded from active headline pipeline.",
            "Q1 promise baseline remains qualified and should not be stated as a clean commitment number.",
        ],
        panel=LIGHT_BLUE_PANEL,
        accent=PRIMARY_BLUE,
    )
    content_panel(
        slide,
        left=6.8,
        top=1.95,
        width=5.9,
        height=4.6,
        title="Source and hierarchy notes",
        lines=[
            "Regional deck built from validated director snapshots and regional rollup logic.",
            "Middle East & Africa remains under EMEA in the forecast hierarchy.",
            "Commercial approval rate is a simple average of component-book approval rates.",
            "Churn slide remains intentionally incomplete until Finance overlay exists.",
        ],
        panel=LIGHT_PANEL_2,
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Regional snapshot, validated fact pack, and hierarchy notes.")


def build_gold_deck(
    *,
    region_snapshot_path: Path,
    output_path: Path,
    master_template_path: Path = DEFAULT_MASTER_TEMPLATE_PATH,
    shell_contract_path: Path = DEFAULT_SHELL_CONTRACT_PATH,
) -> dict[str, Any]:
    snapshot = load_json(region_snapshot_path)
    shell = load_json(shell_contract_path)
    prs = Presentation(str(master_template_path))
    clear_slides(prs)

    add_cover(prs, region_name=snapshot["region_name"], snapshot_date=snapshot["snapshot_date"])
    add_agenda(prs, shell, region_name=snapshot["region_name"], snapshot_date=snapshot["snapshot_date"])
    add_exec_summary(prs, snapshot)
    add_q1_review(prs, snapshot)
    add_quarterly_pipeline(prs, snapshot)
    add_book_breakdown(prs, snapshot)
    add_pipeline_intel(prs, snapshot)
    add_approval_overview(prs, snapshot)
    add_missing_approvals(prs, snapshot)
    add_renewals(prs, snapshot)
    add_slipped(prs, snapshot)
    add_churn(prs, snapshot)
    add_appendix(prs, snapshot)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return {
        "deck_path": str(output_path),
        "slide_count": len(prs.slides),
        "region_name": snapshot["region_name"],
        "snapshot_date": snapshot["snapshot_date"],
        "source_snapshot_path": str(region_snapshot_path),
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--region-snapshot-path", type=Path, required=True)
    parser.add_argument("--output-path", type=Path, required=True)
    parser.add_argument("--master-template-path", type=Path, default=DEFAULT_MASTER_TEMPLATE_PATH)
    parser.add_argument("--shell-contract-path", type=Path, default=DEFAULT_SHELL_CONTRACT_PATH)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    result = build_gold_deck(
        region_snapshot_path=args.region_snapshot_path,
        output_path=args.output_path,
        master_template_path=args.master_template_path,
        shell_contract_path=args.shell_contract_path,
    )
    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
