#!/usr/bin/env python3
"""Build the Sales Global Summary shell using the native SimCorp template path."""

from __future__ import annotations

import argparse
import json
import os
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MASTER_TEMPLATE_PATH = (
    Path.home()
    / "archive"
    / "simcorp-deck-agent-backup"
    / "reference-decks"
    / "SimCorp_PPT_Template.pptx"
)
DEFAULT_SHELL_CONTRACT_PATH = REPO_ROOT / "config" / "sales_global_summary_shell.json"
DEFAULT_OUTPUT_ROOT = REPO_ROOT / "output" / "sales_global_summary_shells"
DEFAULT_JS_BUILDER_PATH = REPO_ROOT / "scripts" / "build_sales_global_summary_shell_v1.js"
DEFAULT_NODE_MODULES_PATH = REPO_ROOT / "output" / "sales_director_monthly_deck_2026-03-31" / "node_modules"

NAVY = RGBColor(0x01, 0x19, 0x46)
PRIMARY_BLUE = RGBColor(0x0E, 0x37, 0x88)
AQUA = RGBColor(0x6F, 0xCC, 0xDD)
AQUA_LIGHT = RGBColor(0xCB, 0xF5, 0xF3)
LIGHT_BLUE_PANEL = RGBColor(0xE6, 0xEE, 0xFE)
LIGHT_PANEL = RGBColor(0xF5, 0xF8, 0xFD)
MAGENTA = RGBColor(0x9D, 0x2E, 0x7B)
MAGENTA_LIGHT = RGBColor(0xF8, 0xE8, 0xF1)
AMBER = RGBColor(0xFB, 0x9B, 0x2A)
GREEN_OK = RGBColor(0x33, 0xD8, 0xCE)
GREY_TEXT = RGBColor(0x5C, 0x74, 0x82)
DIVIDER_GREY = RGBColor(0xD7, 0xE2, 0xE8)

SLIDE_EYEBROWS = {
    "global-executive-summary": "GLOBAL EXECUTIVE VIEW",
    "apac-region-summary": "APAC OPERATING VIEW",
    "emea-region-summary": "EMEA OPERATING VIEW",
    "north-america-region-summary": "NORTH AMERICA OPERATING VIEW",
    "global-commercial-approval-overview": "COMMERCIAL APPROVAL",
    "global-appendix": "APPENDIX",
}

SLIDE_ACCENTS = {
    "global-executive-summary": AQUA,
    "apac-region-summary": AQUA,
    "emea-region-summary": PRIMARY_BLUE,
    "north-america-region-summary": MAGENTA,
    "global-commercial-approval-overview": AMBER,
    "global-appendix": PRIMARY_BLUE,
}


def load_shell_contract(path: Path = DEFAULT_SHELL_CONTRACT_PATH) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def merge_fill_payload_with_shell_contract(
    fill_payload: dict[str, Any],
    shell_contract: dict[str, Any],
) -> dict[str, Any]:
    payload_by_id = {
        slide["id"]: slide
        for slide in fill_payload.get("slides", [])
        if isinstance(slide, dict) and slide.get("id")
    }
    merged_slides: list[dict[str, Any]] = []
    for slide_def in shell_contract.get("slides", []):
        payload_slide = payload_by_id.get(slide_def.get("id"), {})
        merged_slide = dict(slide_def)
        for key, value in payload_slide.items():
            if key == "slots":
                continue
            merged_slide[key] = value
        merged_slide["slots"] = payload_slide.get("slots", {})
        if not merged_slide.get("support_level"):
            merged_slide["support_level"] = (
                slide_def.get("data_contract", {}) or {}
            ).get("support_level")
        merged_slides.append(merged_slide)

    merged_payload = dict(fill_payload)
    merged_payload["template_name"] = shell_contract.get(
        "template_name", fill_payload.get("template_name")
    )
    merged_payload["slides"] = merged_slides
    return merged_payload


def base_fill_payload(*, snapshot_date: str) -> dict[str, Any]:
    return {
        "template_name": "Sales Global Summary",
        "snapshot_date": snapshot_date,
        "slides": [],
    }


def get_layout(prs: Presentation, name: str):
    lowered = name.lower()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name or layout.name.lower() == lowered:
                return layout
    raise KeyError(f"Layout not found: {name}")


def clear_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        rel_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]


def blank_slide(prs: Presentation):
    try:
        layout = get_layout(prs, "Blank")
    except KeyError:
        layout = get_layout(prs, "Title only")
    slide = prs.slides.add_slide(layout)
    for placeholder in list(slide.placeholders):
        placeholder._element.getparent().remove(placeholder._element)
    return slide


def txt(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = NAVY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft Sans Serif"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def multi_paragraph(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    size: int = 12,
    color: RGBColor = NAVY,
    bullet: bool = False,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.clear()
    if not lines:
        return tb
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {line}" if bullet else line
        run.font.name = "Microsoft Sans Serif"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def rounded_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def header(
    slide,
    *,
    eyebrow: str,
    title: str,
    narrative: str,
    accent: RGBColor,
) -> None:
    rect(slide, 0.6, 0.42, 0.08, 0.7, accent)
    txt(slide, 0.82, 0.38, 12.0, 0.24, eyebrow, size=10, bold=True, color=PRIMARY_BLUE)
    txt(slide, 0.82, 0.62, 12.0, 0.55, title, size=28, bold=True, color=NAVY)
    txt(slide, 0.82, 1.25, 12.1, 0.4, narrative, size=13, color=GREY_TEXT)


def source_footer(slide, text: str) -> None:
    txt(slide, 0.6, 7.05, 12.1, 0.22, text, size=9, color=GREY_TEXT)


def metric_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    value: str,
    context: str,
    panel: RGBColor,
    accent: RGBColor,
) -> None:
    rounded_card(slide, left, top, width, height, panel)
    rect(slide, left, top, 0.08, height, accent)
    txt(slide, left + 0.22, top + 0.16, width - 0.35, 0.2, label, size=9, bold=True, color=PRIMARY_BLUE)
    txt(slide, left + 0.22, top + 0.46, width - 0.35, 0.42, value, size=24, bold=True, color=NAVY)
    txt(slide, left + 0.22, top + 0.96, width - 0.35, height - 1.05, context, size=11, color=GREY_TEXT)


def content_panel(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    lines: list[str],
    panel: RGBColor,
    accent: RGBColor,
) -> None:
    rounded_card(slide, left, top, width, height, panel)
    rect(slide, left, top, 0.08, height, accent)
    txt(slide, left + 0.22, top + 0.16, width - 0.35, 0.22, title, size=12, bold=True, color=NAVY)
    multi_paragraph(slide, left + 0.22, top + 0.5, width - 0.4, height - 0.65, lines, size=11, color=GREY_TEXT)


def as_text(value: Any, fallback: str = "—") -> str:
    if value is None:
        return fallback
    if isinstance(value, str):
        cleaned = value.strip()
        return cleaned or fallback
    return str(value)


def list_lines(value: Any) -> list[str]:
    if not value:
        return []
    if isinstance(value, list):
        lines: list[str] = []
        for item in value:
            if isinstance(item, dict):
                lines.append(
                    " | ".join(
                        part
                        for part in [
                            as_text(item.get("region_name"), ""),
                            as_text(item.get("opportunity"), ""),
                            as_text(item.get("owner"), ""),
                            as_text(item.get("stage"), ""),
                            as_text(item.get("arr_eur"), ""),
                        ]
                        if part
                    )
                )
            else:
                lines.append(as_text(item))
        return [line for line in lines if line and line != "—"]
    return [as_text(value)]


def payload_slide_by_id(fill_payload: dict[str, Any], slide_id: str) -> dict[str, Any]:
    for slide in fill_payload.get("slides", []):
        if slide.get("id") == slide_id:
            return slide
    return {}


def shell_slide_title(slide_def: dict[str, Any]) -> str:
    if slide_def["id"] == "global-appendix":
        return "Appendix and Guardrails"
    return slide_def["title"]


def populated_slide_title(slide_def: dict[str, Any], slots: dict[str, Any]) -> str:
    slide_id = slide_def["id"]
    if slide_id == "global-executive-summary":
        return "Global pipeline and leadership actions"
    if slide_id.endswith("region-summary"):
        region_name = as_text(slots.get("region_name"), "Region")
        return f"{region_name}: pipeline, approvals, and renewals"
    if slide_id == "global-commercial-approval-overview":
        return "Commercial approval exposure and follow-up"
    if slide_id == "global-appendix":
        return "Appendix and Guardrails"
    return slide_def["title"]


def add_cover_slide(prs: Presentation, *, snapshot_date: str) -> None:
    slide = blank_slide(prs)
    rect(slide, 0.6, 0.42, 0.12, 6.2, AQUA)
    txt(slide, 0.95, 0.85, 9.0, 0.55, "Sales Global Summary", size=28, bold=True, color=NAVY)
    txt(slide, 0.95, 1.48, 9.6, 0.35, f"Snapshot date: {snapshot_date}", size=14, color=GREY_TEXT)
    txt(
        slide,
        0.95,
        2.15,
        8.2,
        0.7,
        "Deterministic global leadership cut built from validated regional rollups.",
        size=20,
        bold=True,
        color=PRIMARY_BLUE,
    )
    content_panel(
        slide,
        left=0.95,
        top=3.15,
        width=5.8,
        height=1.45,
        title="Run standard",
        lines=[
            "Salesforce and workbook contracts feed the regional snapshots first.",
            "Global output is publishable only when audits, fonts, and canonical promotion all clear.",
        ],
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    content_panel(
        slide,
        left=7.0,
        top=3.15,
        width=5.7,
        height=1.45,
        title="Operating posture",
        lines=[
            "Use this deck for leadership rollup only.",
            "Director books remain the operating source for territory detail.",
        ],
        panel=LIGHT_PANEL,
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Native SimCorp template render — deterministic global baseline.")


def add_agenda_slide(prs: Presentation, shell_contract: dict[str, Any]) -> None:
    slide = blank_slide(prs)
    header(
        slide,
        eyebrow="GLOBAL MONTHLY AGENDA",
        title="Agenda",
        narrative="Fixed leadership sequence: executive readout, three regional operating views, approvals, and appendix notes.",
        accent=PRIMARY_BLUE,
    )
    agenda_lines = [
        "1. Global pipeline and leadership actions",
        "2. APAC operating view",
        "3. EMEA operating view",
        "4. North America operating view",
        "5. Commercial approval exposure and follow-up",
        "6. Appendix and Guardrails",
    ]
    content_panel(
        slide,
        left=0.8,
        top=2.0,
        width=6.0,
        height=4.4,
        title="Sequence",
        lines=agenda_lines,
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    support_lines = [
        f"{slide['title']} ({(slide.get('data_contract') or {}).get('support_level', 'unknown')})"
        for slide in shell_contract.get("slides", [])
    ]
    content_panel(
        slide,
        left=7.0,
        top=2.0,
        width=5.7,
        height=4.4,
        title="Source contract",
        lines=support_lines,
        panel=LIGHT_PANEL,
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Agenda follows the validated global shell contract.")


def add_global_executive_summary_slide(
    slide,
    *,
    slide_def: dict[str, Any],
    slots: dict[str, Any],
    populated: bool,
) -> None:
    title = populated_slide_title(slide_def, slots) if populated else shell_slide_title(slide_def)
    narrative = (
        "Global ARR, renewals, and the single escalation leadership should act on now."
        if populated
        else slide_def.get("subtitle", "")
    )
    header(
        slide,
        eyebrow=SLIDE_EYEBROWS[slide_def["id"]],
        title=title,
        narrative=narrative,
        accent=SLIDE_ACCENTS[slide_def["id"]],
    )
    metric_card(
        slide,
        left=0.8,
        top=1.95,
        width=3.8,
        height=1.7,
        label="Q2 active ARR",
        value=as_text(slots.get("global_pipeline_arr_q2"), "Replace with validated ARR"),
        context="Validated global active pipeline rolled up from regional snapshots.",
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    metric_card(
        slide,
        left=4.8,
        top=1.95,
        width=3.8,
        height=1.7,
        label="Q2 renewal ACV",
        value=as_text(slots.get("global_renewal_acv_q2"), "Replace with validated ACV"),
        context="Renewal exposure remains ACV and stays separate from pipeline ARR.",
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    metric_card(
        slide,
        left=8.8,
        top=1.95,
        width=3.9,
        height=1.7,
        label="Missing approvals",
        value=as_text(slots.get("global_missing_approval_count"), "0"),
        context="Count of unresolved commercial-approval candidates across the three regions.",
        panel=LIGHT_PANEL,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=0.8,
        top=4.0,
        width=5.9,
        height=2.1,
        title="Top risk",
        lines=[as_text(slots.get("global_top_risk"), "Add the leading cross-region risk here.")],
        panel=LIGHT_PANEL,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=6.95,
        top=4.0,
        width=5.75,
        height=2.1,
        title="Top action",
        lines=[as_text(slots.get("global_top_action"), "Add the single global action here.")],
        panel=LIGHT_PANEL,
        accent=GREEN_OK,
    )
    source_footer(slide, "Global executive summary is a deterministic rollup from validated regional snapshots.")


def add_region_summary_slide(
    slide,
    *,
    slide_def: dict[str, Any],
    slots: dict[str, Any],
    populated: bool,
) -> None:
    region_name = as_text(slots.get("region_name"), slide_def["title"].replace(" Regional Summary", ""))
    title = populated_slide_title(slide_def, slots) if populated else shell_slide_title(slide_def)
    quarter_label = as_text(slots.get("quarterly_pipeline_label"), "Q2")
    quarter_title = as_text(slots.get("quarterly_pipeline_title"), "Q2 2026")
    display_reason = as_text(slots.get("quarterly_pipeline_display_reason"), "current_quarter")
    narrative = (
        f"{quarter_title} control view for pipeline, approvals, and renewals."
        if populated
        else slide_def.get("subtitle", "")
    )
    if populated and display_reason == "forward_quarter_fallback":
        narrative = f"{quarter_title} forward-quarter fallback because the current quarter is empty."
    header(
        slide,
        eyebrow=SLIDE_EYEBROWS[slide_def["id"]],
        title=title,
        narrative=narrative,
        accent=SLIDE_ACCENTS[slide_def["id"]],
    )
    metric_card(
        slide,
        left=0.8,
        top=1.95,
        width=2.9,
        height=1.8,
        label=f"{quarter_label} active ARR",
        value=as_text(slots.get("headline_pipeline_arr_q2"), "Add validated ARR"),
        context=f"Region active pipeline for {quarter_title}.",
        panel=LIGHT_BLUE_PANEL,
        accent=AQUA,
    )
    metric_card(
        slide,
        left=3.95,
        top=1.95,
        width=2.7,
        height=1.8,
        label="Commit",
        value=as_text(slots.get("q2_commit_arr"), "—"),
        context=f"{quarter_title} commit ARR.",
        panel=LIGHT_PANEL,
        accent=PRIMARY_BLUE,
    )
    metric_card(
        slide,
        left=6.9,
        top=1.95,
        width=2.7,
        height=1.8,
        label="Best Case",
        value=as_text(slots.get("q2_best_case_arr"), "—"),
        context=f"{quarter_title} best-case ARR.",
        panel=LIGHT_PANEL,
        accent=GREEN_OK,
    )
    metric_card(
        slide,
        left=9.85,
        top=1.95,
        width=2.85,
        height=1.8,
        label="Open renewals ACV",
        value=as_text(slots.get("renewal_open_acv"), "—"),
        context="Renewal exposure stays ACV.",
        panel=MAGENTA_LIGHT,
        accent=MAGENTA,
    )
    content_panel(
        slide,
        left=0.8,
        top=4.0,
        width=3.4,
        height=1.65,
        title="Approval and omitted view",
        lines=[
            f"Approval rate stage 3+: {as_text(slots.get('approval_rate_stage3_plus'), '—')}",
            f"Omitted ARR: {as_text(slots.get('q2_omitted_arr'), '—')}",
        ],
        panel=LIGHT_PANEL,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=4.45,
        top=4.0,
        width=4.0,
        height=1.65,
        title="Top risk",
        lines=[as_text(slots.get("top_risk"), f"Add the top risk for {region_name}.")],
        panel=LIGHT_PANEL,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=8.7,
        top=4.0,
        width=4.0,
        height=1.65,
        title="Top action",
        lines=[as_text(slots.get("top_action"), f"Add the top action for {region_name}.")],
        panel=LIGHT_PANEL,
        accent=PRIMARY_BLUE,
    )
    footnote = as_text(slots.get("quarterly_pipeline_footnote"), "")
    if populated and footnote != "—":
        txt(slide, 0.8, 6.1, 11.8, 0.3, footnote, size=10, color=GREY_TEXT)
    source_footer(slide, f"{region_name} slide is populated only from its validated regional snapshot.")


def add_commercial_approval_slide(
    slide,
    *,
    slide_def: dict[str, Any],
    slots: dict[str, Any],
    populated: bool,
) -> None:
    title = populated_slide_title(slide_def, slots) if populated else shell_slide_title(slide_def)
    header(
        slide,
        eyebrow=SLIDE_EYEBROWS[slide_def["id"]],
        title=title,
        narrative="Approved exposure, missing-approval exposure, and the largest unresolved candidates by region.",
        accent=SLIDE_ACCENTS[slide_def["id"]],
    )
    approved_lines = [
        f"{as_text(row.get('region_name'))}: {as_text(row.get('arr_eur'))} ({as_text(row.get('deal_count'))} deals)"
        for row in (slots.get("approved_2026_by_region") or [])
    ] or ["Add approved 2026 exposure by region."]
    missing_lines = [
        f"{as_text(row.get('region_name'))}: {as_text(row.get('arr_eur'))} ({as_text(row.get('candidate_count'))} candidates)"
        for row in (slots.get("missing_approval_by_region") or [])
    ] or ["Add missing-approval exposure by region."]
    candidate_lines = []
    for row in (slots.get("largest_global_missing_candidates") or [])[:6]:
        candidate_lines.append(
            " | ".join(
                part
                for part in [
                    as_text(row.get("region_name"), ""),
                    as_text(row.get("opportunity"), ""),
                    as_text(row.get("owner"), ""),
                    as_text(row.get("stage"), ""),
                    as_text(row.get("arr_eur"), ""),
                ]
                if part
            )
        )
    if not candidate_lines:
        candidate_lines = ["Add the largest unresolved approval candidates here."]
    content_panel(
        slide,
        left=0.8,
        top=1.95,
        width=3.7,
        height=4.75,
        title="Approved 2026 by region",
        lines=approved_lines,
        panel=LIGHT_BLUE_PANEL,
        accent=GREEN_OK,
    )
    content_panel(
        slide,
        left=4.75,
        top=1.95,
        width=3.7,
        height=4.75,
        title="Missing approval by region",
        lines=missing_lines,
        panel=LIGHT_PANEL,
        accent=AMBER,
    )
    content_panel(
        slide,
        left=8.7,
        top=1.95,
        width=4.0,
        height=4.75,
        title="Largest unresolved candidates",
        lines=candidate_lines,
        panel=LIGHT_PANEL,
        accent=PRIMARY_BLUE,
    )
    source_footer(slide, "Commercial approval overview rolls up validated regional approval exposure and missing candidates.")


def add_appendix_slide(
    slide,
    *,
    slide_def: dict[str, Any],
    slots: dict[str, Any],
) -> None:
    header(
        slide,
        eyebrow=SLIDE_EYEBROWS[slide_def["id"]],
        title="Appendix and Guardrails",
        narrative="Metric definitions, rollup rules, and known limitations for the global cut.",
        accent=SLIDE_ACCENTS[slide_def["id"]],
    )
    content_panel(
        slide,
        left=0.8,
        top=1.95,
        width=3.75,
        height=4.75,
        title="Metric definition notes",
        lines=list_lines(slots.get("metric_definition_notes")) or ["Document ARR versus ACV rules here."],
        panel=LIGHT_BLUE_PANEL,
        accent=PRIMARY_BLUE,
    )
    content_panel(
        slide,
        left=4.8,
        top=1.95,
        width=3.75,
        height=4.75,
        title="Region rollup notes",
        lines=list_lines(slots.get("region_rollup_notes")) or ["Document forecast hierarchy and region inclusion notes here."],
        panel=LIGHT_PANEL,
        accent=AQUA,
    )
    content_panel(
        slide,
        left=8.8,
        top=1.95,
        width=3.9,
        height=4.75,
        title="Known gaps",
        lines=list_lines(slots.get("known_gaps")) or ["List known limitations explicitly instead of hiding them in narration."],
        panel=LIGHT_PANEL,
        accent=AMBER,
    )
    source_footer(slide, "Appendix notes document metric discipline and known global rollup limitations.")


def build_native_shell_deck(
    *,
    snapshot_date: str,
    output_path: Path,
    master_template_path: Path,
    shell_contract_path: Path,
    fill_payload_path: Path | None,
) -> dict[str, Any]:
    if fill_payload_path is not None and not fill_payload_path.exists():
        raise FileNotFoundError(f"Fill payload not found: {fill_payload_path}")
    if not master_template_path.exists():
        raise FileNotFoundError(f"Master template not found: {master_template_path}")

    shell = load_shell_contract(shell_contract_path)
    fill_payload = (
        json.loads(fill_payload_path.read_text(encoding="utf-8"))
        if fill_payload_path
        else base_fill_payload(snapshot_date=snapshot_date)
    )
    merged_payload = merge_fill_payload_with_shell_contract(fill_payload, shell)
    populated = bool(fill_payload_path)

    prs = Presentation(str(master_template_path))
    clear_slides(prs)
    add_cover_slide(prs, snapshot_date=snapshot_date)
    add_agenda_slide(prs, shell)

    for slide_def in shell.get("slides", []):
        slide = blank_slide(prs)
        payload_slide = payload_slide_by_id(merged_payload, slide_def["id"])
        slots = payload_slide.get("slots") or {}
        if slide_def["id"] == "global-executive-summary":
            add_global_executive_summary_slide(
                slide,
                slide_def=slide_def,
                slots=slots,
                populated=populated,
            )
        elif slide_def["id"].endswith("region-summary"):
            add_region_summary_slide(
                slide,
                slide_def=slide_def,
                slots=slots,
                populated=populated,
            )
        elif slide_def["id"] == "global-commercial-approval-overview":
            add_commercial_approval_slide(
                slide,
                slide_def=slide_def,
                slots=slots,
                populated=populated,
            )
        elif slide_def["id"] == "global-appendix":
            add_appendix_slide(slide, slide_def=slide_def, slots=slots)
        else:  # pragma: no cover - defensive against contract drift
            raise ValueError(f"Unsupported shell slide id: {slide_def['id']}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return {
        "deck_path": str(output_path),
        "slide_count": len(prs.slides),
        "master_template_path": str(master_template_path),
        "shell_contract_path": str(shell_contract_path),
        "template_version": shell.get("template_version"),
        "builder": "simcorp-native" if fill_payload_path else "simcorp-native-shell",
        "publish_safe": bool(fill_payload_path),
        "js_builder_path": str(DEFAULT_JS_BUILDER_PATH),
        "fill_payload_path": str(fill_payload_path) if fill_payload_path else None,
    }


def build_shell_deck(
    *,
    snapshot_date: str,
    output_path: Path,
    master_template_path: Path = DEFAULT_MASTER_TEMPLATE_PATH,
    shell_contract_path: Path = DEFAULT_SHELL_CONTRACT_PATH,
    js_builder_path: Path = DEFAULT_JS_BUILDER_PATH,
    node_modules_path: Path = DEFAULT_NODE_MODULES_PATH,
    fill_payload_path: Path | None = None,
    allow_legacy_js_builder: bool = False,
) -> dict[str, Any]:
    if fill_payload_path is not None or not allow_legacy_js_builder:
        return build_native_shell_deck(
            snapshot_date=snapshot_date,
            output_path=output_path,
            master_template_path=master_template_path,
            shell_contract_path=shell_contract_path,
            fill_payload_path=fill_payload_path,
        )

    if not js_builder_path.exists():
        raise FileNotFoundError(f"JS shell builder not found: {js_builder_path}")
    if not node_modules_path.exists():
        raise FileNotFoundError(
            f"Node modules path not found for shell builder dependencies: {node_modules_path}"
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    env = os.environ.copy()
    existing_node_path = env.get("NODE_PATH", "")
    env["NODE_PATH"] = (
        f"{node_modules_path}:{existing_node_path}" if existing_node_path else str(node_modules_path)
    )
    env["SD_GLOBAL_SHELL_VALIDATE_LAYOUT"] = "1"

    cmd = [
        "node",
        str(js_builder_path),
        "--snapshot-date",
        snapshot_date,
        "--contract",
        str(shell_contract_path),
        "--output",
        str(output_path),
    ]
    subprocess.run(cmd, check=True, env=env, cwd=REPO_ROOT)

    shell = load_shell_contract(shell_contract_path)
    return {
        "deck_path": str(output_path),
        "slide_count": len(shell.get("slides", [])) + 2,
        "master_template_path": str(master_template_path),
        "shell_contract_path": str(shell_contract_path),
        "template_version": shell.get("template_version"),
        "builder": "js-v1",
        "publish_safe": False,
        "js_builder_path": str(js_builder_path),
        "fill_payload_path": None,
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--snapshot-date", required=True)
    parser.add_argument("--output-path", type=Path, required=True)
    parser.add_argument("--master-template-path", type=Path, default=DEFAULT_MASTER_TEMPLATE_PATH)
    parser.add_argument("--shell-contract-path", type=Path, default=DEFAULT_SHELL_CONTRACT_PATH)
    parser.add_argument("--js-builder-path", type=Path, default=DEFAULT_JS_BUILDER_PATH)
    parser.add_argument("--node-modules-path", type=Path, default=DEFAULT_NODE_MODULES_PATH)
    parser.add_argument("--fill-payload-path", type=Path, default=None)
    parser.add_argument(
        "--allow-legacy-js-builder",
        action="store_true",
        help="Explicitly allow the non-publish-safe JS global shell renderer.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    result = build_shell_deck(
        snapshot_date=args.snapshot_date,
        output_path=args.output_path,
        master_template_path=args.master_template_path,
        shell_contract_path=args.shell_contract_path,
        js_builder_path=args.js_builder_path,
        node_modules_path=args.node_modules_path,
        fill_payload_path=args.fill_payload_path,
        allow_legacy_js_builder=args.allow_legacy_js_builder,
    )
    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
