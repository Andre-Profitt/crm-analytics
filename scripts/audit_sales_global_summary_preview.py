#!/usr/bin/env python3
"""Audit a populated Sales Global Summary preview deck for shell leakage and title quality."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


DEFAULT_SHELL_CONTRACT = Path(__file__).resolve().parents[1] / "config" / "sales_global_summary_shell.json"

PLACEHOLDER_TOKENS = [
    "€x.xM",
    "xx%",
    "Opportunity A",
    "Owner A",
]

SHELL_STYLE_TOKENS = [
    "Meeting rules",
    "Month in view",
    "Leadership use",
]


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def extract_slide_texts(deck_path: Path) -> list[list[str]]:
    prs = Presentation(str(deck_path))
    slides: list[list[str]] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                texts.append(text.strip())
        slides.append(texts)
    return slides


def first_meaningful_text(texts: list[str]) -> str:
    for text in texts:
        cleaned = " ".join(text.split())
        if cleaned:
            return cleaned
    return ""


def analyze_preview(
    *,
    shell_contract: dict[str, Any],
    fill_payload: dict[str, Any],
    slide_texts: list[list[str]],
) -> dict[str, Any]:
    findings: list[dict[str, Any]] = []
    body_slides = shell_contract.get("slides", [])
    payload_by_id = {slide["id"]: slide for slide in fill_payload.get("slides", [])}
    quarterly_pipeline_disclosures: list[dict[str, Any]] = []

    cover_text = " ".join(slide_texts[0]) if slide_texts else ""
    if "shell" in cover_text.lower():
        findings.append(
            {
                "severity": "error",
                "type": "cover_shell_language",
                "slide_number": 1,
                "message": "Cover still uses shell language.",
            }
        )

    for index, slide_def in enumerate(body_slides, start=3):
        texts = slide_texts[index - 1] if len(slide_texts) >= index else []
        joined = "\n".join(texts)
        slide_id = slide_def["id"]
        payload_slide = payload_by_id.get(slide_id, {})
        slots = payload_slide.get("slots") or {}
        title = first_meaningful_text(texts)
        if slide_id != "global-appendix" and title == slide_def["title"]:
            findings.append(
                {
                    "severity": "warn",
                    "type": "title_not_rewritten",
                    "slide_number": index,
                    "slide_id": slide_id,
                    "message": f"Slide title still matches shell title `{slide_def['title']}`.",
                }
            )
        leaked = [token for token in PLACEHOLDER_TOKENS if token in joined]
        if leaked:
            findings.append(
                {
                    "severity": "warn",
                    "type": "placeholder_leak",
                    "slide_number": index,
                    "slide_id": slide_id,
                    "message": f"Slide still contains placeholder tokens: {', '.join(leaked)}.",
                }
            )
        shell_style = [token for token in SHELL_STYLE_TOKENS if token in joined]
        if shell_style:
            findings.append(
                {
                    "severity": "warn",
                    "type": "shell_style_label",
                    "slide_number": index,
                    "slide_id": slide_id,
                    "message": f"Slide still contains shell-style labels: {', '.join(shell_style)}.",
                }
            )

        display_reason = str(slots.get("quarterly_pipeline_display_reason") or "").strip()
        region_name = str(slots.get("region_name") or slide_id).strip()
        quarter_title = str(slots.get("quarterly_pipeline_title") or "").strip()
        footnote = str(slots.get("quarterly_pipeline_footnote") or "").strip()
        if display_reason:
            disclosure = {
                "slide_id": slide_id,
                "slide_number": index,
                "region_name": region_name,
                "display_reason": display_reason,
                "quarterly_pipeline_title": quarter_title,
                "quarterly_pipeline_footnote": footnote or None,
            }
            quarterly_pipeline_disclosures.append(disclosure)
            if display_reason == "forward_quarter_fallback":
                if not (footnote and footnote in joined):
                    findings.append(
                        {
                            "severity": "error",
                            "type": "forward_quarter_fallback_hidden",
                            "slide_number": index,
                            "slide_id": slide_id,
                            "message": (
                                f"{region_name} is using a forward-quarter fallback but the required footnote "
                                "is not visible on the rendered slide."
                            ),
                        }
                    )

    ok = not any(item["severity"] == "error" for item in findings)
    return {
        "ok": ok,
        "slide_count": len(slide_texts),
        "finding_count": len(findings),
        "findings": findings,
        "quarterly_pipeline_disclosures": quarterly_pipeline_disclosures,
    }


def audit_preview(deck_path: Path, fill_payload_path: Path, shell_contract_path: Path = DEFAULT_SHELL_CONTRACT) -> dict[str, Any]:
    shell_contract = load_json(shell_contract_path)
    fill_payload = load_json(fill_payload_path)
    slide_texts = extract_slide_texts(deck_path)
    return analyze_preview(
        shell_contract=shell_contract,
        fill_payload=fill_payload,
        slide_texts=slide_texts,
    )


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--deck", type=Path, required=True)
    parser.add_argument("--fill-payload", type=Path, required=True)
    parser.add_argument("--shell-contract", type=Path, default=DEFAULT_SHELL_CONTRACT)
    parser.add_argument("--output", type=Path, default=None)
    args = parser.parse_args()

    report = audit_preview(args.deck, args.fill_payload, args.shell_contract)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(json.dumps(report, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report, indent=2))


if __name__ == "__main__":
    main()
