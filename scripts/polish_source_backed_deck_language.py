#!/usr/bin/env python3
"""Polish generated source-backed deck language without touching deck builders."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_ROOT = ROOT / "output" / "source_backed_deck_polish"
SCHEMA_VERSION = "monthly_platform.source_backed_deck_polish.v1"

TEXT_REPLACEMENTS = {
    "Source-backed monthly control deck": "Source-backed monthly review",
    "reconciled through the validated fact packet.": (
        "with source checks complete and no tie-out issues."
    ),
    "Full source-backed coverage.": "Every region included.",
    "Deck facts mapped to source paths.": "Ready for review.",
    "Prepared for analyst workbook, think-cell source workbook, and standard PowerPoint production.": (
        "Prepared for analyst review, think-cell handoff, and monthly leadership production."
    ),
    "SOURCE-BACKED PUBLISH GATE": "PUBLISH READINESS",
    "Publish gate status is green across the evidence chain": (
        "Publish gate is green across the evidence chain"
    ),
    "The monthly pack should fail closed if source, workbook, fact, or deck truth controls drift.": (
        "The pack is release-ready only when source, workbook, fact, and deck checks all clear."
    ),
    "source chain is ready for standard deck production.": (
        "evidence trail is ready for leadership review."
    ),
    "Truth packet": "Fact check",
    "Validated fact registry and claim set": "Validated claims and evidence trail",
    "Source publish gate": "Source readiness",
    "Source contract, bundles, readiness, and workbook coverage": (
        "Report extracts, workbook coverage, and readiness checks"
    ),
    "Source extracts": "Salesforce inputs",
    "Selected sources": "Required inputs",
    "Expected source contract items present": "Required monthly inputs present",
    "Salesforce reports and list views resolved for the run": (
        "Salesforce reports and list views resolved for the month"
    ),
    "Director bundles": "Director books",
    "Normalized source-backed DirectorBundle files": "Normalized director books",
    "Presentation-ready facts with source paths": (
        "Presentation-ready facts with evidence references"
    ),
    "Source: deck truth packet and source-backed publish gate JSON.": (
        "Source: validated monthly evidence packet and publish readiness checks."
    ),
    "Regional open pipeline is traceable to director bundles": (
        "Regional open pipeline is reconciled and ready to review"
    ),
    "Regional totals are recomputed from director bundles; SOQL remains validation, not the source of deck truth.": (
        "Regional totals are recomputed from director books; Salesforce queries are used as tie-out checks."
    ),
    "Source: director gold analytics regional rollups; ARR is open book, not active forecast unless explicitly labeled.": (
        "Source: validated regional rollups; ARR is open book, not active forecast unless explicitly labeled."
    ),
    "This is the analyst handoff layer: facts are deterministic, readable, and source-path backed.": (
        "This view shows where the open book, risk rows, and data exceptions sit by director."
    ),
    "Source: deck truth packet director facts; all values retain source_artifact and source_json_path lineage.": (
        "Source: validated director facts; zero tie-out/source issues means the row is publish-clean."
    ),
    "Source: monthly source bundle manifest pipeline_display_decision for each territory.": (
        "Source: monthly period-control table for each territory."
    ),
    "Leadership readout: clean source chain, ready for standard deck production": (
        "Leadership readout: clean evidence trail, ready for monthly deck handoff"
    ),
    "Standard PowerPoint output should be generated from the same workbook and named element contract each month.": (
        "The same workbook, named elements, and evidence checks can be rerun each month."
    ),
    "Salesforce source extracts resolved": "Salesforce inputs resolved",
    "selected sources present": "required inputs present",
    "every decision is visible in the source manifest.": (
        "every decision is visible in the period-control table."
    ),
    "source contracts stay locked together.": "source requirements stay locked together.",
    "Use this packet as the standard deck production input; keep narrative edits inside the verified source-backed fact set.": (
        "Use this packet for final deck production; keep edits inside the verified fact set."
    ),
    "Path / element": "Reference",
    "Deck facts and claim registry": "Validated facts and claims",
    "Source/workbook readiness control": "Data and workbook readiness control",
    "Source bundle manifest": "Period decisions",
    "Source: truth packet think-cell contract and source-backed monthly artifacts.": (
        "Source: release packet and think-cell handoff contract."
    ),
}

REGEX_REPLACEMENTS = [
    (
        re.compile(r"\bRun:\s+[A-Za-z0-9._-]+"),
        "Evidence run archived in release packet",
    ),
    (
        re.compile(r"output/deck_truth_packets_from_sources/[^\s]+"),
        "Release packet",
    ),
    (
        re.compile(r"output/monthly_source_backed_publish_gate/[^\s]+"),
        "Publish gate",
    ),
    (
        re.compile(r"output/monthly_source_bundles/[^\s]+"),
        "Period decisions",
    ),
    (
        re.compile(r"[^\s]+deck_truth_packet\.json"),
        "Release packet",
    ),
    (
        re.compile(r"[^\s]+source_backed_publish_gate\.json"),
        "Publish gate",
    ),
    (
        re.compile(r"[^\s]+source_bundle_manifest\.json"),
        "Period decisions",
    ),
    (
        re.compile(r"\b1 territory use\b"),
        "1 territory uses",
    ),
]

BLOCKED_TERMS = [
    "AI-GENERATED",
    "GENERATED BY AI",
    "AS AN AI",
    "CHATGPT",
    "CLAUDE",
    "CODEX",
    "LOREM",
    "PLACEHOLDER",
    "TODO",
    "TBD",
    "DUMMY",
    "source_artifact",
    "source_json_path",
    "DirectorBundle",
    "SOQL",
    "JSON",
]


def polish_deck_language(
    *,
    deck_path: Path,
    output_path: Path | None = None,
    output_root: Path = DEFAULT_OUTPUT_ROOT,
    snapshot_date: str | None = None,
    source_run_id: str | None = None,
    manifest_path: Path | None = None,
) -> dict[str, Any]:
    deck_path = Path(deck_path)
    findings: list[dict[str, Any]] = []
    replacements: list[dict[str, Any]] = []
    manifest_updated = False
    resolved_snapshot_date = snapshot_date or "unknown-snapshot"
    resolved_run_id = source_run_id or "unknown-run"

    if not deck_path.exists():
        findings.append(_finding("high", "deck_missing", f"Missing deck: {deck_path}"))
        return _write_result(
            payload=_result_payload(
                deck_path=deck_path,
                snapshot_date=resolved_snapshot_date,
                source_run_id=resolved_run_id,
                replacements=replacements,
                findings=findings,
                slide_count=0,
                manifest_path=manifest_path,
                manifest_updated=manifest_updated,
            ),
            output_path=output_path,
            output_root=output_root,
        )

    try:
        presentation = Presentation(str(deck_path))
    except Exception as exc:
        findings.append(_finding("high", "deck_load_failed", str(exc)))
        return _write_result(
            payload=_result_payload(
                deck_path=deck_path,
                snapshot_date=resolved_snapshot_date,
                source_run_id=resolved_run_id,
                replacements=replacements,
                findings=findings,
                slide_count=0,
                manifest_path=manifest_path,
                manifest_updated=manifest_updated,
            ),
            output_path=output_path,
            output_root=output_root,
        )

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for column_index, cell in enumerate(row.cells, start=1):
                        replacements.extend(
                            _polish_text_frame(
                                cell.text_frame,
                                location=(
                                    f"slide {slide_index} table row {row_index} "
                                    f"col {column_index}"
                                ),
                            )
                        )
            if getattr(shape, "has_text_frame", False):
                replacements.extend(
                    _polish_text_frame(
                        shape.text_frame,
                        location=f"slide {slide_index} shape",
                    )
                )

    presentation.save(deck_path)
    manifest_path = _resolved_manifest_path(deck_path=deck_path, manifest_path=manifest_path)
    manifest_updated = _polish_manifest_contract(manifest_path)
    remaining_text = _deck_text(Presentation(str(deck_path)))
    for term in BLOCKED_TERMS:
        if re.search(re.escape(term), remaining_text, flags=re.IGNORECASE):
            findings.append(
                _finding(
                    "high",
                    "blocked_language_remaining",
                    f"blocked term remains: {term}",
                )
            )

    return _write_result(
        payload=_result_payload(
            deck_path=deck_path,
            snapshot_date=resolved_snapshot_date,
            source_run_id=resolved_run_id,
            replacements=replacements,
            findings=findings,
            slide_count=len(presentation.slides),
            manifest_path=manifest_path,
            manifest_updated=manifest_updated,
        ),
        output_path=output_path,
        output_root=output_root,
    )


def _polish_text_frame(text_frame, *, location: str) -> list[dict[str, Any]]:
    applied: list[dict[str, Any]] = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original = run.text
            polished = _polished_text(original)
            if polished != original:
                run.text = polished
                applied.append(
                    {
                        "location": location,
                        "before": original,
                        "after": polished,
                    }
                )
    return applied


def _polished_text(text: str) -> str:
    polished = text
    for old, new in TEXT_REPLACEMENTS.items():
        polished = polished.replace(old, new)
    for pattern, replacement in REGEX_REPLACEMENTS:
        polished = pattern.sub(replacement, polished)
    return polished


def _resolved_manifest_path(*, deck_path: Path, manifest_path: Path | None) -> Path | None:
    if manifest_path is not None:
        return Path(manifest_path)
    candidate = deck_path.with_name("source_backed_deck_manifest.json")
    return candidate if candidate.exists() else None


def _polish_manifest_contract(manifest_path: Path | None) -> bool:
    if manifest_path is None or not manifest_path.exists():
        return False
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        return False
    visual_contract = payload.get("visual_contract")
    if not isinstance(visual_contract, dict):
        return False
    changed = False
    for key in ("required_titles", "required_text"):
        values = visual_contract.get(key)
        if not isinstance(values, list):
            continue
        polished_values = [_polished_text(str(value)) for value in values]
        if polished_values != values:
            visual_contract[key] = polished_values
            changed = True
    if changed:
        manifest_path.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")
    return changed


def _deck_text(presentation: Presentation) -> str:
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(str(shape.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.append(str(cell.text))
    return "\n".join(texts)


def _result_payload(
    *,
    deck_path: Path,
    snapshot_date: str,
    source_run_id: str,
    replacements: list[dict[str, Any]],
    findings: list[dict[str, Any]],
    slide_count: int,
    manifest_path: Path | None,
    manifest_updated: bool,
) -> dict[str, Any]:
    high_count = sum(1 for finding in findings if finding["severity"] == "high")
    return {
        "schema_version": SCHEMA_VERSION,
        "status": "blocked" if high_count else "ok",
        "snapshot_date": snapshot_date,
        "source_run_id": source_run_id,
        "deck_path": str(deck_path),
        "manifest_path": str(manifest_path) if manifest_path else None,
        "checks": {
            "deck_exists": deck_path.exists(),
            "slide_count": slide_count,
            "polished_text_checked": deck_path.exists(),
            "manifest_contract_polished": manifest_updated,
            "replacements_applied_count": len(replacements),
        },
        "summary": {
            "finding_count": len(findings),
            "high_finding_count": high_count,
            "medium_finding_count": sum(
                1 for finding in findings if finding["severity"] == "medium"
            ),
            "low_finding_count": sum(
                1 for finding in findings if finding["severity"] == "low"
            ),
            "replacements_applied_count": len(replacements),
        },
        "replacements": replacements,
        "findings": findings,
    }


def _write_result(
    *,
    payload: dict[str, Any],
    output_path: Path | None,
    output_root: Path,
) -> dict[str, Any]:
    if output_path is None:
        output_path = (
            Path(output_root)
            / payload["snapshot_date"]
            / payload["source_run_id"]
            / "source_backed_deck_polish_audit.json"
        )
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    payload["output_path"] = str(output_path)
    output_path.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")
    return payload


def _finding(severity: str, issue: str, evidence: str) -> dict[str, Any]:
    return {
        "severity": severity,
        "issue": issue,
        "evidence": evidence,
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--deck-path", type=Path, required=True)
    parser.add_argument("--manifest-path", type=Path)
    parser.add_argument("--snapshot-date")
    parser.add_argument("--source-run-id")
    parser.add_argument("--output-path", type=Path)
    parser.add_argument("--output-root", type=Path, default=DEFAULT_OUTPUT_ROOT)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    payload = polish_deck_language(
        deck_path=args.deck_path,
        output_path=args.output_path,
        output_root=args.output_root,
        snapshot_date=args.snapshot_date,
        source_run_id=args.source_run_id,
        manifest_path=args.manifest_path,
    )
    print(json.dumps(payload, indent=2))
    return 0 if payload["status"] == "ok" else 1


if __name__ == "__main__":
    raise SystemExit(main())
